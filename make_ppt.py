"""
Generate GaussTune STMM method design presentation.
Covers: problem, DB2 STMM algorithm, BRBE extension, Proactive BRBE, experimental design.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Colour palette ────────────────────────────────────────────────────────────
BLUE   = RGBColor(0x1F, 0x49, 0x7D)   # deep blue
LBLUE  = RGBColor(0x41, 0x72, 0xC4)   # medium blue
TEAL   = RGBColor(0x00, 0x70, 0x70)   # teal accent
ORANGE = RGBColor(0xC0, 0x55, 0x14)   # orange accent
RED    = RGBColor(0xC0, 0x00, 0x00)
GREEN  = RGBColor(0x37, 0x86, 0x23)
GREY   = RGBColor(0x40, 0x40, 0x40)
LGREY  = RGBColor(0xD9, 0xD9, 0xD9)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
BLACK  = RGBColor(0x00, 0x00, 0x00)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]   # completely blank


def add_slide():
    return prs.slides.add_slide(BLANK)


def rect(slide, l, t, w, h, fill=None, line=None):
    from pptx.util import Inches
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def textbox(slide, text, l, t, w, h,
            size=18, bold=False, color=BLACK, align=PP_ALIGN.LEFT,
            wrap=True, italic=False):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb


def header_bar(slide, title, subtitle=None):
    rect(slide, 0, 0, 13.33, 1.1, fill=BLUE)
    textbox(slide, title, 0.3, 0.08, 10, 0.55,
            size=28, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        textbox(slide, subtitle, 0.3, 0.65, 12, 0.40,
                size=16, color=LGREY, align=PP_ALIGN.LEFT)


def bullet_list(slide, items, l, t, w, h,
                size=16, color=BLACK, indent=0.25, spacing=0.42):
    for i, item in enumerate(items):
        prefix = "•  " if not item.startswith("–") else item[:3]
        text   = item if not item.startswith("–") else item[3:]
        x = l + (indent if item.startswith("–") else 0)
        y = t + i * spacing
        textbox(slide, prefix + text if not item.startswith("–") else "    –  " + text,
                x, y, w, 0.38, size=size, color=color)


# ── Slide 0: Title ─────────────────────────────────────────────────────────────
sl = add_slide()
rect(sl, 0, 0, 13.33, 7.5, fill=BLUE)
rect(sl, 0, 2.8, 13.33, 2.0, fill=LBLUE)
textbox(sl, "GaussTune STMM", 0.6, 1.0, 12, 1.0,
        size=44, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
textbox(sl, "Self-Tuning Memory Manager for Mixed TP+AP Workloads",
        0.6, 2.1, 12, 0.6, size=22, color=LGREY, align=PP_ALIGN.LEFT)
textbox(sl, "Method Design  ·  Algorithm Reference  ·  Python Implementation",
        0.6, 2.95, 12, 0.5, size=18, color=WHITE, align=PP_ALIGN.LEFT)
textbox(sl, "Based on: DB2 STMM (Zilio et al. 2006)  +  BRBE Extension  +  Proactive BRBE",
        0.6, 3.55, 12, 0.4, size=14, italic=True, color=LGREY, align=PP_ALIGN.LEFT)

# ── Slide 1: Problem Statement ────────────────────────────────────────────────
sl = add_slide()
header_bar(sl, "Problem: TP+AP Co-location Degrades Both Workloads",
           "Without memory management, OLTP and OLAP compete for shared resources")

# Two columns
rect(sl, 0.3, 1.25, 5.9, 5.8, fill=RGBColor(0xFF,0xF0,0xF0), line=RED)
rect(sl, 6.5, 1.25, 6.5, 5.8, fill=RGBColor(0xF0,0xF8,0xFF), line=LBLUE)

textbox(sl, "❶  Without Tuning", 0.55, 1.35, 5.5, 0.4, size=16, bold=True, color=RED)
textbox(sl, "work_mem = 64 MB  (default)", 0.55, 1.85, 5.5, 0.35, size=14, color=GREY)
textbox(sl, "shared_buffers = 128 MB  (default)", 0.55, 2.2, 5.5, 0.35, size=14, color=GREY)

for i, line in enumerate([
    "AP sort/hash >64 MB → spills to disk",
    "4 AP workers × 64 MB = 256 MB sort space",
    "Sort temp I/O competes with TP buffer I/O",
    "Buffer pool too small → frequent eviction",
    "TP queries re-read evicted pages → slow",
    "AP TPS:  low (sort spill overhead)",
    "TP TPS drop:  10–40% under heavy AP",
]):
    textbox(sl, "• " + line, 0.55, 2.65 + i*0.43, 5.5, 0.4, size=13, color=GREY)

textbox(sl, "❷  With DBA Expert Tuning", 6.75, 1.35, 5.8, 0.4, size=16, bold=True, color=GREEN)
textbox(sl, "work_mem = 512 MB  (per sort/hash)", 6.75, 1.85, 5.8, 0.35, size=14, color=GREY)
textbox(sl, "shared_buffers = 6144 MB  (6 GB)", 6.75, 2.2, 5.8, 0.35, size=14, color=GREY)

for i, line in enumerate([
    "400K-row sort (72 MB) fits in WM → no spill",
    "AP queries complete 1.5× faster",
    "Large buffer pool caches TP + AP data",
    "Eviction rate drops → TP re-read eliminated",
    "TP TPS drop: <5% during AP",
    "BUT: requires domain expertise",
    "AND: static — wrong for dynamic workloads",
]):
    textbox(sl, "✓ " + line, 6.75, 2.65 + i*0.43, 5.8, 0.4, size=13, color=GREY)

textbox(sl, "Goal: match expert performance automatically, starting from defaults",
        0.3, 7.05, 12.7, 0.38, size=15, bold=True, color=BLUE, align=PP_ALIGN.CENTER)


# ── Slide 2: DB2 STMM Overview ────────────────────────────────────────────────
sl = add_slide()
header_bar(sl, "DB2 STMM Algorithm Overview",
           "Zilio et al. (2006): OD + MIMO dual-loop controller for shared memory")

# Left: OD/MIMO description
rect(sl, 0.3, 1.2, 7.8, 5.9, fill=RGBColor(0xF5,0xF5,0xFF), line=LBLUE)
textbox(sl, "Core Algorithm  (stmm_controller.py: STMMController)",
        0.5, 1.3, 7.5, 0.4, size=15, bold=True, color=LBLUE)

code_lines = [
    "def tick(blks_hit, blks_read, temp_bytes, n_ap):",
    "  # 1. Compute benefit signals",
    "  wm_ben = temp_bytes_spilled × DISK_COST / wm_mb",
    "  sb_ben = blks_read × PAGE_MB × DISK_COST / sb_mb",
    "",
    "  # 2. MIMO: regression-based optimal size",
    "  if wm_model.valid:   Δwm = model.delta_size(wm_ben)",
    "",
    "  # 3. OD: explore when MIMO not yet valid",
    "  else:   Δwm = od_step(wm_ben)   # +/- probe",
    "",
    "  # 4. RECOVER: shrink toward WM_MIN when idle",
    "  if zero_benefit ≥ 3:  Δwm = -(wm - WM_MIN) × 0.2",
    "",
    "  new_wm = clamp(wm + Δwm, WM_MIN=64, WM_MAX=1024)",
    "  return new_wm, suggest_sb",
]
for i, line in enumerate(code_lines):
    textbox(sl, line, 0.5, 1.8 + i*0.27, 7.5, 0.28,
            size=11, color=BLACK if not line.startswith("  #") else TEAL,
            bold=not line.startswith("  "))

# Right: signal definitions
rect(sl, 8.3, 1.2, 4.8, 5.9, fill=RGBColor(0xF5,0xFF,0xF5), line=GREEN)
textbox(sl, "Benefit Signals", 8.5, 1.3, 4.5, 0.4, size=15, bold=True, color=GREEN)

signals = [
    ("wm_ben", "temp_bytes_delta × disk_cost / wm_mb", "MB of spill saved per extra MB of WM"),
    ("sb_ben", "blks_read × page_MB × disk_cost / sb_mb", "I/O seconds saved per extra MB of SB"),
]
y = 1.8
for name, formula, desc in signals:
    textbox(sl, name, 8.5, y, 4.5, 0.3, size=13, bold=True, color=BLUE)
    textbox(sl, formula, 8.5, y+0.28, 4.5, 0.3, size=11, color=BLACK)
    textbox(sl, desc, 8.5, y+0.55, 4.5, 0.3, size=11, italic=True, color=GREY)
    y += 0.95

textbox(sl, "Constants", 8.5, 3.85, 4.5, 0.35, size=13, bold=True, color=GREY)
consts = [
    "WM_MIN_MB = 64    WM_MAX_MB = 1024",
    "SB_INIT_MB = 2048  SB_MAX_MB = 8000",
    "WM_STEP_MIN = 64  WM_STEP_FINE = 8",
    "DISK_READ_COST ≈ 0.004 s/MB",
    "RECOVERY_INTS = 3  MAX_DEC_RATIO = 0.20",
]
for i, c in enumerate(consts):
    textbox(sl, c, 8.5, 4.25 + i*0.3, 4.5, 0.28, size=11, color=GREY)


# ── Slide 3: BRBE Extension ────────────────────────────────────────────────────
sl = add_slide()
header_bar(sl, "BRBE Extension: Joint WM + SB Trade-off",
           "BRBEController (stmm_controller.py) — weighted marginal benefit comparison")

rect(sl, 0.3, 1.2, 12.7, 1.55, fill=RGBColor(0xE8,0xF4,0xFF), line=LBLUE)
textbox(sl, "Why BRBE?  Memory is shared: WM and SB compete for the same RAM. "
        "Assigning more MB to SB means less for WM. "
        "BRBE computes marginal benefit of each and picks the best allocation.",
        0.5, 1.28, 12.3, 0.5, size=14, color=GREY)

textbox(sl, "BRBE Marginal Benefits  (common currency: saved seconds per MB)",
        0.5, 2.85, 12.3, 0.38, size=15, bold=True, color=LBLUE)
textbox(sl, "mb_wm = α × wm_ben      (α = spill reducibility ∈ [0,1])",
        0.5, 3.28, 6.0, 0.35, size=14, color=BLACK)
textbox(sl, "mb_sb = β × sb_ben      (β = read reducibility ∈ [0,1])",
        0.5, 3.65, 6.0, 0.35, size=14, color=BLACK)
textbox(sl, "α decays when WM grows but spill persists (hardware-bound I/O)",
        0.5, 4.05, 12.3, 0.35, size=13, italic=True, color=GREY)
textbox(sl, "β decays when SB grows but reads persist (data > addressable SB)",
        0.5, 4.4, 12.3, 0.35, size=13, italic=True, color=GREY)

# Two sub-boxes
rect(sl, 0.3, 4.85, 5.9, 2.3, fill=RGBColor(0xFF,0xF8,0xE8), line=ORANGE)
textbox(sl, "SB Suggestion Logic", 0.5, 4.92, 5.5, 0.35, size=14, bold=True, color=ORANGE)
for i, line in enumerate([
    "if mb_sb > mb_wm × 1.5:   grow SB by SB_STEP_MB",
    "if mb_wm > mb_sb × 1.5:   no SB growth this interval",
    "if TP hit_ratio < 0.95:   extrapolate working set",
    "if blks_read drops after SB↑:  β recovers toward 1",
]):
    textbox(sl, "• " + line, 0.5, 5.35 + i*0.38, 5.7, 0.35, size=12, color=GREY)

rect(sl, 6.5, 4.85, 6.5, 2.3, fill=RGBColor(0xF5,0xFF,0xF0), line=GREEN)
textbox(sl, "Code Reference", 6.7, 4.92, 6.2, 0.35, size=14, bold=True, color=GREEN)
for i, line in enumerate([
    "class BRBEController(STMMController):   # line 423",
    "  _update_alpha(spill_mb)               # line 465",
    "  _update_beta(blks_read_delta)         # line 476",
    "  _sb_benefit_brbe(blks_read, n_ap)     # line 489",
    "  _brbe_suggest_sb(wm_ben, sb_ben, n_ap)# line 519",
    "  tick(...)  → overrides STMMController # line 575",
]):
    textbox(sl, line, 6.7, 5.35 + i*0.38, 6.2, 0.35, size=12, color=GREY)

# Alpha/beta bar
rect(sl, 0.3, 2.45, 12.7, 0.35, fill=RGBColor(0xDD,0xEE,0xFF))
textbox(sl,
    "α=1 (spill is reducible)  →  DECAY when WM↑ but spill persists  →  α→0 (hardware-bound)    "
    "β=1 (reads are SB-reducible)  →  DECAY when SB↑ but reads persist  →  β→0 (data > RAM)",
    0.5, 2.48, 12.3, 0.3, size=11, color=GREY)


# ── Slide 4: Proactive BRBE ────────────────────────────────────────────────────
sl = add_slide()
header_bar(sl, "Proactive BRBE: Pre-AP WM and SB Prediction",
           "ProactiveBRBEController — uses EXPLAIN output to set WM/SB before AP starts")

rect(sl, 0.3, 1.2, 12.7, 1.45, fill=RGBColor(0xF0,0xF8,0xFF), line=LBLUE)
textbox(sl, "Key insight (PropoSQL §Sort model):  The exact work_mem threshold to avoid sort "
        "spill is computable from EXPLAIN output: input_bytes = rows × (width + 24B overhead).  "
        "Apply the right WM before AP starts → avoid the 300s reactive ramp-up period.",
        0.5, 1.28, 12.3, 0.5, size=13, italic=True, color=GREY)

# Left: WM prediction
rect(sl, 0.3, 2.75, 6.0, 4.35, fill=RGBColor(0xF5,0xF5,0xFF), line=LBLUE)
textbox(sl, "WM Threshold Model", 0.5, 2.82, 5.7, 0.38, size=15, bold=True, color=LBLUE)
textbox(sl, "_wm_threshold_for_sort(rows, width_b):", 0.5, 3.28, 5.7, 0.3, size=12, bold=True)
for i, line in enumerate([
    "# SORT_TUPLE_OVERHEAD_B = 24",
    "input_bytes = rows × (width + 24)",
    "wm_needed = ceil(input_bytes / 1MB)",
    "return clamp(wm_needed, 64, 1024)",
    "",
    "Example: 400K rows × 210B = 82 MB",
    "  → WM_rec = max(64, 82) = 128 MB",
    "  → rounds to next power of 2",
    "  → set WM=128MB before AP inject",
]):
    textbox(sl, line, 0.5, 3.65 + i*0.3, 5.7, 0.28,
            size=11, color=TEAL if line.startswith("#") else
            (ORANGE if line.startswith("Example") else BLACK))

# Right: SB prediction
rect(sl, 6.5, 2.75, 6.5, 4.35, fill=RGBColor(0xF5,0xFF,0xF0), line=GREEN)
textbox(sl, "SB Threshold Model", 6.7, 2.82, 6.2, 0.38, size=15, bold=True, color=GREEN)
textbox(sl, "_sb_threshold_for_ap(tp_blks_hit, tp_blks_read, ap_rows, ...):",
        6.7, 3.28, 6.2, 0.3, size=12, bold=True)
for i, line in enumerate([
    "tp_ws_mb = min(blks_hit×8KB/MB, sb_mb)",
    "if hit_ratio < 0.95:",
    "  tp_ws_mb = sb_mb / hit_ratio  # extrapolate",
    "",
    "ap_press = n_workers × rows × width",
    "           / 1MB × (1 − 0.8_ring_bypass)",
    "",
    "sb_needed = tp_ws_mb + ap_press",
    "return round_up(sb_needed, SB_STEP=1024)",
]):
    textbox(sl, line, 6.7, 3.65 + i*0.3, 6.2, 0.28,
            size=11, color=TEAL if "# " in line else BLACK)

# Timeline
rect(sl, 0.3, 7.08, 12.7, 0.35, fill=RGBColor(0xDD,0xEE,0xFF))
textbox(sl,
    "PRE phase (60s):  STMM calibrates on TP   →   EXPLAIN + predict_pre_ap() called at t=60s   "
    "→   WM/SB applied   →   AP injected at t=60s with optimal WM from day 0",
    0.5, 7.1, 12.3, 0.3, size=11, color=BLUE)


# ── Slide 5: Experimental Design ──────────────────────────────────────────────
sl = add_slide()
header_bar(sl, "Experimental Design",
           "2 workloads × 5 configs × fair warmup (scaled by SB size)")

rect(sl, 0.3, 1.2, 6.0, 5.9, fill=RGBColor(0xF5,0xF5,0xFF), line=LBLUE)
textbox(sl, "Workloads  (stmm_test.py: WORKLOADS)", 0.5, 1.3, 5.7, 0.38, size=15, bold=True, color=LBLUE)

textbox(sl, "SORT — WM sensitivity", 0.5, 1.8, 5.7, 0.35, size=14, bold=True, color=BLUE)
for i, line in enumerate([
    "SELECT k,c,pad FROM sbtest1",
    "WHERE id ≤ 400K ORDER BY c DESC",
    "Sort input: 400K × 186B = 72MB",
    "WM=64MB → spills to disk (~75MB)",
    "WM=256MB+ → in-memory quicksort",
    "Metric: AP query throughput",
]):
    textbox(sl, "  " + line, 0.5, 2.2 + i*0.33, 5.7, 0.3, size=12, color=GREY)

textbox(sl, "JOIN — SB sensitivity", 0.5, 4.3, 5.7, 0.35, size=14, bold=True, color=TEAL)
for i, line in enumerate([
    "sbtest1 JOIN sbtest2 WHERE id ≤ 500K",
    "ORDER BY s1.c, s2.pad",
    "Scans 2×93MB = 186MB per iteration",
    "4 workers × 186MB = 744MB AP footprint",
    "SB=2048MB: AP evicts 36% TP pages",
    "SB=6144MB: AP fits alongside TP data",
]):
    textbox(sl, "  " + line, 0.5, 4.7 + i*0.33, 5.7, 0.3, size=12, color=GREY)

rect(sl, 6.5, 1.2, 6.5, 5.9, fill=RGBColor(0xF5,0xFF,0xF0), line=GREEN)
textbox(sl, "Configs per workload", 6.7, 1.3, 6.2, 0.38, size=15, bold=True, color=GREEN)

configs = [
    ("1", "Static-Default", "WM=64MB  SB=2048MB", "baseline: no tuning"),
    ("2", "Static-Expert-WM", "WM=512MB  SB=2048MB", "optimal WM, baseline SB"),
    ("3", "Static-Expert-Full", "WM=512MB  SB=6144MB", "DBA-tuned upper bound"),
    ("4", "STMM+BRBE", "WM=64MB→adapt  SB=adapt", "reactive OD+MIMO controller"),
    ("5", "STMM+ProactiveBRBE", "WM=128MB→adapt  SB=adapt", "predict from EXPLAIN"),
]
for i, (n, name, params, desc) in enumerate(configs):
    y = 1.8 + i*1.1
    rect(sl, 6.7, y, 6.0, 0.95, fill=RGBColor(0xEE,0xF7,0xEE) if i%2==0 else WHITE, line=LGREY)
    textbox(sl, f"❶{n}  {name}", 6.85, y+0.05, 5.8, 0.3, size=13, bold=True, color=BLUE)
    textbox(sl, params, 6.85, y+0.32, 5.8, 0.28, size=12, color=BLACK)
    textbox(sl, desc, 6.85, y+0.58, 5.8, 0.28, size=11, italic=True, color=GREY)

textbox(sl,
    "Warmup = max(120s, 120s × SB/2048MB)  →  fair buffer pool fill fraction for all SB sizes",
    0.3, 7.1, 12.7, 0.32, size=12, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)


# ── Slide 6: Measurement Framework ────────────────────────────────────────────
sl = add_slide()
header_bar(sl, "Measurement & Timeline",
           "PRE→AP→POST three-phase measurement (stmm_test.py: run_sysbench_phase)")

# Timeline bar
rect(sl, 0.3, 1.25, 12.7, 0.9, fill=RGBColor(0xE8,0xEE,0xFF), line=LBLUE)
phases = [
    (0.3,  1.5,  "PRE\n60s",    LBLUE,  "TP only"),
    (2.1,  1.5,  "AP INJECT",   ORANGE, "AP starts"),
    (3.0,  1.5,  "AP Phase\n360s", RED, "TP + 4 AP workers"),
    (8.5,  1.5,  "AP KILL",     ORANGE, "AP ends"),
    (9.2,  1.5,  "POST\n180s",  GREEN,  "TP recovery"),
]
for lp, tp, label, color, _ in phases:
    textbox(sl, label, lp, tp, 1.5, 0.55, size=11, bold=True, color=color, align=PP_ALIGN.CENTER)

textbox(sl, "← pre_tps measured here →", 0.3, 1.88, 1.9, 0.3, size=9, color=LBLUE, align=PP_ALIGN.CENTER)
textbox(sl, "←————————  ap_tps measured here  ————————→", 3.0, 1.88, 5.5, 0.3, size=9, color=RED, align=PP_ALIGN.CENTER)
textbox(sl, "← post_tps →", 9.2, 1.88, 1.8, 0.3, size=9, color=GREEN, align=PP_ALIGN.CENTER)

# Left: metrics
rect(sl, 0.3, 2.3, 5.9, 4.8, fill=RGBColor(0xF5,0xF5,0xFF), line=LBLUE)
textbox(sl, "Metrics Computed", 0.5, 2.38, 5.5, 0.35, size=15, bold=True, color=LBLUE)
metrics = [
    ("pre_tps",  "avg TP TPS during PRE (60s)"),
    ("ap_tps",   "avg TP TPS during AP (360s)"),
    ("post_tps", "avg TP TPS during POST (180s)"),
    ("drop%",    "(pre - ap) / pre × 100"),
    ("recovery%","(post - pre) / pre × 100"),
    ("ap_queries","# AP completions in 360s  [TODO run 8]"),
]
for i, (name, desc) in enumerate(metrics):
    textbox(sl, name, 0.5, 2.85 + i*0.55, 1.7, 0.35, size=13, bold=True, color=BLUE)
    textbox(sl, desc, 2.3, 2.85 + i*0.55, 3.8, 0.35, size=12, color=GREY)

# Right: STMM thread
rect(sl, 6.5, 2.3, 6.5, 4.8, fill=RGBColor(0xF5,0xFF,0xF0), line=GREEN)
textbox(sl, "STMM Thread  (15s poll interval)", 6.7, 2.38, 6.2, 0.35, size=15, bold=True, color=GREEN)
for i, line in enumerate([
    "PRE: OD probe → calibrate baseline",
    "    → Proactive: EXPLAIN → predict WM/SB",
    "    → Apply WM change (instant, no restart)",
    "    → Apply SB change (DB restart + warmup)",
    "AP:  tick() every 15s → OD/MIMO/RECOVER",
    "    → WM applied immediately via ALTER",
    "    → SB suggestion queued (pending)",
    "AP KILL: pending SB → POST-AP restart",
    "    n_ap=0: RECOVER WM toward WM_MIN",
    "POST: SB change applied once (n_ap>0 guard)",
    "     Further changes blocked until next AP",
]):
    textbox(sl, "  " + line, 6.7, 2.85 + i*0.38, 6.2, 0.35, size=11,
            color=TEAL if ":" in line[:12] else GREY)


# ── Slide 7: Key Code Pointers ─────────────────────────────────────────────────
sl = add_slide()
header_bar(sl, "Key Code Locations",
           "stmm_controller.py + stmm_test.py — quick navigation reference")

entries = [
    # (file, class/func, line, description)
    ("stmm_controller.py", "STMMController",               "L55–405",  "Base OD+MIMO controller: WM+SB benefit, OD step, MIMO regression, RECOVER"),
    ("stmm_controller.py", "_wm_benefit()",                "L215–230", "temp_bytes × disk_cost / wm_mb — sort spill benefit signal"),
    ("stmm_controller.py", "_sb_benefit()",                "L232–245", "blks_read × page_MB × disk_cost / sb_mb — SB miss benefit"),
    ("stmm_controller.py", "_od_step_wm()",                "L248–268", "OD perturbation: +step if ben>0 / oscillation detect / flip direction"),
    ("stmm_controller.py", "_apply_transfer()",            "L270–290", "Clamp delta to [max_dec, max_inc], round to WM_STEP granularity"),
    ("stmm_controller.py", "STMMController.tick()",        "L295–380", "Main control loop: compute ben → OD/MIMO/RECOVER → return new_wm, suggest_sb"),
    ("stmm_controller.py", "BRBEController",               "L423–650", "Adds α/β reducibility, BRBE marginal comparison, SB suggestion logic"),
    ("stmm_controller.py", "BRBEController.tick()",        "L575–637", "Override: BRBE sb_benefit, update α/β, RECOVER with floor rounding"),
    ("stmm_controller.py", "ProactiveBRBEController",      "L652–785", "predict_pre_ap(), _wm_threshold_for_sort(), _sb_threshold_for_ap()"),
    ("stmm_test.py",       "run_sysbench_phase()",         "L460–610", "Phase runner: inject AP, parse TPS, handle SB changes post-AP"),
    ("stmm_test.py",       "Proactive block",              "L741–771", "EXPLAIN AP_SQL → predict_pre_ap() → apply WM/SB before AP injection"),
    ("stmm_test.py",       "stmm_thread_v2()",             "L661–694", "STMM poll loop: get_db_stats → tick → apply WM → queue SB if n_ap>0"),
    ("stmm_test.py",       "reset_between_runs()",         "L384–402", "DB restart + scaled warmup (120s×SB/2048MB) for fair buffer pool state"),
    ("stmm_test.py",       "run_workload()",               "L926–1059","Run all 5 configs per workload, print summary, save JSON"),
]

rows = entries
ncols = 4
col_w = [1.7, 2.0, 1.0, 7.5]
col_x = [0.3, 2.05, 4.1, 5.15]
headers = ["File", "Symbol", "Lines", "Description"]
header_y = 1.2
for j, (hdr, x, w) in enumerate(zip(headers, col_x, col_w)):
    rect(sl, x, header_y, w-0.05, 0.35, fill=BLUE)
    textbox(sl, hdr, x+0.05, header_y+0.03, w-0.15, 0.3, size=12, bold=True, color=WHITE)

for i, (file, sym, lines, desc) in enumerate(rows):
    y = header_y + 0.38 + i*0.37
    bg = RGBColor(0xF0,0xF5,0xFF) if i%2==0 else WHITE
    for j, (val, x, w) in enumerate(zip([file, sym, lines, desc], col_x, col_w)):
        rect(sl, x, y, w-0.05, 0.35, fill=bg, line=LGREY)
        textbox(sl, val, x+0.05, y+0.03, w-0.15, 0.3, size=10,
                color=TEAL if j == 1 else (GREY if j == 0 else BLACK))


# ── Slide 8: Known Issues & Run 8 Plan ─────────────────────────────────────────
sl = add_slide()
header_bar(sl, "Current Issues & Run 8 Fix Plan",
           "Issues identified from run 6 / run 7 analysis — fixes already coded for run 8")

rect(sl, 0.3, 1.2, 6.0, 5.9, fill=RGBColor(0xFF,0xF5,0xF0), line=RED)
textbox(sl, "Issues Found", 0.5, 1.28, 5.7, 0.38, size=15, bold=True, color=RED)

issues = [
    ("RECOVER stuck at WM=80MB",
     "BRBEController.tick() used round() with fine=True\n(WM_STEP_FINE=8): -3.2 rounds back to 0\n→ WM never reaches 64MB to detect spill",
     "Fixed: floor() with WM_STEP_FINE=8"),
    ("Proactive SB restart unfair comparison",
     "After SB 2048→6144MB restart, Phase 2 sysbench\nstarts fresh with AP at t=0 (no ramp-up)\n→ ap_tps includes cold-start penalty",
     "Run 8: disable proactive SB; WM-only proactive"),
    ("Insufficient warmup for large SB",
     "Expert-Full (6144MB) used 120s warmup:\nbuffer only 30% full → pre_tps=79 (unfair)\n→ Expert-Full appeared worse than Default",
     "Fixed: warmup = 120s × SB/2048MB\n= 360s for SB=6144MB"),
    ("POST-phase SB oscillation",
     "STMM kept queuing SB changes after AP\nbecause n_ap=0 check was missing\n→ 5+ DB restarts in POST phase",
     "Fixed: `if n_ap > 0` guard in stmm_thread"),
    ("Wrong metric for WM claim",
     "Sort spill hurts AP query throughput, not TP TPS\nCurrent code only measures TP TPS → no WM\ndifferentiation visible in sort workload",
     "Run 8: track AP query count in stmm_ap.sh"),
]
for i, (title, problem, fix) in enumerate(issues):
    y = 1.75 + i*1.07
    textbox(sl, f"❌  {title}", 0.5, y, 5.7, 0.32, size=12, bold=True, color=RED)
    textbox(sl, problem, 0.55, y+0.3, 5.5, 0.5, size=10, italic=True, color=GREY)

rect(sl, 6.5, 1.2, 6.5, 5.9, fill=RGBColor(0xF0,0xFF,0xF0), line=GREEN)
textbox(sl, "Run 8 Plan", 6.7, 1.28, 6.2, 0.38, size=15, bold=True, color=GREEN)

for i, (title, problem, fix) in enumerate(issues):
    y = 1.75 + i*1.07
    textbox(sl, f"✓  {fix}", 6.7, y+0.15, 6.2, 0.55, size=12, color=GREEN)


# ── Save ──────────────────────────────────────────────────────────────────────
out = "/home/node/GaussTune/GaussTune_STMM_Design.pptx"
prs.save(out)
print(f"Saved: {out}")
print(f"Slides: {len(prs.slides)}")
