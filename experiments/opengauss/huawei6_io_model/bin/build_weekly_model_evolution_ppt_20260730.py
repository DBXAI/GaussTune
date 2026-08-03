#!/usr/bin/env python3
"""Build a model-evolution presentation from hit replay to mixed TP/AP control."""

from __future__ import annotations

import csv
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches

from build_memory_autonomy_paper_ppt import (
    BLUE, GRAY, GREEN, INK, LIGHT, MID, ORANGE, PALE_GREEN, PALE_ORANGE,
    PALE_RED, PALE_TEAL, PURPLE, RED, TEAL, WHITE, add_arrow, add_bullets,
    add_callout, add_flow_box, add_kpi, add_notes, add_rect, add_table,
    add_text, add_title, set_bg,
)

ROOT = Path(__file__).resolve().parents[1]
LATEST = ROOT / "artifacts/00_latest"
OUT = LATEST / "Huawei5_model_evolution_acceptance_aligned_20260730.pptx"
OUTLINE = LATEST / "Huawei5_model_evolution_acceptance_aligned_20260730_outline.md"

RAW_REC = ROOT / "artifacts/90_archive/ppt_assets/ppt_raw_recommendation_best.png"
S5_HIT_TPS = ROOT / "artifacts/02_validation_figures/s5_20260716_17/s5_tp_sb_hit_vs_total_tps_20260716.png"
AP8_NUMERIC = ROOT / "results/saturated32_ap8_tps_prediction_eval_20260717/ap8_trace_tps_prediction_vs_actual.png"
WORKMEM = ROOT / "artifacts/01_current_joint_model/figures/all_query_workmem_prediction_vs_actual.png"
FIVE_STAGE = LATEST / "five_stage_saturated_tps_validation_20260726.png"
V8 = LATEST / "weekly_v8_controller_timeline_20260729.png"
RECOMMENDATIONS = LATEST / "five_stage_saturated_joint_recommendations_20260726.csv"


def footer(slide, page: int, source: str = "") -> None:
    add_rect(slide, 0.62, 7.17, 12.08, 0.012, MID)
    add_text(slide, "Huawei5 · 从命中率 Replay 到 TP/AP 混合性能推荐 · 2026-07-30", 0.62, 7.22, 7.3, 0.18, size=7.5, color=GRAY)
    if source:
        add_text(slide, f"来源：{source}", 7.15, 7.22, 5.0, 0.18, size=7, color=GRAY, align=PP_ALIGN.RIGHT)
    add_text(slide, str(page), 12.28, 7.22, 0.42, 0.18, size=8, color=GRAY, align=PP_ALIGN.RIGHT)


def add_image(slide, path: Path, x: float, y: float, w: float, h: float) -> None:
    with Image.open(path) as image:
        iw, ih = image.size
    scale = min(w / iw, h / ih)
    dw, dh = iw * scale, ih * scale
    slide.shapes.add_picture(str(path), Inches(x + (w - dw) / 2), Inches(y + (h - dh) / 2), width=Inches(dw), height=Inches(dh))


def read_rows(path: Path) -> list[dict[str, str]]:
    with path.open(newline="", encoding="utf-8") as handle:
        return list(csv.DictReader(handle))


def build() -> None:
    for path in [RAW_REC, S5_HIT_TPS, AP8_NUMERIC, WORKMEM, FIVE_STAGE, V8, RECOMMENDATIONS]:
        if not path.exists():
            raise FileNotFoundError(path)
    recs = read_rows(RECOMMENDATIONS)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    outline: list[str] = []

    # 1 Cover
    slide = prs.slides.add_slide(blank); set_bg(slide, INK)
    add_rect(slide, 0, 0, 0.22, 7.5, TEAL)
    add_text(slide, "这一周模型到底升级了什么？", 0.92, 1.05, 10.8, 0.50, size=24, color=PALE_TEAL, bold=True)
    add_text(slide, "从“预测命中率”到“推荐 TPS 最优配置”", 0.90, 1.73, 11.72, 0.90, size=37, color=WHITE, bold=True)
    add_text(slide, "再到 TP/AP 混合负载的在线稳定控制", 0.92, 2.80, 11.3, 0.58, size=23, color=MID, bold=True)
    add_rect(slide, 0.92, 3.76, 6.26, 0.56, TEAL, radius=True)
    add_text(slide, "Trace Replay · Plan/Spill · SB × work_mem · TP SLO", 1.08, 3.92, 5.94, 0.24, size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "核心说明：当前主模型预测“最优区域/配置排序”，并不直接生成每个点的精确 TPS 数值。", 0.94, 5.26, 11.36, 0.52, size=16, color=MID)
    add_text(slide, "模型演进说明版 · 2026-07-30", 0.94, 6.46, 4.50, 0.24, size=10, color=MID)
    add_notes(slide, "开场先纠正术语：推荐 TPS 最优点不等于直接回归 TPS 数值。")
    outline.append("1. 封面：从命中率到 TPS 最优配置，再到混合负载在线控制。")

    # 2 Three meanings
    slide = prs.slides.add_slide(blank); set_bg(slide)
    add_title(slide, "先把“预测 TPS”拆成三件不同的事", "之前的表述把三层能力混在了一起，这是难以理解的主要原因。", "01 · DEFINITION")
    add_callout(slide, "A · 预测 TPS 数值", "输入配置，直接输出 1280 TPS。需要稳定的跨负载 TPS 映射；当前主模型没有做。", 0.68, 1.76, 3.84, 2.06, RED, PALE_RED)
    add_callout(slide, "B · 推荐 TPS 最优配置", "预测 TP miss、AP spill 与内存安全，找出应进入最高 TPS 平台的 SB × work_mem。当前主模型做到的是这一层。", 4.75, 1.76, 3.84, 2.06, TEAL, PALE_TEAL)
    add_callout(slide, "C · 在线稳定 TPS", "不预测绝对数值，用真实 TP TPS 作为传感器，动态调整 AP CPU/I/O、SB 和后续 Query grant。当前控制器做到这一层。", 8.82, 1.76, 3.84, 2.06, GREEN, PALE_GREEN)
    add_rect(slide, 0.78, 4.32, 11.76, 1.25, INK, radius=True)
    add_text(slide, "本周真正的升级", 1.05, 4.63, 1.75, 0.32, size=17, color=WHITE, bold=True)
    add_text(slide, "从 A 的早期尝试转向 B：使用可解释 Replay 推荐最高 TPS 区域；再用 C 对纯 Replay 未覆盖的 CPU、设备队列和路径变化做在线纠偏。", 2.78, 4.52, 9.34, 0.62, size=15, color=WHITE, bold=True)
    add_text(slide, "后续所有页面都会标明：这一层到底预测了什么，TPS 在哪里出现。", 0.86, 6.05, 11.58, 0.42, size=18, color=TEAL, bold=True, align=PP_ALIGN.CENTER)
    footer(slide, 2, "joint_bidirectional_replay.py；TP SLO controller")
    outline.append("2. 定义：TPS 数值预测、TPS 最优配置推荐、在线 TPS 稳定是三层能力。")

    # 3 Last week
    slide = prs.slides.add_slide(blank); set_bg(slide)
    add_title(slide, "上周：单条 1504MB Trace 只能预测缓存命中率", "当时的目标是回答“换一个 SB，SB/OS/combined hit 会怎样”，还没有 AP 算子内存模型。", "02 · LAST WEEK")
    flow = [
        ("真实运行", "SB=1504MB\n五阶段、低 TP 并发", BLUE),
        ("页访问 Trace", "relation / block\ntime / hit / strategy", PURPLE),
        ("候选 SB Replay", "DB buffer 淘汰\nOS two-list cache", TEAL),
        ("输出", "SB hit / OS hit\ncombined hit", ORANGE),
        ("旧推荐", "最大 combined hit\n作为最佳 SB", GREEN),
    ]
    for i, (title, body, color) in enumerate(flow):
        x = 0.55 + i * 2.55
        add_flow_box(slide, title, body, x, 1.86, 2.18, 1.48, color, fill=LIGHT if i != 2 else PALE_TEAL)
        if i < 4:
            add_arrow(slide, x + 2.20, 2.60, x + 2.48, 2.60, color=GRAY, width=1.5)
    add_callout(slide, "输入中没有", "AP Query 的 Plan、Sort/Hash/Join 内存、spill、并发动态峰值，也没有 32-terminal 饱和 TP 行为。", 0.80, 4.08, 5.78, 1.40, RED, PALE_RED)
    add_callout(slide, "输出中没有", "预测 TPS、AP 性能、work_mem 推荐和资源配额；只有缓存层的命中率与磁盘页数。", 6.78, 4.08, 5.78, 1.40, ORANGE, PALE_ORANGE)
    add_text(slide, "原始结果：最佳 SB Top-1 仅 2/5；combined held-out MAE 约 3.29pp。", 0.84, 5.98, 11.70, 0.50, size=19, color=RED, bold=True, align=PP_ALIGN.CENTER)
    footer(slide, 3, "build_model_progress_ppt.py；原始 Trace Replay")
    outline.append("3. 上周模型：1504MB 单轨迹，只输出缓存命中率，Top-1 为 2/5。")

    # 4 Why hit not TPS
    slide = prs.slides.add_slide(blank); set_bg(slide)
    add_title(slide, "为什么“combined 命中率最高”不能直接当成“TPS 最高”", "这一轮失败推动了后面的饱和 TP 和 AP 执行路径建模。", "02 · GAP")
    add_image(slide, RAW_REC, 0.58, 1.58, 6.60, 4.70)
    add_callout(slide, "问题 1 · 曲线抵消", "SB hit 上升、OS conditional hit 下降，combined 可能几乎不变，无法稳定排序候选。", 7.42, 1.72, 5.12, 1.12, ORANGE, PALE_ORANGE)
    add_callout(slide, "问题 2 · OS hit 仍不是 SB hit", "Linux cache 命中仍要走数据库 miss/read 路径；在饱和 TP 下，代价明显高于直接 SB hit。", 7.42, 3.02, 5.12, 1.12, TEAL, PALE_TEAL)
    add_callout(slide, "问题 3 · TP 没有饱和", "S1-S4 只有约 2 terminals、40 TPS，吞吐被限速，SB 不同也看不出最高 TPS。", 7.42, 4.32, 5.12, 1.12, RED, PALE_RED)
    add_text(slide, "结论：必须先让 TP 饱和，再把 AP 的执行、spill 和资源竞争放进同一个候选评价。", 0.84, 6.38, 11.70, 0.42, size=17, color=INK, bold=True, align=PP_ALIGN.CENTER)
    footer(slide, 4, "原始推荐预测 vs 实测最优 SB")
    outline.append("4. 缺口：combined 抵消、OS hit 成本和低 TP 并发使命中率无法直接代表 TPS。")

    # 5 Evolution timeline
    slide = prs.slides.add_slide(blank); set_bg(slide)
    add_title(slide, "这一周的升级路径：不是一步变成 TPS，而是增加四层机制", "每一层都保留上层输出，并增加新的真实输入与可验证行为。", "03 · EVOLUTION")
    stages = [
        ("V0", "缓存命中", "单轨迹\nSB/OS/combined", GRAY),
        ("V1", "饱和 TP", "32 terminals\nTP-SB TPS 平台", BLUE),
        ("V2", "AP 算子", "Plan/work_mem\nspill/动态峰值", ORANGE),
        ("V3", "TP/AP 联合", "mixed trace\nSB × work_mem 物理 I/O", TEAL),
        ("V4", "在线闭环", "真实 TPS/AP 进展\nCPU/I/O/SB 动态控制", GREEN),
    ]
    for i, (code, title, body, color) in enumerate(stages):
        x = 0.55 + i * 2.55
        add_rect(slide, x, 1.76, 2.18, 2.02, LIGHT, line=color, radius=True, width=1.5)
        add_rect(slide, x + 0.70, 1.47, 0.78, 0.52, color, radius=True)
        add_text(slide, code, x + 0.76, 1.61, 0.66, 0.23, size=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, title, x + 0.12, 2.12, 1.94, 0.34, size=16, color=color, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, body, x + 0.16, 2.63, 1.86, 0.82, size=11, color=INK, align=PP_ALIGN.CENTER)
        if i < 4:
            add_arrow(slide, x + 2.20, 2.76, x + 2.48, 2.76, color=GRAY, width=1.5)
    data = [
        ["版本", "TPS 在哪里出现", "该版本实际输出"],
        ["V0", "不出现", "命中率、磁盘 miss、SB 推荐"],
        ["V1", "实测用于确认 TP-SB 平台规律", "饱和 TP 的候选平台，而非 TPS 数值"],
        ["V2/V3", "独立 TPS sweep 用于最终验证", "TP/AP 联合 I/O、内存安全与配置排序"],
        ["V4", "真实 TPS 是在线传感器", "升降档、冻结、回退和安全上界"],
    ]
    add_table(slide, data, 0.72, 4.32, 11.92, 2.14, widths=[1.42, 4.40, 6.10], font_size=10.3)
    footer(slide, 5, "模型版本演进")
    outline.append("5. 演进总图：缓存、饱和 TP、AP 算子、联合 Replay、在线闭环。")

    # 6 Saturated TP
    slide = prs.slides.add_slide(blank); set_bg(slide)
    add_title(slide, "升级一：用饱和 TP Trace 找“最高 TPS 平台”，不再看 combined 最大值", "最终五阶段重新采集 32-terminal、SB=1504MB、work_mem=1024MB 的饱和锚点轨迹。", "04 · SATURATED TP")
    add_image(slide, S5_HIT_TPS, 0.58, 1.54, 7.08, 4.72)
    add_callout(slide, "观察到的规律", "早期 S5 中 TP-SB hit 随 SB 上升，TPS 同步上升并从 1024MB 开始进入约 220 TPS 平台。", 7.86, 1.72, 4.66, 1.38, TEAL, PALE_TEAL)
    add_callout(slide, "为什么只看 TP-SB", "OS hit 虽避免物理盘，但仍经历数据库 buffer miss 路径；饱和 TP 的平台与直接 SB hit 更一致。", 7.86, 3.30, 4.66, 1.38, BLUE, LIGHT)
    add_callout(slide, "新规则", "保留 TP-SB hit 距本阶段最大值不超过 0.001 的候选；这一步只确定 TPS 平台区域。", 7.86, 4.88, 4.66, 1.38, ORANGE, PALE_ORANGE)
    add_text(slide, "注意：图中的 220 TPS 是早期 S5 12+2 口径；最终 32-terminal 五阶段验证使用新的饱和 trace。", 0.78, 6.45, 11.76, 0.32, size=12.5, color=RED, bold=True, align=PP_ALIGN.CENTER)
    footer(slide, 6, "s5_tp_sb_hit_vs_total_tps_20260716.png")
    outline.append("6. 升级一：32-terminal 饱和 TP，使用 TP-SB hit 识别 TPS 平台。")

    # 7 AP model
    slide = prs.slides.add_slide(blank); set_bg(slide)
    add_title(slide, "升级二：把 AP 从“缓存访问者”升级为“有 Plan 和内存生命周期的查询”", "work_mem 改变 Plan、spill 与动态峰值，这些状态会反向改变 TP 的 OS cache 和 I/O。", "04 · AP MODEL")
    data = [
        ["阶段", "AP Query", "AP 并发", "模型为每条 Query 计算"],
        ["S1", "Q1", "1", "Plan family、Sort/Hash/Agg grant、spill、峰值"],
        ["S2", "Q3", "1", "同上；长时间运行路径单独保留"],
        ["S3", "Q5 + Q7", "2", "两条 Query 生命周期峰值并发聚合"],
        ["S4", "Q9/Q13/Q18/Q21", "4", "四 Query 峰值、动态池硬约束和不可达 no-spill"],
        ["S5", "Q1/Q3/Q5/Q7", "4", "AP spill/cache 污染与高 TP 同时 replay"],
    ]
    add_table(slide, data, 0.66, 1.55, 12.00, 3.06, widths=[1.10, 2.76, 1.30, 6.84], font_size=10.5)
    add_flow_box(slide, "一次真实 Query", "实际 rows/width\n算子开始结束\nobserved spill", 0.76, 5.02, 2.28, 1.34, BLUE, fill=LIGHT)
    add_arrow(slide, 3.10, 5.69, 3.55, 5.69, color=GRAY)
    add_flow_box(slide, "候选 work_mem", "EXPLAIN Plan family\nHash/Sort/Agg 源码容量\n未执行 Plan 合成", 3.62, 5.02, 2.52, 1.34, PURPLE, fill=LIGHT)
    add_arrow(slide, 6.20, 5.69, 6.65, 5.69, color=GRAY)
    add_flow_box(slide, "Query 输出", "Plan 切换区间\nspill I/O\n动态峰值", 6.72, 5.02, 2.28, 1.34, ORANGE, fill=PALE_ORANGE)
    add_arrow(slide, 9.06, 5.69, 9.51, 5.69, color=GRAY)
    add_flow_box(slide, "阶段输出", "并发动态峰值\nAP spill 总 I/O\n内存安全/不可部署", 9.58, 5.02, 2.84, 1.34, GREEN, fill=PALE_GREEN)
    footer(slide, 7, "ONE_SHOT_SOURCE_PLAN_REPLAY.md")
    outline.append("7. 升级二：AP Plan、算子 grant、spill、生命周期和阶段并发。")

    # 8 Mixed mechanics
    slide = prs.slides.add_slide(blank); set_bg(slide)
    add_title(slide, "升级三：TP 和 AP 如何在同一个 Replay 中发生作用", "“混合预测”不是把两条独立曲线相加，而是让 AP 改变 TP 实际可用的缓存和 I/O 环境。", "04 · MIXED REPLAY")
    add_flow_box(slide, "候选 (SB, work_mem)", "例如 4096MB + 1150MB", 0.62, 2.14, 2.18, 1.16, PURPLE, fill=LIGHT)
    add_arrow(slide, 2.86, 2.72, 3.30, 2.72, color=GRAY)
    add_flow_box(slide, "AP 算子 Replay", "Plan / grant / spill\n并发动态峰值", 3.36, 1.66, 2.30, 1.36, ORANGE, fill=PALE_ORANGE)
    add_flow_box(slide, "TP 页 Replay", "候选 SB 淘汰\n只统计 TP hit/miss", 3.36, 3.34, 2.30, 1.36, TEAL, fill=PALE_TEAL)
    add_arrow(slide, 5.72, 2.34, 6.22, 2.34, color=GRAY)
    add_arrow(slide, 5.72, 4.02, 6.22, 4.02, color=GRAY)
    add_flow_box(slide, "OS cache 状态", "SB 与动态峰值扣减容量\nAP bulk/spill 进入 streaming inactive\nTP refault 页进入 active", 6.28, 2.14, 3.04, 1.96, BLUE, fill=LIGHT)
    add_arrow(slide, 9.38, 3.12, 9.88, 3.12, color=GRAY)
    add_flow_box(slide, "联合输出", "TP disk miss / I/O\nAP spill I/O\nMemAvailable / safety", 9.94, 2.14, 2.76, 1.96, GREEN, fill=PALE_GREEN)
    add_callout(slide, "AP 影响 TP 的三条已建模路径", "1. AP 扫描污染 OS cache；2. spill 临时页读写占用 I/O；3. 动态峰值压缩 OS cache 容量。", 0.78, 5.10, 5.74, 1.20, TEAL, PALE_TEAL)
    add_callout(slide, "尚未纯 Replay 的路径", "AP CPU、设备队列延迟和内核调度没有被转换成精确 TPS；这些由运行时真实 TPS 反馈处理。", 6.80, 5.10, 5.74, 1.20, RED, PALE_RED)
    footer(slide, 8, "JOINT_BIDIRECTIONAL_REPLAY.md")
    outline.append("8. 升级三：AP spill/动态峰值/cache 污染反向改变 TP miss。")

    # 9 Scoring
    slide = prs.slides.add_slide(blank); set_bg(slide)
    add_title(slide, "模型怎样从行为预测得到“TPS 最优配置”", "候选排序有四步；输出中没有 tps_pred 字段。", "04 · RECOMMENDATION")
    steps = [
        ("1", "排除不可部署", "动态峰值必须小于 max_dynamic_memory 可用池；MemAvailable 保留量必须达标。", RED),
        ("2", "进入饱和 TP 平台", "保留 TP-SB hit ≥ 本阶段最大值 − 0.001 的候选。", BLUE),
        ("3", "最小化联合物理 I/O", "Predicted I/O = TP 数据盘 I/O + AP spill 读写 I/O。", TEAL),
        ("4", "同等性能选更小内存", "依次比较动态峰值、SB + dynamic 总占用，保留运行时余量。", GREEN),
    ]
    y = 1.60
    for code, title, body, color in steps:
        add_rect(slide, 0.74, y, 0.66, 0.56, color, radius=True)
        add_text(slide, code, 0.79, y + 0.16, 0.56, 0.22, size=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, title, 1.70, y + 0.01, 2.60, 0.34, size=16, color=color, bold=True)
        add_text(slide, body, 4.30, y, 8.08, 0.58, size=13.3)
        y += 0.92
    add_rect(slide, 0.76, 5.38, 11.78, 0.94, INK, radius=True)
    add_text(slide, "输出示例字段", 1.02, 5.68, 1.74, 0.28, size=15, color=WHITE, bold=True)
    add_text(slide, "tp_sb_hit_rate · tp_disk_misses · spill_io_mb · dynamic_peak_mb · predicted_memavailable_mb · predicted_physical_io_mb", 2.72, 5.61, 9.45, 0.42, size=12.5, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "因此准确表述是：模型预测性能机制并推荐 TPS 最高平台配置；TPS 数值由独立实验验证。", 0.82, 6.54, 11.70, 0.34, size=15.5, color=TEAL, bold=True, align=PP_ALIGN.CENTER)
    footer(slide, 9, "joint_bidirectional_replay.py::recommend")
    outline.append("9. 候选排序：内存安全、TP-SB 平台、联合 I/O、最小内存。")

    # 10 Numeric branch
    slide = prs.slides.add_slide(blank); set_bg(slide)
    add_title(slide, "曾经做过“直接 TPS 数值预测”，但没有作为当前主模型", "这条 AP8 支线依赖历史 AP4 TPS 曲线和目标负载的 1504MB 实测锚点。", "05 · NUMERIC TPS BRANCH")
    add_image(slide, AP8_NUMERIC, 0.58, 1.52, 7.52, 4.92)
    add_callout(slide, "转换方法", "用 AP4 数据拟合 TPS = a + b·[-log(1−TP-SB hit)]，再用 AP8@1504MB 的真实 TPS 缩放整条曲线。", 8.28, 1.72, 4.28, 1.42, BLUE, LIGHT)
    add_kpi(slide, "17.05%", "7 个留出点 MAPE", 8.32, 3.52, 1.94, color=ORANGE)
    add_kpi(slide, "0.984", "曲线相关系数", 10.42, 3.52, 1.94, color=GREEN)
    add_callout(slide, "为什么没有采用", "它需要目标负载 TPS 锚点，且跨负载数学映射可能过拟合；只能说明曲线形状相关，不能称为纯 trace 预测。", 8.28, 5.04, 4.28, 1.26, RED, PALE_RED)
    add_text(slide, "结果虽找对 4096MB 最优点，但主线选择了更可解释的“机制预测 + 配置排序 + 独立 TPS 验证”。", 0.74, 6.53, 11.84, 0.30, size=13.5, color=INK, bold=True, align=PP_ALIGN.CENTER)
    footer(slide, 10, "evaluate_ap8_tps_prediction.py；AP8 metrics.json")
    outline.append("10. TPS 数值支线：AP4 拟合 + AP8 单锚点，MAPE 17.05%，未纳入主模型。")

    # 11 Operator validation
    slide = prs.slides.add_slide(blank); set_bg(slide)
    add_title(slide, "AP 子模型效果：多数 Query 的最小 no-spill work_mem 可被定位", "这一步提供 AP spill 和动态峰值，不直接预测 AP QPS。", "06 · EFFECT 1")
    add_image(slide, WORKMEM, 0.60, 1.50, 12.14, 4.78)
    add_callout(slide, "当前结果", "可观测边界 5/6；Plan 7/7、spill 6/7。", 0.80, 6.02, 5.72, 0.84, GREEN, PALE_GREEN)
    add_callout(slide, "边界", "Q5 切换 Plan；Q18/Q21 无可部署 no-spill 点。", 6.80, 6.02, 5.72, 0.84, ORANGE, PALE_ORANGE)
    footer(slide, 11, "all_query_workmem_prediction_vs_actual.png")
    outline.append("11. AP 子模型效果：Plan/spill/work_mem 边界。")

    # 12 Five stage walkthrough
    slide = prs.slides.add_slide(blank); set_bg(slide)
    add_title(slide, "五组静态挑战负载：用于验证推荐排序，不是验收方案的五个控制阶段", "P1-P5 分别独立启动、独立结束；下表不能解释阶段间应该执行的增减动作。", "06 · EFFECT 2")
    data = [["负载组", "TP 输入", "AP 输入", "推荐 SB", "推荐 work_mem", "Replay 预测量", "独立实测 TPS"]]
    ap = ["Q1×1", "Q3×1", "Q5+Q7×2", "Q9/Q13/Q18/Q21×4", "Q1/Q3/Q5/Q7×4"]
    actual = ["1105.90", "1207.15", "1272.86", "1236.88", "1323.87"]
    for i, row in enumerate(recs):
        data.append([
            f"P{i+1}", "32 terminals 饱和 trace", ap[i], f"{row['recommended_sb_mb']}MB",
            f"{row['recommended_work_mem_mb']}MB", "TP miss + AP spill + 峰值", actual[i],
        ])
    add_table(slide, data, 0.44, 1.58, 12.46, 3.54, widths=[0.70, 2.14, 2.26, 1.22, 1.50, 2.86, 1.78], font_size=8.7)
    add_callout(slide, "P1-P5 是局部目标点", "模型在 280 个 SB × work_mem 候选上计算 TP hit/miss、AP spill、动态峰值和内存安全，并冻结推荐。", 0.70, 5.52, 5.82, 1.08, TEAL, PALE_TEAL)
    add_callout(slide, "不能直接当作 S1-S5 动作", "这些 TPS 只验证各负载组的配置排序；它们没有验证 SB 与动态池在连续阶段间如何让渡。", 6.80, 5.52, 5.82, 1.08, RED, PALE_RED)
    footer(slide, 12, "stage_joint_recommendations.csv；saturated TPS validation")
    outline.append("12. 五组静态挑战负载：验证配置排序，不等于原验收五阶段动作。")

    # 13 Recommendation result
    slide = prs.slides.add_slide(blank); set_bg(slide)
    add_title(slide, "静态效果：P1-P5 推荐点均落在各自已测 TPS 最高点的 5% 内", "这证明局部配置排序有效，不证明连续五阶段动作已经通过。", "06 · EFFECT 3")
    add_image(slide, FIVE_STAGE, 0.58, 1.50, 12.18, 4.92)
    add_kpi(slide, "5/5", "阶段推荐通过", 0.90, 6.18, 2.58, color=GREEN)
    add_kpi(slide, "4.55%", "最大 TPS regret", 3.55, 6.18, 2.58, color=ORANGE)
    add_text(slide, "验证定义：regret = (同阶段已测最高 TPS − 推荐配置实测 TPS) / 已测最高 TPS。", 6.40, 6.32, 5.92, 0.34, size=13, color=INK, bold=True, align=PP_ALIGN.CENTER)
    footer(slide, 13, "five_stage_saturated_tps_validation_20260726.png")
    outline.append("13. 静态推荐效果：五组挑战负载的配置 regret 均小于 5%。")

    # 14 Acceptance state machine
    slide = prs.slides.add_slide(blank); set_bg(slide)
    add_title(slide, "原验收方案的 S1-S5：是一条连续压力轨迹中的五种控制状态", "AP 慢 SQL 从 0 持续增加，运行中 Query 跨状态保留；动作方向由内存压力和 TP 负载变化决定。", "07 · ACCEPTANCE STATE MACHINE")
    data = [
        ["状态", "触发条件", "必须发生的算法动作", "要保护的指标"],
        ["S1 内存富裕", "总内存未到 max", "逐步增加 AP 动态内存", "减少 spill、提升 AP 效率"],
        ["S2 触及上限", "SB + 动态池达到 max", "按 granule 降 SB，转给动态池", "buffer hit 降幅可控"],
        ["S3 保护基准", "AP 压力继续增加", "停止降 SB；降低每会话 AP grant", "守住低负载 TP"],
        ["S4 反压排队", "已无安全内存可让渡", "新 AP 入队；存量 AP 自然执行", "不冲击 TP 和存量 AP"],
        ["S5 基准突增", "TP 从低负载跃升到高负载", "阻塞新 AP；优雅降运行 AP；提高 SB", "TP 波动受控、0 cancel"],
    ]
    add_table(slide, data, 0.50, 1.55, 12.34, 4.34, widths=[1.52, 2.56, 4.70, 3.56], font_size=9.7)
    add_rect(slide, 0.72, 6.12, 11.88, 0.60, INK, radius=True)
    add_text(slide, "统一约束：SB + 实际动态内存 ≤ memory_target_max；在线扩缩单次 TPS 抖动 ≤3%；阶段 TP 保持率 ≥95%。", 0.92, 6.30, 11.48, 0.25, size=13.4, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    footer(slide, 14, "内存池动态调整方案验证.pptx")
    outline.append("14. 原验收五阶段：连续压力状态机及每阶段必须执行的动作。")

    # 15 Gap audit
    slide = prs.slides.add_slide(blank); set_bg(slide)
    add_title(slide, "差距审计：当前控制器只有部分机制，尚未按原场景完成动作验收", "问题来自实验语义变化，不是静态推荐表中的数字精度。", "07 · GAP AUDIT")
    data = [
        ["状态", "当前实现/实验", "与原要求是否一致"],
        ["S1 增 AP 内存", "静态 P1 反而选择 work_mem=1MB；闭环仅在 TP 健康时恢复 grant", "不一致"],
        ["S2 降 SB 让渡", "P1/P2 都推荐 4096MB；没有连续的 SB→动态池迁移", "不一致"],
        ["S3 停降 SB + 降 grant", "有安全 grant 与 graceful debt；但 P2→P3 的 SB 是 4096→8192", "机制部分具备，场景未验"],
        ["S4 新 AP 排队", "有 admission queue；基础 AP=4 时未必触发，压力 sweep 到 2×/3×才排队", "机制具备，基准场景未验"],
        ["S5 TP 突增 + 提 SB", "当前五组都用固定 800 TPS；P4/P5 都是 8192MB，没有真实低→高跳变", "不一致"],
    ]
    add_table(slide, data, 0.48, 1.52, 12.38, 4.52, widths=[2.10, 7.12, 3.16], font_size=9.5)
    add_callout(slide, "根因", "同一组 S1-S5 名字同时表示“固定 Query 集合”和“在线控制状态”，把局部最优点误写成了阶段动作。", 0.72, 6.20, 5.78, 0.66, RED, PALE_RED)
    add_callout(slide, "修正", "以后 P1-P5 只表示静态挑战负载；S1-S5 专用于连续验收状态机。", 6.78, 6.20, 5.78, 0.66, TEAL, PALE_TEAL)
    footer(slide, 15, "driver/controller/source scheme audit")
    outline.append("15. 差距审计：S1/S2/S5 不一致，S3/S4 只有部分机制、场景未验。")

    # 16 Online control
    slide = prs.slides.add_slide(blank); set_bg(slide)
    add_title(slide, "在线闭环是执行器，但还需要原方案的五阶段状态约束", "Replay 给安全候选；状态机规定动作方向；真实 TP TPS 与 AP 进展决定动作幅度和回退。", "08 · ONLINE CONTROL")
    add_image(slide, V8, 0.56, 1.48, 8.10, 4.98)
    add_callout(slide, "在线输入", "15 秒 TP TPS、AP CPU/I/O/wait、Query 进展、SB active/target、freezer 恢复结果。", 8.86, 1.72, 3.70, 1.24, BLUE, LIGHT)
    add_callout(slide, "在线输出", "AP CPU/I/O 档位、freeze、SB target，以及仅在 Query 边界生效的 work_mem grant。", 8.86, 3.18, 3.70, 1.24, TEAL, PALE_TEAL)
    add_callout(slide, "为什么必须闭环", "AP 执行路径会随时间变化；纯 Replay 尚不能精确覆盖 CPU 和设备队列，因此要由真实 TPS 兜底。", 8.86, 4.64, 3.70, 1.24, ORANGE, PALE_ORANGE)
    add_text(slide, "v8 已识别 80MiB/s 后期不安全并回退，但全生命周期仍有 10/222 个窗口低于 95%，尚未验收。", 0.72, 6.52, 11.86, 0.30, size=12.8, color=RED, bold=True, align=PP_ALIGN.CENTER)
    footer(slide, 16, "V8_RESULT.md；controller timeline")
    outline.append("16. 在线闭环：Replay、验收状态机和真实反馈三者分工。")

    # 17 What can/cannot
    slide = prs.slides.add_slide(blank); set_bg(slide)
    add_title(slide, "当前模型到底能预测什么、不能预测什么", "用这页作为报告中的统一口径，避免再次把“推荐”写成“精确 TPS 回归”。", "08 · CAPABILITY")
    data = [
        ["能力", "当前状态", "含义"],
        ["不同 SB 下 TP-SB/OS/combined hit", "可以", "页级 Trace Replay；held-out MAE 已验证"],
        ["不同 work_mem 下 Plan 与 spill 边界", "大部分可以", "Plan/源码/trace；新算子和全局 grant 仍有限制"],
        ["TP/AP 混合后的缓存、spill、动态峰值与物理 I/O", "可以", "AP 状态会反向改变 TP 的 OS cache 和 miss"],
        ["推荐接近最高 TPS 的 SB × work_mem", "已测网格可以", "五阶段最大 regret 4.55%"],
        ["每个候选配置的精确 TP TPS 数值", "当前不可以", "AP8 锚点支线 MAPE 17.05%，未进入主线"],
        ["AP 的精确完成时间/QPS", "当前不可以", "目前预测 spill/峰值，在线测量真实进展"],
        ["完整五阶段 TP 全生命周期保持率 ≥95%", "尚未通过", "S2 v8 仍有 10/222 越界窗口"],
    ]
    table = add_table(slide, data, 0.58, 1.54, 12.18, 4.94, widths=[4.12, 2.18, 5.88], font_size=9.7)
    add_text(slide, "准确结论：我们已经从“命中率预测”升级到“TP/AP 执行机制联合预测 + TPS 最优配置推荐 + 在线 TPS 反馈控制”。", 0.80, 6.62, 11.74, 0.32, size=14, color=TEAL, bold=True, align=PP_ALIGN.CENTER)
    footer(slide, 17, "当前能力边界")
    outline.append("17. 能力边界：静态推荐已验证，连续五阶段动作尚未完整验收。")

    # 18 Next step
    slide = prs.slides.add_slide(blank); set_bg(slide)
    add_title(slide, "下一步：用连续负载重新验收五阶段动作，而不是继续比较五个静态点", "保留纯 Replay 的候选生成，同时补上动作方向、阶段触发和跨阶段运行态。", "09 · NEXT")
    add_callout(slide, "1 · 连续压力输入", "S1-S4 保持低 TP，AP 到达率持续上升且 Query 跨阶段运行；S5 再将 TP offered load 阶跃到高档。", 0.72, 1.68, 3.76, 1.70, TEAL, PALE_TEAL)
    add_callout(slide, "2 · 显式状态机", "按 headroom、TP retention、队列和 TP 负载触发 S1-S5，强制检查增 grant、降 SB、冻结 SB、排队和反向让渡的顺序。", 4.78, 1.68, 3.76, 1.70, BLUE, LIGHT)
    add_callout(slide, "3 · 双重验收", "每次 granule 迁移抖动 ≤3%；各阶段 TP retention ≥95%；全部 AP 自然完成，0 cancel、0 restart。", 8.84, 1.68, 3.76, 1.70, GREEN, PALE_GREEN)
    add_rect(slide, 0.76, 4.02, 11.82, 1.54, INK, radius=True)
    add_text(slide, "最终研究目标", 1.05, 4.36, 1.64, 0.32, size=17, color=WHITE, bold=True)
    add_text(slide, "一次负载 + 源码/trace 生成候选 → 联合预测 TP miss 与 AP 路径代价 → 在线真实 TPS/AP 进展纠偏 → 五阶段自然完成且 TP 保持率 ≥95%。", 2.66, 4.22, 9.40, 0.72, size=16, color=WHITE, bold=True)
    add_text(slide, "验收仍然以完整自然生命周期为准，不能只截取前 180 秒，也不能取消未完成 AP Query。", 0.84, 6.13, 11.68, 0.44, size=16, color=RED, bold=True, align=PP_ALIGN.CENTER)
    footer(slide, 18, "原方案对齐后的连续验收计划")
    outline.append("18. 下一步：连续压力、显式五阶段状态机和双重 TPS 验收。")

    prs.save(OUT)
    OUTLINE.write_text(
        "# Huawei5 模型演进与验收方案对齐提纲（2026-07-30）\n\n"
        + "\n".join(f"- {item}" for item in outline)
        + "\n\n## 核心口径\n\n"
        + "- 当前主模型不直接生成每个配置的 TPS 数值。\n"
        + "- 当前主模型预测 TP hit/miss、AP Plan/spill、动态峰值和联合物理 I/O，据此推荐 TPS 最高平台配置。\n"
        + "- 推荐点 TPS 由冻结预测后的独立实验验证；在线控制使用真实 TPS 反馈。\n"
        + "- AP8 的直接 TPS 数值预测依赖 AP4 映射和目标 1504MB 锚点，MAPE 17.05%，不属于当前纯 Replay 主线。\n",
        encoding="utf-8",
    )
    print(OUT)
    print(OUTLINE)


if __name__ == "__main__":
    build()
