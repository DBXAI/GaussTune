#!/usr/bin/env python3
"""Build the Huawei5 weekly progress presentation for 2026-07-23..29."""

from __future__ import annotations

import csv
import json
from pathlib import Path

import matplotlib.pyplot as plt
import numpy as np
from matplotlib import font_manager
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches

from build_memory_autonomy_paper_ppt import (
    BLUE,
    FONT,
    GRAY,
    GREEN,
    INK,
    LIGHT,
    MID,
    ORANGE,
    PALE_GREEN,
    PALE_ORANGE,
    PALE_RED,
    PALE_TEAL,
    PURPLE,
    RED,
    TEAL,
    WHITE,
    add_arrow,
    add_bullets,
    add_callout,
    add_flow_box,
    add_kpi,
    add_notes,
    add_rect,
    add_table,
    add_text,
    add_title,
    configure_plot_font,
    rgb_hex,
    set_bg,
)


ROOT = Path(__file__).resolve().parents[1]
LATEST = ROOT / "artifacts" / "00_latest"
LATEST.mkdir(parents=True, exist_ok=True)

OUT = LATEST / "Huawei5_weekly_progress_20260729.pptx"
OUTLINE = LATEST / "Huawei5_weekly_progress_20260729_outline.md"
CHART_VALIDATION = LATEST / "weekly_tp_validation_20260729.png"
CHART_TIMELINE = LATEST / "weekly_v8_controller_timeline_20260729.png"
CHART_MODEL = LATEST / "paper_model_validation_summary_20260727.png"

TP_ONLY = ROOT / "results" / "tp_only_long_v6_sb4096_rate800_20260729"
V8 = ROOT / "results" / "tp_slo_dynamic_ap_resource_v8_stage2_rate800_20260729"


def footer(slide, page: int, source: str = "") -> None:
    add_rect(slide, 0.62, 7.17, 12.08, 0.012, MID)
    add_text(
        slide,
        "Huawei5 · 一周进展 · Trace Replay 与 TP 优先在线闭环 · 2026-07-29",
        0.62,
        7.22,
        7.2,
        0.18,
        size=7.5,
        color=GRAY,
    )
    if source:
        add_text(slide, f"来源：{source}", 7.25, 7.22, 4.85, 0.18, size=7, color=GRAY, align=PP_ALIGN.RIGHT)
    add_text(slide, str(page), 12.28, 7.22, 0.42, 0.18, size=8, color=GRAY, align=PP_ALIGN.RIGHT)


def load_json(path: Path) -> dict:
    return json.loads(path.read_text(encoding="utf-8"))


def configure_weekly_plot_font() -> None:
    font_path = "/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc"
    font_manager.fontManager.addfont(font_path)
    name = font_manager.FontProperties(fname=font_path).get_name()
    plt.rcParams["font.sans-serif"] = [name]
    plt.rcParams["axes.unicode_minus"] = False


def make_validation_chart() -> None:
    configure_weekly_plot_font()
    tp = load_json(TP_ONLY / "summary.json")
    v8 = load_json(V8 / "summary.json")
    stage = v8["stage_results"]["stage2_reach_limit"]

    labels = ["TP-only\n1小时", "v8 前180秒\n准入窗口", "v8 全生命周期\n自然完成"]
    means = [100 * tp["mean_retention"], 100 * stage["control_window_mean_retention"], 100 * stage["full_lifecycle_mean_retention"]]
    minimums = [100 * tp["minimum_retention"], 100 * min(
        float(r["tp_retention_ratio"])
        for r in csv.DictReader((V8 / "ap_resource_actions.csv").open(encoding="utf-8"))
        if r["phase"] == "admission_window"
    ), 100 * stage["full_lifecycle_min_retention"]]
    violations = [tp["violating_windows"], stage["violating_control_windows"], stage["full_lifecycle_violating_control_windows"]]

    x = np.arange(3)
    fig, axes = plt.subplots(1, 2, figsize=(12.2, 4.2), gridspec_kw={"width_ratios": [1.7, 1]})
    axes[0].axhspan(95, 105, color="#e6f2ec", alpha=0.9)
    axes[0].axhline(95, color=rgb_hex(RED), linestyle="--", linewidth=1.5, label="95% 验收线")
    axes[0].bar(x - 0.18, means, 0.36, color=rgb_hex(TEAL), label="平均保持率")
    axes[0].bar(x + 0.18, minimums, 0.36, color=rgb_hex(ORANGE), label="最低保持率")
    axes[0].set_xticks(x, labels)
    axes[0].set_ylim(80, 110)
    axes[0].set_ylabel("TP 保持率（%）")
    axes[0].set_title("只看前 180 秒会掩盖自然收尾期失败", fontweight="bold")
    axes[0].legend(frameon=False, ncol=3, loc="upper center")
    axes[0].grid(axis="y", alpha=0.18)
    for i, (m, mn) in enumerate(zip(means, minimums)):
        axes[0].text(i - 0.18, m + 0.55, f"{m:.2f}", ha="center", fontsize=9)
        axes[0].text(i + 0.18, mn + 0.55, f"{mn:.2f}", ha="center", fontsize=9)

    colors = [rgb_hex(GREEN), rgb_hex(GREEN), rgb_hex(RED)]
    bars = axes[1].bar(labels, violations, color=colors, width=0.58)
    axes[1].set_title("低于 95% 的 15 秒窗口", fontweight="bold")
    axes[1].set_ylabel("窗口数")
    axes[1].grid(axis="y", alpha=0.18)
    for bar, value in zip(bars, violations):
        axes[1].text(bar.get_x() + bar.get_width() / 2, value + 0.25, str(value), ha="center", fontweight="bold")
    for ax in axes:
        ax.spines[["top", "right"]].set_visible(False)
    fig.tight_layout(w_pad=2.0)
    fig.savefig(CHART_VALIDATION, dpi=190, bbox_inches="tight")
    plt.close(fig)


def make_timeline_chart() -> None:
    configure_weekly_plot_font()
    with (V8 / "ap_resource_actions.csv").open(newline="", encoding="utf-8") as handle:
        rows = [r for r in csv.DictReader(handle) if r["phase"] != "stage_start"]
    with (V8 / "controller_actions.csv").open(newline="", encoding="utf-8") as handle:
        controls = list(csv.DictReader(handle))

    x = np.arange(1, len(rows) + 1)
    retention = np.array([100 * float(r["tp_retention_ratio"]) for r in rows])
    io_level = np.array([float(r["observed_io_mib_per_second"]) for r in rows])
    frozen = np.array([r["observed_ap_frozen"].lower() == "true" for r in rows])
    sb = np.array([float(r["sb_mb"]) for r in controls[: len(rows)]])

    fig, axes = plt.subplots(3, 1, figsize=(12.4, 6.2), sharex=True, gridspec_kw={"height_ratios": [1.55, 1, 0.9]})
    axes[0].axhspan(95, 105, color="#e7f2ec", alpha=0.9)
    axes[0].axhline(95, color=rgb_hex(RED), linestyle="--", linewidth=1.3)
    axes[0].plot(x, retention, color=rgb_hex(BLUE), linewidth=1.45)
    axes[0].set_ylabel("TP 保持率 (%)")
    axes[0].set_title("v8 探索运行：执行路径变化后，原安全档位不再安全", fontweight="bold")

    axes[1].step(x, io_level, where="post", color=rgb_hex(ORANGE), linewidth=1.8, label="实际生效 I/O 档位")
    axes[1].set_ylabel("AP I/O (MiB/s)")
    axes[1].set_yticks([20, 40, 80])
    axes[1].legend(frameon=False, loc="upper right")

    axes[2].step(x, sb, where="post", color=rgb_hex(TEAL), linewidth=1.8)
    axes[2].set_ylabel("SB (MB)")
    axes[2].set_xlabel("15 秒控制窗口（自然收尾期也计入）")

    for i, is_frozen in enumerate(frozen, start=1):
        if is_frozen:
            for ax in axes:
                ax.axvspan(i - 0.5, i + 0.5, color="#c8cdd1", alpha=0.42, linewidth=0)
    for ax in axes:
        ax.grid(alpha=0.17)
        ax.spines[["top", "right"]].set_visible(False)
    axes[0].annotate("80 后期越界", xy=(138, retention[137]), xytext=(103, 88), arrowprops={"arrowstyle": "->", "color": rgb_hex(RED)}, color=rgb_hex(RED), fontsize=10)
    axes[1].annotate("冻结确认因果后\n80→40→20", xy=(148, 20), xytext=(164, 63), arrowprops={"arrowstyle": "->", "color": rgb_hex(TEAL)}, color=rgb_hex(TEAL), fontsize=10)
    fig.tight_layout(h_pad=0.5)
    fig.savefig(CHART_TIMELINE, dpi=190, bbox_inches="tight")
    plt.close(fig)


def build() -> None:
    make_validation_chart()
    make_timeline_chart()
    tp = load_json(TP_ONLY / "summary.json")
    v8 = load_json(V8 / "summary.json")
    stage = v8["stage_results"]["stage2_reach_limit"]
    ap = v8["ap_stage_results"]["stage2_reach_limit"]

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    outline: list[str] = []

    # 1 Cover
    slide = prs.slides.add_slide(blank)
    set_bg(slide, INK)
    add_rect(slide, 0, 0, 0.22, 7.5, TEAL)
    add_text(slide, "从静态联合推荐到 TP 优先在线闭环", 0.92, 1.05, 10.8, 0.48, size=23, color=PALE_TEAL, bold=True)
    add_text(slide, "Huawei5 一周研发进展", 0.92, 1.70, 11.1, 0.75, size=38, color=WHITE, bold=True)
    add_text(slide, "Trace Replay · Plan/Spill · SB/work_mem · 动态配额 · 全生命周期验收", 0.94, 2.73, 11.4, 0.42, size=17, color=MID)
    add_rect(slide, 0.92, 3.56, 6.20, 0.56, TEAL, radius=True)
    add_text(slide, "Huawei5 五阶段 TP/AP 混合负载", 1.05, 3.72, 5.94, 0.25, size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "核心目标：AP 必须自然完成，同时五阶段 TP 相对波动控制在 ±5%", 0.94, 5.26, 11.1, 0.40, size=15, color=MID)
    add_text(slide, "周报范围：2026-07-23 — 2026-07-29", 0.94, 6.44, 5.4, 0.22, size=9, color=MID)
    add_notes(slide, "本周不是继续做一个静态配置表，而是把已有预测模型推进成在线执行和验收闭环。")
    outline.append("1. 封面：从静态联合推荐到 TP 优先在线闭环。")

    # 2 Executive summary
    slide = prs.slides.add_slide(blank); set_bg(slide)
    add_title(slide, "一周结论先行", "模型链路、内核执行和控制器均有实质进展；但五阶段全生命周期验收尚未通过。", "01 · EXECUTIVE SUMMARY")
    add_kpi(slide, "97/97", "自动化测试通过", 0.72, 1.70, 2.55, color=GREEN, note="控制器、replay、内核接口")
    add_kpi(slide, "0", "AP SQL 被取消", 3.45, 1.70, 2.55, color=GREEN, note="Q3 保持同一 session/plan")
    add_kpi(slide, "5535.6s", "Q3 自然完成", 6.18, 1.70, 2.55, color=BLUE, note="SF85，32 terminals + 800 TPS")
    add_kpi(slide, "10", "全生命周期越界窗口", 8.91, 1.70, 2.55, color=RED, note="v8：低于 95%")
    add_callout(slide, "已经完成", "一次负载/源码 replay、SB+work_mem 联合候选、运行态 SB、动态 CPU/I/O 档位、freezer 因果验证、自然结束验收。", 0.82, 3.62, 5.75, 1.52, TEAL, PALE_TEAL)
    add_callout(slide, "尚未通过", "v8 前 180 秒无越界，但自然收尾最低仅 84.50%；因此不能宣称五阶段 TP 波动已小于 5%。", 6.78, 3.62, 5.75, 1.52, RED, PALE_RED)
    add_text(slide, "本周最重要的认识：在线探索本身会污染 page cache；生产运行必须复用已学习的安全路径上界。", 0.72, 5.62, 11.9, 0.58, size=18, color=INK, bold=True, align=PP_ALIGN.CENTER)
    footer(slide, 2, "v8 summary.json、TP-only summary.json")
    add_notes(slide, "先讲边界：功能做出来了，探索也找到原因了，但最终验收没有通过。")
    outline.append("2. 摘要：本周完成项、Q3 自然结束与全生命周期未通过。")

    # 3 Goal evolution
    slide = prs.slides.add_slide(blank); set_bg(slide)
    add_title(slide, "目标从“推荐配置”升级为“实时稳定 TP”", "静态最优点只是初始条件；AP 执行路径变化时，资源竞争也会变化。", "02 · PROBLEM")
    add_flow_box(slide, "过去：离线推荐", "给出每阶段\nSB + work_mem", 0.72, 1.82, 2.52, 1.55, BLUE, fill=LIGHT)
    add_arrow(slide, 3.35, 2.60, 4.15, 2.60, color=GRAY)
    add_flow_box(slide, "本周：在线观测", "TP TPS / AP 进展\nSB / CPU / I/O / freeze", 4.26, 1.82, 2.80, 1.55, TEAL, fill=PALE_TEAL)
    add_arrow(slide, 7.18, 2.60, 7.98, 2.60, color=GRAY)
    add_flow_box(slide, "最终：闭环执行", "TP ≥ 95%\nAP 不饿死且自然结束", 8.10, 1.82, 3.04, 1.55, GREEN, fill=PALE_GREEN)
    add_callout(slide, "验收约束", "32 terminals、固定 offered 800 TPS；五阶段使用同一参考口径；AP 严重影响 TP 时允许限速或暂停，但不能取消 SQL。", 0.82, 4.10, 5.62, 1.42, ORANGE, PALE_ORANGE)
    add_callout(slide, "为什么不能只看命中率", "命中率可定位 SB/OS cache，但 TPS 还受 I/O 队列、算子内存、spill、CPU、并发和执行路径阶段影响。", 6.76, 4.10, 5.62, 1.42, PURPLE, LIGHT)
    footer(slide, 3, "项目验收目标与本周实验口径")
    add_notes(slide, "说明项目意义：实时把资源从 AP 让给 TP，但又不能让 AP 永久不完成。")
    outline.append("3. 目标演进：从静态 SB/work_mem 推荐到 TP 优先在线控制。")

    # 4 Existing model
    slide = prs.slides.add_slide(blank); set_bg(slide)
    add_title(slide, "已有模型：不是 TPS 回归，而是执行行为 replay", "预测输入来自页访问、Plan、算子生命周期和源码容量规则；实测 TPS 不回写训练。", "03 · MODEL")
    boxes = [
        ("一次锚点负载", "页访问 + 算子 trace", TEAL),
        ("Plan/源码 replay", "不同 work_mem 的 plan 与 spill", BLUE),
        ("双向联合搜索", "每个 SB × work_mem 重算", PURPLE),
        ("约束推荐", "TP I/O + AP spill + 内存安全", GREEN),
    ]
    x0 = 0.66
    for i, (title, sub, color) in enumerate(boxes):
        x = x0 + i * 3.12
        add_flow_box(slide, title, sub, x, 1.75, 2.55, 1.44, color, fill=LIGHT if i % 2 else PALE_TEAL)
        if i < len(boxes) - 1:
            add_arrow(slide, x + 2.58, 2.47, x + 3.02, 2.47, color=GRAY)
    add_bullets(slide, [
        "缓存层：模拟 shared_buffers 淘汰，以及 Linux active/inactive、refault 对 TP 热页的保护。",
        "算子层：HashJoin / HashAggregate / Sort 生命周期；预测最小不 spill work_mem 与动态峰值。",
        "联合层：work_mem 改变 spill 与 OS cache，SB 又改变剩余动态内存，两者必须双向重算。",
        "输出：不是单一命中率，而是候选配置、置信度、内存边界、物理 I/O 与执行约束。",
    ], 0.82, 3.72, 11.55, 2.38, size=14, spacing=10)
    footer(slide, 4, "ONE_SHOT_REPLAY_V2、JOINT_BIDIRECTIONAL_REPLAY")
    add_notes(slide, "强调模型基于行为和代码，不用测试集 TPS 去拟合推荐点。")
    outline.append("4. 已有模型：一次锚点、源码 Plan/Spill replay 与 SB+work_mem 双向联合搜索。")

    # 5 Model evidence
    slide = prs.slides.add_slide(blank); set_bg(slide)
    add_title(slide, "预测层已有证据与边界", "本周在线控制是在既有 replay 结果之上增加执行层，不替换原模型。", "03 · MODEL")
    slide.shapes.add_picture(str(CHART_MODEL), Inches(0.72), Inches(1.55), width=Inches(11.92), height=Inches(3.35))
    add_callout(slide, "可以支撑", "缓存 held-out 误差、Plan family 和大部分 spill 分类已通过留出验证；可用于生成候选和解释误差来源。", 0.82, 5.15, 5.62, 1.18, GREEN, PALE_GREEN)
    add_callout(slide, "不能夸大", "扩大未见 Query 后 spill 分类曾降到 15/21；五阶段推荐仅证明已测网格，不是连续空间全局最优。", 6.78, 5.15, 5.62, 1.18, ORANGE, PALE_ORANGE)
    footer(slide, 5, "2026-07-27 论文式汇报中的 held-out / unseen 结果")
    add_notes(slide, "这一页用于承接上周成果，避免听众以为本周从零开始。")
    outline.append("5. 预测效果：缓存误差、Plan/Spill 留出结果与泛化边界。")

    # 6 Experimental rigor
    slide = prs.slides.add_slide(blank); set_bg(slide)
    add_title(slide, "本周首先修正实验口径", "如果基线、并发和自然收尾口径不严，控制器会看起来“很好”，但结论无效。", "04 · METHODOLOGY")
    data = [
        ["问题", "旧风险", "本周修正"],
        ["TP 参考值", "用偏低冷启动 TPS 放宽门槛", "固定 offered 800 TPS"],
        ["冷缓存", "AP 在 TP 尚未稳定时启动", "连续 3×15s ≥784 才提交 AP"],
        ["阶段结束", "180 秒直接取消未完成 SQL", "无限自然 drain；0 个 SQL 取消"],
        ["验收区间", "只统计前 180 秒", "从 AP 提交到全部自然完成"],
        ["控制归因", "SB 与 I/O 同窗口变化", "执行器串行仲裁，观测/决策档位分列"],
    ]
    add_table(slide, data, 0.72, 1.62, 11.92, 3.80, widths=[2.15, 4.15, 5.62], font_size=11.4)
    add_text(slide, "严格口径直接推翻了“前 180 秒通过 = 项目通过”的旧结论。", 0.82, 5.78, 11.5, 0.55, size=20, color=RED, bold=True, align=PP_ALIGN.CENTER)
    footer(slide, 6, "tp_slo_query_boundary_driver.py")
    add_notes(slide, "重点讲自然结束：阶段计时只关闭准入，不取消正在执行的查询。")
    outline.append("6. 实验口径：固定 800、稳定基线、自然结束与全生命周期验收。")

    # 7 Controller architecture
    slide = prs.slides.add_slide(blank); set_bg(slide)
    add_title(slide, "在线闭环：预测负责候选，反馈负责安全执行", "控制器不直接预测 TPS；它使用 replay 约束候选，再用真实 TP/AP 反馈在线选择档位。", "05 · CONTROLLER")
    add_flow_box(slide, "Replay 候选", "SB / work_mem\nPlan / spill / peak", 0.62, 1.78, 2.32, 1.55, PURPLE, fill=LIGHT)
    add_arrow(slide, 3.02, 2.56, 3.66, 2.56, color=GRAY)
    add_flow_box(slide, "TP SLO 控制", "800 TPS 参考\n95% floor / 98% guard", 3.76, 1.78, 2.38, 1.55, TEAL, fill=PALE_TEAL)
    add_arrow(slide, 6.22, 2.56, 6.86, 2.56, color=GRAY)
    add_flow_box(slide, "资源执行器", "动态 SB\nCPU / I/O cgroup / freezer", 6.96, 1.78, 2.52, 1.55, BLUE, fill=LIGHT)
    add_arrow(slide, 9.56, 2.56, 10.18, 2.56, color=GRAY)
    add_flow_box(slide, "真实反馈", "TP TPS\nAP CPU / I/O / wait", 10.26, 1.78, 2.42, 1.55, GREEN, fill=PALE_GREEN)
    add_callout(slide, "内存层", "shared_buffers 运行态 target；work_mem 只在 Query 边界生效；已运行算子形成 graceful debt，不能伪装成已回收。", 0.80, 4.12, 5.72, 1.52, BLUE, LIGHT)
    add_callout(slide, "资源层", "AP backend LWTID 进入 CPU/blkio/freezer cgroup；暂停保留同一 SQL、Plan 和算子状态，不执行 cancel。", 6.80, 4.12, 5.72, 1.52, TEAL, PALE_TEAL)
    footer(slide, 7, "动态 SB 内核 + TP SLO controller + AP cgroup")
    add_notes(slide, "预测和在线反馈各自负责不同问题：预测缩小搜索空间，反馈防止环境偏差。")
    outline.append("7. 控制架构：Replay 候选、TP SLO、动态 SB 与 AP cgroup 执行器。")

    # 8 State machine
    slide = prs.slides.add_slide(blank); set_bg(slide)
    add_title(slide, "动态档位不是固定 0.25 core + 5MiB/s", "档位由在线实验确定，但必须防止探索本身伤害 TP。", "05 · CONTROLLER")
    steps = [
        ("01", "稳定确认", "连续 2 个窗口 ≥98% 才允许升档", TEAL),
        ("02", "90 秒试用", "AP 进展至少提升 10%，期间 TP 不越界", BLUE),
        ("03", "串行干预", "SB 变化窗口不同时学习 CPU/I/O 因果", PURPLE),
        ("04", "冻结因果", "暂停同一 SQL；TP 恢复才确认 AP 是原因", ORANGE),
        ("05", "降档恢复", "80→40→20；高档记为不安全上界", GREEN),
        ("06", "有界暂停", "120 秒仍不恢复就让 SQL 自然完成释放", RED),
    ]
    for i, (code, title, body, color) in enumerate(steps):
        col, row = i % 3, i // 3
        x, y = 0.78 + col * 4.15, 1.68 + row * 2.22
        add_rect(slide, x, y, 0.68, 0.55, color, radius=True)
        add_text(slide, code, x + 0.05, y + 0.15, 0.58, 0.22, size=10, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, title, x + 0.88, y - 0.01, 2.96, 0.34, size=16, color=color, bold=True)
        add_text(slide, body, x + 0.88, y + 0.43, 2.98, 0.72, size=11.5, color=INK)
    add_text(slide, "新增：初始档位与完整回退集合分离，可从 20 快速开始，但仍保留 10/5 作为后期路径回退。", 0.84, 6.12, 11.65, 0.52, size=16, color=INK, bold=True, align=PP_ALIGN.CENTER)
    footer(slide, 8, "tp_slo_ap_resource_controller.py，21 个专项测试")
    add_notes(slide, "解释为什么这不是固定隔离：它测收益、看 TP、记录上界，并随路径阶段切换。")
    outline.append("8. 状态机：90 秒试用、串行归因、冻结因果、降档恢复和 5/10 回退。")

    # 9 control baseline
    slide = prs.slides.add_slide(blank); set_bg(slide)
    add_title(slide, "先证明 800 TPS 本身可持续", "否则 AP 运行时的下降可能只是 TP 自己的冷缓存或磁盘抖动。", "06 · EXPERIMENT")
    add_kpi(slide, f"{tp['mean_retention']*100:.2f}%", "1 小时平均保持率", 0.82, 1.70, 2.55, color=TEAL)
    add_kpi(slide, f"{tp['minimum_retention']*100:.2f}%", "最低 15 秒保持率", 3.57, 1.70, 2.55, color=GREEN)
    add_kpi(slide, str(tp["violating_windows"]), "低于 95% 窗口", 6.32, 1.70, 2.55, color=GREEN)
    add_kpi(slide, str(tp["control_windows"]), "完整控制窗口", 9.07, 1.70, 2.55, color=BLUE)
    add_callout(slide, "实验条件", "冷 cache、SB=4096MB、32 terminals、固定 800 TPS、无 AP；基线稳定后连续测量 3600 秒。", 0.90, 3.65, 5.52, 1.45, BLUE, LIGHT)
    add_callout(slide, "结论", "800 TPS 在无 AP 时没有低于 95% 的窗口。因此 v8 后期下降是 AP 运行态干扰，不是 offered load 本身不可持续。", 6.88, 3.65, 5.52, 1.45, GREEN, PALE_GREEN)
    add_text(slide, "冷启动达到稳定基线曾耗时约 733 秒：预热门控是实验的一部分，不是可省略的等待。", 0.94, 5.65, 11.4, 0.48, size=15, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
    footer(slide, 9, "tp_only_long_v6_sb4096_rate800_20260729")
    add_notes(slide, "这一页用于回答：是不是 32 并发、800 TPS 本来就跑不稳。答案是否。")
    outline.append("9. TP-only 对照：800 TPS 一小时可持续，0 个越界窗口。")

    # 10 timeline
    slide = prs.slides.add_slide(blank); set_bg(slide)
    add_title(slide, "v8：控制器能识别路径变化并逐级回退", "Q3 不取消；灰色区域表示 freezer 暂停同一个 SQL。", "06 · EXPERIMENT")
    slide.shapes.add_picture(str(CHART_TIMELINE), Inches(0.66), Inches(1.46), width=Inches(12.02), height=Inches(5.35))
    footer(slide, 10, "v8 ap_resource_actions.csv + controller_actions.csv")
    add_notes(slide, "先讲成功点：80 前期安全、后期不安全；冻结证明确实是 AP；恢复时逐级降档，没有回到旧高档。")
    outline.append("10. v8 时间线：80 后期越界，freezer 因果确认后 80→40→20。")

    # 11 acceptance result
    slide = prs.slides.add_slide(blank); set_bg(slide)
    add_title(slide, "但 v8 仍未通过全生命周期验收", "前 180 秒 0 越界；继续等 Q3 自然结束后，问题才完整暴露。", "06 · RESULT")
    slide.shapes.add_picture(str(CHART_VALIDATION), Inches(0.72), Inches(1.48), width=Inches(11.90), height=Inches(4.15))
    add_text(slide, f"Q3 自然完成 {float(str(ap['query_completion_seconds']).split('=')[1]):.1f}s · 全生命周期平均 {stage['full_lifecycle_mean_retention']*100:.2f}% · 最低 {stage['full_lifecycle_min_retention']*100:.2f}% · 10/222 越界", 0.80, 5.87, 11.75, 0.50, size=16, color=RED, bold=True, align=PP_ALIGN.CENTER)
    footer(slide, 11, "TP-only 与 v8 summary.json")
    add_notes(slide, "这里要明确说失败。平均值接近 100% 不代表稳定，最低值和越界窗口才决定验收。")
    outline.append("11. 验收结果：v8 全生命周期最低 84.50%，10/222 越界，判定失败。")

    # 12 cause
    slide = prs.slides.add_slide(blank); set_bg(slide)
    add_title(slide, "为什么事后回退仍然来不及", "根因是状态具有滞后：配额降低不等于已经恢复缓存和释放算子内存。", "07 · DIAGNOSIS")
    causes = [
        ("Page cache 污染不可瞬间逆转", "Q3 顺序扫描驱逐 TP 热页；降到 10MiB/s 或冻结 120 秒后，TP 仍可能低于 95%。", RED, PALE_RED),
        ("运行中 work_mem 不能热缩", "控制器可以降低下一条 Query 的 grant，但当前 HashAgg/Join 已持有内存，只能形成 graceful debt。", ORANGE, PALE_ORANGE),
        ("同一 Query 的瓶颈会变", "扫描阶段 I/O-bound，HashAgg 阶段可能 CPU/内存-bound；一个固定档位无法覆盖完整执行路径。", PURPLE, LIGHT),
        ("探索也有代价", "80MiB/s 经过 90 秒仍可能在更晚路径变得不安全；生产不能每次重新学习同一个错误。", TEAL, PALE_TEAL),
    ]
    for i, (title, body, color, fill) in enumerate(causes):
        col, row = i % 2, i // 2
        add_callout(slide, title, body, 0.78 + col * 6.0, 1.65 + row * 2.28, 5.62, 1.72, color, fill)
    add_text(slide, "策略转变：探索阶段学习路径上界；正式运行直接使用已学习的低档，预防污染，而不是污染后再补救。", 0.86, 6.12, 11.55, 0.46, size=16, color=INK, bold=True, align=PP_ALIGN.CENTER)
    footer(slide, 12, "v7/v8 诊断记录与 Linux cache 行为")
    add_notes(slide, "这页是核心技术解释：为什么模型不是简单把 I/O 再减半就完成。")
    outline.append("12. 根因：page cache 滞后、运行中 work_mem 不可热缩、路径阶段变化与探索代价。")

    # 13 status matrix
    slide = prs.slides.add_slide(blank); set_bg(slide)
    add_title(slide, "截至本周：哪些已经完成，哪些仍是缺口", "所有结论按真实验证范围陈述。", "08 · STATUS")
    data = [
        ["模块", "当前状态", "证据 / 边界"],
        ["缓存 + 算子 Trace Replay", "已完成", "held-out / unseen 验证；扩大未见 Query 仍有边界"],
        ["SB + work_mem 双向联合候选", "已完成", "五阶段已测网格 regret <5%；非连续全局最优证明"],
        ["动态 shared_buffers 内核", "已实现", "运行态 target 生效；生产恢复流程已验证"],
        ["AP CPU/I/O/freezer 闭环", "已实现并探索验证", "80→40→20、同一 SQL 自然完成、0 cancel"],
        ["S2 全生命周期 TP ≤5%", "未通过", "v8：最低 84.50%，10/222 越界"],
        ["五阶段全生命周期 ≤5%", "尚未验证", "必须先完成 S2 预防式低档验收"],
    ]
    add_table(slide, data, 0.68, 1.52, 12.00, 4.88, widths=[3.02, 2.54, 6.44], font_size=10.8)
    footer(slide, 13, "本周代码、实验与结果审计")
    add_notes(slide, "这页适合答辩时直接用：避免别人追问哪些是真的完成。")
    outline.append("13. 状态审计：预测、内核和闭环已实现；S2 与五阶段全生命周期仍未通过。")

    # 14 next steps
    slide = prs.slides.add_slide(blank); set_bg(slide)
    add_title(slide, "下一步：从“探索控制器”转为“已学习策略验收”", "先让 S2 干净通过，再扩展五阶段；不能跳过。", "09 · NEXT")
    next_steps = [
        ("01", "持久化路径安全上界", "保存阶段/Query/路径对应的 CPU、I/O 和 freeze 结论；生产运行不重复试探 80。", TEAL),
        ("02", "S2 从低档主动启动", "完整集合保留 5/10/20…；从已学习安全档开始，Q3 全程自然完成，验收所有窗口。", RED),
        ("03", "联合 Query 边界 work_mem", "运行前选择 grant；运行中只记录 graceful debt，不假装 work_mem 已经缩小。", ORANGE),
        ("04", "扩展到五阶段", "每阶段固定 800 TPS 参考；比较归一化保持率，要求跨阶段和阶段内均满足 ±5%。", GREEN),
    ]
    y = 1.62
    for code, title, body, color in next_steps:
        add_rect(slide, 0.82, y, 0.72, 0.58, color, radius=True)
        add_text(slide, code, 0.88, y + 0.16, 0.60, 0.22, size=10.5, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, title, 1.82, y - 0.01, 3.26, 0.34, size=16, color=color, bold=True)
        add_text(slide, body, 5.02, y - 0.02, 7.20, 0.65, size=12.5, color=INK)
        y += 1.16
    add_callout(slide, "最终验收门槛", "AP 查询全部自然完成；任何阶段不能只截取前 180 秒；五阶段 TP 归一化波动 ≤5%；未通过时保留失败证据而不是校准到答案。", 0.84, 6.18, 11.62, 0.72, RED, PALE_RED)
    footer(slide, 14, "下一轮实验计划")
    add_notes(slide, "收尾时强调顺序：先持久化策略，再做 S2 exploit-only，最后才是五阶段。")
    outline.append("14. 下一步：持久化安全上界、S2 低档预防式验收、Query 边界 grant、五阶段扩展。")

    prs.save(OUT)
    OUTLINE.write_text(
        "# Huawei5 一周研发进展提纲（2026-07-23—29）\n\n"
        + "\n".join(f"- {item}" for item in outline)
        + "\n\n## 汇报口径\n\n"
        + "- v8 是探索运行，不是最终配置推荐或五阶段验收通过证明。\n"
        + "- Q3 自然完成，未取消；手工排空边界后的数据不作为端到端验收证据。\n"
        + "- 前 180 秒准入窗口与完整自然收尾必须分开报告。\n"
        + "- 预测模型不使用实测 TPS 最优点作为训练标签。\n",
        encoding="utf-8",
    )
    print(OUT)
    print(OUTLINE)


if __name__ == "__main__":
    build()
