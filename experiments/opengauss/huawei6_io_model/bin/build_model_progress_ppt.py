#!/usr/bin/env python3
"""Build the Huawei5 cache-model progress presentation."""

from __future__ import annotations

import csv
from pathlib import Path

import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
ARTIFACTS = ROOT / "artifacts"
ARTIFACTS.mkdir(parents=True, exist_ok=True)

OUT = ARTIFACTS / "Huawei5_SB_prediction_model_progress_20260715.pptx"
OUTLINE = ARTIFACTS / "Huawei5_SB_prediction_model_progress_20260715_outline.md"

FONT = "Microsoft YaHei"
MONO = "Consolas"

INK = RGBColor(27, 38, 48)
TEAL = RGBColor(29, 111, 122)
GREEN = RGBColor(76, 149, 108)
ORANGE = RGBColor(217, 130, 43)
RED = RGBColor(198, 83, 79)
BLUE = RGBColor(54, 112, 166)
LIGHT = RGBColor(244, 246, 247)
MID = RGBColor(214, 221, 225)
GRAY = RGBColor(101, 112, 120)
WHITE = RGBColor(255, 255, 255)


def rgb_hex(color: RGBColor) -> str:
    return f"#{color[0]:02x}{color[1]:02x}{color[2]:02x}"


def set_bg(slide, color=WHITE):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(
    slide,
    text,
    x,
    y,
    w,
    h,
    *,
    size=20,
    color=INK,
    bold=False,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    font=FONT,
    margin=0.05,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    box.text_frame.clear()
    box.text_frame.margin_left = Inches(margin)
    box.text_frame.margin_right = Inches(margin)
    box.text_frame.margin_top = Inches(margin)
    box.text_frame.margin_bottom = Inches(margin)
    box.text_frame.vertical_anchor = valign
    p = box.text_frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_rect(slide, x, y, w, h, fill, *, line=None, radius=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line or fill
    if radius:
        shape.adjustments[0] = 0.08
    return shape


def add_title(slide, title, subtitle=None, section=None):
    add_rect(slide, 0, 0, 13.333, 0.12, TEAL)
    if section:
        add_text(slide, section.upper(), 0.62, 0.28, 2.8, 0.32, size=10, color=TEAL, bold=True)
    add_text(slide, title, 0.62, 0.58, 12.0, 0.55, size=27, bold=True)
    if subtitle:
        add_text(slide, subtitle, 0.64, 1.13, 11.9, 0.38, size=12, color=GRAY)


def add_footer(slide, page):
    add_rect(slide, 0.62, 7.18, 12.08, 0.012, MID)
    add_text(slide, "Huawei5 shared_buffers prediction | 2026-07-15", 0.62, 7.22, 5.8, 0.2, size=8, color=GRAY)
    add_text(slide, str(page), 12.2, 7.22, 0.5, 0.2, size=8, color=GRAY, align=PP_ALIGN.RIGHT)


def add_bullets(slide, items, x, y, w, h, *, size=17, color=INK, spacing=9):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.05)
    for index, item in enumerate(items):
        p = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        p.text = f"• {item}"
        p.level = 0
        p.font.name = FONT
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(spacing)
    return box


def add_kpi(slide, value, label, x, y, w, *, color=TEAL, note=None):
    add_text(slide, value, x, y, w, 0.52, size=30, color=color, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, label, x, y + 0.55, w, 0.34, size=12, color=INK, bold=True, align=PP_ALIGN.CENTER)
    if note:
        add_text(slide, note, x, y + 0.9, w, 0.42, size=9, color=GRAY, align=PP_ALIGN.CENTER)


def add_flow_box(slide, text, x, y, w, h, color, *, sub=None):
    add_rect(slide, x, y, w, h, LIGHT, line=color, radius=True)
    add_text(slide, text, x + 0.08, y + 0.13, w - 0.16, 0.38, size=16, color=color, bold=True, align=PP_ALIGN.CENTER)
    if sub:
        add_text(slide, sub, x + 0.1, y + 0.55, w - 0.2, h - 0.62, size=10, color=GRAY, align=PP_ALIGN.CENTER)


def add_arrow(slide, x1, y1, x2, y2, color=GRAY, width=2.0):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    line.line.color.rgb = color
    line.line.width = Pt(width)
    line.line.end_arrowhead = True
    return line


def add_table(slide, rows, cols, data, x, y, w, h, widths=None, font_size=11):
    table = slide.shapes.add_table(rows, cols, Inches(x), Inches(y), Inches(w), Inches(h)).table
    if widths:
        for idx, width in enumerate(widths):
            table.columns[idx].width = Inches(width)
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = str(data[r][c])
            cell.margin_left = Inches(0.05)
            cell.margin_right = Inches(0.05)
            cell.margin_top = Inches(0.03)
            cell.margin_bottom = Inches(0.03)
            cell.fill.solid()
            cell.fill.fore_color.rgb = INK if r == 0 else (LIGHT if r % 2 else WHITE)
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER
                for run in p.runs:
                    run.font.name = FONT
                    run.font.size = Pt(font_size if r else font_size - 1)
                    run.font.bold = r == 0
                    run.font.color.rgb = WHITE if r == 0 else INK
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    return table


def make_recommendation_chart(path: Path):
    stages = ["S1", "S2", "S3", "S4", "S5"]
    pred = [1024, 1504, 2048, 512, 1504]
    actual = [256, 1504, 2048, 128, 256]
    x = range(len(stages))
    fig, ax = plt.subplots(figsize=(8.8, 3.4))
    ax.bar([i - 0.18 for i in x], pred, width=0.36, color=rgb_hex(ORANGE), label="raw prediction")
    ax.bar([i + 0.18 for i in x], actual, width=0.36, color=rgb_hex(TEAL), label="actual best")
    ax.set_yscale("log", base=2)
    ax.set_xticks(list(x), stages)
    ax.set_ylabel("shared_buffers (MB, log2)")
    ax.grid(axis="y", alpha=0.2)
    ax.legend(frameon=False, ncol=2, loc="upper left")
    for i, value in enumerate(pred):
        ax.text(i - 0.18, value * 1.07, str(value), ha="center", fontsize=8)
    for i, value in enumerate(actual):
        ax.text(i + 0.18, value * 1.07, str(value), ha="center", fontsize=8)
    fig.tight_layout()
    fig.savefig(path, dpi=180, transparent=False)
    plt.close(fig)


def make_mae_chart(path: Path):
    labels = ["raw", "holdout residual", "path-state diagnostic"]
    values = [3.289, 0.694, 0.495]
    colors = [rgb_hex(RED), rgb_hex(ORANGE), rgb_hex(TEAL)]
    fig, ax = plt.subplots(figsize=(7.6, 3.2))
    bars = ax.bar(labels, values, color=colors, width=0.58)
    ax.set_ylabel("combined hit-rate MAE (pp)")
    ax.set_ylim(0, 3.7)
    ax.grid(axis="y", alpha=0.2)
    for bar, value in zip(bars, values):
        ax.text(bar.get_x() + bar.get_width() / 2, value + 0.08, f"{value:.3f}", ha="center", fontsize=10)
    fig.tight_layout()
    fig.savefig(path, dpi=180)
    plt.close(fig)


def add_status_row(slide, y, status, title, detail, color):
    add_rect(slide, 0.75, y, 1.12, 0.42, color, radius=True)
    add_text(slide, status, 0.76, y + 0.07, 1.1, 0.25, size=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, title, 2.05, y - 0.01, 3.4, 0.34, size=16, bold=True)
    add_text(slide, detail, 5.35, y - 0.01, 7.2, 0.55, size=12, color=GRAY)


def build() -> None:
    chart_best = ARTIFACTS / "ppt_raw_recommendation_best.png"
    chart_mae = ARTIFACTS / "ppt_model_mae_comparison.png"
    make_recommendation_chart(chart_best)
    make_mae_chart(chart_mae)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    pages = []

    # 1. Cover
    slide = prs.slides.add_slide(blank)
    set_bg(slide, INK)
    add_rect(slide, 0, 0, 0.18, 7.5, TEAL)
    add_text(slide, "Huawei5 五阶段 shared_buffers 预测模型", 0.9, 1.35, 11.5, 0.9, size=34, color=WHITE, bold=True)
    add_text(slide, "现有框架、验证结果与多锚点升级进展", 0.92, 2.35, 10.5, 0.55, size=22, color=RGBColor(190, 218, 220))
    add_rect(slide, 0.92, 3.38, 3.05, 0.48, TEAL, radius=True)
    add_text(slide, "阶段性汇报 · 快速验证进行中", 1.02, 3.49, 2.85, 0.25, size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "数据环境：openGauss · TPC-H SF85 · TPC-C 250 warehouses", 0.92, 5.75, 8.2, 0.35, size=13, color=RGBColor(194, 201, 206))
    add_text(slide, "2026-07-15", 0.92, 6.25, 3.0, 0.3, size=12, color=RGBColor(194, 201, 206))
    pages.append("封面：说明这是阶段性汇报，快速验证尚未完成。")

    # 2. Problem
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, "问题定义：同一工作负载，不同阶段需要不同 SB", "目标不是寻找一个全局固定值，而是为五个阶段分别推荐配置。", "01 / Problem")
    stage_names = [
        ("Stage 1", "Memory rich", "Q1 · AP×1", GREEN),
        ("Stage 2", "Reach limit", "Q3 · AP×1", BLUE),
        ("Stage 3", "Protect TP", "Q5/Q7 · AP×2", TEAL),
        ("Stage 4", "Backpressure", "Q9/Q13/Q18/Q21 · AP×4", ORANGE),
        ("Stage 5", "TP surge", "Q1/Q3/Q5/Q7 · AP×4 + TP surge", RED),
    ]
    x = 0.68
    for idx, (name, desc, load, color) in enumerate(stage_names):
        add_rect(slide, x, 1.78, 2.35, 1.23, LIGHT, line=color, radius=True)
        add_text(slide, name, x + 0.1, 1.92, 2.15, 0.28, size=13, color=color, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, desc, x + 0.1, 2.24, 2.15, 0.28, size=15, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, load, x + 0.08, 2.61, 2.19, 0.24, size=9, color=GRAY, align=PP_ALIGN.CENTER)
        if idx < 4:
            add_arrow(slide, x + 2.35, 2.38, x + 2.52, 2.38, MID, 1.5)
        x += 2.52
    add_rect(slide, 0.68, 3.45, 7.48, 2.55, LIGHT, line=MID)
    add_text(slide, "预测目标", 0.94, 3.72, 2.0, 0.35, size=19, color=TEAL, bold=True)
    add_bullets(slide, [
        "给定阶段与候选 shared_buffers，预测 SB hit、OS conditional hit 和 combined hit",
        "在可行配置中选择命中率最高或综合性能最优的 SB",
        "最终还需要加入 TPS、AP 延迟、临时文件和失败风险，形成性能推荐",
    ], 0.94, 4.18, 6.92, 1.5, size=15)
    add_rect(slide, 8.48, 3.45, 4.18, 2.55, INK)
    add_text(slide, "combined = 1 − (1−SB) × (1−OS)", 8.7, 4.05, 3.75, 0.5, size=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER, font=MONO)
    add_text(slide, "当前阶段重点：先验证缓存命中预测，再扩展到 TPS/延迟。", 8.85, 4.86, 3.45, 0.65, size=13, color=RGBColor(205, 216, 220), align=PP_ALIGN.CENTER)
    add_footer(slide, 2)
    pages.append("问题定义：五个阶段的最优 SB 不同，当前先解决命中率预测。")

    # 3. Original flow
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, "现有原始模型：单轨迹 Trace Replay", "使用 1504MB 运行轨迹，模拟其他候选 SB 下的缓存行为。", "02 / Baseline")
    flow = [
        ("真实运行", "SB=1504MB\n五阶段负载", GREEN),
        ("采集 trace", "relation/block\ntime/hit/strategy", BLUE),
        ("阶段切分", "边界连续\n缓存状态不重置", TEAL),
        ("DB Buffer", "bulk ring +\nshared table", ORANGE),
        ("OS Cache", "two-list\n容量随 SB 调整", RED),
        ("推荐", "最大化\ncombined hit", INK),
    ]
    x = 0.52
    for idx, (title, sub, color) in enumerate(flow):
        add_flow_box(slide, title, x, 2.05, 1.78, 1.45, color, sub=sub)
        if idx < len(flow) - 1:
            add_arrow(slide, x + 1.8, 2.78, x + 2.08, 2.78, color, 1.8)
        x += 2.1
    add_rect(slide, 0.72, 4.15, 12.0, 1.68, LIGHT, line=MID)
    add_text(slide, "核心假设", 0.95, 4.38, 1.5, 0.32, size=18, color=TEAL, bold=True)
    add_bullets(slide, [
        "改变 SB 后，页面访问序列和 SQL 执行路径基本不变",
        "SB 增大占用的内存会从 Linux page cache 中扣除",
        "缓存命中率足以作为阶段配置推荐的主要依据",
    ], 2.38, 4.3, 9.95, 1.2, size=15, spacing=6)
    add_footer(slide, 3)
    pages.append("介绍单轨迹 replay 流程，强调三条核心假设。")

    # 4. Cache details
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, "缓存模拟细节", "原模型是机制模拟，不是神经网络训练。", "02 / Baseline")
    add_rect(slide, 0.68, 1.72, 5.95, 4.85, LIGHT, line=TEAL)
    add_text(slide, "Shared Buffer 模型", 0.98, 2.02, 4.8, 0.42, size=21, color=TEAL, bold=True)
    add_bullets(slide, [
        "全局 shared-buffer 页面表",
        "普通访问使用近似 clock sweep",
        "BulkRead 使用每 backend / strategy 的私有 ring 选择 victim",
        "采样率 1/64，同时缩放模拟容量",
        "输出 SB hits、misses 和淘汰页面流",
    ], 0.98, 2.63, 5.1, 2.65, size=15)
    add_text(slide, "原 trace 的缺陷：记录的是 gaussdb PID，而非执行线程，导致不同 backend 的 ring 可能混在一起。", 0.98, 5.55, 5.05, 0.68, size=12, color=RED, bold=True)
    add_rect(slide, 6.82, 1.72, 5.83, 4.85, LIGHT, line=ORANGE)
    add_text(slide, "Linux Page Cache 模型", 7.12, 2.02, 4.8, 0.42, size=21, color=ORANGE, bold=True)
    add_bullets(slide, [
        "SB miss 进入 OS cache 模拟",
        "Two-list active / inactive 近似",
        "候选 OS 容量 = 当前容量 + 1504 − candidate SB",
        "代表性参数：readahead=0，OS scale=0.75",
        "输出 OS conditional hit 与磁盘页数",
    ], 7.12, 2.63, 5.0, 2.65, size=15)
    add_text(slide, "风险：combined 可能因为 SB 与 OS 的误差相互抵消，看起来准确但机制实际错误。", 7.12, 5.55, 5.0, 0.68, size=12, color=RED, bold=True)
    add_footer(slide, 4)
    pages.append("讲清 DB buffer 和 OS cache 两层模拟，以及误差抵消风险。")

    # 5. Workload
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, "实验数据与五阶段负载", "稳定种子、固定 SQL、查询边界和冷 OS cache，减少不同 SB 运行之间的漂移。", "03 / Data")
    add_kpi(slide, "≈143 GiB", "数据库规模", 0.7, 1.68, 2.2, color=TEAL, note="大于全部有效 SB 配置")
    add_kpi(slide, "SF85", "TPC-H", 3.1, 1.68, 1.7, color=BLUE)
    add_kpi(slide, "250", "TPC-C warehouses", 5.0, 1.68, 2.1, color=GREEN)
    add_kpi(slide, "1024 MB", "AP work_mem", 7.35, 1.68, 2.1, color=ORANGE)
    add_kpi(slide, "8 points", "已完成 SB 实测", 9.7, 1.68, 2.2, color=RED, note="128MB–8192MB")
    data = [
        ["阶段", "AP 查询", "AP 并发", "TP 状态", "目的"],
        ["S1", "Q1", "1", "低负载 ≈40 TPS", "内存宽裕"],
        ["S2", "Q3", "1", "低负载", "接近容量边界"],
        ["S3", "Q5 + Q7", "2", "低负载", "保护 TP"],
        ["S4", "Q9/Q13/Q18/Q21", "4", "低负载", "AP backpressure"],
        ["S5", "Q1/Q3/Q5/Q7", "4", "TP surge ≈220 TPS 目标", "混合压力"],
    ]
    add_table(slide, 6, 5, data, 0.72, 3.23, 11.9, 2.95, widths=[1.05, 2.45, 1.25, 2.7, 4.45], font_size=12)
    add_footer(slide, 5)
    pages.append("说明数据规模、负载固定方式和五阶段定义。")

    # 6. Baseline result
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, "原始模型效果：最佳 SB 推荐命中 2/5", "8 个有效实测点：128、256、512、1024、1504、2048、4096、8192MB。", "04 / Results")
    add_kpi(slide, "40%", "Top-1 推荐命中率", 0.7, 1.62, 2.35, color=RED, note="2 / 5 stages")
    add_kpi(slide, "3.29 pp", "combined MAE", 3.05, 1.62, 2.35, color=ORANGE, note="严格留出集合")
    add_kpi(slide, "2 / 5", "正确阶段", 5.4, 1.62, 2.15, color=GREEN, note="Stage 2、Stage 3")
    slide.shapes.add_picture(str(chart_best), Inches(0.72), Inches(3.05), width=Inches(7.1), height=Inches(2.85))
    rec_data = [
        ["阶段", "Raw 推荐", "实际最优", "结果"],
        ["S1", "1024", "256", "×"],
        ["S2", "1504", "1504", "✓"],
        ["S3", "2048", "2048", "✓"],
        ["S4", "512", "128", "×"],
        ["S5", "1504", "256", "×"],
    ]
    table = add_table(slide, 6, 4, rec_data, 8.15, 1.72, 4.48, 4.55, widths=[1.0, 1.25, 1.25, 0.98], font_size=12)
    for r in range(1, 6):
        cell = table.cell(r, 3)
        for p in cell.text_frame.paragraphs:
            for run in p.runs:
                run.font.color.rgb = GREEN if r in (2, 3) else RED
                run.font.bold = True
                run.font.size = Pt(18)
    add_footer(slide, 6)
    pages.append("原始模型正式结果：Top-1 2/5，combined MAE 约 3.29pp。")

    # 7. Flat curves
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, "为什么 Pred 曲线几乎没有波动？", "SB hit 上升与 OS hit 下降在 combined 公式中发生抵消。", "05 / Diagnosis")
    ranges = [("S1", 0.033, GREEN), ("S2", 0.060, BLUE), ("S3", 0.074, TEAL), ("S4", 0.795, ORANGE), ("S5", 7.073, RED)]
    add_text(slide, "候选 SB 范围内 predicted combined 的波动幅度", 0.78, 1.72, 6.0, 0.36, size=17, bold=True)
    y = 2.25
    for stage, value, color in ranges:
        add_text(slide, stage, 0.82, y + 0.03, 0.6, 0.26, size=12, color=color, bold=True)
        width = 0.08 + 5.15 * (value / 7.073)
        add_rect(slide, 1.45, y, width, 0.32, color)
        add_text(slide, f"{value:.3f} pp", 6.72, y + 0.02, 1.0, 0.25, size=11, color=INK, bold=True, align=PP_ALIGN.RIGHT)
        y += 0.68
    add_rect(slide, 8.05, 1.72, 4.55, 4.28, INK)
    add_text(slide, "误差抵消", 8.32, 2.08, 4.0, 0.4, size=22, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "SB ↑", 8.45, 2.85, 1.15, 0.42, size=23, color=GREEN, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "+", 9.62, 2.85, 0.45, 0.42, size=23, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "OS ↓", 10.05, 2.85, 1.15, 0.42, size=23, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "≈", 11.18, 2.85, 0.45, 0.42, size=23, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "Combined 不变", 8.78, 3.65, 3.1, 0.5, size=20, color=RGBColor(199, 224, 226), bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "因此 Stage 1–3 的曲线无法提供稳定的最优点排序，即使某些点的 combined 误差不大。", 8.52, 4.65, 3.65, 0.82, size=13, color=RGBColor(211, 219, 223), align=PP_ALIGN.CENTER)
    add_footer(slide, 7)
    pages.append("解释 Pred 曲线平坦：SB 与 OS 误差抵消，无法可靠排序。")

    # 8. Stage 4/5 failures
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, "Stage 4 / 5 的主要失效机制", "这两类误差不是简单调整一个系数就能解决。", "05 / Diagnosis")
    add_rect(slide, 0.68, 1.7, 5.95, 4.95, LIGHT, line=ORANGE)
    add_text(slide, "Stage 4：高并发 AP backpressure", 0.98, 2.02, 5.2, 0.42, size=20, color=ORANGE, bold=True)
    s4 = [
        ["SB", "actual SB", "pred SB", "actual combined", "pred combined"],
        ["128", "77.59%", "83.81%", "85.13%", "88.40%"],
        ["512", "78.02%", "86.27%", "81.71%", "88.44%"],
    ]
    add_table(slide, 3, 5, s4, 0.96, 2.68, 5.42, 1.32, widths=[0.7, 1.18, 1.18, 1.18, 1.18], font_size=9)
    add_bullets(slide, [
        "模型高估增大 SB 对 Q9/Q13/Q18/Q21 并发的收益",
        "没有模拟执行重叠、临时文件、脏页和 backpressure",
        "真实最优为 128MB，原模型推荐 512MB",
    ], 0.98, 4.35, 5.15, 1.55, size=13, spacing=6)
    add_rect(slide, 6.82, 1.7, 5.83, 4.95, LIGHT, line=RED)
    add_text(slide, "Stage 5：小 SB 的 Linux page cache 补偿", 7.12, 2.02, 5.0, 0.42, size=20, color=RED, bold=True)
    add_text(slide, "SB = 256MB", 7.15, 2.65, 2.0, 0.35, size=15, bold=True)
    add_kpi(slide, "83.89%", "实际 OS hit", 7.08, 3.12, 2.15, color=GREEN)
    add_kpi(slide, "20.68%", "原模型 OS hit", 9.35, 3.12, 2.15, color=RED)
    add_text(slide, "actual combined 98.06%  vs  predicted 90.06%", 7.2, 4.45, 4.95, 0.42, size=15, color=INK, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "SB hit 本身较准；主要错误来自 OS cache 补偿被严重低估。", 7.2, 5.18, 4.95, 0.6, size=13, color=GRAY, align=PP_ALIGN.CENTER)
    add_footer(slide, 8)
    pages.append("Stage4 是并发路径问题，Stage5 是 Linux page cache 补偿问题。")

    # 9. Calibration caution
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, "校准实验：数值误差下降，但不能替代预测", "训练点拟合得好不等于能推荐未见 SB。", "06 / Validation")
    slide.shapes.add_picture(str(chart_mae), Inches(0.72), Inches(1.82), width=Inches(6.6), height=Inches(2.8))
    add_rect(slide, 7.65, 1.8, 4.98, 2.82, LIGHT, line=ORANGE)
    add_text(slide, "严格留出：128/512/1504/4096 训练", 7.92, 2.08, 4.45, 0.42, size=15, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "256/1024/2048/8192 测试", 7.92, 2.53, 4.45, 0.4, size=15, bold=True, align=PP_ALIGN.CENTER)
    add_kpi(slide, "1/5 → 2/5", "留出点 Top-1", 8.15, 3.14, 3.9, color=ORANGE)
    add_rect(slide, 0.72, 5.0, 11.92, 1.35, INK)
    add_text(slide, "结论", 1.0, 5.32, 1.1, 0.35, size=18, color=WHITE, bold=True)
    add_text(slide, "Residual calibration 可以降低 held-out MAE，但最佳 SB 排序仍不可靠；in-sample 5/5 属于数据泄漏式诊断，不能作为模型效果。", 2.08, 5.22, 10.05, 0.65, size=14, color=RGBColor(211, 222, 226), bold=True)
    add_footer(slide, 9)
    pages.append("强调校准的边界：留出 MAE 改善，但推荐只从 1/5 到 2/5。")

    # 10. Upgraded architecture
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, "升级框架：多锚点、执行路径感知的 Replay", "核心变化：目标 SB 不再只依赖 1504MB 的固定访问序列。", "07 / Upgrade")
    anchors = [("128MB", 0.65, GREEN), ("512MB", 0.65, BLUE), ("4096MB", 0.65, ORANGE)]
    for idx, (label, _, color) in enumerate(anchors):
        y = 1.7 + idx * 1.22
        add_flow_box(slide, label, 0.68, y, 1.72, 0.88, color, sub="完整 page trace")
        add_arrow(slide, 2.42, y + 0.44, 3.08, y + 0.44, color)
    add_rect(slide, 3.08, 1.65, 2.25, 3.55, LIGHT, line=TEAL, radius=True)
    add_text(slide, "路径标识", 3.32, 1.95, 1.75, 0.4, size=20, color=TEAL, bold=True, align=PP_ALIGN.CENTER)
    add_bullets(slide, ["thread ID", "query ID → SQL", "strategy / ring", "BufFile temp I/O", "Hash spill", "OS read"], 3.32, 2.48, 1.75, 2.25, size=12, spacing=4)
    add_arrow(slide, 5.35, 3.42, 6.02, 3.42, TEAL)
    add_rect(slide, 6.02, 1.65, 2.55, 3.55, LIGHT, line=ORANGE, radius=True)
    add_text(slide, "目标 SB Replay", 6.25, 1.95, 2.1, 0.4, size=20, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
    add_bullets(slide, ["每个锚点轨迹分别重放", "thread-aware bulk ring", "DB + OS 双层缓存", "保留阶段连续状态", "输出路径级命中/IO"], 6.25, 2.52, 2.05, 2.0, size=12, spacing=6)
    add_arrow(slide, 8.6, 3.42, 9.28, 3.42, ORANGE)
    add_rect(slide, 9.28, 1.65, 3.32, 3.55, INK, radius=True)
    add_text(slide, "路径感知组合", 9.55, 1.95, 2.8, 0.4, size=20, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_bullets(slide, ["相邻锚点权重", "执行路径切换", "spill / 回读阈值", "不确定性与失败风险", "生成阶段 SB 推荐"], 9.68, 2.52, 2.45, 2.0, size=12, color=RGBColor(220, 228, 231), spacing=6)
    add_rect(slide, 0.68, 5.62, 11.92, 0.72, RGBColor(232, 243, 239), line=GREEN)
    add_text(slide, "验证原则：锚点只用于建模；256MB、1504MB 等 held-out 点只用于评估，禁止用实际结果反向修正预测。", 0.92, 5.83, 11.42, 0.3, size=14, color=INK, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, 10)
    pages.append("真正升级框架：多个完整 page trace、路径标识、逐锚点 replay 和 held-out。")

    # 11. Engineering status
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, "当前工程进展", "快速验证由独立后台服务运行，不依赖当前会话。", "08 / Progress")
    add_status_row(slide, 1.72, "DONE", "Trace v2", "记录 thread、query ID、strategy、页面访问和路径事件。", GREEN)
    add_status_row(slide, 2.42, "DONE", "路径映射", "query ID 与实际 TPC-H SQL hash/text 每 2 秒关联。", GREEN)
    add_status_row(slide, 3.12, "DONE", "阶段隔离", "预热/测量分离；阶段结束强制清理 tpch_ap 后端。", GREEN)
    add_status_row(slide, 3.82, "RUNNING", "快速采集", "anchors=128/512/4096；held-out=256/1504。", ORANGE)
    add_status_row(slide, 4.52, "NEXT", "Multi-trace replay", "相邻锚点轨迹重放、路径组合与严格留出评估。", TEAL)
    add_status_row(slide, 5.22, "NEXT", "性能模型", "加入 TPS、AP latency、temp bytes、失败风险，超越命中率目标。", BLUE)
    add_rect(slide, 0.75, 6.12, 11.8, 0.48, LIGHT, line=MID)
    add_text(slide, "后台服务：huawei5-quick-path-validation-20260715.service  ·  预计总时长 2–3 小时", 0.9, 6.23, 11.45, 0.24, size=12, color=GRAY, align=PP_ALIGN.CENTER)
    add_footer(slide, 11)
    pages.append("当前进展：采集器和隔离逻辑完成，快速验证运行中，multi-trace replay 待数据。")

    # 12. Evidence levels
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, "阶段性效果：必须按证据等级解读", "避免把诊断性结果包装成正式预测能力。", "09 / Conclusion")
    columns = [
        (0.68, GREEN, "已证明", ["原始 Top-1 = 2/5", "Stage 4/5 是主要失败点", "固定轨迹假设不足", "单纯 residual 校准不可靠"]),
        (4.55, ORANGE, "诊断性证据", ["路径状态插值 MAE：3.289→0.495pp", "Top-1：1/5→3/5", "Stage4/5 数值明显改善", "但不是 page-level multi-trace replay"]),
        (8.42, TEAL, "正在验证", ["真实多锚点 page trace", "thread/query/spill 路径感知", "256/1504 严格留出", "完成后才能报告正式升级效果"]),
    ]
    for x, color, title, items in columns:
        add_rect(slide, x, 1.72, 3.55, 4.6, LIGHT, line=color)
        add_rect(slide, x, 1.72, 3.55, 0.65, color)
        add_text(slide, title, x + 0.12, 1.88, 3.31, 0.3, size=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_bullets(slide, items, x + 0.25, 2.68, 3.02, 2.85, size=14, spacing=10)
    add_text(slide, "现阶段结论：trace replay 有机制价值，但单轨迹不能独立承担最终 SB 推荐；多锚点执行路径 replay 是当前合理升级方向。", 0.88, 6.52, 11.55, 0.42, size=15, color=INK, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, 12)
    pages.append("总结证据等级：已证明、诊断性、正在验证三层。")

    # 13. Discussion
    slide = prs.slides.add_slide(blank)
    set_bg(slide, INK)
    add_text(slide, "明天讨论建议聚焦 3 个决策", 0.82, 0.72, 11.7, 0.62, size=30, color=WHITE, bold=True)
    decisions = [
        ("01", "最终优化目标", "只优化 combined hit，还是直接优化 TPS + AP latency + 失败风险？", TEAL),
        ("02", "锚点成本", "每个阶段允许采集多少 SB 锚点？3 点、4 点还是自适应增点？", ORANGE),
        ("03", "验收标准", "建议：combined MAE <1pp、Top-2 ≥4/5、相对默认配置无明显回退。", GREEN),
    ]
    y = 1.72
    for num, title, detail, color in decisions:
        add_text(slide, num, 0.9, y, 1.0, 0.55, size=26, color=color, bold=True)
        add_text(slide, title, 2.05, y, 2.45, 0.42, size=20, color=WHITE, bold=True)
        add_text(slide, detail, 4.6, y, 7.55, 0.65, size=15, color=RGBColor(208, 218, 222))
        add_rect(slide, 0.9, y + 0.82, 11.3, 0.012, RGBColor(74, 88, 98))
        y += 1.45
    add_rect(slide, 0.9, 6.3, 11.3, 0.55, TEAL, radius=True)
    add_text(slide, "建议会议表述：原模型验证了 trace replay 的可行性，也暴露了固定执行路径假设；当前工作正在补齐多锚点与路径变化。", 1.05, 6.44, 11.0, 0.28, size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    pages.append("讨论页：最终目标、锚点成本和验收标准。")

    prs.save(OUT)
    outline_lines = [
        "# Huawei5 SB Prediction Model Progress - Speaker Outline",
        "",
    ]
    for idx, text in enumerate(pages, 1):
        outline_lines.append(f"{idx}. {text}")
    outline_lines += [
        "",
        "## Key numbers",
        "",
        "- Raw recommendation Top-1: 2/5 on all completed points.",
        "- Strict held-out raw combined MAE: 3.289 pp.",
        "- Residual holdout: 0.694 pp and 2/5, not sufficient.",
        "- Path-state interpolation diagnostic: 0.495 pp and 3/5; not formal multi-trace replay evidence.",
        "- Quick validation split: anchors 128/512/4096; held-out 256/1504.",
    ]
    OUTLINE.write_text("\n".join(outline_lines) + "\n", encoding="utf-8")
    print(OUT)
    print(OUTLINE)


if __name__ == "__main__":
    build()
