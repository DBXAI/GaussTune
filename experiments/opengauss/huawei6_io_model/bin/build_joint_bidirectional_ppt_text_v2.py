#!/usr/bin/env python3
"""Build a text-first Huawei5 joint replay presentation with explicit accuracy."""

from __future__ import annotations

import csv
from pathlib import Path

import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches

from build_model_progress_ppt import (
    BLUE,
    GRAY,
    GREEN,
    INK,
    LIGHT,
    MID,
    ORANGE,
    RED,
    TEAL,
    WHITE,
    add_bullets,
    add_kpi,
    add_rect,
    add_table,
    add_text,
    add_title,
    rgb_hex,
    set_bg,
)


ROOT = Path(__file__).resolve().parents[1]
ARTIFACTS = ROOT / "artifacts"
LATEST = ARTIFACTS / "00_latest"
CURRENT = ARTIFACTS / "01_current_joint_model"
FIGURES = CURRENT / "figures"
RESULTS = ROOT / "results"
OUT = LATEST / "Huawei5_five_stage_simple_results_v6_20260722.pptx"
OUTLINE = LATEST / "Huawei5_five_stage_simple_results_v6_20260722_outline.md"
HIT_CHART = FIGURES / "joint_v2_hit_accuracy.png"
S5_CHART = FIGURES / "joint_v3_s5_recommendation_actual_tps.png"
COMBO_CHART = FIGURES / "joint_v4_s5_sb_workmem_intersection.png"
SIMPLE_FIGURES = FIGURES / "five_stage_simple"


def read_csv(path: Path) -> list[dict[str, str]]:
    with path.open(newline="", encoding="utf-8") as fh:
        return list(csv.DictReader(fh))


def footer(slide, page: int) -> None:
    add_rect(slide, 0.62, 7.17, 12.08, 0.012, MID)
    add_text(slide, "Huawei5 | Joint SB/work_mem trace replay | 2026-07-22", 0.62, 7.22, 7.0, 0.18, size=8, color=GRAY)
    add_text(slide, str(page), 12.2, 7.22, 0.5, 0.18, size=8, color=GRAY, align=PP_ALIGN.RIGHT)


def conclusion_bar(slide, text: str, y: float = 1.58, color=TEAL) -> None:
    add_rect(slide, 0.68, y, 11.96, 0.64, color, radius=True)
    add_text(slide, text, 0.92, y + 0.16, 11.48, 0.28, size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)


def status_box(slide, title: str, body: str, x: float, y: float, w: float, color) -> None:
    add_rect(slide, x, y, w, 1.32, LIGHT, line=color, radius=True)
    add_text(slide, title, x + 0.16, y + 0.18, w - 0.32, 0.32, size=17, color=color, bold=True)
    add_text(slide, body, x + 0.16, y + 0.61, w - 0.32, 0.52, size=12, color=INK)


def make_hit_accuracy_chart() -> None:
    labels = ["SB hit", "OS conditional hit", "Combined hit"]
    values = [0.605933, 1.821841, 0.495384]
    colors = [rgb_hex(TEAL), rgb_hex(ORANGE), rgb_hex(GREEN)]
    fig, ax = plt.subplots(figsize=(8.6, 3.1))
    bars = ax.barh(labels[::-1], values[::-1], color=colors[::-1], height=0.55)
    ax.set_xlabel("Held-out MAE (percentage points; lower is better)")
    ax.set_xlim(0, 2.15)
    ax.grid(axis="x", alpha=0.2)
    for bar, value in zip(bars, values[::-1]):
        ax.text(value + 0.04, bar.get_y() + bar.get_height() / 2, f"{value:.2f} pp", va="center", fontsize=12, fontweight="bold")
    for spine in ["top", "right", "left"]:
        ax.spines[spine].set_visible(False)
    fig.tight_layout()
    fig.savefig(HIT_CHART, dpi=190)
    plt.close(fig)


def make_s5_chart() -> None:
    rows = [
        row for row in read_csv(RESULTS / "tp_only_performance_alignment_20260716/tp_only_performance_points.csv")
        if row["stage"] == "stage5_tp_surge"
    ]
    rows.sort(key=lambda row: int(row["sb_mb"]))
    sbs = [int(row["sb_mb"]) for row in rows]
    tps = [float(row["total_tp_tps"]) for row in rows]
    max_tps = max(tps)
    x = list(range(len(sbs)))
    fig, ax = plt.subplots(figsize=(10.0, 4.0))
    ax.plot(x, tps, marker="o", markersize=8, linewidth=2.8, color=rgb_hex(BLUE), label="Actual measured TPS")
    ax.axhspan(max_tps * 0.99, max_tps + 2.0, color=rgb_hex(GREEN), alpha=0.12, label="99% maximum TPS region")
    selected = sbs.index(1024)
    ax.scatter([selected], [tps[selected]], marker="*", s=330, color=rgb_hex(RED), edgecolor="white", linewidth=1.2, zorder=5, label="Model recommended SB")
    ax.annotate(
        "Predicted recommendation: 1024MB\nActual TPS: 220.04 (measured maximum)",
        xy=(selected, tps[selected]),
        xytext=(selected + 0.45, 187),
        arrowprops={"arrowstyle": "->", "color": rgb_hex(RED), "linewidth": 1.6},
        color=rgb_hex(RED),
        fontsize=11,
        fontweight="bold",
    )
    for index, value in enumerate(tps):
        ax.text(index, value + 2.2, f"{value:.1f}", ha="center", fontsize=9, color=rgb_hex(INK))
    ax.set_xticks(x, [str(value) for value in sbs])
    ax.set_ylim(120, 226)
    ax.set_xlabel("shared_buffers (MB)")
    ax.set_ylabel("Actual total TP TPS")
    ax.grid(alpha=0.2)
    ax.legend(frameon=False, loc="lower right", fontsize=9)
    fig.tight_layout()
    fig.savefig(S5_CHART, dpi=190)
    plt.close(fig)


def make_combo_chart() -> None:
    performance = [
        row for row in read_csv(RESULTS / "tp_only_performance_alignment_20260716/tp_only_performance_points.csv")
        if row["stage"] == "stage5_tp_surge"
    ]
    performance.sort(key=lambda row: int(row["sb_mb"]))
    sbs = [int(row["sb_mb"]) for row in performance]
    tps = [float(row["total_tp_tps"]) for row in performance]

    candidates = read_csv(RESULTS / "joint_bidirectional_replay_20260722/replay/joint_bidirectional_candidates.csv")
    supported_work_mem = [256, 1137, 1150, 1208]
    by_work_mem = {
        int(row["work_mem_mb"]): float(row["spill_io_mb"]) / 1024
        for row in candidates
        if row["stage"] == "stage5_tp_surge"
        and int(row["sb_mb"]) == 1024
        and row["plan_supported"].lower() == "true"
    }
    spill = [by_work_mem[value] for value in supported_work_mem]

    fig, axes = plt.subplots(1, 2, figsize=(12.0, 3.75))
    x = list(range(len(sbs)))
    axes[0].plot(x, tps, marker="o", linewidth=2.6, color=rgb_hex(BLUE))
    axes[0].axhspan(max(tps) * 0.99, max(tps) + 2, color=rgb_hex(GREEN), alpha=0.13)
    selected_sb = sbs.index(1024)
    axes[0].scatter([selected_sb], [tps[selected_sb]], marker="*", s=280, color=rgb_hex(RED), edgecolor="white", zorder=5)
    axes[0].set_xticks(x, [str(value) for value in sbs], rotation=28)
    axes[0].set_ylim(120, 226)
    axes[0].set_xlabel("shared_buffers (MB)")
    axes[0].set_ylabel("Actual TP TPS")
    axes[0].set_title("1. First maximum-TPS SB: 1024MB", fontweight="bold")
    axes[0].grid(alpha=0.2)

    colors = [rgb_hex(ORANGE), rgb_hex(ORANGE), rgb_hex(GREEN), rgb_hex(GRAY)]
    bars = axes[1].bar(range(len(supported_work_mem)), spill, color=colors, width=0.62)
    axes[1].scatter([2], [0.25], marker="*", s=280, color=rgb_hex(RED), edgecolor="white", zorder=5)
    axes[1].set_xticks(range(len(supported_work_mem)), [str(value) for value in supported_work_mem])
    axes[1].set_xlabel("work_mem (MB)")
    axes[1].set_ylabel("Predicted spill I/O (GiB)")
    axes[1].set_title("2. First zero-spill work_mem: 1150MB", fontweight="bold")
    axes[1].grid(axis="y", alpha=0.2)
    for bar, value in zip(bars, spill):
        label_y = value + 0.35 if value else 0.85
        axes[1].text(bar.get_x() + bar.get_width() / 2, label_y, f"{value:.1f}", ha="center", fontsize=10, fontweight="bold")
    axes[1].set_ylim(0, max(spill) * 1.18)
    fig.tight_layout()
    fig.savefig(COMBO_CHART, dpi=190)
    plt.close(fig)


def build() -> None:
    LATEST.mkdir(parents=True, exist_ok=True)
    FIGURES.mkdir(parents=True, exist_ok=True)
    make_hit_accuracy_chart()
    make_s5_chart()
    recs = read_csv(RESULTS / "joint_bidirectional_replay_20260722/replay/stage_joint_recommendations.csv")

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    outline: list[str] = []

    # 1 Cover
    slide = prs.slides.add_slide(blank)
    set_bg(slide, INK)
    add_rect(slide, 0, 0, 0.18, 7.5, TEAL)
    add_text(slide, "Huawei5 联合配置预测模型", 0.94, 1.35, 10.9, 0.72, size=38, color=WHITE, bold=True)
    add_text(slide, "模型是什么、怎样预测、目前到底准不准", 0.96, 2.34, 10.8, 0.60, size=27, color=RGBColor(181, 219, 222), bold=True)
    add_text(slide, "shared_buffers × work_mem 双向 Trace Replay", 0.96, 3.34, 8.0, 0.38, size=18, color=WHITE)
    add_rect(slide, 0.96, 4.12, 4.05, 0.52, TEAL, radius=True)
    add_text(slide, "文字说明版 · 准确率单独解释", 1.08, 4.25, 3.80, 0.25, size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "2026-07-22", 0.96, 6.23, 2.0, 0.26, size=12, color=RGBColor(194, 201, 206))
    outline.append("1. 封面。")

    # 2 First answer
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, "先说结论：现在能预测什么？", "不要把三个不同问题混成一个“总准确率”。", "01 / Summary")
    conclusion_bar(slide, "模型已经能给出五阶段 SB/work_mem 推荐，但只有部分结果完成了真实最优验证。")
    status_box(slide, "能直接预测", "每个候选配置下的 TP-SB/OS 命中率、spill I/O、动态内存峰值和内存余量。", 0.76, 2.62, 3.75, TEAL)
    status_box(slide, "能给出推荐", "在计划路径有 trace、内存安全的候选中，联合选择 SB 和 work_mem。", 4.79, 2.62, 3.75, GREEN)
    status_box(slide, "尚不能证明", "五个阶段的推荐值全部等于真实全局最优；S4 仍缺新计划路径 trace。", 8.82, 2.62, 3.75, RED)
    add_bullets(slide, ["S5：预测与实际 TPS 平台一致", "S3：现有证据基本一致", "S1/S2：TPS 限速，无法识别唯一最优", "S4：当前推荐是 trace 覆盖范围内的暂定值"], 1.02, 4.45, 11.15, 1.68, size=16, spacing=10)
    footer(slide, 2)
    outline.append("2. 当前能力与限制：能预测和推荐，但尚未证明五阶段全局最优。")

    # 3 Plain workflow
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, "模型怎样工作：一组配置需要走完 6 个步骤", "输入一个 (SB, work_mem)，输出该配置的行为和风险。", "02 / Method")
    steps = [
        ("1", "检查执行计划", "确认当前 work_mem 下属于哪个计划族；没有同路径 trace 时不进入正式推荐。"),
        ("2", "重放算子内存", "模拟 Hash Join、HashAggregate、Sort 的 grant、batch、spill 和生命周期。"),
        ("3", "计算系统内存", "把 SB 和并发动态内存峰值换算成剩余 MemAvailable 与 OS cache 容量。"),
        ("4", "重放 mixed trace", "TP 和 AP 访问共同改变 SB/OS 状态，但最终只统计 TP 页命中。"),
        ("5", "模拟 Linux 保护", "active/inactive/refault 保护高频 TP 页，优先回收 AP/spill streaming 页。"),
        ("6", "比较二维候选", "同时比较 TP miss、spill 物理 I/O、内存占用，选择 Pareto 候选。"),
    ]
    for index, (number, title, body) in enumerate(steps):
        col = index % 2
        row = index // 2
        x = 0.78 + col * 6.15
        y = 1.75 + row * 1.65
        color = [BLUE, ORANGE, TEAL, GREEN, RED, INK][index]
        add_rect(slide, x, y, 0.58, 0.58, color, radius=True)
        add_text(slide, number, x + 0.01, y + 0.14, 0.56, 0.24, size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, title, x + 0.78, y - 0.01, 4.95, 0.34, size=18, color=color, bold=True)
        add_text(slide, body, x + 0.78, y + 0.45, 4.95, 0.72, size=13, color=INK)
    footer(slide, 3)
    outline.append("3. 六步预测流程。")

    # 4 Experiment design
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, "实验怎样做：五阶段负载 + 两类 trace", "预测数据与实际 TPS 数据分开采集、分开使用。", "03 / Experiment")
    stage_table = [
        ["阶段", "AP 查询", "AP 并发", "TP 负载"],
        ["S1", "Q1", "1", "2 terminals / target 40"],
        ["S2", "Q3", "1", "2 terminals / target 40"],
        ["S3", "Q5, Q7", "2", "2 terminals / target 40"],
        ["S4", "Q9, Q13, Q18, Q21", "4", "2 terminals / target 40"],
        ["S5", "Q1, Q3, Q5, Q7", "4", "+12 terminals / target 180"],
    ]
    add_table(slide, 6, 4, stage_table, 0.72, 1.72, 7.45, 3.90, widths=[1.05, 2.45, 1.15, 2.80], font_size=12)
    add_rect(slide, 8.54, 1.72, 3.95, 3.90, LIGHT, line=TEAL, radius=True)
    add_text(slide, "采集内容", 8.83, 1.98, 3.35, 0.38, size=21, color=TEAL, bold=True, align=PP_ALIGN.CENTER)
    add_bullets(slide, ["mixed TP/AP 页面访问 trace", "8 个 AP 查询算子内存 trace", "多个 work_mem 的计划族", "8 个 SB 实测档位", "最后单独采集 TPS 与 P95"], 8.91, 2.58, 3.20, 2.42, size=14, spacing=10)
    add_rect(slide, 1.08, 6.05, 11.12, 0.55, INK, radius=True)
    add_text(slide, "预测器不读取实际 TPS 和实际最优配置；这些值只在预测完成后做验证。", 1.32, 6.19, 10.64, 0.25, size=15, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    footer(slide, 4)
    outline.append("4. 五阶段实验配置与 trace/TPS 数据隔离。")

    # 5 Concrete example
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, "具体算例：为什么 S5 推荐 1024MB + 1150MB？", "用一组实际输出数字看完整预测过程。", "04 / Example")
    conclusion_bar(slide, "1024MB 已让 TP-SB 命中进入平台；1150MB 刚好消除关键 AP spill，再增加内存收益很小。", color=GREEN)
    example = [
        ["预测步骤", "S5 候选结果", "含义"],
        ["计划路径", "Q1/Q3/Q5/Q7 均有同族 trace", "该点可进入正式推荐"],
        ["work_mem replay", "动态峰值 4415MB；spill 0", "1150MB 已覆盖关键算子需求"],
        ["SB replay", "TP-SB hit = 99.39%", "从 1024MB 开始接近最大值"],
        ["Linux cache replay", "TP-OS conditional hit = 86.00%", "剩余 SB miss 多数仍被 OS 吸收"],
        ["最终结果", "combined = 99.915%；物理 I/O 2895MB", "继续增大 SB/work_mem 不再降低预测 I/O"],
    ]
    add_table(slide, 6, 3, example, 0.72, 2.50, 11.90, 3.72, widths=[2.15, 4.25, 5.50], font_size=12)
    footer(slide, 5)
    outline.append("5. S5 1024/1150 的完整数值算例。")

    # 6 Accuracy overview
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, "目前的准确率：必须分成三项看", "这三项使用不同验证数据，不能相加成一个总分。", "05 / Accuracy")
    conclusion_bar(slide, "当前最强证据：命中率误差低、5/6 可观测内存边界正确、S5 TPS 平台精确命中。")
    status_box(slide, "① 命中率预测", "Held-out MAE：SB 0.61pp；OS 1.82pp；combined 0.50pp。", 0.76, 2.65, 3.75, TEAL)
    status_box(slide, "② 动态内存边界", "6 个可观测查询中 5 个 operational 最小值正确，即 83.3%。", 4.79, 2.65, 3.75, ORANGE)
    status_box(slide, "③ 配置推荐", "唯一非限速阶段 S5：预测和实际 TPS 99% 平台都为 1024MB。", 8.82, 2.65, 3.75, GREEN)
    add_bullets(slide, ["命中率 MAE 来自多锚点 path replay 的 held-out SB 点", "83.3% 不包含无可实现 no-spill 配置的 Q18/Q21", "S1-S4 的 TPS 被目标速率限制，不能计算“最优点命中率”", "因此目前不能声称五阶段联合最优准确率为 100%"], 1.02, 4.45, 11.10, 1.65, size=15, spacing=10)
    footer(slide, 6)
    outline.append("6. 三类准确率总览：命中率、内存边界、配置推荐。")

    # 7 Hit accuracy
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, "准确率 1：缓存命中率预测误差", "单位为百分点（pp），越低越好。", "06 / Accuracy")
    conclusion_bar(slide, "在 held-out SB 点上，combined 命中率平均误差 0.50pp；OS 子模型仍是主要误差来源。")
    slide.shapes.add_picture(str(HIT_CHART), Inches(0.82), Inches(2.50), width=Inches(7.85), height=Inches(2.83))
    add_rect(slide, 9.05, 2.50, 3.35, 2.83, LIGHT, line=TEAL, radius=True)
    add_text(slide, "如何理解", 9.35, 2.79, 2.75, 0.34, size=20, color=TEAL, bold=True, align=PP_ALIGN.CENTER)
    add_bullets(slide, ["0.50pp 不是 50%", "例如真实 99.0%，平均预测约在 ±0.5pp 范围", "OS 误差 1.82pp，仍需更多真实 refault 数据"], 9.33, 3.37, 2.75, 1.55, size=13, spacing=9)
    add_text(slide, "说明：这是已有多锚点数据上的 held-out 结果，不是新采集的完全盲测集。", 1.00, 5.83, 11.35, 0.45, size=13, color=RED, bold=True, align=PP_ALIGN.CENTER)
    footer(slide, 7)
    outline.append("7. 命中率子模型 MAE：SB 0.61pp、OS 1.82pp、combined 0.50pp。")

    # 8 Dynamic memory accuracy
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, "准确率 2：最小 no-spill work_mem 是否预测正确？", "可观测的 6 个查询中，5 个 operational boundary 正确。", "07 / Accuracy")
    boundary = [
        ["查询", "预测最小值", "实测最小值", "结论"],
        ["Q1", "1MB", "1MB", "正确"],
        ["Q3", "1150MB", "1150MB", "逐 MB 精确"],
        ["Q5", "997MB", "305MB", "错误：305MB 切换计划"],
        ["Q7", "1083MB", "1083MB", "数值精确；边界计划变化"],
        ["Q9", "5707MB", "5707MB", "逐 MB 精确"],
        ["Q13", "1174MB", "1174MB", "逐 MB 精确"],
        ["Q18", "16539MB", "主机无法提供", "不计入准确率"],
        ["Q21", "16732MB", "引擎分配不可行", "不计入准确率"],
    ]
    add_table(slide, 9, 4, boundary, 0.62, 1.70, 8.05, 4.95, widths=[1.00, 1.75, 1.95, 3.35], font_size=11)
    add_kpi(slide, "5 / 6", "可观测边界正确", 9.17, 1.93, 2.65, color=GREEN, note="83.3%")
    add_kpi(slide, "3 / 3", "同计划边界逐 MB 精确", 9.17, 3.45, 2.65, color=TEAL, note="Q3 / Q9 / Q13")
    add_rect(slide, 9.05, 5.05, 3.00, 1.03, LIGHT, line=RED, radius=True)
    add_text(slide, "主要失败：Q5", 9.26, 5.24, 2.58, 0.28, size=16, color=RED, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "原因不是内存公式，而是优化器在 305MB 换了执行路径。", 9.26, 5.61, 2.58, 0.38, size=11, align=PP_ALIGN.CENTER)
    footer(slide, 8)
    outline.append("8. work_mem 边界验证：5/6 正确，Q5 因计划切换失败。")

    # 9 S5 TPS
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, "最重要的结果图：模型推荐点的真实 TPS 是不是最高？", "S5 是唯一没有被 40 TPS 目标速率限制、可以直接验证最高点的阶段。", "08 / Accuracy")
    conclusion_bar(slide, "是：模型推荐 1024MB；真实测试在 1024MB 得到 220.04 TPS，也是实测最高值。", color=GREEN)
    slide.shapes.add_picture(str(S5_CHART), Inches(0.72), Inches(2.36), width=Inches(9.10), height=Inches(3.64))
    add_kpi(slide, "1024MB", "模型事先推荐", 9.92, 2.62, 2.45, color=TEAL)
    add_kpi(slide, "220.04", "该点真实 TPS", 9.92, 4.02, 2.45, color=GREEN)
    add_text(slide, "它不是预测 TPS 数值，而是预测“最高 TPS 平台中的最小配置点”。", 0.98, 6.18, 11.35, 0.42, size=14, color=INK, bold=True, align=PP_ALIGN.CENTER)
    footer(slide, 9)
    outline.append("9. S5：预测 1024MB 拐点精确对应真实 TPS 平台，Pearson 0.990。")

    # 10 Joint combination
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, "SB 和 work_mem 到底是什么关系？", "不是两个独立答案取交集，而是逐个重放二维组合。", "09 / Joint Combination")
    conclusion_bar(slide, "二者共同消耗内存，并通过 OS cache 和 spill 相互影响；模型比较的是每一个 (SB, work_mem) 点。", color=GREEN)
    status_box(slide, "SB 增大", "TP-SB hit 提高；但固定缓存占用增加，MemAvailable 和 OS cache 余量下降。", 0.72, 2.48, 3.72, TEAL)
    status_box(slide, "work_mem 增大", "Join/Agg/Sort spill 减少；但动态内存峰值增加，也会压缩 OS cache。", 4.80, 2.48, 3.72, ORANGE)
    status_box(slide, "联合结果", "重新得到 TP miss、spill I/O、内存安全和总内存占用，再选择 Pareto 点。", 8.88, 2.48, 3.72, GREEN)
    combo_table = [
        ["SB + work_mem", "TP-SB hit", "spill I/O", "预测物理 I/O", "总内存占用", "结论"],
        ["512 + 1137", "97.06%", "1.67GiB", "4.50GiB", "4.80GiB", "TP 未到平台"],
        ["1024 + 1137", "99.39%", "1.67GiB", "4.50GiB", "5.30GiB", "仍有 spill"],
        ["1024 + 1150", "99.39%", "0", "2.83GiB", "5.31GiB", "推荐"],
        ["2048 + 1150", "99.80%", "0", "2.83GiB", "6.31GiB", "I/O 不降，多占 1GB"],
    ]
    add_table(slide, 5, 6, combo_table, 0.64, 4.18, 12.05, 2.28, widths=[1.78, 1.55, 1.50, 1.85, 1.85, 3.52], font_size=10)
    footer(slide, 10)
    outline.append("10. 联合关系：SB/work_mem 共同影响 OS cache、spill 和内存余量，逐点比较二维候选。")

    # 11 Recommendations
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, "当前五阶段推荐：哪些可信，哪些只是暂定？", "推荐值后面必须同时给出验证状态。", "10 / Recommendation")
    data = [["阶段", "推荐 SB", "推荐 work_mem", "当前证据"]]
    names = {"stage1_memory_rich": "S1", "stage2_reach_limit": "S2", "stage3_protect_tp": "S3", "stage4_backpressure": "S4", "stage5_tp_surge": "S5"}
    evidence = {
        "S1": "TPS 限速；work_mem floor 已验证",
        "S2": "TPS 限速；Q3 1150MB 精确",
        "S3": "最低 P95 在 512MB；Q7 1083MB 精确",
        "S4": "SB 从 256MB 达目标 TPS；work_mem 暂定",
        "S5": "1024MB TPS 平台精确；1150MB 无 spill",
    }
    for row in recs:
        short = names[row["stage"]]
        data.append([short, f"{row['recommended_sb_mb']}MB", f"{row['recommended_work_mem_mb']}MB", evidence[short]])
    add_table(slide, 6, 4, data, 0.70, 1.78, 11.92, 4.25, widths=[1.10, 1.75, 2.25, 6.82], font_size=13)
    add_rect(slide, 0.92, 6.27, 11.50, 0.48, INK, radius=True)
    add_text(slide, "不能写成“五阶段全部准确”：S5 已证实；S3 基本一致；S1/S2 不可识别；S4 待补计划 trace。", 1.14, 6.39, 11.06, 0.24, size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    footer(slide, 11)
    outline.append("11. 推荐表及每阶段证据强度。")

    # 12 S4 limitation
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, "为什么 S4 还不能说找到了真实最优？", "当前 512MB 是“有同计划 trace 的候选中最优”，不是全局最优。", "11 / Limitation")
    conclusion_bar(slide, "Q21 在 1024MB 以上切换到新计划族；缺少该路径的算子 trace，模型主动停止正式外推。", color=RED)
    add_bullets(slide, ["128/256/512MB：Q9/Q13/Q18/Q21 都有同计划族 trace，可以正式比较", "1024/1174MB：Q21 计划改变，当前数值只能作为 provisional，不能进入推荐", "Q18 受 max_dynamic_memory 限制；Q21 的全 no-spill bucket 分配超过 MaxAllocSize", "因此 S4 不存在简单的“把 work_mem 加到不 spill”方案"], 1.05, 2.58, 11.15, 2.25, size=17, spacing=13)
    add_rect(slide, 1.02, 5.30, 11.25, 0.82, LIGHT, line=TEAL, radius=True)
    add_text(slide, "下一步验证：补采 Q21 新计划族 trace → 扩展 S4 work_mem 网格 → 对推荐点及相邻点做真实二维复跑。", 1.36, 5.55, 10.58, 0.31, size=16, color=TEAL, bold=True, align=PP_ALIGN.CENTER)
    footer(slide, 12)
    outline.append("12. S4 限制与下一步验证。")

    # 13-17 Five simple stage figures
    simple_slides = [
        ("S1", "推荐 256MB SB + 1MB work_mem；TPS 被限速，SB 最优无法识别；1MB 已足够无 spill。"),
        ("S2", "推荐 256MB SB + 1150MB work_mem；TPS 被限速；Q3 的 1150MB no-spill 边界逐 MB 验证。"),
        ("S3", "推荐 512MB SB + 1083MB work_mem；TPS 被限速但最低 P95 在 512MB；work_mem 边界验证通过。"),
        ("S4", "推荐 256MB SB + 512MB work_mem；TPS 被限速；1024MB 以上灰柱缺 Q21 同计划 trace，当前值暂定。"),
        ("S5", "推荐 1024MB SB + 1150MB work_mem；1024MB 命中实测最高 TPS；1150MB 为首个零 spill 支持点。"),
    ]
    for page, (short, explanation) in enumerate(simple_slides, start=13):
        slide = prs.slides.add_slide(blank)
        set_bg(slide)
        add_title(slide, f"五阶段统一结果图：{short}", "左边是真实 TPS，右边是预测 spill；红星表示模型推荐点。", "Result")
        image = SIMPLE_FIGURES / f"{short.lower()}_simple_tps_workmem_recommendation.png"
        slide.shapes.add_picture(str(image), Inches(0.62), Inches(1.52), width=Inches(12.10), height=Inches(4.28))
        add_rect(slide, 0.90, 6.08, 11.54, 0.58, INK, radius=True)
        add_text(slide, explanation, 1.14, 6.23, 11.06, 0.27, size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        footer(slide, page)
        outline.append(f"{page}. {short} 统一 TPS/spill 结果图。")

    prs.save(OUT)
    OUTLINE.write_text("# 文字说明版 PPT 大纲\n\n" + "\n".join(f"- {line}" for line in outline) + "\n", encoding="utf-8")
    print(OUT)
    print(OUTLINE)


if __name__ == "__main__":
    build()
