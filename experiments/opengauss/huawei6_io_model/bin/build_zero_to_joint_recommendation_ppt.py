#!/usr/bin/env python3
"""Build a newcomer-oriented deck for the complete Huawei6 model workflow."""

from __future__ import annotations

import csv
import json
import shutil
import statistics
from pathlib import Path

import matplotlib.pyplot as plt
import numpy as np
from matplotlib import font_manager
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches

from build_memory_autonomy_paper_ppt import (
    BLUE,
    GRAY,
    GREEN,
    INK,
    LIGHT,
    MID,
    ORANGE,
    PALE_GREEN,
    PALE_ORANGE,
    PALE_RED,
    PALE_TEAL,
    RED,
    TEAL,
    WHITE,
    add_arrow,
    add_bullets,
    add_callout,
    add_flow_box,
    add_kpi,
    add_rect,
    add_table,
    add_text,
    add_title,
    configure_plot_font,
    rgb_hex,
    set_bg,
)


ROOT = Path(__file__).resolve().parents[1]
LATEST = ROOT / "artifacts" / "00_latest"
LATEST.mkdir(parents=True, exist_ok=True)
OUT = LATEST / "Huawei6_从零启动_双向联合配置预测全流程_20260802.pptx"
ROOT_COPY = Path("/root/Huawei6_双向联合预测_从零到推荐全流程.pptx")
OUTLINE = LATEST / "Huawei6_从零启动_双向联合配置预测全流程_20260802_讲稿提纲.md"

SURFACE_CSV = ROOT / "results/mixed_storage_surface_sync_strict_20260802/train/mixed_storage_surface.csv"
SURFACE_FROZEN = ROOT / "results/mixed_storage_surface_sync_strict_20260802/frozen/frozen_surface.json"
FORMULA_REPORT = ROOT / "results/tp_io_causal_validation_20260802/final_buffered_path_holdout_v1_evaluation/online_tps_holdout_report.json"
FORMULA_ROWS = ROOT / "results/tp_io_causal_validation_20260802/final_buffered_path_holdout_v1_evaluation/online_tps_holdout_comparisons.csv"
DECISIONS = ROOT / "results/huawei6_observation_driven_joint_prediction_20260802_final_v2/summary.json"
STAGE_REPORT = ROOT / "results/huawei6_observation_driven_five_stage_equal_tps_20260802/validation_report.json"

CHART_SURFACE = LATEST / "portable_new_machine_storage_surface_20260802.png"
CHART_FORMULA = LATEST / "portable_formula_holdout_explained_20260802.png"
CHART_STAGE = LATEST / "portable_five_stage_actions_tps_20260802.png"


def read_json(path: Path):
    return json.loads(path.read_text(encoding="utf-8"))


def read_csv(path: Path):
    with path.open(newline="", encoding="utf-8") as handle:
        return list(csv.DictReader(handle))


def footer(slide, page: int, source: str = "") -> None:
    add_rect(slide, 0.62, 7.17, 12.08, 0.012, MID)
    add_text(slide, "Huawei6 · 从零标定到双向联合推荐 · 2026-08-02", 0.62, 7.22, 6.2, 0.18, size=7.5, color=GRAY)
    if source:
        add_text(slide, source, 6.4, 7.22, 5.75, 0.18, size=6.8, color=GRAY, align=PP_ALIGN.RIGHT)
    add_text(slide, str(page), 12.28, 7.22, 0.42, 0.18, size=8, color=GRAY, align=PP_ALIGN.RIGHT)


def make_charts() -> None:
    configure_plot_font()
    cjk_font = "/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc"
    if Path(cjk_font).exists():
        font_manager.fontManager.addfont(cjk_font)
        plt.rcParams["font.sans-serif"] = [
            font_manager.FontProperties(fname=cjk_font).get_name()
        ]

    rows = read_csv(SURFACE_CSV)
    by_qd: dict[int, list[float]] = {}
    for row in rows:
        by_qd.setdefault(int(row["ap_queue_depth"]), []).append(float(row["tp_await_ms"]))
    qd = sorted(by_qd)
    actual = [statistics.fmean(by_qd[item]) for item in qd]
    baseline = actual[0]
    frozen = read_json(SURFACE_FROZEN)
    predicted = [baseline + float(frozen["tp_added_await_ms_by_ap_queue_depth"][str(item)]) for item in qd]
    fig, ax = plt.subplots(figsize=(10.8, 4.35))
    ax.plot(qd, actual, "o-", linewidth=2.8, markersize=7, color=rgb_hex(BLUE), label="本机实测 TP 8KiB await")
    ax.plot(qd, predicted, "s--", linewidth=2.0, markersize=5, color=rgb_hex(ORANGE), label="冻结设备曲面")
    for x, y in zip(qd, actual):
        ax.annotate(f"{y:.1f}", (x, y), xytext=(0, 8), textcoords="offset points", ha="center", fontsize=9)
    ax.set_title("新机器必须先测：AP 队列压力如何抬高 TP 物理 I/O 延迟", fontsize=15, fontweight="bold")
    ax.set_xlabel("AP 128KiB 随机读队列深度 qAP")
    ax.set_ylabel("TP 8KiB 请求延迟 (ms)")
    ax.grid(alpha=0.2)
    ax.spines[["top", "right"]].set_visible(False)
    ax.legend(frameon=False, loc="upper left")
    fig.tight_layout()
    fig.savefig(CHART_SURFACE, dpi=190, bbox_inches="tight")
    plt.close(fig)

    formula_rows = [row for row in read_csv(FORMULA_ROWS) if int(row["external_queue_depth"]) > 0]
    depths = sorted({int(row["external_queue_depth"]) for row in formula_rows})
    pred_lat = [statistics.fmean(float(row["predicted_tp_request_await_ms"]) for row in formula_rows if int(row["external_queue_depth"]) == depth) for depth in depths]
    real_lat = [statistics.fmean(float(row["actual_tp_request_await_ms"]) for row in formula_rows if int(row["external_queue_depth"]) == depth) for depth in depths]
    pred_tps = [statistics.fmean(float(row["predicted_tp_tps"]) for row in formula_rows if int(row["external_queue_depth"]) == depth) for depth in depths]
    real_tps = [statistics.fmean(float(row["actual_tp_tps"]) for row in formula_rows if int(row["external_queue_depth"]) == depth) for depth in depths]
    fig, axes = plt.subplots(1, 2, figsize=(11.3, 4.1))
    x = np.arange(len(depths))
    width = 0.34
    axes[0].bar(x - width / 2, pred_lat, width, color=rgb_hex(BLUE), label="预测")
    axes[0].bar(x + width / 2, real_lat, width, color="#aab4ba", label="实测")
    axes[0].set_xticks(x, [f"QD{item}" for item in depths])
    axes[0].set_title("TP I/O latency", fontweight="bold")
    axes[0].set_ylabel("ms")
    axes[0].legend(frameon=False)
    axes[1].bar(x - width / 2, pred_tps, width, color=rgb_hex(TEAL), label="预测")
    axes[1].bar(x + width / 2, real_tps, width, color="#aab4ba", label="实测")
    axes[1].set_xticks(x, [f"QD{item}" for item in depths])
    axes[1].set_title("TP TPS", fontweight="bold")
    axes[1].set_ylabel("TPS")
    axes[1].legend(frameon=False)
    for ax in axes:
        ax.grid(axis="y", alpha=0.18)
        ax.spines[["top", "right"]].set_visible(False)
    fig.suptitle("冻结后未见压力点：先预测，再注入 I/O 压力", fontsize=15, fontweight="bold")
    fig.tight_layout()
    fig.savefig(CHART_FORMULA, dpi=190, bbox_inches="tight")
    plt.close(fig)

    report = read_json(STAGE_REPORT)
    stages = report["stages"]
    labels = [row["stage"] for row in stages]
    protected = [float(row["protected_tp_tps"]) for row in stages]
    total = [float(row["protected_tp_tps"]) + float(row["surge_tp_tps"]) for row in stages]
    fig, ax = plt.subplots(figsize=(11.4, 4.35))
    x = np.arange(len(labels))
    stable_mean = statistics.fmean(protected)
    ax.fill_between([-0.3, 4.3], stable_mean * 0.95, stable_mean * 1.05, color=rgb_hex(GREEN), alpha=0.14, label="S1-S5 保护 TPS ±5%")
    ax.plot(x, protected, "o-", linewidth=2.8, markersize=8, color=rgb_hex(BLUE), label="实测保护 TP TPS")
    ax.plot(x, total, "o--", linewidth=2.0, markersize=6, color=rgb_hex(ORANGE), label="总 TP TPS（S5 含新增需求）")
    for index, value in enumerate(protected):
        ax.annotate(f"{value:.0f}", (index, value), xytext=(0, 9), textcoords="offset points", ha="center", fontsize=9, fontweight="bold")
    ax.set_xticks(x, labels)
    ax.set_ylabel("TPS")
    ax.set_ylim(0, 5100)
    variation = (max(protected) - min(protected)) / statistics.fmean(protected) * 100
    ax.set_title(f"推荐配置执行后：S1-S5 保护 TP TPS 波动 {variation:.2f}%", fontsize=15, fontweight="bold")
    ax.grid(axis="y", alpha=0.18)
    ax.spines[["top", "right"]].set_visible(False)
    ax.legend(frameon=False, loc="upper left")
    fig.tight_layout()
    fig.savefig(CHART_STAGE, dpi=190, bbox_inches="tight")
    plt.close(fig)


def cover(prs, outline):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, INK)
    add_rect(slide, 0, 0, 0.22, 7.5, TEAL)
    add_text(slide, "Huawei6", 0.92, 0.88, 4.0, 0.38, size=18, color=TEAL, bold=True)
    add_text(slide, "从零开始运行双向联合预测", 0.90, 1.42, 11.6, 0.70, size=35, color=WHITE, bold=True)
    add_text(slide, "新机器如何标定，Trace 如何描述 TP/AP，\n以及 1→2→3 / 2→1→3 如何找到 SB、work_mem 与 AP 准入", 0.93, 2.45, 10.9, 1.03, size=21, color=WHITE)
    add_rect(slide, 0.94, 4.16, 8.55, 0.58, TEAL, radius=True)
    add_text(slide, "机器标定  →  候选回放  →  双向搜索  →  I/O/TPS 校正  →  真实验证", 1.06, 4.32, 8.30, 0.25, size=13.5, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_callout(slide, "核心原则", "推荐过程不读取候选配置的真实混合 TPS；真实运行只用于冻结推荐之后的独立验证。", 0.94, 5.33, 9.28, 0.86, TEAL, PALE_TEAL)
    add_text(slide, "项目执行说明 · 面向第一次接触本项目的读者", 0.95, 6.66, 6.2, 0.25, size=11, color=GRAY)
    outline.append("1. 封面：从新机器标定到双向联合推荐的完整执行流程。")


def problem(prs, outline):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(slide)
    add_title(slide, "项目要解决什么问题？", "TP 要稳定，但 AP、TP 和数据库内存共享同一台机器，三个参数形成闭环。", "01 问题")
    items = [
        ("shared_buffers", "大：TP miss 少\n小：给 AP 留更多内存", BLUE, PALE_TEAL),
        ("逐 Query work_mem", "大：AP 少 spill\n小：动态内存占用低", ORANGE, PALE_ORANGE),
        ("AP 准入 / cap", "高：AP 进度快\n低：减少 I/O 和内存竞争", RED, PALE_RED),
    ]
    for index, (title, body, color, fill) in enumerate(items):
        add_flow_box(slide, title, body, 0.76 + index * 4.18, 1.82, 3.72, 1.56, color, fill=fill)
    add_arrow(slide, 2.61, 3.52, 2.61, 4.23, BLUE, 2.0)
    add_arrow(slide, 6.79, 3.52, 6.79, 4.23, ORANGE, 2.0)
    add_arrow(slide, 10.97, 3.52, 10.97, 4.23, RED, 2.0)
    add_flow_box(slide, "共同结果：TP 物理 I/O + AP spill I/O → NVMe 排队 → TP 延迟与 TPS", "仅预测缓存命中率不够；相同 hit ratio 下，AP spill 仍可能把 TP TPS 压低。", 1.24, 4.35, 10.84, 1.20, TEAL, fill=PALE_GREEN)
    add_callout(slide, "最终输出", "每个阶段推荐一个完整组合：(SB，逐 Query work_mem，AP-cap / 是否阻塞新 AP)，目标是 TP 稳定且 AP 尽量前进。", 1.24, 5.93, 10.84, 0.74, GREEN, PALE_GREEN)
    footer(slide, 2)
    outline.append("2. 问题：SB、work_mem、AP 准入共同影响内存、spill I/O 与 TP TPS。")


def lifecycle(prs, outline):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(slide)
    add_title(slide, "完整生命周期：三类工作不能混在一起", "一次性机器标定、每次负载建模、运行时/阶段级选择各自有不同输入和频率。", "02 总览")
    blocks = [
        ("A. 新机器一次性标定", "设备曲面 qAP→ΔLdevice\n数据库路径系数 kpath\nTP-only 基线 X0、L0、n0\n冻结后用未见 QD 验证", BLUE, PALE_TEAL),
        ("B. 新负载 / Trace 建模", "TP：SB→miss / I/O per tx\nAP：plan×work_mem→动态峰值、spill\nAP mix / cap→qAP 与 utility\n生成候选 CSV", ORANGE, PALE_ORANGE),
        ("C. 每阶段双向选择", "TP-first：1→2→3\nAP-first：2→1→3\n同一公式比较完整候选\n重启应用 SB，运行并监控", GREEN, PALE_GREEN),
    ]
    for index, (title, body, color, fill) in enumerate(blocks):
        x = 0.66 + index * 4.22
        add_flow_box(slide, title, body, x, 1.78, 3.72, 2.45, color, fill=fill)
        if index < 2:
            add_arrow(slide, x + 3.77, 3.00, x + 4.15, 3.00, TEAL, 2.1)
    add_callout(slide, "为什么分层", "机器不同，I/O 延迟规律必须重测；负载不同，plan、spill 和 TP miss 必须重放；阶段不同，只需在已生成的候选中执行双向联合搜索。", 0.92, 4.82, 11.50, 0.91, TEAL, PALE_TEAL)
    add_text(slide, "标定数据不等于测试答案；候选实际 TPS 不参与模型冻结。", 1.35, 6.18, 10.6, 0.34, size=17, color=RED, bold=True, align=PP_ALIGN.CENTER)
    footer(slide, 3)
    outline.append("3. 生命周期：机器标定、负载 Trace 建模和阶段级双向选择分离。")


def bootstrap(prs, outline):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(slide)
    add_title(slide, "从零启动：一份配置 + 一条命令", "控制器检查依赖、记录状态、允许中断续跑，并在缺少候选时拒绝伪装成完整运行。", "03 新机器")
    add_flow_box(slide, "填写 machine.json", "NVMe 设备名、openGauss 路径\nTP 命令与 terminals\n标定文件目录和 CPU mask\n候选 replay 命令", 0.68, 1.72, 3.30, 2.15, BLUE, fill=PALE_TEAL)
    add_flow_box(slide, "执行启动器", "export 数据库密码\nrun_portable_model.sh machine.json\nworkspace 加互斥锁\n所有输出写 modelctl.log", 5.02, 1.72, 3.30, 2.15, TEAL, fill=PALE_GREEN)
    add_flow_box(slide, "状态与恢复", "state.json 保存每阶段状态\n校验产物大小与哈希\n失败后同一命令续跑\n数据库在 finally 中恢复", 9.36, 1.72, 3.30, 2.15, GREEN, fill=PALE_GREEN)
    add_arrow(slide, 4.04, 2.80, 4.94, 2.80, TEAL, 2.2); add_arrow(slide, 8.38, 2.80, 9.28, 2.80, TEAL, 2.2)
    add_callout(slide, "完整顺序", "doctor → prepare-files → storage calibration → TP path anchors → freeze model → unseen holdout → candidate replay → TP/AP 双向推荐", 0.86, 4.49, 11.66, 0.78, TEAL, PALE_TEAL)
    add_text(slide, "bin/run_portable_model.sh /absolute/path/to/machine.json", 1.33, 5.70, 10.67, 0.44, size=17, color=INK, bold=True, align=PP_ALIGN.CENTER, font="Consolas")
    add_callout(slide, "原版 openGauss", "不要求修改内核；shared_buffers 作为阶段级静态配置，阶段自然结束后重启应用。", 2.34, 6.20, 8.65, 0.56, GREEN, PALE_GREEN)
    footer(slide, 4, "huawei6_modelctl.py / run_portable_model.sh")
    outline.append("4. 一键启动：配置、锁、日志、状态、哈希和数据库恢复。")


def storage_calibration(prs, outline):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(slide)
    add_title(slide, "机器标定 1：先测 NVMe 的混合 I/O 延迟曲面", "数据库暂时停止，避免数据库后台 I/O 污染设备本身的响应曲线。", "04 设备曲面")
    slide.shapes.add_picture(str(CHART_SURFACE), Inches(0.72), Inches(1.55), width=Inches(7.65))
    add_callout(slide, "TP 代理", "8KiB、同步阻塞读；线程数与 TP terminals 一致，模拟数据库页 miss 的关键路径。", 8.68, 1.70, 3.88, 1.02, BLUE, PALE_TEAL)
    add_callout(slide, "AP 代理", "128KiB 随机读；扫描 QD 0/2/4/8/16/32，模拟 spill / 扫描造成的队列压力。", 8.68, 2.95, 3.88, 1.02, ORANGE, PALE_ORANGE)
    add_callout(slide, "冻结与留出", "训练曲面后才运行 QD6/12/24 留出点；当前设备曲面留出 MAPE 3.74%。", 8.68, 4.20, 3.88, 1.02, GREEN, PALE_GREEN)
    add_callout(slide, "得到的函数", "ΔLdevice(qAP)：在这台机器上，给定 AP 队列深度，TP 物理请求会额外等待多少毫秒。", 8.68, 5.45, 3.88, 0.95, TEAL, PALE_TEAL)
    footer(slide, 5, "mixed_storage_surface_sync_strict_20260802")
    outline.append("5. 设备标定：用混合 I/O 实验得到 qAP 到 TP 增量延迟的本机曲面。")


def path_calibration(prs, outline):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(slide)
    add_title(slide, "机器标定 2：把裸设备曲面迁移到真实 openGauss TP 路径", "设备延迟不能直接等同数据库延迟；必须用 BPF 跟踪 openGauss LWTID 的实际块 I/O。", "05 数据库路径")
    steps = [
        ("运行 TP-only 基线", "采集 X0：TPS\nL0：TP 请求 await\nn0：关键 I/O / transaction", BLUE, PALE_TEAL),
        ("注入路径锚点", "在 QD6/12/24 下施压\n仅使用前后 TP 请求 latency\n拟合 kpath(q)", ORANGE, PALE_ORANGE),
        ("冻结模型包", "设备曲面 + kpath\n机器 inventory + domain\n明确禁止 TPS 拟合系数", TEAL, PALE_GREEN),
        ("未见 QD 盲测", "先写 prediction.json\n再注入 QD9/QD18\n检查 latency/TPS 门槛", GREEN, PALE_GREEN),
    ]
    for i, (title, body, color, fill) in enumerate(steps):
        x = 0.56 + i * 3.20
        add_flow_box(slide, title, body, x, 1.78, 2.70, 1.85, color, fill=fill)
        if i < 3: add_arrow(slide, x + 2.73, 2.70, x + 3.14, 2.70, TEAL, 2.0)
    add_callout(slide, "路径迁移公式", "Lpred = L0 + kpath(qAP) × ΔLdevice(qAP)", 1.10, 4.18, 5.20, 0.84, TEAL, PALE_TEAL)
    add_callout(slide, "防止结果泄漏", "kpath 只读取请求 latency 字段；压力期 TPS 被保留用于审计，但不参与路径系数冻结。", 7.02, 4.18, 5.20, 0.84, RED, PALE_RED)
    add_callout(slide, "模型有效域", "当前 v1：openGauss buffered blocking TP 读（约 8KiB）+ 128KiB random-read AP 压力；超出请求类型或 QD 范围直接拒绝。", 1.10, 5.43, 11.12, 0.88, BLUE, PALE_TEAL)
    footer(slide, 6, "portable_tp_path_probe.py / portable_joint_model.py")
    outline.append("6. 路径标定：用 BPF 锚点得到 kpath，冻结后再做未见 QD 验证。")


def workload_features(prs, outline):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(slide)
    add_title(slide, "负载如何进入模型：不输入阶段答案，只输入可观测事实", "机器公式回答“这些 I/O 会造成多大延迟”；Trace Replay 回答“某个配置会产生多少 I/O”。", "06 负载输入")
    add_flow_box(slide, "TP 侧", "TP terminals / offered rate\n每个 SB 的 TP-only TPS0\nL0 与 n0\ncache replay 的 n_candidate(B)", 0.70, 1.70, 3.45, 2.20, BLUE, fill=PALE_TEAL)
    add_flow_box(slide, "AP 侧", "活跃 / 到达 Query ID\nplan family × work_mem\ndynamic_peak(q,w)\nspill_io(q,w)、service time", 4.94, 1.70, 3.45, 2.20, ORANGE, fill=PALE_ORANGE)
    add_flow_box(slide, "运行环境", "可用总内存\n当前 SB 与动态内存\nAP-cap / queued AP\nBPF 或 replay 给出的 qAP", 9.18, 1.70, 3.45, 2.20, GREEN, fill=PALE_GREEN)
    add_callout(slide, "候选行的含义", "一行不是“单独一个 SB”或“单独一个 work_mem”，而是一个完整状态：stage + SB + 每条 Query 的 work_mem + AP-cap + TP I/O 特征 + AP qAP。", 0.94, 4.48, 11.50, 0.92, TEAL, PALE_TEAL)
    add_callout(slide, "计划变化", "不同 work_mem 可选择不同 plan family；只有 replay 支持且内存安全的 plan 候选进入搜索。未知 plan 的 spill→qAP 仍需 operator/source replay 或额外锚点。", 0.94, 5.68, 11.50, 0.86, RED, PALE_RED)
    footer(slide, 7, "candidate_source.command → candidates.csv")
    outline.append("7. 负载输入：TP cache/I/O 特征、AP plan/spill 特征和当前资源状态。")


def coupling(prs, outline):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(slide)
    add_title(slide, "为什么 SB 与 work_mem 必须联合枚举？", "两者通过总内存和 I/O 竞争耦合，不能分别取各自最优再简单相加。", "07 候选空间")
    add_callout(slide, "内存可行性", "M(B,{wq}) = B + Σ active dynamic_peak(q,wq) + reserve\n若 M 超过可用内存，候选立即淘汰。", 0.72, 1.72, 5.62, 1.06, BLUE, PALE_TEAL)
    add_callout(slide, "性能闭环", "SB↓ 可能让 AP work_mem↑、spill↓、qAP↓；但也可能让 TP n_candidate↑。最终必须比较两者共同作用后的 TPS。", 6.98, 1.72, 5.62, 1.06, ORANGE, PALE_ORANGE)
    headers = ["候选", "SB", "work_mem / Query", "AP-cap", "TP I/O/tx", "qAP", "结论"]
    data = [headers,
            ["A", "8192", "Q18=1150", "1", "低", "低", "TP 富裕"],
            ["B", "4096", "Q18/Q21=1150", "2", "中", "更低", "为 AP 让 SB"],
            ["C", "4096", "Q9/Q13/Q18/Q21=256", "4", "中", "高", "限制 AP 内存"],
            ["D", "8192", "Q18/Q21=256", "2", "低", "中", "TP 突增保护"]]
    add_table(slide, data, 0.78, 3.34, 11.82, 2.30, widths=[0.8, 1.0, 3.35, 1.0, 1.45, 0.75, 2.20], font_size=9.3)
    add_text(slide, "搜索不会把所有组合都真实运行；Trace/公式先筛选，最终推荐才进入真实阶段验证。", 1.12, 6.17, 11.1, 0.34, size=15, color=TEAL, bold=True, align=PP_ALIGN.CENTER)
    footer(slide, 8)
    outline.append("8. 耦合：每个候选同时包含 SB、逐 Query work_mem 和 AP-cap，并先过内存门槛。")


def tp_first(prs, outline):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(slide)
    add_title(slide, "路线 A：TP-first，执行 1→2→3", "适合 TP 已接近容量边界时：先确认 TP 最少需要多少 SB，再让 AP 使用剩余资源。", "08 TP-first")
    add_flow_box(slide, "1. 找 TP 的 SB 性能下界", "比较各 SB 的 TP-only TPS0 与 miss/I/O 曲线。\n得到满足 TP 目标的最小值 B_TP_min；\n可行条件是 SB ≥ B_TP_min。", 0.66, 1.75, 3.55, 2.04, BLUE, fill=PALE_TEAL)
    add_flow_box(slide, "2. 剩余内存分给 AP", "固定 TP 的 SB 下界后，比较各 Query 的 work_mem。\n只保留 memory_safe / plan_supported，\n优先 AP utility 高的 grant。", 4.89, 1.75, 3.55, 2.04, ORANGE, fill=PALE_ORANGE)
    add_flow_box(slide, "3. 用 I/O 重新检查 TPS", "把该组合的 TP n_candidate 与 AP qAP\n代入本机冻结公式。\n若 TP SLO 不满足，收紧 AP 或回到步骤 1。", 9.12, 1.75, 3.55, 2.04, GREEN, fill=PALE_GREEN)
    add_arrow(slide, 4.25, 2.77, 4.81, 2.77, TEAL, 2.2); add_arrow(slide, 8.48, 2.77, 9.04, 2.77, TEAL, 2.2)
    add_callout(slide, "得到的不是单独 SB", "TP-first 最终输出完整候选 A=(SB, {wq}, AP-cap)，并带 predicted TPS、await、AP utility 和是否满足 SLO。", 1.10, 4.50, 11.12, 0.86, BLUE, PALE_TEAL)
    add_callout(slide, "典型动作", "TP 饱和时保留 SB；随后缩小新 AP 的 work_mem，必要时阻塞新 AP，但已运行 AP 自然结束。", 1.10, 5.67, 11.12, 0.82, GREEN, PALE_GREEN)
    footer(slide, 9, "portable_joint_model.choose_recommendations")
    outline.append("9. TP-first：先 TP SB 下界，再分 AP 内存，最后以 I/O/TPS 复核并迭代。")


def ap_first(prs, outline):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(slide)
    add_title(slide, "路线 B：AP-first，执行 2→1→3", "适合 TP 仍有余量而 AP 到达时：先避免 AP 因内存不足产生大量 spill，再反查 TP 可接受的 SB。", "09 AP-first")
    add_flow_box(slide, "2. 先找 AP 的合适 grant", "从 plan/operator replay 选择 AP utility 前沿：\n动态峰值可控、spill 少、service time 合理，\n并保留不同 Query 的独立 work_mem。", 0.66, 1.75, 3.55, 2.04, ORANGE, fill=PALE_ORANGE)
    add_flow_box(slide, "1. 得到 SB 的内存上界", "先扣除 AP grant 的动态峰值与系统预留：\nB_AP_max = M_total - M_AP - reserve。\n可行条件是 SB ≤ B_AP_max。", 4.89, 1.75, 3.55, 2.04, BLUE, fill=PALE_TEAL)
    add_flow_box(slide, "3. 同一公式复核", "计算完整候选 B 的 qAP、Lpred、TPSpred。\n若 AP 少 spill 带来的收益大于 SB 减少的代价，\n该路径可能胜出。", 9.12, 1.75, 3.55, 2.04, GREEN, fill=PALE_GREEN)
    add_arrow(slide, 4.25, 2.77, 4.81, 2.77, TEAL, 2.2); add_arrow(slide, 8.48, 2.77, 9.04, 2.77, TEAL, 2.2)
    add_callout(slide, "S2 的例子", "S1 的 SB=8192MB 很充裕；新 AP 到达后，AP-first 发现把 SB 降到 4096MB 可保留 Q18/Q21=1150MB grant，并减少 AP spill。", 1.10, 4.50, 11.12, 0.91, ORANGE, PALE_ORANGE)
    add_callout(slide, "为什么两条路线都要跑", "TP-first 防止为了 AP 牺牲 TP；AP-first 防止为了缓存保留过大的 SB，反而让 AP spill I/O 把 TP 拖慢。", 1.10, 5.71, 11.12, 0.79, TEAL, PALE_TEAL)
    footer(slide, 10)
    outline.append("10. AP-first：先 AP grant，再得到 SB 内存上界，最后用同一 I/O/TPS 公式复核。")


def bounds_intersection(prs, outline):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(slide)
    add_title(slide, "两条路线的几何意义：一个给下界，一个给上界", "联合配置首先必须落在上下界交集中；I/O→TPS 公式再从交集内选择最终点。", "10 上下界交集")
    add_flow_box(slide, "TP-first：性能下界", "B_TP_min\nTP 至少需要多大的 SB，\n才能满足 miss / TPS / SLO。\n\nSB ≥ B_TP_min", 0.78, 1.70, 3.30, 2.35, BLUE, fill=PALE_TEAL)
    add_flow_box(slide, "AP-first：内存上界", "B_AP_max = M_total - M_AP - reserve\n保留选定 AP grant 后，\n最多还能给 SB 多少内存。\n\nSB ≤ B_AP_max", 9.25, 1.70, 3.30, 2.35, ORANGE, fill=PALE_ORANGE)
    add_arrow(slide, 4.14, 2.86, 5.00, 2.86, BLUE, 2.2)
    add_arrow(slide, 9.18, 2.86, 8.32, 2.86, ORANGE, 2.2)
    add_flow_box(slide, "可行交集", "B_TP_min ≤ SB ≤ B_AP_max\n\n若交集非空：在其中比较 TPS 与 AP utility。\n若交集为空：降低 work_mem / AP-cap，\n或阻塞新 AP 后重新计算。", 5.08, 1.58, 3.16, 2.62, TEAL, fill=PALE_GREEN)
    add_callout(slide, "S2 的离散候选例子", "候选 SB={4096,8192}MB；TP 下界 B_TP_min=4096MB。Q18/Q21=1150MB grant 的动态峰值约 4402MB，在 10500MB 预算下 B_AP_max≈6098MB，因此离散交集只有 SB=4096MB。", 0.88, 4.63, 11.56, 1.05, TEAL, PALE_TEAL)
    add_callout(slide, "不要混淆", "TP-first 找的是“至少多少才够”，因此是下界；AP-first 找的是“保留 AP 后最多还能给多少”，因此是上界。TP 拐点附近收益变小，不代表它自动成为硬上界。", 0.88, 5.95, 11.56, 0.70, RED, PALE_RED)
    footer(slide, 11)
    outline.append("11. 上下界：TP-first 给性能下界，AP-first 给内存上界，最终在交集中选择。")


def formula(prs, outline):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(slide)
    add_title(slide, "共同的第 3 步：I/O latency 如何反哺 TPS", "这是两条路径的共同裁判；不存在一条路径用命中率、另一条路径用 TPS 的情况。", "11 公式")
    add_flow_box(slide, "候选输入", "qAP：AP plan/spill replay\nn：TP critical I/O / tx\nX0、L0、n0：该 SB 的 TP-only 基线\nN：TP terminals", 0.62, 1.67, 2.72, 2.12, BLUE, fill=PALE_TEAL)
    add_flow_box(slide, "预测 I/O 延迟", "ΔL = surface(qAP)\nLpred = L0 + kpath(qAP)×ΔL", 3.62, 1.67, 2.72, 2.12, ORANGE, fill=PALE_ORANGE)
    add_flow_box(slide, "拆出非 I/O 时间", "Rbase = N×1000/X0\nRnonio = max(0, Rbase - n0×L0)", 6.62, 1.67, 2.72, 2.12, TEAL, fill=PALE_GREEN)
    add_flow_box(slide, "得到候选 TPS", "Rpred = Rnonio + n×Lpred + extra_non_io\nTPSpred = N×1000/Rpred", 9.62, 1.67, 3.05, 2.12, GREEN, fill=PALE_GREEN)
    for x in (3.38, 6.38, 9.38): add_arrow(slide, x, 2.73, x + 0.18, 2.73, TEAL, 2.0)
    add_callout(slide, "不同配置如何进入公式", "SB 改变 X0/L0/n0/n；work_mem、plan mix 与 AP-cap 改变 qAP；CPU/锁等独立代价进入 extra_non_io。机器变化则重新测 surface 与 kpath。", 0.94, 4.40, 11.50, 0.94, TEAL, PALE_TEAL)
    add_callout(slide, "不是 TPS 回归", "TPS 换算系数固定为 1.0，直接来自 closed-loop response-time law；模型包明确记录 fitted_tps_coefficient=false。", 0.94, 5.63, 11.50, 0.83, RED, PALE_RED)
    footer(slide, 12, "portable_joint_model.py:169-267")
    outline.append("12. 公式：qAP→Lpred，再将物理 I/O 等待加入事务响应时间得到 TPSpred。")


def merge(prs, outline):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(slide)
    add_title(slide, "两条路线如何合并成一个动作", "每条路线都先完成 1/2 选择和第 3 步校正，再比较完整候选，而不是混搭两个局部最优。", "12 决策")
    add_flow_box(slide, "候选 A：TP-first", "完整 (SB,{wq},cap)\nTP SLO、await、AP utility", 0.72, 1.86, 3.05, 1.42, BLUE, fill=PALE_TEAL)
    add_flow_box(slide, "候选 B：AP-first", "完整 (SB,{wq},cap)\nTP SLO、await、AP utility", 0.72, 4.14, 3.05, 1.42, ORANGE, fill=PALE_ORANGE)
    add_arrow(slide, 3.82, 2.56, 4.73, 3.16, TEAL, 2.0); add_arrow(slide, 3.82, 4.84, 4.73, 3.78, TEAL, 2.0)
    add_flow_box(slide, "统一过滤与排序", "1. memory_safe / plan_supported\n2. 优先满足 TP SLO\n3. AP utility 更高\n4. TPS / await 打破平局", 4.87, 2.44, 3.30, 2.02, RED, fill=PALE_RED)
    add_arrow(slide, 8.23, 3.45, 9.05, 3.45, TEAL, 2.2)
    add_flow_box(slide, "最终阶段动作", "推荐 SB\n逐 Query work_mem\nAP-cap / block_new_ap\n预测 TPS 与原因", 9.18, 2.44, 3.45, 2.02, GREEN, fill=PALE_GREEN)
    add_callout(slide, "阶段状态约束", "TP 饱和后限制 AP grant；新 AP 到达且资源不足时阻塞；TP demand 比保护目标高 5% 以上时要求 SB 增大。已运行 AP 不被取消。", 1.20, 5.91, 10.92, 0.70, TEAL, PALE_TEAL)
    footer(slide, 13)
    outline.append("13. 合并：分别得到完整候选后，以可行性、TP SLO、AP utility 和 TPS 排序。")


def five_stage(prs, outline):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(slide)
    add_title(slide, "五阶段输出：模型从观测推导出与验收要求一致的动作", "推荐文件不包含阶段名或实际混合 TPS；下表是推荐完成后按 S1-S5 对齐展示。", "13 推荐结果")
    report = read_json(STAGE_REPORT)
    data = [["阶段", "观测变化", "推荐 SB", "推荐 work_mem", "AP 动作", "实测保护 TPS"]]
    descriptions = ["基准 TP，内存富余", "新增 AP", "TP 触及容量边界", "饱和且新 AP 到达", "TP demand 突增"]
    action_names = ["保持富裕", "降低 SB 给 AP", "降低 AP 内存", "阻塞新 AP", "提高 SB"]
    for desc, action, row in zip(descriptions, action_names, report["stages"]):
        work_mem = "；".join(f"{key.upper()}={value}MB" for key, value in row["recommended_work_mem"].items())
        data.append([row["stage"], desc, f"{row['recommended_sb_mb']}MB", work_mem, action, f"{row['protected_tp_tps']:.1f}"])
    add_table(slide, data, 0.50, 1.62, 12.33, 3.88, widths=[0.65, 1.75, 1.15, 4.30, 2.08, 1.60], font_size=8.6)
    add_callout(slide, "动作序列", "S1 保持 → S2 SB 8192→4096 给 AP → S3 将 AP work_mem 降至 256MB → S4 阻塞新 AP → S5 TP 突增，SB 4096→8192。", 0.78, 5.78, 11.80, 0.73, GREEN, PALE_GREEN)
    footer(slide, 14, "huawei6_observation_driven_five_stage_validation_20260802")
    outline.append("14. 五阶段：模型得到降低 SB、限制 AP 内存、阻塞新 AP和 S5 提高 SB 的动作。")


def formula_evidence(prs, outline):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(slide)
    add_title(slide, "公式验证：冻结后未见 QD 的 latency 与 TPS 都能预测", "预测文件在 I/O 注入之前写出；QD9/QD18 不参与冻结。", "14 公式证据")
    slide.shapes.add_picture(str(CHART_FORMULA), Inches(0.63), Inches(1.55), width=Inches(7.55))
    metrics = read_json(FORMULA_REPORT)["metrics"]
    add_kpi(slide, f"{metrics['pressure_latency_mape_pct']:.2f}%", "I/O latency MAPE", 8.55, 1.72, 3.65, color=BLUE, note="未见 QD9/QD18")
    add_kpi(slide, f"{metrics['post_intervention_tps_mape_pct']:.2f}%", "端到端 TPS MAPE", 8.55, 3.08, 3.65, color=TEAL, note="目标 ≤5%")
    add_kpi(slide, f"{metrics['tps_conversion_with_measured_latency_mape_pct']:.2f}%", "实测 latency→TPS MAPE", 8.55, 4.44, 3.65, color=GREEN, note="单独验证 TPS 换算")
    add_callout(slide, "方向性", "6/6 个实质压力点，TPS 升降方向全部正确。", 8.58, 5.72, 3.60, 0.82, GREEN, PALE_GREEN)
    footer(slide, 15, "FORMULA_VALIDATION_REPORT.md")
    outline.append("15. 公式证据：未见 QD latency MAPE 1.53%，TPS MAPE 1.81%。")


def stage_evidence(prs, outline):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(slide)
    add_title(slide, "最终验证：配置动作一致，TP 稳定性满足 5% 目标", "所有 AP 自然完成；S4/S5 只阻塞新 AP，不取消已运行 SQL。", "15 五阶段证据")
    slide.shapes.add_picture(str(CHART_STAGE), Inches(0.62), Inches(1.52), width=Inches(8.20))
    checks = read_json(STAGE_REPORT)["checks"]
    add_kpi(slide, f"{checks['protected_tp_variation_s1_s5_percent']:.2f}%", "S1-S5 保护 TPS 波动", 9.12, 1.72, 3.10, color=GREEN, note="验收目标 <5%")
    add_callout(slide, "动作检查", "✓ S2 降低 SB\n✓ S3 降低 AP work_mem\n✓ S4 阻塞新 AP\n✓ S5 提高 SB", 9.08, 3.12, 3.25, 1.55, TEAL, PALE_TEAL)
    add_callout(slide, "执行检查", "推荐配置全部应用；所有 AP 正常自然结束；无 AP failure。", 9.08, 4.93, 3.25, 0.88, BLUE, PALE_TEAL)
    add_callout(slide, "解释 S5", "保护 TP 4308→4416 TPS；新增 demand 约 298 TPS。提高 SB 用于承接突增。", 9.08, 5.99, 3.25, 0.82, ORANGE, PALE_ORANGE)
    footer(slide, 16, "validation_report.json")
    outline.append(f"16. 五阶段证据：动作全部匹配，S1-S5 保护 TPS 波动 {checks['protected_tp_variation_s1_s5_percent']:.2f}%。")


def operations(prs, outline):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(slide)
    add_title(slide, "执行结束后会得到什么，以及如何判断可用", "每一步都有可审计产物；没有通过 holdout 或缺少候选时，流程不会输出“可用模型”结论。", "16 交付与边界")
    data = [["产物", "作用", "验收点"],
            ["machine_inventory.json", "记录机器、内存、CPU、NVMe 与工具", "确认模型绑定当前机器"],
            ["storage_surface/frozen_surface.json", "设备 qAP→增量延迟曲面", "设备留出必须通过"],
            ["tp_path_anchors/anchors.json", "openGauss TP 路径锚点", "不使用压力 TPS 拟合"],
            ["model/frozen_model.json", "可移植机器模型包", "含 domain 与源文件哈希"],
            ["model_holdout/holdout_report.json", "未见 QD latency/TPS 验证", "latency≤10%，TPS≤5%"],
            ["predictions/recommendations.csv", "TP-first/AP-first 与最终联合推荐", "含 SB/work_mem/AP-cap"]]
    add_table(slide, data, 0.65, 1.53, 12.05, 3.90, widths=[3.05, 5.30, 3.70], font_size=9.0)
    add_callout(slide, "已完成的范围", "8KiB buffered TP + 128KiB random-read AP 压力下，新机器自动标定、公式冻结、盲测、候选预测与双向推荐已经代码化。", 0.84, 5.72, 5.70, 0.83, GREEN, PALE_GREEN)
    add_callout(slide, "仍需明确的边界", "新 plan 的 spill→qAP 仍依赖 operator/source replay 或额外锚点；写型、顺序型或不同块大小 I/O 必须新增设备曲面，不能静默外推。", 6.81, 5.72, 5.70, 0.83, RED, PALE_RED)
    footer(slide, 17, "docs/PORTABLE_MODEL_BOOTSTRAP.md")
    outline.append("17. 产物与边界：自动流程、验收门槛和当前支持域。")


def build() -> None:
    make_charts()
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    outline: list[str] = []
    for builder in (
        cover, problem, lifecycle, bootstrap, storage_calibration,
        path_calibration, workload_features, coupling, tp_first, ap_first,
        bounds_intersection, formula, merge, five_stage, formula_evidence,
        stage_evidence, operations,
    ):
        builder(prs, outline)
    prs.core_properties.title = "Huawei6 从零启动双向联合配置预测全流程"
    prs.core_properties.subject = "TP-first 1→2→3 与 AP-first 2→1→3 的机器标定、Trace Replay、联合推荐和验证"
    prs.core_properties.author = "Huawei6 project"
    prs.save(OUT)
    shutil.copy2(OUT, ROOT_COPY)
    OUTLINE.write_text("# Huawei6 从零到双向联合推荐：讲稿提纲\n\n" + "\n".join(f"- {line}" for line in outline) + "\n", encoding="utf-8")
    print(json.dumps({"pptx": str(OUT), "root_copy": str(ROOT_COPY), "slides": len(prs.slides), "outline": str(OUTLINE)}, ensure_ascii=False, indent=2))


if __name__ == "__main__":
    build()
