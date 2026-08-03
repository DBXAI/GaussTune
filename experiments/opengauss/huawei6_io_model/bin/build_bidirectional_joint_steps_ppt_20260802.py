#!/usr/bin/env python3
"""Build a focused Huawei6 deck explaining bidirectional joint prediction."""

from __future__ import annotations

import json
from pathlib import Path

import matplotlib.pyplot as plt
import numpy as np
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches

from build_memory_autonomy_paper_ppt import (
    BLUE,
    GRAY,
    GREEN,
    INK,
    LIGHT,
    MID,
    ORANGE,
    PALE_GREEN,
    PALE_ORANGE,
    PALE_RED,
    PALE_TEAL,
    RED,
    TEAL,
    WHITE,
    add_arrow,
    add_bullets,
    add_callout,
    add_flow_box,
    add_kpi,
    add_rect,
    add_table,
    add_text,
    add_title,
    configure_plot_font,
    rgb_hex,
    set_bg,
)


ROOT = Path(__file__).resolve().parents[1]
LATEST = ROOT / "artifacts" / "00_latest"
LATEST.mkdir(parents=True, exist_ok=True)
REPORT = ROOT / "results" / "huawei6_observation_driven_five_stage_validation_20260802" / "validation_report.json"
OUT = LATEST / "Huawei6_双向联合预测步骤与五阶段验证_20260802.pptx"
OUTLINE = LATEST / "Huawei6_双向联合预测步骤与五阶段验证_20260802_提纲.md"
CHART = LATEST / "Huawei6_双向联合预测_五阶段验证_20260802.png"


def add_footer(slide, page: int, source: str = "") -> None:
    add_rect(slide, 0.62, 7.17, 12.08, 0.012, MID)
    add_text(slide, "Huawei6 · 双向联合预测 · 2026-08-02", 0.62, 7.22, 5.8, 0.18, size=7.5, color=GRAY)
    if source:
        add_text(slide, source, 6.4, 7.22, 5.8, 0.18, size=7.0, color=GRAY, align=PP_ALIGN.RIGHT)
    add_text(slide, str(page), 12.28, 7.22, 0.42, 0.18, size=8, color=GRAY, align=PP_ALIGN.RIGHT)


def work_mem_text(values: dict[str, int]) -> str:
    if not values:
        return "无 AP"
    return "，".join(f"{query}={value}MB" for query, value in values.items())


def make_validation_chart(report: dict) -> None:
    configure_plot_font()
    stages = report["stages"]
    labels = [item["stage"] for item in stages]
    protected = [item["protected_tp_tps"] for item in stages]
    total = [item["protected_tp_tps"] + item["surge_tp_tps"] for item in stages]
    stable_indices = np.array([2, 3, 4])
    stable_mean = float(np.mean(np.array(protected)[stable_indices]))
    lower, upper = stable_mean * 0.95, stable_mean * 1.05

    fig, ax = plt.subplots(figsize=(11.4, 4.7))
    fig.patch.set_facecolor("white")
    x = np.arange(len(labels))
    ax.axvspan(1.5, 4.5, color=rgb_hex(PALE_GREEN), alpha=0.62, zorder=0)
    ax.fill_between([1.7, 4.3], lower, upper, color=rgb_hex(GREEN), alpha=0.16, label="S3-S5 protected TPS +/-5% target")
    ax.plot(x, protected, marker="o", markersize=8, linewidth=2.8, color=rgb_hex(BLUE), label="Measured protected TP TPS")
    ax.plot(x, total, marker="o", markersize=6, linewidth=2.0, linestyle="--", color=rgb_hex(ORANGE), label="Measured total TP TPS (S5 includes surge)")
    for idx, value in enumerate(protected):
        ax.annotate(f"{value:.0f}", (idx, value), xytext=(0, 10), textcoords="offset points", ha="center", fontsize=10, color=rgb_hex(BLUE), fontweight="bold")
    ax.annotate(f"{total[-1]:.0f}", (4, total[-1]), xytext=(0, 10), textcoords="offset points", ha="center", fontsize=10, color=rgb_hex(ORANGE), fontweight="bold")
    ax.set_xticks(x, labels)
    ax.set_ylabel("TPS")
    ax.set_title("Independent validation: S3-S5 protected TP TPS variation = 2.48%", fontsize=15, fontweight="bold", pad=15)
    ax.set_ylim(0, 5100)
    ax.grid(axis="y", alpha=0.18)
    ax.spines[["top", "right"]].set_visible(False)
    ax.legend(frameon=False, loc="upper left", fontsize=9)
    fig.tight_layout()
    fig.savefig(CHART, dpi=200, bbox_inches="tight")
    plt.close(fig)


def slide_cover(prs: Presentation, outline: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, INK)
    add_rect(slide, 0, 0, 0.22, 7.5, TEAL)
    add_text(slide, "Huawei6", 0.92, 1.05, 4.0, 0.40, size=19, color=TEAL, bold=True)
    add_text(slide, "双向联合预测", 0.90, 1.58, 11.4, 0.76, size=40, color=WHITE, bold=True)
    add_text(slide, "如何用 1→2→3 与 2→1→3 联合推荐\nshared_buffers、逐 Query work_mem 与 AP 准入", 0.93, 2.60, 10.8, 0.86, size=22, color=WHITE)
    add_rect(slide, 0.94, 4.10, 7.0, 0.60, TEAL, radius=True)
    add_text(slide, "Trace Replay + 机器观测 + I/O 队列延迟反哺 TPS", 1.10, 4.27, 6.68, 0.24, size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_callout(slide, "边界", "候选配置的选择不读取真实混合 TPS；真实 TPS 只用于推荐之后的独立验证。", 0.94, 5.34, 8.8, 0.88, TEAL, PALE_TEAL)
    add_text(slide, "汇报材料 · 2026-08-02", 0.95, 6.65, 4.2, 0.25, size=11, color=GRAY)
    outline.append("1. 封面：双向联合预测的输入、输出和验证边界。")


def slide_problem(prs: Presentation, outline: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "问题：SB 与 work_mem 不是两个独立参数", "必须把 AP spill I/O 对 TP 磁盘等待和 TPS 的反向影响纳入同一次推导。", "问题定义")
    boxes = [
        ("1. shared_buffers", "SB 变大：TP 缓存 miss 下降\n但可分给 AP 动态内存减少", 0.65, BLUE, PALE_TEAL),
        ("2. AP work_mem", "work_mem 变小：AP 动态峰值下降\n但 spill I/O 上升", 3.95, ORANGE, PALE_ORANGE),
        ("3. NVMe 排队", "AP spill 与 TP miss 同时争抢 I/O\n相同命中率也可能有不同 TPS", 7.25, RED, PALE_RED),
        ("4. TP TPS", "I/O await 增加会拉长每笔事务耗时\n最终决定 TP 能否稳定", 10.55, GREEN, PALE_GREEN),
    ]
    for title, sub, x, color, fill in boxes:
        add_flow_box(slide, title, sub, x, 2.0, 2.12, 1.55, color, fill=fill)
    for index in range(3):
        add_arrow(slide, 2.78 + index * 3.30, 2.77, 3.88 + index * 3.30, 2.77, TEAL, 2.1)
    add_arrow(slide, 11.58, 3.62, 11.58, 4.70, GREEN, 2.0)
    add_arrow(slide, 11.58, 5.16, 1.70, 5.16, GREEN, 2.0)
    add_arrow(slide, 1.70, 5.16, 1.70, 3.63, GREEN, 2.0)
    add_text(slide, "闭环：不能只比较缓存命中率；需要把 replay 的 AP spill 变成 I/O 压力，再将 I/O await 修正到 TP TPS。", 0.78, 5.73, 11.9, 0.38, size=16, color=INK, bold=True, align=PP_ALIGN.CENTER)
    add_callout(slide, "联合输出", "每个控制窗输出：(SB，逐 Query work_mem，是否阻塞新 AP)。", 3.06, 6.10, 7.2, 0.72, TEAL, PALE_TEAL)
    add_footer(slide, 2)
    outline.append("2. 问题：SB、AP 内存、spill I/O、I/O await 和 TP TPS 构成闭环。")


def slide_inputs(prs: Presentation, outline: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "输入与输出：机器观测负责当前状态，Trace Replay 负责候选后果", "控制器不输入阶段名、期望动作或候选配置的真实混合 TPS。", "模型输入")
    add_flow_box(slide, "实时机器观测", "SB、已用/峰值动态内存\n活跃/到达 AP Query ID、排队数\nTP terminals、offered rate、TP-only Ctp\nCPU 与 NVMe IOPS", 0.72, 1.82, 3.45, 2.35, TEAL, fill=PALE_TEAL)
    add_flow_box(slide, "历史 Trace Replay", "每个 Query × work_mem：\ndynamic_peak(q,w)\nspill_io_mb(q,w)\nplan_confidence(q,w)\n另有 TP miss_sb(B) 曲线", 4.94, 1.82, 3.45, 2.35, ORANGE, fill=PALE_ORANGE)
    add_flow_box(slide, "联合推荐", "shared_buffers B\n每条 AP 的 work_mem wq\n是否阻塞新 AP\n并给出导致该动作的信号", 9.16, 1.82, 3.45, 2.35, GREEN, fill=PALE_GREEN)
    add_arrow(slide, 4.23, 2.98, 4.87, 2.98, TEAL, 2.2)
    add_arrow(slide, 8.45, 2.98, 9.10, 2.98, TEAL, 2.2)
    add_callout(slide, "内存约束", "M(B,{wq}) = B + Σ dynamic_peak(q,wq)\n只有 M 不超过可用总内存的候选，才会进入 I/O / TPS 比较。", 0.98, 4.85, 5.25, 1.00, BLUE, PALE_TEAL)
    add_callout(slide, "验证边界", "推荐完成后才重启到对应 SB / work_mem 运行完整阶段。\n因此验证 TPS 不可能反过来参与候选选择。", 7.12, 4.85, 5.25, 1.00, RED, PALE_RED)
    add_text(slide, "“逐 Query work_mem”代表不同 AP 可得到不同的内存档位，而不是对所有 AP 使用同一个固定值。", 1.03, 6.28, 11.2, 0.30, size=14, color=INK, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, 3)
    outline.append("3. 输入：实时机器状态 + 历史 Trace；输出为 SB、逐 Query work_mem 和 AP 准入。")


def slide_tp_first(prs: Presentation, outline: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "路径 A：TP-first，按 1→2→3 推导", "先守住 TP 的缓存和容量边界，再在剩余内存中选择 AP grant。", "双向路径 A")
    add_flow_box(slide, "1. 找 TP 的 SB 下界", "对每个 B 回放 miss_sb(B)。\n找到 TP miss 的拐点 / 满足 TP 容量的最小 SB。", 0.68, 1.82, 3.58, 1.70, BLUE, fill=PALE_TEAL)
    add_flow_box(slide, "2. 在剩余内存选 AP", "仅保留内存安全候选。\n在这些候选中最大化 AP utility，\n并保留 Query 对应的 trace grant。", 4.88, 1.82, 3.58, 1.70, ORANGE, fill=PALE_ORANGE)
    add_flow_box(slide, "3. I/O→TPS 校正", "把 spill 转成 AP IOPS，\n与 TP miss IOPS 一起进入队列等待，\n预测保护 TP TPS。", 9.08, 1.82, 3.58, 1.70, GREEN, fill=PALE_GREEN)
    add_arrow(slide, 4.30, 2.67, 4.80, 2.67, TEAL, 2.2)
    add_arrow(slide, 8.50, 2.67, 9.00, 2.67, TEAL, 2.2)
    add_callout(slide, "步骤 1 的含义", "不是盲目取最大 SB，而是得到“TP 已经足够”的 SB 范围；剩余内存才可被 AP 使用。", 0.88, 4.23, 3.40, 1.00, BLUE, PALE_TEAL)
    add_callout(slide, "步骤 2 的含义", "AP utility 同时惩罚较慢服务时间和较大的 replay spill，避免只看动态内存而漏掉 I/O 代价。", 4.96, 4.23, 3.40, 1.00, ORANGE, PALE_ORANGE)
    add_callout(slide, "适用状态", "TP 已接近或达到容量边界时，TP-first 优先保护 TP，再压缩 AP grant 或阻塞新到达 AP。", 9.04, 4.23, 3.40, 1.00, GREEN, PALE_GREEN)
    add_text(slide, "TP-first 的选型顺序：TP miss 曲线 → 内存安全 AP grant → I/O 排队与 TPS。", 1.30, 6.20, 10.7, 0.36, size=17, color=TEAL, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, 4)
    outline.append("4. TP-first：1 找 TP SB 下界，2 选内存安全 AP grant，3 用 I/O/TPS 校正。")


def slide_ap_first(prs: Presentation, outline: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "路径 B：AP-first，按 2→1→3 推导", "当低 TP 下出现 AP 到达，先保留 AP Trace 中验证过的合适 grant，再寻找能容纳它的 SB。", "双向路径 B")
    add_flow_box(slide, "2. 先选 AP trace grant", "对当前/到达 Query 选历史 utility 最高、\nplan 置信度足够的 work_mem。\n目标是减少不必要 spill。", 0.68, 1.82, 3.58, 1.70, ORANGE, fill=PALE_ORANGE)
    add_flow_box(slide, "1. 再找可容纳它的 SB", "在保持 TP miss 下界的前提下，\n选能够容纳 AP dynamic_peak 的\n最强 SB。", 4.88, 1.82, 3.58, 1.70, BLUE, fill=PALE_TEAL)
    add_flow_box(slide, "3. 同一 I/O→TPS 校正", "对完整候选计算 spill IOPS、\nNVMe await 和保护 TPS。\n不安全候选被淘汰。", 9.08, 1.82, 3.58, 1.70, GREEN, fill=PALE_GREEN)
    add_arrow(slide, 4.30, 2.67, 4.80, 2.67, TEAL, 2.2)
    add_arrow(slide, 8.50, 2.67, 9.00, 2.67, TEAL, 2.2)
    add_callout(slide, "S2 的因果规则", "若当前富裕 SB 装不下 AP-first grant，而减小一档 SB 能装下，则降低 SB，而不是静默把 AP work_mem 降低。", 0.88, 4.22, 5.35, 1.12, ORANGE, PALE_ORANGE)
    add_callout(slide, "为什么需要这条路径", "仅从 TP 出发会把所有内存留给 SB，遗漏“适度让出 SB 能避免 AP spill，最终也降低 TP I/O 等待”的机会。", 7.10, 4.22, 5.35, 1.12, BLUE, PALE_TEAL)
    add_text(slide, "AP-first 的选型顺序：Trace grant / plan → 能容纳 grant 的 SB → I/O 排队与 TPS。", 1.22, 6.22, 10.9, 0.34, size=17, color=TEAL, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, 5)
    outline.append("5. AP-first：2 保留 Trace grant，1 找容纳它的 SB，3 用同一 I/O/TPS 校正。")


def slide_queue_derivation(prs: Presentation, outline: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "I/O→TPS 子模型：不是黑盒 TPS 回归，而是可展开的四步计算", "候选配置只提供 B 和 {wq}；每个候选从 Trace 中推导 TP/AP I/O，再用排队近似求增量时延。", "子模型推导")
    add_flow_box(slide, "A. TP miss 频率", "TP cache replay 给 miss_sb(B)。\nTP-only 标定得到 23.129682\nlogical pages / transaction。\nmp(B)=23.129682×miss_sb(B)", 0.62, 1.70, 2.86, 1.75, BLUE, fill=PALE_TEAL)
    add_flow_box(slide, "B. AP spill 频率", "Trace 给 spill_io_mb(q,w)。\n除以该 Query 的 anchor service time，\n得到 spill bytes/s；再除以 128KiB。", 3.72, 1.70, 2.86, 1.75, ORANGE, fill=PALE_ORANGE)
    add_flow_box(slide, "C. NVMe 排队", "λ = TPS×mp(B) + AP_IOPS\nρ = min(0.985, λ×S / Q)\nawait = S / (1-ρ)", 6.82, 1.70, 2.86, 1.75, RED, fill=PALE_RED)
    add_flow_box(slide, "D. TPS 不动点", "tx_ms = base_tx_ms + α×mp(B)\n× max(0, await - await_no_AP)\nTPS = min(offered, terminals×1000/tx_ms)\n迭代至 TPS 收敛。", 9.92, 1.70, 2.78, 1.75, GREEN, fill=PALE_GREEN)
    for x in (3.52, 6.62, 9.72):
        add_arrow(slide, x, 2.57, x + 0.16, 2.57, TEAL, 2.0)
    add_callout(slide, "实际的 AP IOPS", "AP_IOPS = min(spill bytes/s ÷ 8KiB，spill bytes/s ÷ 128KiB)。128KiB 是独立 I/O 实验得到的写回物化粒度；单 Query 直接测到的物理 IOPS 只做诊断，不作为“零 I/O”覆盖。", 0.92, 4.35, 11.50, 0.93, ORANGE, PALE_ORANGE)
    add_callout(slide, "为什么计算 await_no_AP", "TPS 修正只采用 AP 额外引入的等待 await-await_no_AP，避免把 TP 自己已有的 I/O 延迟重复惩罚。", 0.92, 5.56, 11.50, 0.68, TEAL, PALE_TEAL)
    add_footer(slide, 6, "实现：huawei6_bidirectional_joint_predictor.py:177-249")
    outline.append("6. I/O→TPS 子模型：TP miss 与 AP spill 推导 IOPS，用队列公式和不动点迭代得到 TPS。")


def slide_queue_calibration(prs: Presentation, outline: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "子模型参数如何得到：独立训练、按 profile 留出验证", "这一步会使用历史训练 profile 的 TPS/await 来标定常数；但当前五阶段候选的实际 TPS 从不参与推荐。", "训练与验证")
    add_flow_box(slide, "BPF 争用矩阵训练", "基线：baseline_sb4096\n训练：highmem cap8、lowmem cap8\n采集：块层 TP/AP/other IOPS、await、TPS\n拟合 S、Q、α", 0.72, 1.72, 3.72, 2.10, BLUE, fill=PALE_TEAL)
    add_flow_box(slide, "最终在线常数", "S = 0.572646 ms\nQ = 12 effective queues\nα = 10.861136\n来自 BPF 训练集，而非五阶段候选。", 4.80, 1.72, 3.72, 2.10, TEAL, fill=PALE_GREEN)
    add_flow_box(slide, "独立留出 profile", "holdout：lowmem cap4\n8 个稳定窗口\nawait MAE = 0.0147 ms\nTPS MAPE = 1.67%", 8.88, 1.72, 3.72, 2.10, GREEN, fill=PALE_GREEN)
    add_arrow(slide, 4.48, 2.77, 4.70, 2.77, TEAL, 2.0)
    add_arrow(slide, 8.56, 2.77, 8.78, 2.77, TEAL, 2.0)
    add_callout(slide, "AP spill 物化比例的独立实验", "I/O latency matrix：无 AP 基线 + 3 个 AP 训练 profile，留出 lowmem cap4。由 temp_bytes / device write I/O 得到 128KiB / I/O；13 个留出窗口的 TPS MAPE 为 2.03%。", 0.94, 4.35, 7.32, 1.05, ORANGE, PALE_ORANGE)
    add_callout(slide, "严格边界", "训练/留出 profile 只标定机器 I/O 规律。当前五阶段仅输入 B、TP demand、AP Query ID/arrival 与历史 trace。实际阶段 TPS 在推荐写出后才比较。", 8.62, 4.35, 3.78, 1.05, RED, PALE_RED)
    add_text(slide, "所以它不是“看到 S1-S5 结果后拟合一个 TPS 曲线”，而是用独立争用实验确定 I/O→等待→TPS 的可解释映射。", 0.90, 6.18, 11.6, 0.34, size=15, color=INK, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, 7, "来源：bpf_contention_matrix_20260731 与 io_latency_matrix_20260731")
    outline.append("7. 参数训练：BPF 争用矩阵与 I/O 物化矩阵独立留出验证，非五阶段候选拟合。")


def slide_io_tps(prs: Presentation, outline: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "共同的第 3 步：将 spill I/O 显式反哺到 TP TPS", "候选配置相同的缓存命中率，可能因 AP spill 造成不同 NVMe 等待，从而有不同 TPS。", "I/O 与 TPS")
    add_flow_box(slide, "Replay 输出", "spill_io_mb(q,w)\nTP miss_sb(B)\n每个候选均可得到", 0.70, 1.82, 2.25, 1.38, ORANGE, fill=PALE_ORANGE)
    add_flow_box(slide, "转换为 IOPS", "AP_IOPS = spill bytes/s ÷ 131072\nTP_IOPS 由 TP miss / tx 得到", 3.55, 1.82, 2.55, 1.38, BLUE, fill=PALE_TEAL)
    add_flow_box(slide, "排队等待", "await = service_ms / (1 - utilization)\nutilization 由 TP_IOPS + AP_IOPS 得到", 6.70, 1.82, 2.55, 1.38, RED, fill=PALE_RED)
    add_flow_box(slide, "保护 TPS", "TPS = terminals / 每笔事务时间\n取 min(offered, capacity)", 9.85, 1.82, 2.55, 1.38, GREEN, fill=PALE_GREEN)
    for x in (3.05, 6.20, 9.35):
        add_arrow(slide, x, 2.51, x + 0.42, 2.51, TEAL, 2.1)
    add_callout(slide, "实际使用的公式", "await = service_ms / [1 - min(0.985, (TP_IOPS + AP_IOPS) × service_ms / queues / 1000)]", 1.23, 4.06, 10.88, 0.70, RED, PALE_RED)
    add_callout(slide, "TPS 修正", "TPS = min(offered, terminals×1000 / [terminals×1000/Ctp + weight×tp_miss_per_tx×(await - await_no_ap)])", 1.23, 5.03, 10.88, 0.70, GREEN, PALE_GREEN)
    add_text(slide, "参数来源：spill / dynamic_peak / plan 来自历史 trace；TP miss 来自缓存回放；实际候选混合 TPS 不进入公式。", 0.82, 6.24, 11.7, 0.30, size=13.5, color=GRAY, align=PP_ALIGN.CENTER)
    add_footer(slide, 8)
    outline.append("8. 第 3 步摘要：以 Trace spill 与 TP miss 估计 IOPS、await，再修正保护 TPS。")


def slide_selection(prs: Presentation, outline: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "两条路径如何合并：以可行性和保护 TPS 选择最终动作", "不是分别给一个 SB 和一个 work_mem；每条路径都产生完整的 (SB, {work_mem}, AP-cap) 候选。", "联合决策")
    add_flow_box(slide, "TP-first 1→2→3", "优先 TP SB 下界\n得到完整候选 A", 0.70, 1.87, 2.55, 1.35, BLUE, fill=PALE_TEAL)
    add_flow_box(slide, "AP-first 2→1→3", "优先 AP trace grant\n得到完整候选 B", 0.70, 4.08, 2.55, 1.35, ORANGE, fill=PALE_ORANGE)
    add_arrow(slide, 3.33, 2.54, 4.23, 3.09, TEAL, 2.0)
    add_arrow(slide, 3.33, 4.75, 4.23, 3.62, TEAL, 2.0)
    add_flow_box(slide, "共同过滤", "1. 内存 M(B,{wq}) 安全\n2. TP 容量 / miss 下界\n3. I/O await 不越界\n4. AP arrival 与队列规则", 4.35, 2.36, 3.15, 1.95, RED, fill=PALE_RED)
    add_arrow(slide, 7.60, 3.34, 8.45, 3.34, TEAL, 2.0)
    add_flow_box(slide, "最终推荐", "最高保护 TPS 的可行候选\n并输出原因：\n降低/提升 SB、降低 AP work_mem、或阻塞新 AP", 8.57, 2.36, 3.95, 1.95, GREEN, fill=PALE_GREEN)
    add_callout(slide, "状态不是标签输入", "低 TP + AP 到达触发 AP-first 机会；TP 饱和触发 TP-first 保护；新 AP 到达且驻留集已保护时进入阻塞；TP surge 时要求提高 SB。", 1.05, 5.55, 11.25, 0.78, TEAL, PALE_TEAL)
    add_footer(slide, 9)
    outline.append("9. 合并：两条路径都给完整配置，经过同一组约束后，以保护 TPS 选择最终动作。")


def slide_actions(prs: Presentation, report: dict, outline: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "最终五阶段：由观测和 replay 推出与验收要求一致的动作", "以下推荐在不输入阶段标签、期望动作或真实混合 TPS 的条件下冻结；随后才执行真实阶段验证。", "五阶段推荐")
    data = [["阶段", "模型推断动作", "推荐 SB", "逐 Query work_mem", "新 AP"]]
    actions = {
        "keep_rich_memory": "保留富裕内存",
        "yield_sb_to_ap": "降低 SB 给 AP",
        "reduce_ap_work_mem": "降低 AP work_mem",
        "block_new_ap": "阻塞新 AP",
        "raise_sb_for_tp_surge": "TP 突增，提高 SB",
    }
    for item in report["stages"]:
        data.append([
            item["stage"],
            actions[item["inferred_action"]],
            f"{item['recommended_sb_mb']}MB",
            work_mem_text(item["recommended_work_mem"]),
            "阻塞（排队 7）" if item["recommend_block_new_ap"] else "允许",
        ])
    add_table(slide, data, 0.58, 1.72, 12.15, 3.55, widths=[0.75, 2.45, 1.30, 5.10, 2.55], font_size=11.0)
    add_callout(slide, "关键转换", "S1→S2：8192MB → 4096MB，保留 Q18/Q21=1150MB 的 Trace grant；S5：TP surge 后 4096MB → 8192MB，同时仅保留内存安全 AP grant。", 0.82, 5.80, 11.72, 0.70, TEAL, PALE_TEAL)
    add_footer(slide, 10, "来源：最终冻结推荐 + 独立真实验证报告")
    outline.append("10. 五阶段推荐：S2 降低 SB 给 AP，S3 降低 AP work_mem，S4 阻塞新 AP，S5 提升 SB。")


def slide_validation(prs: Presentation, report: dict, outline: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "真实验证：五个推荐均被实际应用，S3-S5 保护 TPS 波动 2.48%", "阶段边界均采用 restart；AP 查询自然结束，不取消正在执行的 AP。", "独立验证")
    slide.shapes.add_picture(str(CHART), Inches(0.72), Inches(1.55), width=Inches(8.25), height=Inches(4.62))
    checks = report["checks"]
    add_kpi(slide, "5 / 5", "PPT 要求动作匹配", 9.18, 1.80, 3.05, color=GREEN, note="S2 降 SB，S3 降 AP 内存，S4 阻塞，S5 提升 SB")
    add_kpi(slide, f"{checks['protected_tp_variation_s3_s5_percent']:.2f}%", "S3-S5 保护 TPS 波动", 9.18, 3.30, 3.05, color=TEAL, note="满足不超过 5% 的目标")
    add_kpi(slide, "0", "AP 失败 / 强制取消", 9.18, 4.80, 3.05, color=BLUE, note="所有已运行 AP 自然完成")
    add_text(slide, "结论：双向步骤能将 TP 保护、AP 内存、spill I/O 和 AP 准入放入同一决策；最终动作序列与目标五阶段一致。", 0.92, 6.40, 11.55, 0.36, size=15, color=INK, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, 11, "来源：validation_report.json（post-decision real validation）")
    outline.append("11. 独立验证：5/5 动作匹配，S3-S5 保护 TPS 波动 2.48%，AP 均自然结束。")


def build() -> None:
    report = json.loads(REPORT.read_text(encoding="utf-8"))
    make_validation_chart(report)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    outline: list[str] = []
    slide_cover(prs, outline)
    slide_problem(prs, outline)
    slide_inputs(prs, outline)
    slide_tp_first(prs, outline)
    slide_ap_first(prs, outline)
    slide_queue_derivation(prs, outline)
    slide_queue_calibration(prs, outline)
    slide_io_tps(prs, outline)
    slide_selection(prs, outline)
    slide_actions(prs, report, outline)
    slide_validation(prs, report, outline)
    prs.save(OUT)
    OUTLINE.write_text("# Huawei6 双向联合预测步骤与五阶段验证\n\n" + "\n".join(outline) + "\n", encoding="utf-8")
    print(OUT)
    print(CHART)


if __name__ == "__main__":
    build()
