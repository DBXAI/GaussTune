#!/usr/bin/env python3
"""Build the Huawei5 joint SB/work_mem replay presentation."""

from __future__ import annotations

import csv
import json
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches

from build_model_progress_ppt import (
    BLUE,
    FONT,
    GRAY,
    GREEN,
    INK,
    LIGHT,
    MID,
    ORANGE,
    RED,
    TEAL,
    WHITE,
    add_arrow,
    add_bullets,
    add_flow_box,
    add_kpi,
    add_rect,
    add_table,
    add_text,
    add_title,
    set_bg,
)


ROOT = Path(__file__).resolve().parents[1]
RESULTS = ROOT / "results/joint_bidirectional_replay_20260722"
REPLAY = RESULTS / "replay"
VALIDATION = RESULTS / "validation/validation_summary.json"
SOURCE_FIGURE = REPLAY / "joint_bidirectional_five_stage.png"
RECOMMENDATIONS = REPLAY / "stage_joint_recommendations.csv"
ARTIFACTS = ROOT / "artifacts/01_current_joint_model"
FIGURES = ARTIFACTS / "figures"
OUT = ARTIFACTS / "Huawei5_SB_workmem_joint_bidirectional_replay_20260722.pptx"
OUTLINE = ARTIFACTS / "Huawei5_SB_workmem_joint_bidirectional_replay_20260722_outline.md"
TOP_FIGURE = FIGURES / "joint_bidirectional_effect_s1_s2_20260722.png"
MIDDLE_FIGURE = FIGURES / "joint_bidirectional_effect_s3_s4_20260722.png"
BOTTOM_FIGURE = FIGURES / "joint_bidirectional_effect_s5_20260722.png"


def read_csv(path: Path) -> list[dict[str, str]]:
    with path.open(newline="", encoding="utf-8") as fh:
        return list(csv.DictReader(fh))


def add_footer(slide, page: int) -> None:
    add_rect(slide, 0.62, 7.17, 12.08, 0.012, MID)
    add_text(
        slide,
        "Huawei5 | SB × work_mem bidirectional trace replay | 2026-07-22",
        0.62,
        7.22,
        7.5,
        0.18,
        size=8,
        color=GRAY,
    )
    add_text(slide, str(page), 12.2, 7.22, 0.5, 0.18, size=8, color=GRAY, align=PP_ALIGN.RIGHT)


def crop_effect_figure() -> None:
    image = Image.open(SOURCE_FIGURE)
    width, height = image.size
    first_split = int(height * 0.40)
    middle_end = int(height * 0.793)
    bottom_start = int(height * 0.80)
    image.crop((0, 0, width, first_split)).save(TOP_FIGURE)
    image.crop((0, first_split, width, middle_end)).save(MIDDLE_FIGURE)
    image.crop((0, bottom_start, width, height)).save(BOTTOM_FIGURE)


def add_note(slide, text: str, x: float, y: float, w: float, color=TEAL) -> None:
    add_rect(slide, x, y, 0.08, 0.52, color)
    add_text(slide, text, x + 0.18, y - 0.01, w - 0.18, 0.55, size=12, color=GRAY)


def build() -> None:
    FIGURES.mkdir(parents=True, exist_ok=True)
    crop_effect_figure()
    recommendations = read_csv(RECOMMENDATIONS)
    validation = json.loads(VALIDATION.read_text(encoding="utf-8"))

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    outline: list[str] = []

    # 1. Cover
    slide = prs.slides.add_slide(blank)
    set_bg(slide, INK)
    add_rect(slide, 0, 0, 0.18, 7.5, TEAL)
    add_text(slide, "Huawei5 SB × work_mem", 0.92, 1.25, 10.8, 0.65, size=37, color=WHITE, bold=True)
    add_text(slide, "双向联合 Trace Replay 预测模型", 0.94, 2.02, 11.1, 0.75, size=34, color=WHITE, bold=True)
    add_text(
        slide,
        "模型机制 · 实验设计 · 五阶段预测效果 · 实际验证",
        0.96,
        3.08,
        9.8,
        0.42,
        size=20,
        color=RGBColor(179, 218, 221),
    )
    add_rect(slide, 0.96, 4.02, 3.45, 0.48, TEAL, radius=True)
    add_text(slide, "不使用实际 TPS 标签训练", 1.06, 4.13, 3.25, 0.24, size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "openGauss · TPC-H SF85 · TPC-C 250 warehouses", 0.96, 5.82, 8.2, 0.3, size=13, color=RGBColor(194, 201, 206))
    add_text(slide, "2026-07-22", 0.96, 6.26, 2.0, 0.25, size=12, color=RGBColor(194, 201, 206))
    outline.append("1. 封面：Huawei5 SB × work_mem 双向联合 Trace Replay。")

    # 2. What the model is
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, "模型是什么：从两个独立参数变成一个闭环", "目标是预测配置导致的执行行为，而不是用实测最优点拟合答案。", "01 / Model")
    add_flow_box(slide, "shared_buffers", 0.72, 1.82, 2.45, 1.15, TEAL, sub="决定 TP/AP 数据页进入 SB 还是 OS")
    add_flow_box(slide, "work_mem", 0.72, 3.35, 2.45, 1.15, ORANGE, sub="决定 Join/Agg/Sort grant 与 spill")
    add_arrow(slide, 3.25, 2.38, 4.15, 2.38, TEAL)
    add_arrow(slide, 3.25, 3.92, 4.15, 3.92, ORANGE)
    add_rect(slide, 4.18, 1.78, 3.15, 2.8, LIGHT, line=INK, radius=True)
    add_text(slide, "同一个系统状态", 4.38, 2.08, 2.75, 0.38, size=20, bold=True, align=PP_ALIGN.CENTER)
    add_bullets(slide, ["动态内存峰值", "spill 临时页与 I/O", "Linux active/inactive/refault", "TP-only disk miss"], 4.55, 2.63, 2.45, 1.65, size=13, spacing=7)
    add_arrow(slide, 7.42, 3.15, 8.30, 3.15, GREEN)
    add_flow_box(slide, "联合推荐", 8.35, 2.40, 2.25, 1.48, GREEN, sub="在命中率、物理 I/O 与内存占用之间选择 Pareto 点")
    add_rect(slide, 10.92, 1.90, 1.68, 2.48, INK, radius=True)
    add_text(slide, "输出", 11.08, 2.14, 1.36, 0.3, size=15, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "SB hit\nOS hit\ncombined\nrefault\nspill I/O\n内存余量", 11.08, 2.62, 1.36, 1.55, size=12, color=WHITE, align=PP_ALIGN.CENTER)
    add_note(slide, "关键区别：work_mem 会改变 OS cache，SB 也会挤占动态内存余量，二者不再独立。", 0.78, 5.55, 11.9)
    add_footer(slide, 2)
    outline.append("2. 模型定义：work_mem 与 SB 通过动态内存、spill 和 OS cache 双向耦合。")

    # 3. Experiment inputs and isolation
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, "实验输入：预测数据与验证数据严格隔离", "预测器不会读取实际 TPS 或实际最优配置。", "02 / Experiment")
    add_rect(slide, 0.7, 1.72, 5.75, 4.65, LIGHT, line=TEAL, radius=True)
    add_text(slide, "预测输入", 0.98, 1.98, 2.2, 0.4, size=22, color=TEAL, bold=True)
    add_bullets(
        slide,
        [
            "597MB mixed TP/AP 二进制访问轨迹（sample=64）",
            "Q1/Q3/Q5/Q7/Q9/Q13/Q18/Q21 算子内存 trace",
            "Hash Join / HashAggregate / Sort 生命周期与分配事件",
            "当前数据库在多个 work_mem 点的 EXPLAIN 计划族",
            "8 个 SB 候选与系统内存消耗系数",
        ],
        0.98,
        2.58,
        5.0,
        3.05,
        size=15,
        spacing=11,
    )
    add_rect(slide, 6.85, 1.72, 5.75, 4.65, LIGHT, line=ORANGE, radius=True)
    add_text(slide, "仅用于最后验证", 7.13, 1.98, 3.2, 0.4, size=22, color=ORANGE, bold=True)
    add_bullets(
        slide,
        [
            "不同 SB 下真实 TP TPS 与 P95",
            "逐查询 no-spill 边界实测",
            "AP8 二维短跑与长期内存压力实验",
            "实际命中率和实际最优点",
        ],
        7.13,
        2.58,
        4.85,
        2.5,
        size=15,
        spacing=13,
    )
    add_rect(slide, 7.25, 5.30, 4.85, 0.58, WHITE, line=RED, radius=True)
    add_text(slide, "这些文件不由 joint predictor 读取", 7.43, 5.45, 4.5, 0.26, size=13, color=RED, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, 3)
    outline.append("3. 实验输入：trace/EXPLAIN 用于预测，TPS 与实际边界只用于 held-out 验证。")

    # 4. Full workflow
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, "完整预测流程：每个二维点都重新 replay", "候选点不是线性插值，而是经过执行路径、算子和缓存三层模型。", "03 / Pipeline")
    boxes = [
        ("1 计划族", "扫描 EXPLAIN\n匹配同路径 trace", BLUE),
        ("2 算子 replay", "grant / batch / merge\n动态内存时间线", ORANGE),
        ("3 spill 回灌", "临时页 + 读写 I/O\nstreaming pollution", RED),
        ("4 SB replay", "mixed TP/AP 访问\nBulkRead ring", TEAL),
        ("5 Linux cache", "active/inactive\nrefault 保护", GREEN),
    ]
    x_positions = [0.55, 3.08, 5.61, 8.14, 10.67]
    for (title, sub, color), x in zip(boxes, x_positions):
        add_flow_box(slide, title, x, 2.05, 2.08, 1.55, color, sub=sub)
    for x in [2.67, 5.20, 7.73, 10.26]:
        add_arrow(slide, x, 2.82, x + 0.34, 2.82, GRAY)
    add_arrow(slide, 11.70, 3.72, 11.70, 4.50, GREEN)
    add_arrow(slide, 11.70, 4.50, 2.10, 4.50, GREEN)
    add_arrow(slide, 2.10, 4.50, 2.10, 4.04, GREEN)
    add_text(slide, "OS miss/refault 结果反馈到联合物理 I/O 与内存约束", 4.25, 4.58, 5.7, 0.35, size=13, color=GREEN, bold=True, align=PP_ALIGN.CENTER)
    add_rect(slide, 2.05, 5.25, 9.2, 0.82, LIGHT, line=INK, radius=True)
    add_text(slide, "输出：TP-SB hit · TP-OS conditional hit · combined · refault · spill I/O · MemAvailable", 2.30, 5.49, 8.7, 0.3, size=15, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, 4)
    outline.append("4. 流程：计划族、算子 replay、spill 回灌、SB replay、Linux cache、联合指标。")

    # 5. Operator replay
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, "动态内存预测：从算子执行路径得到 work_mem 需求", "同一查询计划切换后必须使用新的 trace 锚点，不能跨路径偷算。", "04 / Operator Replay")
    operator_data = [
        ["算子", "replay 状态", "低于需求时的结果"],
        ["Hash Join", "tuple bytes / bucket / batch", "batch 增加，临时表读写"],
        ["HashAggregate", "group 数 / 每组分配", "未容纳 group spill"],
        ["Sort", "tuple chunk / memtuple / pass", "external merge 多轮 I/O"],
    ]
    add_table(slide, 4, 3, operator_data, 0.72, 1.80, 7.55, 2.6, widths=[1.65, 2.95, 2.95], font_size=13)
    add_rect(slide, 8.62, 1.80, 3.95, 2.6, LIGHT, line=ORANGE, radius=True)
    add_text(slide, "并发峰值", 8.92, 2.08, 3.35, 0.38, size=20, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "按 operator start/end 时间线\n计算查询峰值，再按阶段并发叠加", 9.00, 2.70, 3.18, 0.90, size=15, align=PP_ALIGN.CENTER)
    add_text(slide, "grant = min(work_mem, operator requirement)", 8.92, 3.78, 3.35, 0.25, size=11, color=GRAY, align=PP_ALIGN.CENTER)
    add_kpi(slide, "3 / 8", "同计划逐 MB 精确边界", 0.95, 4.92, 2.55, color=TEAL, note="Q3 / Q9 / Q13")
    add_kpi(slide, "5 / 8", "operational boundary 通过", 3.72, 4.92, 2.55, color=GREEN, note="验证文件不参与预测")
    add_kpi(slide, "Q18", "主机动态内存上限", 6.49, 4.92, 2.55, color=ORANGE, note="无法实现全 no-spill")
    add_kpi(slide, "Q21", "MaxAllocSize 上限", 9.26, 4.92, 2.55, color=RED, note="单次 4GB bucket 分配不可行")
    add_footer(slide, 5)
    outline.append("5. 算子模型：Join/Agg/Sort replay、生命周期峰值和计划路径约束。")

    # 6. Cache coupling
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, "缓存联动：spill 不只增加 I/O，也会改变 TP 的缓存环境", "模型显式模拟 Linux active/inactive/refault 对频繁 TP 页的保护。", "05 / Cache Replay")
    add_flow_box(slide, "SB miss", 0.75, 2.02, 2.05, 1.18, TEAL, sub="mixed trace 中 TP/AP 都改变状态")
    add_arrow(slide, 2.88, 2.62, 3.65, 2.62, TEAL)
    add_flow_box(slide, "OS inactive", 3.72, 1.62, 2.28, 1.22, BLUE, sub="普通页二次命中后可晋升")
    add_flow_box(slide, "Streaming", 3.72, 3.18, 2.28, 1.22, ORANGE, sub="AP bulk read + spill 临时页")
    add_arrow(slide, 6.10, 2.25, 6.90, 2.25, GREEN)
    add_arrow(slide, 6.10, 3.78, 6.90, 3.08, RED)
    add_flow_box(slide, "Active / refault", 6.98, 1.84, 2.40, 1.30, GREEN, sub="短距离 refault 保护 TP 热页")
    add_flow_box(slide, "优先回收", 6.98, 3.42, 2.40, 1.18, RED, sub="streaming → normal → active")
    add_arrow(slide, 9.48, 2.82, 10.18, 2.82, GRAY)
    add_flow_box(slide, "TP-only 计分", 10.25, 2.02, 2.18, 1.80, INK, sub="SB hit / OS hit\nrefault / disk miss")
    add_rect(slide, 0.78, 5.30, 11.72, 0.70, LIGHT, line=MID, radius=True)
    add_text(slide, "MemAvailable = 23546.38 - 0.29220 × SB - 0.41804 × dynamic_peak  (MB)", 1.05, 5.51, 11.18, 0.28, size=17, color=INK, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, 6)
    outline.append("6. 缓存模型：spill streaming 页、Linux refault 保护、TP-only 命中率与内存公式。")

    # 7. Candidate grid and selection
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, "二维实验与推荐：192 个点中只使用有证据的区域", "没有同计划族 trace 的点保留预测值，但不允许进入最终推荐。", "06 / Selection")
    add_kpi(slide, "192", "二维候选点", 0.72, 1.72, 2.45, color=TEAL, note="5 stages × SB × work_mem")
    add_kpi(slide, "144", "同计划族支持点", 3.27, 1.72, 2.45, color=GREEN, note="可进入推荐")
    add_kpi(slide, "8", "SB 测试档位", 5.82, 1.72, 2.45, color=BLUE, note="128MB 至 8192MB")
    add_kpi(slide, "3.2GB", "系统保留量", 8.37, 1.72, 2.45, color=ORANGE, note="低于该值判为不安全")
    steps = [
        ("1", "计划支持 + 内存安全", BLUE),
        ("2", "TP-SB hit 达到最大值的 99%", TEAL),
        ("3", "物理 I/O 与内存占用 Pareto 前沿", GREEN),
        ("4", "近最优 I/O 中选择最小内存", ORANGE),
    ]
    for index, (number, label, color) in enumerate(steps):
        y = 3.42 + index * 0.67
        add_rect(slide, 1.20, y, 0.48, 0.48, color, radius=True)
        add_text(slide, number, 1.21, y + 0.10, 0.46, 0.22, size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, label, 1.92, y + 0.04, 5.15, 0.36, size=16, bold=True)
    add_rect(slide, 7.65, 3.52, 4.25, 2.35, LIGHT, line=RED, radius=True)
    add_text(slide, "不会做的事", 7.98, 3.82, 3.60, 0.35, size=20, color=RED, bold=True, align=PP_ALIGN.CENTER)
    add_bullets(slide, ["读取 TPS 最优点再校准", "跨计划族沿用算子 trace", "把 combined 的平坦最大值当 TPS 最优"], 8.04, 4.35, 3.40, 1.25, size=13, spacing=8)
    add_footer(slide, 7)
    outline.append("7. 选择规则：计划支持、内存安全、TP-SB knee、I/O/内存 Pareto。")

    # 8. Effect plot S1-S2
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, "预测效果图（S1-S2）", "左：TP-only SB hit；中：work_mem 对 spill I/O；右：SB × work_mem 联合物理 I/O。", "07 / Result")
    slide.shapes.add_picture(str(TOP_FIGURE), Inches(1.02), Inches(1.55), width=Inches(11.30), height=Inches(5.42))
    add_footer(slide, 8)
    outline.append("8. 效果图：S1-S2。星号为推荐点，灰色/空白为缺少计划锚点。")

    # 9. Effect plot S3-S4
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, "预测效果图（S3-S4）", "S4 的更高 work_mem 点缺 Q21 新计划族 trace，因此 512MB 只标为覆盖区内暂定。", "08 / Result")
    slide.shapes.add_picture(str(MIDDLE_FIGURE), Inches(1.02), Inches(1.58), width=Inches(11.30), height=Inches(5.33))
    add_footer(slide, 9)
    outline.append("9. 效果图：S3-S4；S4 仍位于计划 trace 覆盖边界。")

    # 10. Effect plot S5
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, "预测效果图（S5）", "S5 的 TP-SB 曲线在 1024MB 进入平台；该点与独立实测 TPS 平台起点一致。", "09 / Result")
    slide.shapes.add_picture(str(BOTTOM_FIGURE), Inches(0.62), Inches(1.48), width=Inches(12.10), height=Inches(2.91))
    add_kpi(slide, "1024MB", "预测 SB 平台起点", 1.05, 4.82, 3.15, color=TEAL)
    add_kpi(slide, "1024MB", "实测 TPS 99% 平台", 5.08, 4.82, 3.15, color=GREEN)
    add_kpi(slide, "0.0%", "TPS regret", 9.10, 4.82, 3.15, color=BLUE)
    add_footer(slide, 10)
    outline.append("10. 效果图：S5；预测和实测 TPS 平台均为 1024MB。")

    # 11. Recommendations and validation
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, "当前推荐与真实验证结论", "完整机制已经运行；验证强度因阶段是否限速、是否具备计划锚点而不同。", "10 / Conclusion")
    table_data = [["阶段", "推荐 SB", "推荐 work_mem", "状态"]]
    short_names = {
        "stage1_memory_rich": "S1",
        "stage2_reach_limit": "S2",
        "stage3_protect_tp": "S3",
        "stage4_backpressure": "S4",
        "stage5_tp_surge": "S5",
    }
    for row in recommendations:
        table_data.append([
            short_names[row["stage"]],
            f"{row['recommended_sb_mb']}MB",
            f"{row['recommended_work_mem_mb']}MB",
            "暂定" if row["coverage_limited"].lower() == "true" else "完整网格",
        ])
    add_table(slide, 6, 4, table_data, 0.68, 1.72, 7.15, 3.75, widths=[1.05, 1.85, 2.05, 2.20], font_size=13)
    add_rect(slide, 8.18, 1.72, 4.42, 3.75, LIGHT, line=TEAL, radius=True)
    add_text(slide, "验证结果", 8.48, 2.02, 3.82, 0.38, size=22, color=TEAL, bold=True, align=PP_ALIGN.CENTER)
    s5_regret = float(validation["s5_tps_regret_pct"])
    add_bullets(
        slide,
        [
            f"S5：1024MB 精确命中 TPS 99% 平台，regret {s5_regret:.1f}%",
            "S3：SB 与最低 P95 点一致，1083MB 边界通过",
            "S1/S2/S4：约 40 TPS 限速，不能证明唯一最优",
            "S4：Q21 高 work_mem 新计划族 trace 尚未覆盖",
        ],
        8.45,
        2.60,
        3.72,
        2.35,
        size=14,
        spacing=10,
    )
    add_rect(slide, 0.76, 5.82, 11.72, 0.62, INK, radius=True)
    add_text(slide, "结论：S5 已验证一致；S3 基本一致；S1/S2 不可识别；S4 仍需补采新计划 trace。", 1.02, 5.99, 11.20, 0.28, size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, 11)
    outline.append("11. 推荐与验证：S5 一致，S3 基本一致，S1/S2 不可识别，S4 待补 trace。")

    prs.save(OUT)
    OUTLINE.write_text("# PPT 大纲\n\n" + "\n".join(f"- {item}" for item in outline) + "\n", encoding="utf-8")
    print(OUT)
    print(OUTLINE)


if __name__ == "__main__":
    build()
