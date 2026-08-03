#!/usr/bin/env python3
"""Build a target-first Huawei5 deck aligned with the original acceptance plan."""

from __future__ import annotations

import csv
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches

from build_memory_autonomy_paper_ppt import (
    BLUE,
    GRAY,
    GREEN,
    INK,
    LIGHT,
    MID,
    ORANGE,
    PALE_GREEN,
    PALE_ORANGE,
    PALE_RED,
    PALE_TEAL,
    PURPLE,
    RED,
    TEAL,
    WHITE,
    add_arrow,
    add_callout,
    add_flow_box,
    add_kpi,
    add_notes,
    add_rect,
    add_table,
    add_text,
    add_title,
    set_bg,
)


ROOT = Path(__file__).resolve().parents[1]
LATEST = ROOT / "artifacts/00_latest"
OUT = LATEST / "Huawei5_final_acceptance_path_20260730.pptx"
OUTLINE = LATEST / "Huawei5_final_acceptance_path_20260730_outline.md"

WORKMEM = ROOT / "artifacts/01_current_joint_model/figures/all_query_workmem_prediction_vs_actual.png"
STATIC_REC_DATA = LATEST / "five_stage_saturated_tps_validation_20260726.csv"
STATIC_REC = LATEST / "p1_p5_static_recommendation_validation_20260730.png"
RATE600 = LATEST / "five_stage_tp_slo_acceptance_20260727.png"
KERNEL_RESIZE = LATEST / "paper_kernel_resize_acceptance_20260727.png"
V8_TIMELINE = LATEST / "weekly_v8_controller_timeline_20260729.png"


def footer(slide, page: int, source: str = "") -> None:
    add_rect(slide, 0.62, 7.17, 12.08, 0.012, MID)
    add_text(slide, "Huawei5 · 内存池动态调整最终验收路径 · 2026-07-30", 0.62, 7.22, 7.2, 0.18, size=7.5, color=GRAY)
    if source:
        add_text(slide, f"来源：{source}", 7.10, 7.22, 5.05, 0.18, size=7, color=GRAY, align=PP_ALIGN.RIGHT)
    add_text(slide, str(page), 12.28, 7.22, 0.42, 0.18, size=8, color=GRAY, align=PP_ALIGN.RIGHT)


def add_image(slide, path: Path, x: float, y: float, w: float, h: float) -> None:
    with Image.open(path) as image:
        iw, ih = image.size
    scale = min(w / iw, h / ih)
    dw, dh = iw * scale, ih * scale
    slide.shapes.add_picture(
        str(path),
        Inches(x + (w - dw) / 2),
        Inches(y + (h - dh) / 2),
        width=Inches(dw),
        height=Inches(dh),
    )


def make_static_recommendation_plot() -> None:
    import matplotlib.pyplot as plt

    with STATIC_REC_DATA.open(newline="", encoding="utf-8") as handle:
        rows = list(csv.DictReader(handle))
    labels = [f"P{i}" for i in range(1, len(rows) + 1)]
    recommended = [float(row["recommended_actual_tps"]) for row in rows]
    best = [float(row["best_challenger_tps"]) for row in rows]
    regrets = [float(row["tps_regret_pct"]) for row in rows]

    x = list(range(len(labels)))
    width = 0.36
    fig, ax = plt.subplots(figsize=(10.8, 5.1))
    ax.bar([v - width / 2 for v in x], recommended, width, color="#2d7fb8", label="Replay recommendation: actual TPS")
    ax.bar([v + width / 2 for v in x], best, width, color="#9da8b1", label="Best measured challenge TPS")
    for i, (top, regret) in enumerate(zip(best, regrets)):
        ax.text(i, top + 18, f"regret {regret:.2f}%", ha="center", va="bottom", color="#25835d", fontsize=10, fontweight="bold")
    ax.set_xticks(x, labels)
    ax.set_ylabel("Saturated TP TPS under AP pressure")
    ax.set_title("P1-P5 static recommendation: every measured-grid TPS regret is below 5%", fontweight="bold")
    ax.set_ylim(0, max(best) * 1.15)
    ax.grid(axis="y", alpha=0.18)
    ax.legend(loc="lower right", frameon=False)
    fig.tight_layout()
    fig.savefig(STATIC_REC, dpi=180)
    plt.close(fig)


def build() -> None:
    for path in [WORKMEM, STATIC_REC_DATA, RATE600, KERNEL_RESIZE, V8_TIMELINE]:
        if not path.exists():
            raise FileNotFoundError(path)
    make_static_recommendation_plot()

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    outline: list[str] = []

    # 1 Cover
    slide = prs.slides.add_slide(blank)
    set_bg(slide, INK)
    add_rect(slide, 0, 0, 0.22, 7.5, TEAL)
    add_text(slide, "Huawei5 内存池动态调整", 0.92, 1.12, 10.8, 0.48, size=23, color=PALE_TEAL, bold=True)
    add_text(slide, "当前有效结果与五阶段最终验收路径", 0.90, 1.86, 11.80, 0.92, size=36, color=WHITE, bold=True)
    add_text(slide, "Trace Replay 负责预测边界 · 状态机负责动作顺序 · 真实 TPS 负责闭环纠偏", 0.94, 3.02, 11.54, 0.42, size=17, color=MID, bold=True)
    add_rect(slide, 0.94, 4.12, 11.52, 1.02, TEAL, radius=True)
    add_text(slide, "最终目标：AP 慢 SQL 持续加压时，动态让渡 SB 与算子内存；TP 突增时优先恢复 TP，同时不取消 AP。", 1.20, 4.42, 11.02, 0.43, size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "验收对齐版 · 2026-07-30", 0.94, 6.48, 4.0, 0.24, size=10, color=MID)
    add_notes(slide, "开场先说明：当前有分项证据，但连续五阶段尚未完成最终验收。")
    outline.append("1. 封面：当前有效结果与五阶段最终验收路径。")

    # 2 Problem and final target
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, "要解决的问题：AP 和 TP 争夺同一份内存，但必须优先保持 TP 平稳", "项目不是寻找一组永久固定参数，而是根据负载变化实时移动内存和控制 AP。", "01 · PROBLEM")
    add_flow_box(slide, "TP / OLTP", "高频小事务\n需要稳定 SB\n关注 TPS", 0.72, 1.78, 2.64, 1.66, BLUE, fill=LIGHT)
    add_arrow(slide, 3.42, 2.61, 4.06, 2.61, color=GRAY)
    add_flow_box(slide, "统一内存目标", "shared_buffers\n+ AP 动态内存\n≤ memory_target_max", 4.12, 1.62, 4.02, 1.98, TEAL, fill=PALE_TEAL)
    add_arrow(slide, 8.20, 2.61, 8.84, 2.61, color=GRAY)
    add_flow_box(slide, "AP / OLAP", "Join / Sort / Agg\n需要算子内存\n不足时 spill", 8.90, 1.78, 3.70, 1.66, ORANGE, fill=PALE_ORANGE)
    add_callout(slide, "正常方向", "AP 压力上升时，先用富余内存，再由 SB 向动态池让渡；达到保护边界后限制或排队 AP。", 0.78, 4.25, 5.74, 1.26, TEAL, PALE_TEAL)
    add_callout(slide, "反向方向", "TP 突增时，停止新 AP、降低 AP 后续 grant，等实际动态内存归还后提高 SB。", 6.80, 4.25, 5.74, 1.26, BLUE, LIGHT)
    add_rect(slide, 0.78, 5.92, 11.76, 0.70, INK, radius=True)
    add_text(slide, "最终不是“预测一个最好参数”，而是“预测安全边界 + 在线决定何时、向哪个方向调整”。", 1.02, 6.14, 11.28, 0.29, size=15, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    footer(slide, 2, "内存池动态调整方案验证.pptx")
    outline.append("2. 问题：统一内存池中动态平衡 TP 的 SB 与 AP 算子内存。")

    # 3 Acceptance definition
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, "最终验收口径：必须同时满足动作正确、TP 稳定和 AP 自然完成", "只看配置 regret、命中率或前 180 秒都不能单独代表验收通过。", "01 · ACCEPTANCE")
    add_kpi(slide, "≤3%", "单次在线调整 TPS 抖动", 0.72, 1.67, 2.76, color=TEAL, note="granule 扩缩全过程")
    add_kpi(slide, "≥95%", "每阶段 TP retention", 3.62, 1.67, 2.76, color=BLUE, note="相对同 offered load 的 no-AP 基线")
    add_kpi(slide, "<5pp", "五阶段归一化差", 6.52, 1.67, 2.76, color=ORANGE, note="阶段间 TP 相对波动")
    add_kpi(slide, "0", "restart / cancel", 9.42, 1.67, 2.76, color=GREEN, note="AP 最终自然完成")
    data = [
        ["验收维度", "必须观察到的证据"],
        ["内存守恒", "SB + 实际 AP 动态内存始终不超过 memory_target_max"],
        ["动作顺序", "S1 增 grant → S2 降 SB → S3 降 grant → S4 排队 → S5 提 SB"],
        ["TP 保护", "逐秒 TPS、15 秒 retention、切换窗口抖动均满足门槛"],
        ["AP 完整性", "记录等待、开始、进展、结束；不使用取消制造阶段结束"],
    ]
    add_table(slide, data, 0.80, 3.67, 11.74, 2.55, widths=[2.44, 9.30], font_size=10.8)
    add_text(slide, "这四类证据必须来自同一次连续运行。", 0.86, 6.52, 11.62, 0.34, size=15, color=RED, bold=True, align=PP_ALIGN.CENTER)
    footer(slide, 3, "原验收约束 + TP 稳定性补充要求")
    outline.append("3. 验收：≤3% 调整抖动、≥95% retention、<5pp 阶段差、0 restart/cancel。")

    # 4 Five-stage state machine
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, "原 PPT 的五阶段：不是五套 Query，而是一条连续压力轨迹", "AP 从 0 持续增加，运行中 Query 跨阶段保留；S5 才发生 TP 低负载到高负载的跃升。", "02 · FIVE STATES")
    stages = [
        ("S1", "内存富裕", "增加 AP grant", "减少 spill", TEAL),
        ("S2", "触及上限", "降低 SB", "让渡 granule", PURPLE),
        ("S3", "保护基准", "冻结 SB 下限\n降低 AP grant", "守住低 TP", BLUE),
        ("S4", "反压排队", "新 AP 入队", "存量自然运行", ORANGE),
        ("S5", "TP 突增", "阻塞 AP + 提 SB", "优先恢复 TP", GREEN),
    ]
    for i, (code, title, action, goal, color) in enumerate(stages):
        x = 0.48 + i * 2.57
        add_rect(slide, x, 1.78, 2.24, 3.30, LIGHT, line=color, radius=True, width=1.5)
        add_rect(slide, x + 0.68, 1.50, 0.88, 0.54, color, radius=True)
        add_text(slide, code, x + 0.74, 1.66, 0.76, 0.22, size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, title, x + 0.12, 2.25, 2.00, 0.34, size=16, color=color, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, "必须动作", x + 0.20, 2.92, 1.84, 0.25, size=10, color=GRAY, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, action, x + 0.18, 3.27, 1.88, 0.55, size=13, color=INK, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, goal, x + 0.18, 4.30, 1.88, 0.30, size=11, color=GRAY, align=PP_ALIGN.CENTER)
        if i < 4:
            add_arrow(slide, x + 2.26, 3.44, x + 2.51, 3.44, color=GRAY, width=1.5)
    add_rect(slide, 0.76, 5.62, 11.82, 0.86, INK, radius=True)
    add_text(slide, "状态触发来自真实观测：内存 headroom、TP retention、AP 队列和 TP offered load；不能由 Query 编号代替。", 1.02, 5.89, 11.30, 0.34, size=14.5, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    footer(slide, 4, "内存池动态调整方案验证.pptx")
    outline.append("4. 原五阶段：连续 AP 加压和最终 TP 突增下的规定动作。")

    # 5 Model architecture
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, "为实现最终验收，模型应分成四层，而不是输出一张静态配置表", "Replay 不读取验证 TPS 最优答案；运行时 TPS 只作为传感器和安全反馈。", "03 · MODEL")
    layers = [
        ("离线 Trace Replay", "TP 页访问\nAP Plan/算子 trace\n源码容量规则", "输出安全候选与代价", PURPLE),
        ("五阶段状态机", "内存 headroom\nTP 负载变化\nAP 排队压力", "决定动作方向", TEAL),
        ("在线反馈控制", "TP TPS/retention\nAP CPU/I/O/进展\n实际动态内存", "决定动作幅度与回退", BLUE),
        ("执行器", "SB granule\nper-query work_mem\n准入/排队/freezer", "真正修改运行状态", ORANGE),
    ]
    for i, (title, inputs, output, color) in enumerate(layers):
        x = 0.58 + i * 3.18
        add_flow_box(slide, title, inputs, x, 1.64, 2.74, 2.12, color, fill=LIGHT if i != 1 else PALE_TEAL)
        add_rect(slide, x + 0.20, 4.10, 2.34, 0.82, color, radius=True)
        add_text(slide, output, x + 0.30, 4.35, 2.14, 0.28, size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        if i < 3:
            add_arrow(slide, x + 2.78, 2.70, x + 3.10, 2.70, color=GRAY)
    add_callout(slide, "Replay 的价值", "预测不同 SB/work_mem 下的缓存 miss、Plan、spill、动态峰值和不可部署区域。", 0.76, 5.48, 5.76, 0.92, TEAL, PALE_TEAL)
    add_callout(slide, "闭环的价值", "处理纯 trace 难以精确覆盖的 CPU、I/O 队列、page cache 污染和执行路径时变。", 6.80, 5.48, 5.76, 0.92, BLUE, LIGHT)
    footer(slide, 5, "当前模型与目标控制结构")
    outline.append("5. 模型：Replay、状态机、在线反馈和执行器四层分工。")

    # 6 Per-query work_mem
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, "逐 Query work_mem：当前 Replay 已输出差异化分配", "统一 work_mem 只用于横向比较；实际控制器在每条 AP Query 的独立 session 启动前应用对应 grant。", "03 · PER-QUERY GRANT")
    data = [
        ["静态负载组", "并发 Query", "逐 Query work_mem 分配", "预测动态峰值", "预测 spill I/O"],
        ["P1", "Q1", "Q1=1MB", "2MB", "0"],
        ["P2", "Q3", "Q3=1150MB", "1400MB", "0"],
        ["P3", "Q5 + Q7", "Q5=1024MB；Q7=1083MB", "2936MB", "0"],
        ["P4", "Q9/Q13/Q18/Q21", "1174 / 1024 / 4096 / 2968MB", "15207MB", "26828MiB"],
        ["P5", "Q1/Q3/Q5/Q7", "256 / 1150 / 1024 / 1137MB", "4392MB", "0"],
    ]
    add_table(slide, data, 0.46, 1.52, 12.42, 3.54, widths=[1.32, 2.08, 4.42, 2.10, 2.50], font_size=9.8)
    add_callout(slide, "P4 为什么必须不同", "统一 2048MB 会把内存给收益较低的 Query，同时 Q18/Q21 仍 spill。逐 Query 搜索把更多 grant 给重算子。", 0.70, 5.43, 5.78, 1.08, ORANGE, PALE_ORANGE)
    add_callout(slide, "当前收益", "P4 预测 spill I/O 从 31895 降至 26828MiB，下降约 15.9%；动态峰值 15207MB，仍低于 15291MB 动态池。", 6.76, 5.43, 5.86, 1.08, TEAL, PALE_TEAL)
    add_text(slide, "限制：grant 只在 Query 启动边界生效；运行中 Query 已占用内存只能 graceful 回收。", 0.82, 6.66, 11.70, 0.27, size=13, color=RED, bold=True, align=PP_ALIGN.CENTER)
    footer(slide, 6, "stage_work_mem_recommendations.csv；runtime controller assignments")
    outline.append("6. 逐 Query work_mem：P1-P5 已输出差异化 session grant，P4 spill 预测降低约15.9%。")

    # 7 Replay evidence
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, "有效结果一：Replay 能预测多数 Query 的 Plan/spill 边界", "这是后续 S1-S4 调整 AP grant 的依据，但不直接证明连续五阶段已经通过。", "04 · VALID EVIDENCE 1")
    add_image(slide, WORKMEM, 0.54, 1.45, 8.60, 5.12)
    add_kpi(slide, "7/7", "Plan family", 9.34, 1.70, 1.42, color=GREEN)
    add_kpi(slide, "6/7", "spill 分类", 10.94, 1.70, 1.42, color=ORANGE)
    add_callout(slide, "能证明", "模型能把 Query 的执行计划、算子内存、spill 和动态峰值转成候选 work_mem 边界。", 9.30, 3.12, 3.10, 1.30, TEAL, PALE_TEAL)
    add_callout(slide, "不能证明", "未知算子和所有 scale factor 都准确，也不能单凭这张图推出五阶段 TP 稳定。", 9.30, 4.72, 3.10, 1.30, RED, PALE_RED)
    footer(slide, 7, "all_query_workmem_prediction_vs_actual.png")
    outline.append("7. 分项证据：Plan/spill 和 Query 最小 work_mem 边界。")

    # 8 Static recommendation evidence
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, "有效结果二：P1-P5 静态推荐接近各自最高 TPS", "它验证五组独立负载的配置排序，不再称作原方案 S1-S5。", "04 · VALID EVIDENCE 2")
    add_image(slide, STATIC_REC, 0.54, 1.44, 9.42, 5.02)
    add_kpi(slide, "5/5", "P1-P5 推荐通过", 10.14, 1.72, 2.26, color=GREEN)
    add_kpi(slide, "4.55%", "最大 TPS regret", 10.14, 3.06, 2.26, color=ORANGE)
    add_callout(slide, "有效结论", "在已经测过的 SB × work_mem 挑战网格中，Replay 的候选排序具有实用价值。", 10.10, 4.44, 2.34, 1.32, TEAL, PALE_TEAL)
    add_text(slide, "不能表述为：原方案连续五阶段已经完成内存让渡和 TPS 验收。", 0.82, 6.53, 11.70, 0.31, size=14, color=RED, bold=True, align=PP_ALIGN.CENTER)
    footer(slide, 8, "P1-P5 relabel of five_stage_saturated_tps_validation_20260726.csv")
    outline.append("8. 分项证据：P1-P5 静态配置推荐最大 TPS regret 4.55%。")

    # 9 Fixed-rate control evidence
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, "当前有效结果三：固定 600 TPS 条件下，闭环具备保护 TP 的能力", "该实验验证反馈、准入和资源限制链路；负载轨迹仍不同于原 PPT。", "04 · VALID EVIDENCE 3")
    add_image(slide, RATE600, 0.56, 1.46, 8.68, 5.02)
    add_kpi(slide, "99.74%", "最低阶段 retention", 9.46, 1.74, 2.76, color=GREEN)
    add_kpi(slide, "1.32pp", "五组最大最小差", 9.46, 3.12, 2.76, color=TEAL)
    add_callout(slide, "边界", "固定 TP offered load、固定 AP 资源隔离；没有验证 AP 持续累积、S2 让渡和 S5 TP 突增。", 9.40, 4.54, 2.88, 1.30, RED, PALE_RED)
    add_text(slide, "所以它是“控制链路可用”的证据，不是最终五阶段验收结果。", 0.82, 6.53, 11.70, 0.31, size=14, color=INK, bold=True, align=PP_ALIGN.CENTER)
    footer(slide, 9, "five_stage_tp_slo_acceptance_20260727.png")
    outline.append("9. 分项证据：固定 600 TPS 闭环稳定，但不是原连续五阶段。")

    # 10 Kernel evidence
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, "当前有效结果四：shared_buffers 在线扩缩的内核执行路径已经跑通", "这证明可以做到不重启和低抖动调整；生产规模与完整五阶段仍需再验。", "04 · VALID EVIDENCE 4")
    add_image(slide, KERNEL_RESIZE, 0.56, 1.45, 8.72, 5.16)
    add_kpi(slide, "0", "数据库重启", 9.48, 1.72, 2.72, color=GREEN)
    add_kpi(slide, "2.65%", "最差单秒下降", 9.48, 3.08, 2.72, color=TEAL)
    add_callout(slide, "有效边界", "128→64MB 隔离原型、8MB/s 读写 TP；验证的是 resize 执行器，不是完整控制策略。", 9.42, 4.48, 2.86, 1.30, ORANGE, PALE_ORANGE)
    add_text(slide, "最终验收还需在 memory_target_max 下同时观察 SB、动态池、AP spill 和 TP TPS。", 0.82, 6.53, 11.70, 0.31, size=14, color=INK, bold=True, align=PP_ALIGN.CENTER)
    footer(slide, 10, "paper_kernel_resize_acceptance_20260727.png")
    outline.append("10. 分项证据：在线 SB 扩缩 0 restart，隔离原型抖动低于3%。")

    # 11 Status audit
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, "当前完成度：组件证据具备，连续五阶段尚未闭合", "下表是现在最准确的项目状态，避免把“组件完成”写成“系统验收完成”。", "05 · STATUS")
    data = [
        ["能力", "状态", "现有证据", "最终还缺什么"],
        ["TP 缓存 Replay", "已验证", "held-out SB/OS/combined 误差", "纳入连续状态的实时工作集"],
        ["AP Plan/spill Replay", "大部分验证", "Plan 7/7、spill 6/7", "补未知算子与更大泛化集"],
        ["P1-P5 静态联合推荐", "已验证已测网格", "最大 TPS regret 4.55%", "不能替代动态动作验收"],
        ["AP 准入/排队/graceful debt", "机制已实现", "控制日志和压力 sweep", "按 S1-S5 正确时刻触发"],
        ["SB 在线扩缩", "原型已验证", "0 restart、局部抖动 <3%", "生产粒度连续迁移"],
        ["原方案连续五阶段", "尚未完成", "当前没有同口径结果", "重构负载并端到端实跑"],
    ]
    add_table(slide, data, 0.46, 1.50, 12.42, 4.86, widths=[2.62, 2.05, 3.45, 4.30], font_size=9.4)
    add_rect(slide, 0.76, 6.48, 11.82, 0.46, RED, radius=True)
    add_text(slide, "当前可以说“具备进入系统级验收的基础”，不能说“五阶段已经验收通过”。", 1.04, 6.61, 11.26, 0.22, size=13.5, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    footer(slide, 11, "现有实验审计")
    outline.append("11. 状态审计：组件级证据充分，连续五阶段尚未完成。")

    # 12 Why mismatch
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, "为什么上一版看起来全都不一致：实验把“状态”改成了“Query 组合”", "这是负载协议没有对齐造成的，不是通过修改推荐数字可以修复。", "05 · ROOT CAUSE")
    data = [
        ["原验收方案", "此前构造的实验", "直接后果"],
        ["一次连续运行", "P1-P5 分别独立运行", "没有跨阶段内存压力"],
        ["AP 从 0 持续增加", "每组固定 Query 集合", "无法自然触发 S1-S4"],
        ["运行 AP 跨阶段保留", "每阶段结束等待 AP 全部完成", "没有 graceful 跨阶段债务"],
        ["S1-S4 低 TP，S5 突增", "五组都固定 800 TPS", "无法验证 S5 反向让渡"],
        ["状态机规定动作方向", "静态优化器选择局部最优", "可能出现 S2 不降 SB、S5 不升 SB"],
    ]
    add_table(slide, data, 0.54, 1.54, 12.24, 4.46, widths=[3.72, 4.10, 4.42], font_size=10.1)
    add_callout(slide, "命名修正", "P1-P5：五组独立静态挑战负载，只验证配置排序。", 0.76, 6.18, 5.76, 0.70, TEAL, PALE_TEAL)
    add_callout(slide, "验收命名", "S1-S5：连续压力状态，只用于最终动态验收。", 6.80, 6.18, 5.76, 0.70, BLUE, LIGHT)
    footer(slide, 12, "driver 与原方案对比")
    outline.append("12. 根因：把连续状态错误地替换为五组独立 Query 组合。")

    # 13 Correct experiment
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, "正确的下一轮实验：一个不中断的 TP/AP 运行，依次制造五种状态", "阶段切换由观测阈值触发；时间上限只关闭新 AP 准入，不取消已经运行的 SQL。", "06 · FINAL EXPERIMENT")
    rows = [
        ["状态", "负载构造", "控制器预期动作", "关键记录"],
        ["S1", "低 TP；AP 从 0 缓慢注入", "提高 AP grant，直到 spill 收益饱和", "grant、spill、动态峰值"],
        ["S2", "继续加 AP，达到总内存 max", "逐 granule 降 SB，转给动态池", "SB active/target、hit、抖动"],
        ["S3", "AP 压力继续增长", "冻结 SB 下限，降低新 Query grant", "TP retention、graceful debt"],
        ["S4", "内存和 I/O 均到保护边界", "新 AP 排队，存量自然运行", "requested/admitted/queued"],
        ["S5", "TP offered load 阶跃到高档", "阻塞新 AP；归还后提高 SB", "逐秒 TPS、回收、扩容顺序"],
    ]
    add_table(slide, rows, 0.46, 1.50, 12.42, 4.64, widths=[1.12, 3.30, 4.45, 3.55], font_size=9.6)
    add_rect(slide, 0.70, 6.30, 11.94, 0.56, INK, radius=True)
    add_text(slide, "必须保持：同一数据库进程、同一 TP 流、运行 AP 跨阶段保留、统一 memory_target_max、0 cancel。", 0.96, 6.47, 11.42, 0.25, size=13.5, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    footer(slide, 13, "重构后的连续验收协议")
    outline.append("13. 最终实验：连续 TP/AP、AP 累积、S5 TP 阶跃和全生命周期记录。")

    # 14 Evidence expected
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, "最终验收输出：用一张时间轴证明动作与 TPS", "最终报告必须能从同一时间轴核对负载、内存、动作和结果之间的因果顺序，而不是再给五个静态配置。", "06 · ACCEPTANCE OUTPUT")
    add_flow_box(slide, "轨道 1 · 负载", "TP offered load\nAP requested/running/queued", 0.70, 1.70, 2.70, 1.32, BLUE, fill=LIGHT)
    add_flow_box(slide, "轨道 2 · 内存", "SB active/target\nactual dynamic/debt", 3.70, 1.70, 2.70, 1.32, TEAL, fill=PALE_TEAL)
    add_flow_box(slide, "轨道 3 · 动作", "grant / admission\nSB granule / freeze", 6.70, 1.70, 2.70, 1.32, ORANGE, fill=PALE_ORANGE)
    add_flow_box(slide, "轨道 4 · 结果", "TP TPS/retention\nAP spill/progress/end", 9.70, 1.70, 2.70, 1.32, GREEN, fill=PALE_GREEN)
    add_arrow(slide, 1.02, 3.58, 12.10, 3.58, color=GRAY, width=2.0)
    labels = ["S1 增 grant", "S2 降 SB", "S3 降 grant", "S4 排队", "S5 提 SB"]
    colors = [TEAL, PURPLE, BLUE, ORANGE, GREEN]
    for i, (label, color) in enumerate(zip(labels, colors)):
        x = 0.78 + i * 2.42
        add_rect(slide, x, 3.34, 0.46, 0.46, color, radius=True)
        add_text(slide, label, x - 0.30, 4.02, 1.80, 0.34, size=11, color=color, bold=True, align=PP_ALIGN.CENTER)
    add_callout(slide, "通过条件", "动作顺序正确；内存不超限；调整抖动 ≤3%；阶段 retention ≥95%；阶段差 <5pp；0 restart/cancel。", 0.78, 5.08, 7.28, 1.22, TEAL, PALE_TEAL)
    add_callout(slide, "失败条件", "任何一个阶段未触发规定动作、AP 被取消、只截取阶段前 180 秒，或者 TP 越过门槛，都必须保留为失败结果。", 8.30, 5.08, 4.24, 1.22, RED, PALE_RED)
    footer(slide, 14, "最终验收证据设计")
    outline.append("14. 最终输出：同一时间轴展示负载、内存、动作和 TP/AP 结果。")

    # 15 Conclusion
    slide = prs.slides.add_slide(blank)
    set_bg(slide, INK)
    add_text(slide, "当前结论", 0.92, 0.92, 2.50, 0.44, size=22, color=PALE_TEAL, bold=True)
    add_text(slide, "项目已经具备进入最终系统验收的基础", 0.90, 1.58, 11.36, 0.72, size=33, color=WHITE, bold=True)
    add_callout(slide, "已经成立", "Trace Replay 能生成 SB/work_mem 安全候选；静态 P1-P5 推荐在已测网格最大 regret 4.55%；在线控制和 SB 执行器已有实测证据。", 0.78, 2.78, 5.78, 1.68, TEAL, PALE_TEAL)
    add_callout(slide, "尚未成立", "原方案的连续 AP 加压、跨阶段运行态、S2 内存让渡、S4 排队和 S5 TP 突增尚未在同一次运行中完整通过。", 6.80, 2.78, 5.78, 1.68, RED, PALE_RED)
    add_rect(slide, 0.78, 5.02, 11.80, 1.08, TEAL, radius=True)
    add_text(slide, "下一项明确工作：重构连续五阶段负载并冻结控制策略，完成一次端到端实跑；结果不通过就保留失败证据并继续修正。", 1.08, 5.34, 11.18, 0.46, size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "不能再用五个静态最优点替代最终动态验收。", 0.92, 6.54, 11.46, 0.30, size=14, color=MID, bold=True, align=PP_ALIGN.CENTER)
    add_notes(slide, "收尾明确短期结果有效，但最终验收仍需连续五阶段实跑。")
    outline.append("15. 结论：组件基础具备；下一步完成连续五阶段端到端实跑。")

    prs.save(OUT)
    OUTLINE.write_text(
        "# Huawei5 内存池动态调整最终验收路径（2026-07-30）\n\n"
        + "\n".join(f"- {item}" for item in outline)
        + "\n\n## 汇报口径\n\n"
        + "- P1-P5 仅表示五组独立静态挑战负载，不再使用 S1-S5 名称。\n"
        + "- S1-S5 仅表示原验收方案的一条连续压力轨迹。\n"
        + "- 当前结果属于预测器、控制链路和执行器的分项验证。\n"
        + "- 只有连续五阶段同一次运行满足全部门槛，才能表述为最终验收通过。\n",
        encoding="utf-8",
    )
    print(OUT)
    print(OUTLINE)


if __name__ == "__main__":
    build()
