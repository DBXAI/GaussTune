#!/usr/bin/env python3
"""Build a clean, narrative defense deck for Huawei6."""

from __future__ import annotations

import json
import shutil
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from build_memory_autonomy_paper_ppt import (
    GRAY,
    INK,
    MID,
    WHITE,
    add_arrow,
    add_rect,
    add_table,
    add_text,
    set_bg,
)
from build_zero_to_joint_recommendation_ppt import (
    CHART_FORMULA,
    CHART_STAGE,
    CHART_SURFACE,
    FORMULA_REPORT,
    STAGE_REPORT,
    make_charts,
)


ROOT = Path(__file__).resolve().parents[1]
LATEST = ROOT / "artifacts/00_latest"
OUT = LATEST / "Huawei6_项目答辩_连续叙事版_20260802.pptx"
ROOT_COPY = Path("/root/Huawei6_项目答辩_连续叙事版.pptx")
SCRIPT = LATEST / "Huawei6_项目答辩_连续叙事版_20260802_讲稿.md"

NAVY = RGBColor(18, 39, 72)
BLUE = RGBColor(38, 96, 208)
CYAN = RGBColor(0, 142, 166)
GREEN = RGBColor(39, 145, 94)
ORANGE = RGBColor(234, 112, 25)
RED = RGBColor(210, 49, 52)
PALE_BLUE = RGBColor(240, 245, 255)
PALE_GREEN = RGBColor(239, 249, 244)
PALE_ORANGE = RGBColor(255, 246, 236)
PALE_RED = RGBColor(254, 241, 241)
FONT = "Microsoft YaHei"


def read_json(path: Path):
    return json.loads(path.read_text(encoding="utf-8"))


def title(slide, section: str, headline: str, subtitle: str = "") -> None:
    add_rect(slide, 0, 0, 13.333, 0.08, BLUE)
    add_rect(slide, 0.55, 0.28, 0.06, 0.42, BLUE)
    add_text(slide, section.upper(), 0.72, 0.29, 2.6, 0.22, size=8.5, color=BLUE, bold=True)
    add_text(slide, headline, 0.72, 0.69, 11.85, 0.58, size=26, color=NAVY, bold=True)
    if subtitle:
        add_text(slide, subtitle, 0.74, 1.27, 11.6, 0.32, size=11.5, color=GRAY)


def footer(slide, page: int) -> None:
    add_rect(slide, 0.62, 7.18, 12.08, 0.012, MID)
    add_text(slide, "Huawei6 · TP/AP 混合负载内存联合预测与稳定控制", 0.62, 7.22, 6.2, 0.18, size=7.2, color=GRAY)
    add_text(slide, str(page), 12.30, 7.22, 0.38, 0.18, size=8, color=GRAY, align=PP_ALIGN.RIGHT)


def takeaway(slide, text: str, *, color=BLUE) -> None:
    add_rect(slide, 0.72, 6.48, 0.08, 0.48, color)
    add_text(slide, text, 0.98, 6.52, 11.55, 0.32, size=14, color=NAVY, bold=True)


def circle(slide, x: float, y: float, d: float, color, text_value: str, size=15):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    shape.fill.solid(); shape.fill.fore_color.rgb = color
    shape.line.color.rgb = color
    add_text(slide, text_value, x, y + d * 0.30, d, d * 0.32, size=size, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    return shape


def divider(slide, x: float, y: float, h: float, color=MID):
    add_rect(slide, x, y, 0.012, h, color)


def cover(prs, script):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(slide, NAVY)
    add_rect(slide, 0, 0, 13.333, 0.12, BLUE)
    add_text(slide, "项目答辩", 0.82, 0.84, 2.0, 0.28, size=14, color=RGBColor(120, 166, 255), bold=True)
    add_text(slide, "面向 TP/AP 混合负载的", 0.80, 1.52, 10.9, 0.55, size=30, color=WHITE, bold=True)
    add_text(slide, "内存联合预测与 TP 稳定控制", 0.80, 2.12, 11.4, 0.72, size=39, color=WHITE, bold=True)
    add_rect(slide, 0.82, 3.34, 6.82, 0.035, BLUE)
    add_text(slide, "Trace Replay  ×  机器因果模型  ×  双向约束搜索", 0.82, 3.62, 9.4, 0.38, size=18, color=RGBColor(205, 219, 241))
    add_text(slide, "从命中率预测，走向可执行配置推荐和 TP TPS 稳定保护", 0.82, 4.33, 10.3, 0.40, size=17, color=RGBColor(120, 166, 255), bold=True)
    add_text(slide, "答辩人：__________", 0.82, 6.34, 3.3, 0.26, size=11, color=RGBColor(164, 180, 202))
    add_text(slide, "2026-08-02", 10.48, 6.34, 1.8, 0.26, size=11, color=RGBColor(164, 180, 202), align=PP_ALIGN.RIGHT)
    script.append("开场：我们的工作从一个简单问题开始：在 AP 与 TP 同机运行时，为什么命中率看起来不差，TP TPS 仍然会下降？")


def s2_problem(prs, script):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(slide)
    title(slide, "01 BACKGROUND", "我们面对的不是“内存够不够”，而是一个自我强化的性能闭环")
    add_text(slide, "shared_buffers", 0.82, 1.88, 2.2, 0.30, size=17, color=BLUE, bold=True)
    add_text(slide, "决定 TP page miss", 0.82, 2.26, 2.5, 0.26, size=12, color=GRAY)
    add_arrow(slide, 2.92, 2.40, 4.12, 2.40, BLUE, 2.3)
    add_text(slide, "AP work_mem", 4.32, 1.88, 2.2, 0.30, size=17, color=ORANGE, bold=True)
    add_text(slide, "决定 plan 与 spill", 4.32, 2.26, 2.5, 0.26, size=12, color=GRAY)
    add_arrow(slide, 6.42, 2.40, 7.62, 2.40, ORANGE, 2.3)
    add_text(slide, "NVMe queue", 7.82, 1.88, 2.2, 0.30, size=17, color=RED, bold=True)
    add_text(slide, "TP miss + AP spill", 7.82, 2.26, 2.6, 0.26, size=12, color=GRAY)
    add_arrow(slide, 9.92, 2.40, 11.12, 2.40, RED, 2.3)
    add_text(slide, "TP TPS", 11.30, 1.88, 1.3, 0.30, size=17, color=GREEN, bold=True)
    add_text(slide, "响应时间", 11.30, 2.26, 1.2, 0.26, size=12, color=GRAY)
    add_arrow(slide, 11.88, 2.83, 11.88, 4.15, GREEN, 2.1)
    add_arrow(slide, 11.88, 4.15, 1.77, 4.15, GREEN, 2.1)
    add_arrow(slide, 1.77, 4.15, 1.77, 2.83, GREEN, 2.1)
    add_text(slide, "闭环", 6.18, 3.90, 1.0, 0.28, size=12, color=GREEN, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "SB 变大能减少 TP miss，却会压缩 AP 动态内存；AP spill 增加后又会抬高 TP I/O latency。", 1.40, 4.75, 10.55, 0.42, size=18, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "因此配置选择必须同时解释“内存如何分”和“I/O 如何回到 TPS”。", 1.58, 5.47, 10.20, 0.34, size=15, color=GRAY, align=PP_ALIGN.CENTER)
    takeaway(slide, "下一步问题：只预测缓存命中率，能否解释这个闭环？")
    footer(slide, 2)
    script.append("背景：SB 和 work_mem 不是两个独立旋钮。SB 影响 TP miss，work_mem 影响 AP spill，二者最后在 NVMe 队列里汇合并改变 TPS。")


def s3_gap(prs, script):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(slide)
    title(slide, "02 GAP", "旧模型止步于命中率，因此无法解释“相同 hit、不同 TPS”")
    add_text(slide, "原有能力", 0.82, 1.78, 2.0, 0.30, size=15, color=GRAY, bold=True)
    add_text(slide, "Trace → SB / OS hit ratio", 0.82, 2.27, 4.6, 0.42, size=25, color=BLUE, bold=True)
    add_text(slide, "能够回答：不同 SB 下，有多少访问命中数据库缓存和 Linux page cache。", 0.82, 2.90, 5.2, 0.58, size=13, color=GRAY)
    divider(slide, 6.28, 1.78, 3.10)
    add_text(slide, "缺失能力", 6.68, 1.78, 2.0, 0.30, size=15, color=RED, bold=True)
    add_text(slide, "hit ratio  ↛  TP TPS", 6.68, 2.27, 4.8, 0.42, size=25, color=RED, bold=True)
    add_text(slide, "无法回答：AP spill 抢占 I/O 后，同样的 TP miss 为什么会等待更久。", 6.68, 2.90, 5.3, 0.58, size=13, color=GRAY)
    add_rect(slide, 0.82, 4.05, 11.55, 0.04, MID)
    add_text(slide, "我们需要补上的推导", 0.82, 4.42, 3.2, 0.30, size=15, color=CYAN, bold=True)
    add_text(slide, "配置", 0.90, 5.08, 1.0, 0.32, size=18, color=NAVY, bold=True)
    add_arrow(slide, 1.88, 5.24, 2.78, 5.24, BLUE, 2.0)
    add_text(slide, "plan / spill", 2.96, 5.08, 1.7, 0.32, size=18, color=ORANGE, bold=True)
    add_arrow(slide, 4.68, 5.24, 5.58, 5.24, ORANGE, 2.0)
    add_text(slide, "I/O latency", 5.76, 5.08, 1.7, 0.32, size=18, color=RED, bold=True)
    add_arrow(slide, 7.48, 5.24, 8.38, 5.24, RED, 2.0)
    add_text(slide, "transaction time", 8.56, 5.08, 2.1, 0.32, size=18, color=CYAN, bold=True)
    add_arrow(slide, 10.70, 5.24, 11.34, 5.24, CYAN, 2.0)
    add_text(slide, "TPS", 11.52, 5.08, 0.8, 0.32, size=18, color=GREEN, bold=True)
    takeaway(slide, "所以，本项目不是再拟合一条 TPS 曲线，而是补齐配置到 TPS 的三条因果链。")
    footer(slide, 3)
    script.append("研究缺口：原模型只到命中率，缺少 AP spill 对 I/O latency 的影响，也缺少 latency 到 TPS 的响应时间换算。")


def s4_solution(prs, script):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(slide)
    title(slide, "03 APPROACH", "我们把问题拆成三条可独立验证、最终可闭环的因果链")
    y = [1.85, 3.22, 4.59]
    labels = [
        ("01", "负载链", "一次 Trace / 源码算子行为", "不同 work_mem 下的 plan、动态峰值与 spill", BLUE),
        ("02", "机器链", "AP I/O 队列压力", "当前机器上的 TP 物理 I/O latency", ORANGE),
        ("03", "控制链", "TP 性能下界 + AP 内存上界", "SB / 逐 Query work_mem / AP-cap", GREEN),
    ]
    for yy, (num, name, left, right, color) in zip(y, labels):
        circle(slide, 0.88, yy, 0.58, color, num, size=12)
        add_text(slide, name, 1.72, yy + 0.06, 1.40, 0.30, size=17, color=color, bold=True)
        add_text(slide, left, 3.26, yy + 0.06, 3.25, 0.30, size=15, color=NAVY)
        add_arrow(slide, 6.48, yy + 0.22, 7.42, yy + 0.22, color, 2.0)
        add_text(slide, right, 7.66, yy + 0.06, 4.65, 0.42, size=15, color=NAVY, bold=True)
        if yy < 4.5: add_rect(slide, 1.72, yy + 0.88, 10.58, 0.012, MID)
    takeaway(slide, "三条链分别回答“候选会发生什么”“机器会慢多少”“最后选哪一个”。")
    footer(slide, 4)
    script.append("总体思路：负载链描述候选后果，机器链计算延迟代价，控制链在可行域中选择配置。")


def s5_replay(prs, script):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(slide)
    title(slide, "04 WORK 1", "一次负载 Trace，可以推演不同 work_mem 下的 plan 与 spill")
    add_text(slide, "一次真实执行", 0.82, 1.83, 2.0, 0.30, size=16, color=BLUE, bold=True)
    circle(slide, 1.18, 2.45, 0.66, BLUE, "SQL", size=11)
    add_arrow(slide, 1.86, 2.78, 3.05, 2.78, BLUE, 2.0)
    add_text(slide, "EXPLAIN / 算子输入规模 / 执行路径 / 访问轨迹", 3.24, 2.56, 4.46, 0.44, size=15, color=NAVY, bold=True)
    add_arrow(slide, 7.72, 2.78, 8.84, 2.78, ORANGE, 2.0)
    circle(slide, 9.02, 2.45, 0.66, ORANGE, "Replay", size=9)
    add_text(slide, "枚举 work_mem", 9.88, 2.56, 2.1, 0.36, size=16, color=ORANGE, bold=True)
    add_rect(slide, 0.82, 3.52, 11.50, 0.025, MID)
    data = [["work_mem", "plan family", "dynamic_peak", "spill I/O", "用途"],
            ["Q18=1150MB", "trace-supported", "较高", "低", "AP 富裕阶段"],
            ["Q18=256MB", "trace-supported", "较低", "增加", "TP 饱和保护"],
            ["Q9/Q13=256MB", "plan-supported", "内存安全", "可计算", "S3/S4 grant"]]
    add_table(slide, data, 0.86, 3.85, 11.42, 1.67, widths=[2.2, 2.55, 2.05, 1.85, 2.77], font_size=9.3)
    add_text(slide, "输出不是“Q18 固定用 1150MB”，而是每个 Query × work_mem 的行为曲线和置信度。", 1.08, 5.83, 11.0, 0.36, size=15, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    takeaway(slide, "负载链给出了 AP 动态内存、spill 和 plan，也给出了 TP 在不同 SB 下的 I/O/tx。")
    footer(slide, 5)
    script.append("工作一：通过 plan-aware operator replay，把一次采集扩展为多个候选配置，而不是把最优 work_mem 写死。")


def s6_machine(prs, script):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(slide)
    title(slide, "05 WORK 2", "新机器只需小规模标定，就能得到 AP I/O 对 TP latency 的影响")
    slide.shapes.add_picture(str(CHART_SURFACE), Inches(0.68), Inches(1.57), width=Inches(8.10))
    add_text(slide, "设备层", 9.18, 1.84, 1.1, 0.28, size=15, color=BLUE, bold=True)
    add_text(slide, "8KiB TP 同步读\n+ 128KiB AP 随机读\n扫描 qAP=0…32", 9.18, 2.25, 3.0, 0.92, size=13, color=NAVY)
    add_rect(slide, 9.18, 3.41, 2.90, 0.025, MID)
    add_text(slide, "数据库路径", 9.18, 3.75, 1.5, 0.28, size=15, color=ORANGE, bold=True)
    add_text(slide, "BPF 按 openGauss LWTID\n采集 TP 请求 latency\n得到 kpath(qAP)", 9.18, 4.17, 3.0, 0.92, size=13, color=NAVY)
    add_rect(slide, 9.18, 5.31, 2.90, 0.025, MID)
    add_text(slide, "冻结函数", 9.18, 5.58, 1.5, 0.28, size=15, color=GREEN, bold=True)
    add_text(slide, "Lpred = L0 + kpath(qAP) × ΔLdevice(qAP)", 9.18, 5.95, 3.20, 0.42, size=12.5, color=NAVY, bold=True)
    takeaway(slide, "机器变化时重测 surface 与 kpath；负载变化时不需要重新拟合公式结构。")
    footer(slide, 6)
    script.append("工作二第一步：先测设备曲面，再用 BPF 把裸设备延迟迁移到真实 openGauss TP 路径。")


def s7_formula(prs, script):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(slide)
    title(slide, "06 WORK 2", "latency 进入事务响应时间后直接得到 TPS，不需要黑盒 TPS 回归")
    add_text(slide, "基线事务时间", 0.84, 1.90, 2.2, 0.28, size=15, color=GRAY, bold=True)
    add_text(slide, "Rbase = N × 1000 / X0", 0.84, 2.37, 3.0, 0.38, size=22, color=NAVY, bold=True)
    add_arrow(slide, 3.48, 2.56, 4.24, 2.56, BLUE, 2.0)
    add_text(slide, "拆出非 I/O 时间", 4.42, 1.90, 2.2, 0.28, size=15, color=GRAY, bold=True)
    add_text(slide, "Rnonio = Rbase - n0 × L0", 4.42, 2.37, 3.2, 0.38, size=21, color=CYAN, bold=True)
    add_arrow(slide, 7.57, 2.56, 8.33, 2.56, CYAN, 2.0)
    add_text(slide, "加入候选 I/O 等待", 8.50, 1.90, 2.4, 0.28, size=15, color=GRAY, bold=True)
    add_text(slide, "Rpred = Rnonio + n × Lpred", 8.50, 2.37, 3.6, 0.38, size=20, color=ORANGE, bold=True)
    add_rect(slide, 0.84, 3.32, 11.36, 0.025, MID)
    add_text(slide, "TPSpred", 1.32, 4.12, 2.1, 0.48, size=30, color=GREEN, bold=True)
    add_text(slide, "=", 3.18, 4.10, 0.50, 0.50, size=30, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "N × 1000", 4.10, 3.78, 2.25, 0.40, size=24, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_rect(slide, 4.05, 4.27, 2.35, 0.025, NAVY)
    add_text(slide, "Rpred", 4.10, 4.43, 2.25, 0.40, size=24, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    divider(slide, 7.10, 3.70, 1.50)
    add_text(slide, "候选差异如何进入", 7.50, 3.75, 2.3, 0.30, size=15, color=BLUE, bold=True)
    add_text(slide, "SB → X0 / L0 / n0 / n\nwork_mem + plan + AP-cap → qAP\nCPU / 锁 → extra_non_io", 7.50, 4.22, 4.4, 0.98, size=13, color=NAVY)
    add_text(slide, "fitted_tps_coefficient = false", 7.50, 5.48, 3.7, 0.30, size=14, color=RED, bold=True)
    takeaway(slide, "现在，负载链给出的 n/qAP 可以通过机器链变成每个候选的 TPSpred。")
    footer(slide, 7)
    script.append("工作二第二步：利用 closed-loop response-time law，把物理 I/O 等待加入事务时间，TPS 系数固定为 1。")


def s8_bounds(prs, script):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(slide)
    title(slide, "07 WORK 3", "TP 给出 SB 性能下界，AP 给出 SB 内存上界")
    add_text(slide, "TP-first", 0.88, 1.84, 1.5, 0.30, size=16, color=BLUE, bold=True)
    add_text(slide, "至少多少才够", 0.88, 2.20, 2.0, 0.28, size=13, color=GRAY)
    add_text(slide, "BTP_min = 4096MB", 0.88, 2.68, 2.7, 0.36, size=20, color=BLUE, bold=True)
    add_text(slide, "AP-first", 9.62, 1.84, 1.5, 0.30, size=16, color=ORANGE, bold=True)
    add_text(slide, "保留 grant 后最多给多少", 9.62, 2.20, 2.8, 0.28, size=13, color=GRAY)
    add_text(slide, "BAP_max ≈ 6098MB", 9.62, 2.68, 2.8, 0.36, size=20, color=ORANGE, bold=True)
    line_y = 4.12
    add_rect(slide, 1.18, line_y, 10.90, 0.04, MID)
    add_rect(slide, 5.02, line_y - 0.04, 2.70, 0.12, GREEN)
    for x, label in [(1.18, "0"), (5.02, "4096"), (7.72, "6098"), (11.18, "8192")]:
        add_rect(slide, x, line_y - 0.18, 0.025, 0.40, NAVY if label in {"0", "8192"} else (BLUE if label == "4096" else ORANGE))
        add_text(slide, label, x - 0.35, line_y + 0.34, 0.72, 0.26, size=11, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "可行交集", 5.45, 3.54, 1.8, 0.28, size=15, color=GREEN, bold=True, align=PP_ALIGN.CENTER)
    circle(slide, 4.84, line_y - 0.30, 0.38, BLUE, "✓", size=11)
    circle(slide, 11.00, line_y - 0.30, 0.38, RED, "×", size=11)
    add_text(slide, "离散候选 {4096, 8192}MB 中，只有 4096MB 落在交集。", 3.10, 5.18, 7.1, 0.36, size=17, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "若交集为空，则降低 AP work_mem / AP-cap，或阻塞新 AP 后重新计算。", 2.32, 5.72, 8.7, 0.34, size=14, color=GRAY, align=PP_ALIGN.CENTER)
    takeaway(slide, "这解释了 S2 为什么不是继续保留 8192MB SB，而是降低到 4096MB 给 AP。")
    footer(slide, 8)
    script.append("工作三的几何解释：TP-first 给下界，AP-first 给上界，最终候选必须落在交集。")


def s9_two_paths(prs, script):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(slide)
    title(slide, "08 ALGORITHM", "1→2→3 与 2→1→3 从两侧逼近，再由同一 TPS 公式裁决")
    add_text(slide, "TP-first", 0.82, 1.84, 1.5, 0.30, size=17, color=BLUE, bold=True)
    steps_top = [("1", "SB 下界"), ("2", "AP grant"), ("3", "I/O→TPS")]
    for i, (num, label) in enumerate(steps_top):
        x = 2.40 + i * 2.30
        circle(slide, x, 1.73, 0.56, BLUE if i == 0 else (ORANGE if i == 1 else GREEN), num, size=12)
        add_text(slide, label, x - 0.35, 2.42, 1.3, 0.26, size=12, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
        if i < 2: add_arrow(slide, x + 0.65, 2.01, x + 2.13, 2.01, MID, 1.8)
    add_arrow(slide, 8.92, 2.01, 10.12, 2.01, BLUE, 2.0)
    add_text(slide, "候选 A", 10.35, 1.86, 1.4, 0.30, size=17, color=BLUE, bold=True)
    add_rect(slide, 0.82, 3.16, 11.44, 0.02, MID)
    add_text(slide, "AP-first", 0.82, 3.62, 1.5, 0.30, size=17, color=ORANGE, bold=True)
    steps_bottom = [("2", "AP grant"), ("1", "SB 上界"), ("3", "I/O→TPS")]
    for i, (num, label) in enumerate(steps_bottom):
        x = 2.40 + i * 2.30
        circle(slide, x, 3.51, 0.56, ORANGE if i == 0 else (BLUE if i == 1 else GREEN), num, size=12)
        add_text(slide, label, x - 0.35, 4.20, 1.3, 0.26, size=12, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
        if i < 2: add_arrow(slide, x + 0.65, 3.79, x + 2.13, 3.79, MID, 1.8)
    add_arrow(slide, 8.92, 3.79, 10.12, 3.79, ORANGE, 2.0)
    add_text(slide, "候选 B", 10.35, 3.64, 1.4, 0.30, size=17, color=ORANGE, bold=True)
    add_arrow(slide, 10.92, 4.30, 10.92, 5.25, GREEN, 2.0)
    add_text(slide, "先满足 TP SLO，再比较 AP utility 与 predicted TPS", 6.45, 5.38, 5.55, 0.34, size=15, color=GREEN, bold=True, align=PP_ALIGN.CENTER)
    takeaway(slide, "双向搜索不是两套答案：每条路径都生成完整配置，最后只输出一个联合推荐。")
    footer(slide, 9)
    script.append("算法：TP-first 和 AP-first 分别得到完整候选 A/B，然后用同一 SLO、AP utility 和 TPS 排序。")


def s10_stages(prs, script):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(slide)
    title(slide, "09 VALIDATION SCENARIO", "五阶段是同一系统压力逐步演化，而不是五次互不相关的调参")
    xs = [1.04, 3.45, 5.86, 8.27, 10.68]
    colors = [GREEN, ORANGE, ORANGE, RED, BLUE]
    labels = [
        ("S1", "基准 TP\n内存富余"),
        ("S2", "新增 AP\n动态内存上升"),
        ("S3", "AP 压力继续增\nTP 触及容量边界"),
        ("S4", "保留运行 AP\n新 AP 到达"),
        ("S5", "TP demand\n进一步突增"),
    ]
    add_rect(slide, 1.34, 2.64, 9.60, 0.04, MID)
    for x, color, (stage, body) in zip(xs, colors, labels):
        circle(slide, x, 2.30, 0.72, color, stage, size=12)
        add_text(slide, body, x - 0.32, 3.34, 1.38, 0.72, size=12.5, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "AP 压力逐步增加", 2.30, 1.84, 5.10, 0.30, size=15, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
    add_arrow(slide, 2.62, 2.13, 7.40, 2.13, ORANGE, 2.2)
    add_text(slide, "TP 饱和后保护", 6.25, 4.48, 3.45, 0.30, size=15, color=GREEN, bold=True, align=PP_ALIGN.CENTER)
    add_arrow(slide, 6.47, 4.88, 10.82, 4.88, GREEN, 2.2)
    add_text(slide, "硬约束：任一阶段 AP 数、TP 利用率或 I/O 压力不满足定义，则整次运行无效。", 1.16, 5.55, 11.0, 0.36, size=15, color=RED, bold=True, align=PP_ALIGN.CENTER)
    takeaway(slide, "因此，模型是否有效，要看它能否从压力演化中推导出正确动作并保持 TP 稳定。")
    footer(slide, 10)
    script.append("验证场景：S1-S4 承载相同基准 TP TPS；S1/S2 仅表示内存仍有余量，S3 随 AP 增长触及容量边界，S4 阻塞新 AP，S5 TP demand 突增。")


def s11_actions(prs, script):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(slide)
    title(slide, "10 DECISION", "模型只看观测信号，依次推导出让内存、限内存、阻塞和扩 SB")
    data = [["压力变化", "模型判断", "推荐动作"],
            ["S1：基准 TP，单 AP", "TP 内存余量充足，AP grant 可容纳", "保持 SB=8192，Q18=1150"],
            ["S2：新增 AP", "AP grant 与 8192MB SB 冲突", "SB 8192→4096，保留 Q18/Q21=1150"],
            ["S3：TP 达饱和", "优先满足 TP 下界和 SLO", "SB=4096，AP work_mem→256"],
            ["S4：饱和且新 AP 到达", "新增候选无安全交集", "阻塞新 AP，运行中 AP 自然完成"],
            ["S5：TP demand 突增", "原 SB 无法承接新增需求", "SB 4096→8192，保持受控 AP grant"]]
    add_table(slide, data, 0.62, 1.58, 12.08, 4.02, widths=[2.35, 4.25, 5.48], font_size=9.2)
    add_text(slide, "动作不是按阶段名写死：输入是 TP demand/capacity、AP 到达、Query trace、内存预算和冻结机器模型。", 1.08, 5.92, 11.15, 0.38, size=14.5, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    takeaway(slide, "到这里，模型已经给出了与验收要求一致的五阶段配置；下一步必须证明它不是事后拟合。")
    footer(slide, 11)
    script.append("决策结果：S2 让出 SB，S3 降 work_mem，S4 阻塞新 AP，S5 提高 SB。")


def s12_protocol(prs, script):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(slide)
    title(slide, "11 EXPERIMENT", "实验严格按“训练 → 冻结 → 盲测 → 真实验证”执行")
    stages = [
        ("机器训练", "设备 QD0/2/4/8/16/32\n路径 QD6/12/24", BLUE),
        ("冻结模型", "surface / kpath\nX0 / L0 / n0", CYAN),
        ("公式盲测", "先写预测\n再注入 QD9/QD18", ORANGE),
        ("冻结推荐", "只读观测与历史 trace\n输出配置", GREEN),
        ("真实五阶段", "应用推荐后运行\nTPS 仅作比较", RED),
    ]
    add_rect(slide, 1.18, 2.57, 10.25, 0.04, MID)
    for i, (name, body, color) in enumerate(stages):
        x = 1.00 + i * 2.48
        circle(slide, x, 2.22, 0.72, color, str(i + 1), size=12)
        add_text(slide, name, x - 0.38, 3.18, 1.48, 0.28, size=14, color=color, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, body, x - 0.62, 3.63, 1.96, 0.70, size=11, color=NAVY, align=PP_ALIGN.CENTER)
    add_rect(slide, 5.56, 1.80, 0.04, 3.05, RED)
    add_text(slide, "冻结边界", 5.08, 5.04, 1.0, 0.26, size=11, color=RED, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "推荐过程中禁止读取：阶段名、期望动作、候选实际混合 TPS、验证阶段实际 await。", 1.18, 5.58, 10.95, 0.36, size=15, color=RED, bold=True, align=PP_ALIGN.CENTER)
    takeaway(slide, "预测产物时间戳早于干预与真实运行，这是结论不发生数据泄漏的前提。")
    footer(slide, 12)
    script.append("实验隔离：训练与验证按时间分开，真实五阶段结果不回灌模型。")


def s13_formula_result(prs, script):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(slide)
    title(slide, "12 RESULT 1", "未见 QD 下，TP latency 误差 1.53%，端到端 TPS 误差 1.81%")
    slide.shapes.add_picture(str(CHART_FORMULA), Inches(0.60), Inches(1.52), width=Inches(8.25))
    metrics = read_json(FORMULA_REPORT)["metrics"]
    add_text(slide, f"{metrics['pressure_latency_mape_pct']:.2f}%", 9.42, 1.78, 2.55, 0.52, size=31, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "latency MAPE", 9.42, 2.36, 2.55, 0.28, size=13, color=GRAY, bold=True, align=PP_ALIGN.CENTER)
    add_rect(slide, 9.65, 2.88, 2.10, 0.02, MID)
    add_text(slide, f"{metrics['post_intervention_tps_mape_pct']:.2f}%", 9.42, 3.20, 2.55, 0.52, size=31, color=CYAN, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "TPS MAPE", 9.42, 3.78, 2.55, 0.28, size=13, color=GRAY, bold=True, align=PP_ALIGN.CENTER)
    add_rect(slide, 9.65, 4.30, 2.10, 0.02, MID)
    add_text(slide, f"{metrics['tps_conversion_with_measured_latency_mape_pct']:.2f}%", 9.42, 4.62, 2.55, 0.52, size=31, color=GREEN, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "实测 latency→TPS", 9.42, 5.20, 2.55, 0.28, size=13, color=GRAY, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "6/6 个实质压力点方向正确", 8.92, 5.86, 3.55, 0.32, size=14, color=GREEN, bold=True, align=PP_ALIGN.CENTER)
    takeaway(slide, "拆分验证表明：误差主要来自 latency 预测，而 latency→TPS 换算本身误差仅 0.57%。")
    footer(slide, 13)
    script.append("公式结果：未见 QD9/QD18 下 latency 和 TPS 都通过门槛，并且 TPS 换算被单独验证。")


def s14_stage_result(prs, script):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(slide)
    checks = read_json(STAGE_REPORT)["checks"]
    variation = checks["protected_tp_variation_s1_s5_percent"]
    title(slide, "13 RESULT 2", f"推荐动作全部命中，S1-S5 保护 TP TPS 波动 {variation:.2f}%")
    slide.shapes.add_picture(str(CHART_STAGE), Inches(0.58), Inches(1.50), width=Inches(8.65))
    add_text(slide, "动作验证", 9.55, 1.80, 1.6, 0.30, size=16, color=BLUE, bold=True)
    add_text(slide, "S2  降低 SB\nS3  降低 AP work_mem\nS4  阻塞新 AP\nS5  提高 SB", 9.55, 2.28, 2.55, 1.30, size=14, color=NAVY)
    add_rect(slide, 9.55, 3.88, 2.42, 0.02, MID)
    add_text(slide, "执行验证", 9.55, 4.22, 1.6, 0.30, size=16, color=GREEN, bold=True)
    add_text(slide, "所有配置实际应用\n所有 AP 自然结束\n无 AP failure", 9.55, 4.70, 2.55, 0.92, size=14, color=NAVY)
    add_text(slide, f"{variation:.2f}%  <  5%", 9.55, 5.85, 2.55, 0.38, size=22, color=GREEN, bold=True, align=PP_ALIGN.CENTER)
    takeaway(slide, "结果证明模型不仅能选到一个配置，还能给出与压力演化一致的控制动作。")
    footer(slide, 14)
    script.append(f"五阶段结果：动作全部匹配，S1-S5 保护 TPS 波动为 {variation:.2f}%，满足 5% 验收目标。")


def s15_contrib(prs, script):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(slide)
    title(slide, "14 CONTRIBUTION", "我们的贡献，是把 Trace、机器与控制串成可验证、可执行的闭环")
    items = [
        ("01", "从一次执行推演候选", "Plan-aware replay 支持逐 Query work_mem、plan 变化和 spill 估计。", BLUE),
        ("02", "从 I/O 因果推导 TPS", "设备曲面 + 数据库路径 + 响应时间定律，不使用 TPS 拟合乘数。", ORANGE),
        ("03", "从双方约束得到配置", "TP 给 SB 下界，AP 给 SB 上界，交集内联合选择 SB/work_mem/AP-cap。", GREEN),
        ("04", "从预测落到稳定控制", "原版 openGauss 上通过阶段重启、逐 Query grant 和 AP 准入完成五阶段验证。", CYAN),
    ]
    for i, (num, headline, body, color) in enumerate(items):
        y = 1.75 + i * 1.08
        add_text(slide, num, 0.82, y, 0.62, 0.32, size=15, color=color, bold=True)
        add_text(slide, headline, 1.68, y - 0.02, 3.25, 0.32, size=17, color=NAVY, bold=True)
        add_text(slide, body, 5.12, y, 6.95, 0.40, size=13.2, color=NAVY)
        if i < 3: add_rect(slide, 1.68, y + 0.72, 10.42, 0.012, MID)
    variation = read_json(STAGE_REPORT)["checks"]["protected_tp_variation_s1_s5_percent"]
    add_text(slide, f"量化证据：latency MAPE 1.53%  ·  TPS MAPE 1.81%  ·  五阶段保护 TPS 波动 {variation:.2f}%", 1.15, 5.98, 11.0, 0.36, size=16, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
    takeaway(slide, "项目完成了从“解释缓存”到“预测配置并控制 TPS”的能力升级。")
    footer(slide, 15)
    script.append("贡献总结：不是一个单点公式，而是候选生成、机器预测、联合搜索和执行控制四部分闭环。")


def s16_limits(prs, script):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(slide)
    title(slide, "15 LIMITATION", "当前结论有明确边界；下一步扩展新 plan 与多 I/O 类型")
    add_text(slide, "已经严格验证", 0.82, 1.78, 2.2, 0.30, size=17, color=GREEN, bold=True)
    add_text(slide, "• 8KiB buffered TP 读\n• 128KiB random-read AP 压力\n• 当前 terminals 与 QD 有效域\n• 未见 QD latency/TPS holdout", 0.82, 2.35, 3.25, 1.52, size=14, color=NAVY)
    divider(slide, 4.38, 1.78, 3.55)
    add_text(slide, "尚未完全解决", 4.80, 1.78, 2.2, 0.30, size=17, color=RED, bold=True)
    add_text(slide, "• 新 plan 的 spill→qAP 泛化\n• 写型、顺序型、不同块大小 I/O\n• CPU / 锁 / OS cache 跨负载影响\n• 更大规模未见 Query 与机器", 4.80, 2.35, 3.40, 1.52, size=14, color=NAVY)
    divider(slide, 8.55, 1.78, 3.55)
    add_text(slide, "下一步", 8.98, 1.78, 2.2, 0.30, size=17, color=BLUE, bold=True)
    add_text(slide, "• 增加多 I/O class 设备曲面\n• 建立新 plan/spill 未见 Query holdout\n• 联合 CPU/锁与 I/O 响应时间\n• 扩大跨机器、跨 scale factor 验证", 8.98, 2.35, 3.48, 1.52, size=14, color=NAVY)
    add_rect(slide, 0.82, 4.55, 11.42, 0.03, MID)
    add_text(slide, "当前可以主张", 0.82, 4.94, 1.8, 0.30, size=15, color=GREEN, bold=True)
    add_text(slide, "在已记录的机器与 I/O domain 内，模型能产生与五阶段要求一致的动作，并将保护 TPS 波动控制在 5% 内。", 2.52, 4.94, 9.44, 0.52, size=14, color=NAVY)
    add_text(slide, "当前不能主张", 0.82, 5.75, 1.8, 0.30, size=15, color=RED, bold=True)
    add_text(slide, "任意新 SQL、任意 I/O 类型和任意机器无需标定即可达到同样误差；超出 domain 的候选会被拒绝。", 2.52, 5.75, 9.44, 0.52, size=14, color=NAVY)
    takeaway(slide, "明确边界不是削弱结论，而是保证项目结果可复现、可继续扩展。")
    footer(slide, 16)
    script.append("边界：当前公式在明确 domain 内通过验证，新 plan 和更多 I/O 类型仍是下一步重点。")


def conclusion(prs, script):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(slide, NAVY)
    add_rect(slide, 0, 0, 13.333, 0.12, BLUE)
    add_text(slide, "CONCLUSION", 0.82, 0.82, 2.0, 0.28, size=11, color=RGBColor(120, 166, 255), bold=True)
    add_text(slide, "我们建立了一条从配置到 TPS 的完整因果链", 0.82, 1.42, 11.4, 0.58, size=32, color=WHITE, bold=True)
    add_text(slide, "Trace 描述候选后果", 0.88, 2.62, 3.1, 0.36, size=18, color=RGBColor(120, 166, 255), bold=True)
    add_arrow(slide, 3.60, 2.82, 4.28, 2.82, RGBColor(120, 166, 255), 2.0)
    add_text(slide, "机器模型计算 latency/TPS", 4.48, 2.62, 3.5, 0.36, size=18, color=RGBColor(255, 173, 101), bold=True)
    add_arrow(slide, 7.96, 2.82, 8.64, 2.82, RGBColor(255, 173, 101), 2.0)
    add_text(slide, "双向搜索产生控制动作", 8.84, 2.62, 3.4, 0.36, size=18, color=RGBColor(113, 213, 160), bold=True)
    add_rect(slide, 0.84, 3.62, 11.50, 0.025, RGBColor(66, 88, 119))
    add_text(slide, "1.53%", 1.28, 4.18, 2.0, 0.52, size=31, color=RGBColor(120, 166, 255), bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "latency MAPE", 1.28, 4.76, 2.0, 0.28, size=12, color=RGBColor(180, 194, 214), align=PP_ALIGN.CENTER)
    add_text(slide, "1.81%", 5.50, 4.18, 2.0, 0.52, size=31, color=RGBColor(255, 173, 101), bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "TPS MAPE", 5.50, 4.76, 2.0, 0.28, size=12, color=RGBColor(180, 194, 214), align=PP_ALIGN.CENTER)
    variation = read_json(STAGE_REPORT)["checks"]["protected_tp_variation_s1_s5_percent"]
    add_text(slide, f"{variation:.2f}%", 9.72, 4.18, 2.0, 0.52, size=31, color=RGBColor(113, 213, 160), bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "保护 TPS 波动", 9.72, 4.76, 2.0, 0.28, size=12, color=RGBColor(180, 194, 214), align=PP_ALIGN.CENTER)
    add_text(slide, "谢谢", 0.82, 6.34, 1.5, 0.38, size=24, color=WHITE, bold=True)
    add_text(slide, "Q & A", 10.60, 6.36, 1.6, 0.34, size=19, color=RGBColor(164, 180, 202), align=PP_ALIGN.RIGHT)
    script.append("结论：我们完成了 Trace、机器模型和控制算法三部分闭环，并在严格盲测与五阶段实验中达到验收目标。")


def build() -> None:
    make_charts()
    prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    script: list[str] = []
    builders = (cover, s2_problem, s3_gap, s4_solution, s5_replay, s6_machine,
                s7_formula, s8_bounds, s9_two_paths, s10_stages, s11_actions,
                s12_protocol, s13_formula_result, s14_stage_result, s15_contrib,
                s16_limits, conclusion)
    for builder in builders:
        builder(prs, script)
    prs.core_properties.title = "Huawei6 项目答辩：连续叙事版"
    prs.core_properties.subject = "从性能闭环到因果建模、双向搜索和五阶段验证"
    prs.core_properties.author = "Huawei6 project"
    prs.save(OUT); shutil.copy2(OUT, ROOT_COPY)
    SCRIPT.write_text("# Huawei6 项目答辩讲稿（连续叙事版）\n\n" + "\n\n".join(f"## 第 {i} 页\n{text}" for i, text in enumerate(script, 1)) + "\n", encoding="utf-8")
    print(json.dumps({"pptx": str(OUT), "root_copy": str(ROOT_COPY), "slides": len(prs.slides), "script": str(SCRIPT)}, ensure_ascii=False, indent=2))


if __name__ == "__main__":
    build()
