#!/usr/bin/env python3
"""Build a concise Chinese deck explaining the statistical PPT policy formula."""
from __future__ import annotations

import json
from pathlib import Path

import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches

from build_memory_autonomy_paper_ppt import (
    BLUE, GRAY, GREEN, INK, LIGHT, MID, ORANGE, PALE_GREEN, PALE_ORANGE,
    PALE_RED, PALE_TEAL, PURPLE, RED, TEAL, WHITE, add_arrow, add_callout,
    add_flow_box, add_kpi, add_rect, add_table, add_text, add_title, set_bg,
)


ROOT = Path(__file__).resolve().parents[1]
RUN = ROOT / "results/restart_five_stage_complex_ap_20260802"
VALIDATION = RUN / "statistical_decision_validation.json"
AUDIT = RUN / "restart_five_stage_steady_audit.json"
OUT_DIR = ROOT / "artifacts/00_latest"
OUT = OUT_DIR / "Huawei6_统计状态机预测公式_20260802.pptx"
TPS_FIG = OUT_DIR / "huawei6_formula_steady_tps_20260802.png"
MEM_FIG = OUT_DIR / "huawei6_formula_memory_trigger_20260802.png"


def footer(slide, page: int, source: str) -> None:
    add_rect(slide, 0.62, 7.17, 12.08, 0.012, MID)
    add_text(slide, "Huawei6 · 五阶段统计状态机预测公式 · 2026-08-02", 0.62, 7.22, 7.4, 0.18, size=7.5, color=GRAY)
    add_text(slide, f"来源：{source}", 7.15, 7.22, 5.0, 0.18, size=7, color=GRAY, align=PP_ALIGN.RIGHT)
    add_text(slide, str(page), 12.28, 7.22, 0.42, 0.18, size=8, color=GRAY, align=PP_ALIGN.RIGHT)


def add_image(slide, path: Path, x: float, y: float, w: float, h: float) -> None:
    from PIL import Image
    with Image.open(path) as image:
        iw, ih = image.size
    scale = min(w / iw, h / ih)
    slide.shapes.add_picture(
        str(path), Inches(x + (w - iw * scale) / 2), Inches(y + (h - ih * scale) / 2),
        width=Inches(iw * scale), height=Inches(ih * scale),
    )


def build_figures(validation: dict[str, object], audit: dict[str, object]) -> None:
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    stages = audit["stages"]
    assert isinstance(stages, list)
    labels = [str(row["stage"]) for row in stages[2:]]
    protected = [float(row["steady_protected_tp_tps"]) for row in stages[2:]]
    total = [float(row["steady_total_tp_tps"]) for row in stages[2:]]
    x = list(range(len(labels)))
    fig, ax = plt.subplots(figsize=(9.5, 4.5))
    ax.plot(x, protected, marker="o", linewidth=2.8, color="#237db2", label="Protected TP TPS")
    ax.plot(x, total, marker="o", linewidth=2.2, color="#d68127", label="Total TP TPS")
    baseline = sum(protected) / len(protected)
    ax.axhspan(baseline * 0.95, baseline * 1.05, color="#cae8d9", alpha=0.75, label="±5% protected band")
    for index, value in enumerate(protected):
        ax.text(index, value + 20, f"{value:.0f}", ha="center", fontsize=10, weight="bold", color="#237db2")
    for index, value in enumerate(total):
        ax.text(index, value - 80, f"{value:.0f}", ha="center", fontsize=9, color="#a95a13")
    ax.set_xticks(x, labels)
    ax.set_ylabel("TPS")
    ax.set_title("S3-S5: protected TP remains stable while S5 adds surge TPS", weight="bold")
    ax.grid(axis="y", alpha=0.2)
    ax.legend(frameon=False, loc="upper left")
    fig.tight_layout()
    fig.savefig(TPS_FIG, dpi=180)
    plt.close(fig)

    limit = float(validation["memory_target_max_mb"])
    s1 = float(validation["s1_projected_managed_memory_mb"])
    s2 = float(validation["s2_pre_action_projected_managed_memory_mb"])
    post_s2 = 4096 + 3311
    fig, ax = plt.subplots(figsize=(9.5, 4.5))
    names = ["S1: 8GB SB + 1 AP", "S2 pre-action: 8GB SB + 2 AP", "S2 recommendation: 4GB SB + 2 AP"]
    values = [s1, s2, post_s2]
    colors = ["#237db2", "#c44742", "#2f8a63"]
    bars = ax.bar(names, values, color=colors, width=0.58)
    ax.axhline(limit, color="#353c44", linewidth=2, linestyle="--", label=f"memory_target_max = {limit:.0f}MB")
    for bar, value in zip(bars, values):
        ax.text(bar.get_x() + bar.get_width() / 2, value + 210, f"{value:.0f}MB", ha="center", weight="bold", fontsize=11)
    ax.set_ylim(0, 12800)
    ax.set_ylabel("Projected managed memory (MB)")
    ax.set_title("S2 trigger: preserve AP grant only after yielding shared_buffers", weight="bold")
    ax.grid(axis="y", alpha=0.2)
    ax.legend(frameon=False, loc="upper right")
    fig.tight_layout()
    fig.savefig(MEM_FIG, dpi=180)
    plt.close(fig)


def build() -> None:
    validation = json.loads(VALIDATION.read_text(encoding="utf-8"))
    audit = json.loads(AUDIT.read_text(encoding="utf-8"))
    build_figures(validation, audit)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # Cover
    slide = prs.slides.add_slide(blank)
    set_bg(slide, INK)
    add_rect(slide, 0, 0, 0.22, 7.5, TEAL)
    add_text(slide, "Huawei6 五阶段内存调度", 0.90, 1.04, 10.8, 0.45, size=23, color=PALE_TEAL, bold=True)
    add_text(slide, "现在如何预测？预测公式是什么？", 0.88, 1.73, 11.8, 0.86, size=37, color=WHITE, bold=True)
    add_text(slide, "Trace/监控估计资源需求  →  统计状态机选择动作  →  实测 TPS 只做后验验证", 0.92, 2.91, 11.45, 0.42, size=17, color=MID, bold=True)
    add_rect(slide, 0.94, 4.10, 11.46, 0.90, TEAL, radius=True)
    add_text(slide, "当前模型不是“直接回归 TPS 数值”；它预测的是在当前资源压力下应采取的 SB、work_mem 和 AP 准入动作。", 1.18, 4.39, 10.98, 0.33, size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "验证负载：复杂 TPC-H AP + 饱和 TP + S5 TP 突增", 0.95, 6.48, 6.2, 0.24, size=10.5, color=MID)

    # Inputs / outputs
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, "1. 预测对象：从运行统计和 Trace 推导“下一步动作”", "它不把当前实际混合 TPS 当成推荐标签；TPS 在决策完成后才用于验收。", "01 · INPUT / OUTPUT")
    add_flow_box(slide, "输入 A：Trace Replay", "算子/Plan\n每条 AP 的动态峰值\nspill 与 I/O 代价", 0.62, 1.65, 2.64, 1.72, PURPLE, fill=LIGHT)
    add_flow_box(slide, "输入 B：在线统计", "当前 SB\n运行/新到达 AP 数\nTP offered rate", 3.54, 1.65, 2.64, 1.72, TEAL, fill=PALE_TEAL)
    add_flow_box(slide, "状态机", "容量超限？\nTP 饱和？\nAP 是否继续到达？\nTP 是否突增？", 6.48, 1.65, 2.64, 1.72, BLUE, fill=LIGHT)
    add_flow_box(slide, "输出配置", "shared_buffers\nAP work_mem\nAP cap\nblock_new_AP", 9.42, 1.65, 2.64, 1.72, ORANGE, fill=PALE_ORANGE)
    for x in (3.28, 6.22, 9.16):
        add_arrow(slide, x, 2.52, x + 0.20, 2.52, color=GRAY)
    data = [
        ["符号", "含义", "来源"],
        ["SB", "当前 shared_buffers", "数据库参数/阶段配置"],
        ["D_hat", "当前 AP 集合预测动态内存需求", "AP trace replay + 当前并发"],
        ["N_run, N_in", "正在运行 AP、刚到达 AP", "AP 调度器"],
        ["T_offer, T_protect", "总 TP 请求速率、需保护 TP 基线", "TP 发生器/容量标定"],
        ["M_max", "统一受控内存上限", "部署策略，当前为 10500MB"],
    ]
    add_table(slide, data, 0.74, 4.20, 11.84, 2.10, widths=[1.35, 5.55, 4.94], font_size=10.5)
    footer(slide, 2, "statistical_ppt_state_machine.py")

    # Formula
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, "2. 核心预测公式：先计算容量，再判断状态转移", "公式中的 D_hat 是由 Trace/在线统计提供的需求估计；公式本身不使用实际混合 TPS。", "02 · FORMULA")
    add_rect(slide, 0.82, 1.52, 11.70, 0.72, INK, radius=True)
    add_text(slide, "M_current = SB_current + D_hat          M_candidate = SB_candidate + D_hat", 1.10, 1.74, 11.12, 0.28, size=21, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_callout(slide, "S2：让渡 SB", "N_in > 0 ∧ M_current > M_max ∧ (SB_floor + D_hat ≤ M_max)\n→ SB = 4096MB，保持高 work_mem", 0.76, 2.70, 5.70, 1.35, PURPLE, PALE_TEAL)
    add_callout(slide, "S3：保护 TP", "T_protect ≥ T_sat ∧ N_in > 0\n→ 停止降低 SB；work_mem = 256MB", 6.86, 2.70, 5.70, 1.35, BLUE, LIGHT)
    add_callout(slide, "S4：反压队列", "state = protect_tp ∧ N_run ≥ AP_cap ∧ N_in > 0\n→ block_new_AP = true；存量 AP 不取消", 0.76, 4.35, 5.70, 1.35, ORANGE, PALE_ORANGE)
    add_callout(slide, "S5：TP 突增反向恢复", "T_offer ≥ T_protect + ΔT，当前 ΔT = 300 TPS\n→ SB = 8192MB；work_mem = 256MB；AP_cap = 2", 6.86, 4.35, 5.70, 1.35, GREEN, PALE_GREEN)
    add_text(slide, "S1 是剩余默认状态：M_current ≤ M_max，AP 不增长到压力边界，保持 SB=8192MB、work_mem=1150MB。", 0.88, 6.34, 11.62, 0.33, size=14, color=INK, bold=True, align=PP_ALIGN.CENTER)
    footer(slide, 3, "当前状态机的实际 if 条件")

    # From trace to demand estimate
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, "3. D_hat 如何得到：按 AP 查询/算子聚合，而不是拍脑袋设置", "当前阶段验证采用动作前探针作为 D_hat；生产路径应由多锚点 Trace Replay 按当前 Plan/并发估计。", "03 · TRACE TO DEMAND")
    add_rect(slide, 0.72, 1.50, 12.00, 0.82, PURPLE, radius=True)
    add_text(slide, "D_hat(stage) = D_base + Σ [ active(q) × peak_dynamic_hat(q, plan(q), work_mem(q)) ]", 0.95, 1.76, 11.54, 0.32, size=19, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_flow_box(slide, "Plan", "EXPLAIN / Plan family\nHash Join / Sort / Agg", 0.78, 3.04, 2.40, 1.45, PURPLE, fill=LIGHT)
    add_arrow(slide, 3.28, 3.76, 3.83, 3.76, color=GRAY)
    add_flow_box(slide, "Replay anchor", "同 Plan trace\n算子峰值与 spill\n或源码结构估计", 3.94, 3.04, 2.40, 1.45, TEAL, fill=PALE_TEAL)
    add_arrow(slide, 6.44, 3.76, 6.99, 3.76, color=GRAY)
    add_flow_box(slide, "并发叠加", "按运行生命周期\n而非简单 work_mem 相加\n得到 D_hat", 7.10, 3.04, 2.40, 1.45, BLUE, fill=LIGHT)
    add_arrow(slide, 9.60, 3.76, 10.15, 3.76, color=GRAY)
    add_flow_box(slide, "容量判断", "与 SB / M_max 组合\n生成 S2-S5 动作", 10.26, 3.04, 2.40, 1.45, ORANGE, fill=PALE_ORANGE)
    add_callout(slide, "当前实验中的 D_hat", "S2 在 8GB SB 下先运行 Q18+Q21 动作前探针，实测峰值 3311MB，作为本次 D_hat。这个探针结果不含混合 TPS。", 0.78, 5.30, 11.82, 0.86, TEAL, PALE_TEAL)
    footer(slide, 4, "trace replay + restart_complex_ap_s2_predecision_20260802")

    # Numerical S2 trigger
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, "4. S2 的实际计算：为什么模型会建议降低 SB", "这一页展示的就是容量公式在本次复杂 AP 负载上的真实输入和输出。", "04 · S2 EXAMPLE")
    add_image(slide, MEM_FIG, 0.56, 1.45, 7.10, 4.90)
    add_callout(slide, "S1：仍可容纳", "8192 + 2097 = 10289MB ≤ 10500MB\n保持 SB=8192MB、work_mem=1150MB", 7.90, 1.72, 4.58, 1.25, BLUE, LIGHT)
    add_callout(slide, "S2 动作前：超限", "8192 + 3311 = 11503MB > 10500MB\nAP 继续到达，触发让渡 SB", 7.90, 3.20, 4.58, 1.25, RED, PALE_RED)
    add_callout(slide, "推荐后：重新可行", "4096 + 3311 = 7407MB ≤ 10500MB\n所以不需要先压低 AP work_mem", 7.90, 4.68, 4.58, 1.25, GREEN, PALE_GREEN)
    footer(slide, 5, "S1/S2 复杂 AP 动作前探针")

    # Five decision trace
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, "5. 五阶段的实际判定与输出配置", "每一行先给出决策时观察到的信号，再给出状态机输出；输出与 PPT 的动作顺序一致。", "05 · DECISION TRACE")
    data = [
        ["状态", "触发观测", "预测输出"],
        ["S1", "1 AP；M=10289≤10500；TP=700", "SB 8192 · work_mem 1150 · AP cap 1"],
        ["S2", "AP 新到达；M_pre=11503>10500", "SB 4096 · work_mem 1150 · AP cap 2"],
        ["S3", "TP protect=4000≥T_sat；AP 继续到达", "SB 4096 · work_mem 256 · AP cap 4"],
        ["S4", "N_run=4=AP cap；又有新 AP", "SB 4096 · work_mem 256 · block new AP"],
        ["S5", "T_offer=4300=4000+300", "SB 8192 · work_mem 256 · AP cap 2 · block"],
    ]
    add_table(slide, data, 0.72, 1.54, 11.92, 3.05, widths=[1.28, 5.44, 5.20], font_size=11.2)
    add_rect(slide, 0.78, 5.05, 11.76, 1.04, INK, radius=True)
    add_text(slide, "关键点：S2 由内存公式驱动；S3/S4 由 TP 饱和与 AP 到达驱动；S5 由 TP 请求增量驱动。五者均不读取实测混合 TPS。", 1.03, 5.37, 11.24, 0.41, size=15.5, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    footer(slide, 6, "statistical_decision_validation.json")

    # Validation
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, "6. 推荐后验证：保护 TP 稳定，S5 增加的是额外 TP 需求", "S5 的总 TPS 上升并不意味着保护目标失效，应看持续存在的 4000 TPS 保护流。", "06 · VALIDATION")
    add_image(slide, TPS_FIG, 0.54, 1.42, 7.34, 4.90)
    add_kpi(slide, "0.825%", "S3-S5 保护 TP TPS 波动", 8.12, 1.72, 4.18, color=GREEN, note="门槛：≤5%")
    add_kpi(slide, "5 / 5", "阶段推荐与 PPT 一致", 8.12, 3.22, 4.18, color=TEAL, note="无阶段标签输入")
    add_kpi(slide, "0", "AP 取消 / 失败", 8.12, 4.72, 4.18, color=BLUE, note="全部自然完成")
    footer(slide, 7, "restart_five_stage_steady_audit.json")

    # Boundary
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, "7. 结论与边界：公式能解释动作，但不是万能 TPS 回归器", "当前结果已经证明给定五阶段负载下，统计信号可以生成与 PPT 一致的配置推荐。", "07 · CONCLUSION")
    add_callout(slide, "已经验证", "复杂 AP + 饱和 TP 下，S1→S5 动作、推荐、自然完成和 TP 稳定性均通过。", 0.76, 1.60, 5.74, 1.38, GREEN, PALE_GREEN)
    add_callout(slide, "公式的含义", "D_hat 决定内存是否需要让渡；TP 饱和与 AP 到达决定是否限额/排队；TP 增量决定是否反向恢复 SB。", 6.82, 1.60, 5.74, 1.38, TEAL, PALE_TEAL)
    add_callout(slide, "当前边界", "原版 openGauss 不能在线修改 shared_buffers。因此 SB 动作在阶段边界重启应用；运行中 AP work_mem 不被伪装成可热缩。", 0.76, 3.42, 5.74, 1.62, ORANGE, PALE_ORANGE)
    add_callout(slide, "下一步工程化", "将 D_hat 从动作前探针替换为持续 Trace Replay 估计，并把状态机接入实时 AP 调度器与监控流。", 6.82, 3.42, 5.74, 1.62, BLUE, LIGHT)
    add_rect(slide, 0.78, 5.68, 11.76, 0.78, INK, radius=True)
    add_text(slide, "一句话：模型预测“资源压力下最合适的控制动作”，而不是用已知 TPS 反推一个看似正确的配置。", 1.05, 5.93, 11.22, 0.29, size=15, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    footer(slide, 8, "Huawei6 final validation")

    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    build()
