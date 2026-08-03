#!/usr/bin/env python3
"""Build a defense-style Huawei6 project presentation."""

from __future__ import annotations

import json
import shutil
from pathlib import Path

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches

from build_memory_autonomy_paper_ppt import (
    BLUE,
    GRAY,
    GREEN,
    INK,
    MID,
    ORANGE,
    PALE_GREEN,
    PALE_ORANGE,
    PALE_RED,
    PALE_TEAL,
    RED,
    TEAL,
    WHITE,
    add_arrow,
    add_callout,
    add_flow_box,
    add_kpi,
    add_rect,
    add_table,
    add_text,
    add_title,
    set_bg,
)
from build_zero_to_joint_recommendation_ppt import (
    CHART_FORMULA,
    CHART_STAGE,
    CHART_SURFACE,
    FORMULA_REPORT,
    STAGE_REPORT,
    make_charts,
)


ROOT = Path(__file__).resolve().parents[1]
LATEST = ROOT / "artifacts/00_latest"
OUT = LATEST / "Huawei6_项目答辩_混合负载内存联合优化_20260802.pptx"
ROOT_COPY = Path("/root/Huawei6_项目答辩_混合负载内存联合优化.pptx")
OUTLINE = LATEST / "Huawei6_项目答辩_混合负载内存联合优化_20260802_讲稿.md"


def read_json(path: Path):
    return json.loads(path.read_text(encoding="utf-8"))


def footer(slide, page: int, source: str = "") -> None:
    add_rect(slide, 0.62, 7.17, 12.08, 0.012, MID)
    add_text(slide, "Huawei6 · 混合负载内存联合预测与 TP 稳定控制 · 项目答辩", 0.62, 7.22, 6.6, 0.18, size=7.4, color=GRAY)
    if source:
        add_text(slide, source, 7.0, 7.22, 5.1, 0.18, size=6.8, color=GRAY, align=PP_ALIGN.RIGHT)
    add_text(slide, str(page), 12.28, 7.22, 0.42, 0.18, size=8, color=GRAY, align=PP_ALIGN.RIGHT)


def badge(slide, number: str, label: str, x: float, y: float, color=TEAL) -> None:
    add_rect(slide, x, y, 0.48, 0.48, color, radius=True)
    add_text(slide, number, x, y + 0.08, 0.48, 0.22, size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, label, x + 0.62, y + 0.08, 2.70, 0.28, size=13, color=INK, bold=True)


def cover(prs, notes):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(slide, INK)
    add_rect(slide, 0, 0, 0.22, 7.5, TEAL)
    add_text(slide, "项目答辩", 0.92, 0.82, 3.0, 0.34, size=17, color=TEAL, bold=True)
    add_text(slide, "面向 TP/AP 混合负载的\n内存联合预测与 TP 稳定控制", 0.90, 1.35, 11.5, 1.45, size=35, color=WHITE, bold=True)
    add_text(slide, "Trace Replay  ·  机器标定因果模型  ·  TP-first / AP-first 双向约束搜索", 0.94, 3.25, 11.1, 0.42, size=17, color=WHITE)
    add_rect(slide, 0.94, 4.25, 8.35, 0.60, TEAL, radius=True)
    add_text(slide, "目标：在 AP 压力变化和 TP 突增下，将保护 TP TPS 波动控制在 5% 内", 1.05, 4.42, 8.12, 0.26, size=13.2, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_callout(slide, "答辩主线", "问题定义 → 三项核心工作 → 严格验证 → 五阶段结果 → 工程落地与边界", 0.94, 5.47, 9.20, 0.78, TEAL, PALE_TEAL)
    add_text(slide, "Huawei6 项目组 · 2026-08-02", 0.95, 6.65, 4.3, 0.25, size=11, color=GRAY)
    notes.append("开场：本项目不是单纯预测命中率，而是面向 TP/AP 竞争场景，联合推荐 SB、逐 Query work_mem 和 AP 准入。")


def background(prs, notes):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(slide)
    add_title(slide, "研究背景：HTAP 混合负载下，内存调优形成闭环", "TP 关注低延迟与稳定 TPS；AP 需要算子内存，内存不足会 spill 并反向争抢 I/O。", "01 背景")
    add_flow_box(slide, "TP：短事务", "依赖 shared_buffers 与 OS cache\nSB 小 → page miss / 物理 I/O 增加\nI/O await 直接拉长事务响应时间", 0.72, 1.78, 3.55, 1.86, BLUE, fill=PALE_TEAL)
    add_flow_box(slide, "AP：复杂查询", "Join / Sort / Agg 消耗 work_mem\nwork_mem 小 → spill / writeback 增加\nAP 并发进一步放大队列压力", 9.06, 1.78, 3.55, 1.86, ORANGE, fill=PALE_ORANGE)
    add_arrow(slide, 4.35, 2.70, 5.14, 2.70, BLUE, 2.1)
    add_arrow(slide, 8.98, 2.70, 8.19, 2.70, ORANGE, 2.1)
    add_flow_box(slide, "共享资源冲突", "总内存固定\nSB 与 AP 动态内存互相挤占\nTP miss I/O 与 AP spill I/O\n共同进入 NVMe 排队", 5.25, 1.65, 2.82, 2.14, RED, fill=PALE_RED)
    add_arrow(slide, 6.66, 3.90, 6.66, 4.56, RED, 2.1)
    add_flow_box(slide, "最终表现：相同命中率，也可能得到不同 TPS", "只预测 hit ratio 无法解释 AP spill 对 TP 物理 I/O latency 的影响。", 2.10, 4.67, 9.12, 1.18, TEAL, fill=PALE_GREEN)
    add_callout(slide, "核心矛盾", "SB 过大压缩 AP 内存；SB 过小增加 TP miss。最优点必须同时考虑内存可行性、AP spill、I/O 排队和 TP TPS。", 1.12, 6.14, 11.08, 0.58, TEAL, PALE_TEAL)
    footer(slide, 2)
    notes.append("背景页：强调两个方向的反馈，解释为什么 SB 和 work_mem 不能独立调优。")


def objective(prs, notes):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(slide)
    add_title(slide, "项目目标：输出可执行配置，并在负载变化中保护 TP", "使用原版 openGauss；SB 可在阶段自然结束后重启生效，已运行 AP 不被强制取消。", "02 目标")
    add_flow_box(slide, "配置输出", "阶段级 shared_buffers\n逐 Query work_mem\nAP-cap / 是否阻塞新 AP", 0.72, 1.82, 3.45, 1.68, BLUE, fill=PALE_TEAL)
    add_flow_box(slide, "性能目标", "AP 与 TP 交织运行\nTP 饱和后保护 TPS\n阶段间相对波动 < 5%", 4.94, 1.82, 3.45, 1.68, GREEN, fill=PALE_GREEN)
    add_flow_box(slide, "验收动作", "S2 降低 SB 给 AP\nS3 限制 AP work_mem\nS4 阻塞新 AP\nS5 TP 突增后提高 SB", 9.16, 1.82, 3.45, 1.68, ORANGE, fill=PALE_ORANGE)
    add_callout(slide, "五阶段负载硬约束", "S1→S3 AP 数量、动态内存和 I/O 压力逐步增加；S3 TP 达真实饱和；S4 保留运行 AP 并阻塞新 AP；S5 在饱和基础上增加 TP demand。", 0.94, 4.28, 11.52, 0.95, TEAL, PALE_TEAL)
    add_callout(slide, "成功标准", "推荐必须来自采集信号、Trace Replay 和冻结公式；实际混合 TPS 只能在推荐写出后用于独立验证。", 0.94, 5.58, 11.52, 0.78, RED, PALE_RED)
    footer(slide, 3)
    notes.append("目标页：先给验收标准，后续所有方法和实验都围绕这三个输出与四个动作展开。")


def challenges(prs, notes):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(slide)
    add_title(slide, "三个关键难点，对应本项目的三项核心工作", "答辩重点不在脚本数量，而在三条因果链是否被建立并验证。", "03 难点与工作")
    blocks = [
        ("难点 1", "配置变化后，\nAP 会换 plan、何时 spill？", "工作一：Plan-aware\noperator Trace Replay", BLUE, PALE_TEAL),
        ("难点 2", "TP miss 与 AP spill 争抢 I/O，\n如何反哺 TPS？", "工作二：机器标定的\nI/O latency→TPS 因果模型", ORANGE, PALE_ORANGE),
        ("难点 3", "SB 与 work_mem 强耦合，\n如何避免组合爆炸？", "工作三：TP/AP 双向约束搜索\n1→2→3 与 2→1→3", GREEN, PALE_GREEN),
    ]
    for index, (tag, question, answer, color, fill) in enumerate(blocks):
        x = 0.64 + index * 4.24
        add_text(slide, tag, x + 1.30, 1.66, 1.0, 0.25, size=10, color=color, bold=True, align=PP_ALIGN.CENTER)
        add_flow_box(slide, question, answer, x, 2.05, 3.72, 2.18, color, fill=fill)
    add_callout(slide, "统一原则", "不以五阶段名称或期望动作作为模型输入；候选配置的真实 TPS 不进入推荐过程。", 1.18, 4.83, 10.96, 0.76, RED, PALE_RED)
    add_text(slide, "Trace 描述候选后果  +  机器模型计算性能代价  +  双向搜索产生控制动作", 1.10, 6.10, 11.2, 0.37, size=17, color=TEAL, bold=True, align=PP_ALIGN.CENTER)
    footer(slide, 4)
    notes.append("贡献预告：三项工作分别解决 plan/spill、I/O/TPS 和联合搜索。")


def architecture(prs, notes):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(slide)
    add_title(slide, "总体技术路线：离线建模、机器标定、在线/阶段级决策", "将负载规律和机器规律分开，避免把某台机器、某组 Query 的经验系数当成通用模型。", "04 总体方案")
    layers = [
        ("负载特征层", "TP cache trace\nAP plan/operator trace\nQuery 到达与并发", BLUE, PALE_TEAL),
        ("Replay 层", "SB→TP I/O/tx\nwork_mem→plan/dynamic/spill\nAP mix→qAP / utility", ORANGE, PALE_ORANGE),
        ("机器因果层", "device surface(qAP)\nkpath(qAP)\nresponse time→TPS", TEAL, PALE_GREEN),
        ("联合决策层", "TP-first 1→2→3\nAP-first 2→1→3\nSLO / admission / restart", GREEN, PALE_GREEN),
    ]
    for index, (title, body, color, fill) in enumerate(layers):
        x = 0.48 + index * 3.22
        add_flow_box(slide, title, body, x, 1.82, 2.75, 2.02, color, fill=fill)
        if index < 3: add_arrow(slide, x + 2.78, 2.84, x + 3.15, 2.84, TEAL, 2.0)
    add_callout(slide, "输出", "每个控制窗输出完整候选：(SB，逐 Query work_mem，AP-cap，block_new_ap，predicted TPS / await，决策原因)。", 0.82, 4.48, 11.72, 0.78, TEAL, PALE_TEAL)
    add_callout(slide, "验证闭环", "先冻结模型与推荐 → 重启应用 SB / 设置 session work_mem → 运行真实阶段 → 仅比较动作和 TPS，不回灌本次结果。", 0.82, 5.57, 11.72, 0.82, RED, PALE_RED)
    footer(slide, 5)
    notes.append("总体方案：负载 replay 和机器标定分层，最后由双向搜索连接。")


def work_replay(prs, notes):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(slide)
    add_title(slide, "核心工作一：Plan-aware 算子内存与 spill Replay", "目标是在不真实运行全部配置的前提下，预测不同 work_mem 下的 plan、动态峰值和 spill。", "05 核心工作一")
    add_flow_box(slide, "一次负载采集", "SQL / EXPLAIN plan\n算子输入规模与执行路径\nSort / HashJoin / HashAgg trace\nTP 访问轨迹", 0.66, 1.74, 3.20, 2.15, BLUE, fill=PALE_TEAL)
    add_arrow(slide, 3.92, 2.82, 4.66, 2.82, TEAL, 2.1)
    add_flow_box(slide, "源码与算子模型重放", "按 Query × work_mem 枚举 plan family\n估计 dynamic_peak(q,w)\n估计 spill_io(q,w)\n输出 plan_confidence", 4.75, 1.74, 3.72, 2.15, ORANGE, fill=PALE_ORANGE)
    add_arrow(slide, 8.53, 2.82, 9.27, 2.82, TEAL, 2.1)
    add_flow_box(slide, "候选特征", "逐 Query work_mem\nplan_supported / memory_safe\nAP utility 与 qAP\nTP n_candidate(B)", 9.36, 1.74, 3.30, 2.15, GREEN, fill=PALE_GREEN)
    data = [["Query", "work_mem", "Replay 输出", "进入联合搜索"],
            ["Q18", "1150MB", "较高动态峰值、低 spill", "AP 富裕阶段"],
            ["Q18", "256MB", "较低动态峰值、spill 增加", "TP 饱和保护"],
            ["Q9/Q13", "256MB", "plan 支持、内存安全", "S3/S4 AP grant"]]
    add_table(slide, data, 1.02, 4.48, 11.28, 1.52, widths=[1.55, 1.55, 4.38, 3.80], font_size=9.4)
    add_callout(slide, "意义", "使用算子行为和 Trace 推导配置，不把 Query ID 对应的最佳 work_mem 写死。", 1.44, 6.08, 10.42, 0.66, TEAL, PALE_TEAL)
    footer(slide, 6, "source_plan_replay / hash_join / sort / hash_agg replay")
    notes.append("工作一：说明输入输出和逐 Query work_mem，不展开所有算子公式。")


def work_causal_model(prs, notes):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(slide)
    add_title(slide, "核心工作二：建立可迁移的 I/O latency→TPS 因果链", "新机器先测设备曲面，再用 BPF 锚定 openGauss 路径；TPS 不使用额外拟合系数。", "06 核心工作二")
    slide.shapes.add_picture(str(CHART_SURFACE), Inches(0.60), Inches(1.50), width=Inches(6.70))
    add_callout(slide, "① 设备曲面", "ΔLdevice(qAP)：AP 队列深度对 TP 8KiB 请求额外延迟的影响。", 7.62, 1.58, 4.90, 0.82, BLUE, PALE_TEAL)
    add_callout(slide, "② 数据库路径迁移", "Lpred = L0 + kpath(qAP) × ΔLdevice(qAP)", 7.62, 2.64, 4.90, 0.82, ORANGE, PALE_ORANGE)
    add_callout(slide, "③ 响应时间换算", "Rnonio = N×1000/X0 - n0×L0\nRpred = Rnonio + n×Lpred + extra_non_io", 7.62, 3.70, 4.90, 1.02, TEAL, PALE_GREEN)
    add_callout(slide, "④ TPS", "TPSpred = N×1000/Rpred；fitted_tps_coefficient=false", 7.62, 4.96, 4.90, 0.82, GREEN, PALE_GREEN)
    add_callout(slide, "可迁移性", "机器变化时重测 surface/kpath；负载变化时更新 n、qAP 和非 I/O 代价。", 7.62, 5.88, 4.90, 0.82, RED, PALE_RED)
    footer(slide, 7, "portable_joint_model.py")
    notes.append("工作二：强调因果分解和新机器重标定，不是看到最终 TPS 后做回归。")


def work_search(prs, notes):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(slide)
    add_title(slide, "核心工作三：TP-first / AP-first 双向约束搜索", "两条路线分别从 TP 与 AP 出发，最终在 SB 上下界交集内用同一 TPS 公式选择。", "07 核心工作三")
    add_flow_box(slide, "TP-first：1→2→3", "1. TP replay 得到 SB 性能下界 B_TP_min\n2. 在剩余内存选择逐 Query AP grant\n3. I/O→TPS 公式复核", 0.66, 1.72, 3.45, 2.12, BLUE, fill=PALE_TEAL)
    add_flow_box(slide, "联合可行域", "B_TP_min ≤ SB ≤ B_AP_max\n\n交集非空：比较 SLO、AP utility、TPS\n交集为空：降低 work_mem/AP-cap\n或阻塞新 AP 后重算", 4.94, 1.58, 3.45, 2.40, TEAL, fill=PALE_GREEN)
    add_flow_box(slide, "AP-first：2→1→3", "2. AP replay 先选合适 grant\n1. 得到 SB 内存上界 B_AP_max\n3. I/O→TPS 公式复核", 9.22, 1.72, 3.45, 2.12, ORANGE, fill=PALE_ORANGE)
    add_arrow(slide, 4.17, 2.78, 4.86, 2.78, BLUE, 2.2); add_arrow(slide, 9.15, 2.78, 8.46, 2.78, ORANGE, 2.2)
    add_callout(slide, "S2 实例", "候选 SB={4096,8192}MB；TP 下界=4096MB。Q18/Q21=1150MB grant 下 B_AP_max≈6098MB，因此离散交集仅有 4096MB，模型选择降低 SB 给 AP。", 0.91, 4.54, 11.50, 0.96, TEAL, PALE_TEAL)
    add_callout(slide, "为什么比全网格更有效", "两条路径先确定性能下界和内存上界，再在交集内计算昂贵的 I/O/TPS，减少无意义组合，同时保留双方视角。", 0.91, 5.80, 11.50, 0.82, GREEN, PALE_GREEN)
    footer(slide, 8)
    notes.append("工作三：明确下界、上界和交集，这是双向搜索最容易被问到的地方。")


def controller(prs, notes):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(slide)
    add_title(slide, "从推荐到控制：观测信号驱动阶段动作", "状态机不读取 S1-S5 标签，而是根据 TP demand、AP 到达、队列和 replay 内存需求推导动作。", "08 控制策略")
    data = [["观测状态", "判断依据", "模型动作"],
            ["TP 低负载，无新增 AP", "TP demand / capacity < 0.70", "保留富裕 SB 与 AP grant"],
            ["TP 低负载，新增 AP", "AP grant 在当前 SB 不可容纳", "降低 SB，为 AP 动态内存让空间"],
            ["TP 达饱和", "TP demand / capacity ≥ 0.70", "保留 SB，降低后续 AP work_mem"],
            ["饱和且新 AP 到达", "内存 / I/O 无剩余安全空间", "阻塞新 AP，运行中 AP 自然完成"],
            ["TP demand 突增", "demand > 1.05×protected", "提高 SB，保留受控 AP grant"]]
    add_table(slide, data, 0.68, 1.55, 11.98, 3.86, widths=[3.00, 4.25, 4.73], font_size=9.3)
    add_callout(slide, "原版 openGauss 的执行方式", "阶段自然结束 → 应用推荐 shared_buffers → 重启 openGauss → warmup → 设置各 AP session work_mem / 准入规则 → 运行下一阶段。", 0.88, 5.72, 11.58, 0.88, BLUE, PALE_TEAL)
    footer(slide, 9, "observation-driven joint controller")
    notes.append("控制页：将模型输出对应到可执行动作，强调不在线修改 SB、不取消 AP。")


def protocol(prs, notes):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(slide)
    add_title(slide, "实验设计：先冻结、后验证，防止“看到结果再预测”", "将机器训练、公式盲测、配置推荐和五阶段真实运行严格按时间隔离。", "09 实验方法")
    steps = [
        ("机器训练", "设备 QD0/2/4/8/16/32\n数据库路径 QD6/12/24", BLUE, PALE_TEAL),
        ("冻结模型", "写入 surface、kpath、X0/L0/n0\n记录输入文件哈希", TEAL, PALE_GREEN),
        ("公式盲测", "先写预测\n再注入未见 QD9/QD18", ORANGE, PALE_ORANGE),
        ("冻结推荐", "只读取机器观测与历史 trace\n输出 SB/work_mem/AP-cap", GREEN, PALE_GREEN),
        ("真实五阶段", "按推荐重启并执行\n实际 TPS 只做最终比较", RED, PALE_RED),
    ]
    for index, (title, body, color, fill) in enumerate(steps):
        x = 0.38 + index * 2.58
        add_flow_box(slide, title, body, x, 1.78, 2.18, 1.82, color, fill=fill)
        if index < 4: add_arrow(slide, x + 2.20, 2.68, x + 2.52, 2.68, TEAL, 1.8)
    add_rect(slide, 5.33, 1.54, 0.04, 2.35, RED)
    add_text(slide, "冻结边界", 4.73, 4.02, 1.25, 0.28, size=11, color=RED, bold=True, align=PP_ALIGN.CENTER)
    add_callout(slide, "不允许进入推荐的字段", "阶段名、期望动作、候选配置的实际混合 TPS、验证阶段的实际 await。", 0.88, 4.57, 5.62, 0.86, RED, PALE_RED)
    add_callout(slide, "允许进入推荐的字段", "TP terminals / demand、SB、活跃与到达 Query、历史 plan/operator trace、TP cache replay、机器冻结模型。", 6.84, 4.57, 5.62, 0.86, GREEN, PALE_GREEN)
    add_text(slide, "结论成立的前提：预测产物时间戳早于干预与真实阶段运行。", 1.18, 6.20, 11.0, 0.34, size=16, color=INK, bold=True, align=PP_ALIGN.CENTER)
    footer(slide, 10)
    notes.append("实验方法页：这是回应过拟合质疑的关键页，强调时间顺序和冻结边界。")


def formula_result(prs, notes):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(slide)
    add_title(slide, "公式验证结果：未见 I/O 压力下，latency 与 TPS 均通过门槛", "QD9/QD18 不参与冻结；预测文件在压力注入之前生成。", "10 实验结果一")
    slide.shapes.add_picture(str(CHART_FORMULA), Inches(0.58), Inches(1.48), width=Inches(7.60))
    metrics = read_json(FORMULA_REPORT)["metrics"]
    add_kpi(slide, f"{metrics['pressure_latency_mape_pct']:.2f}%", "TP latency MAPE", 8.55, 1.64, 3.55, color=BLUE, note="未见 QD9/QD18")
    add_kpi(slide, f"{metrics['post_intervention_tps_mape_pct']:.2f}%", "端到端 TPS MAPE", 8.55, 3.02, 3.55, color=TEAL, note="门槛 ≤5%")
    add_kpi(slide, f"{metrics['tps_conversion_with_measured_latency_mape_pct']:.2f}%", "实测 latency→TPS MAPE", 8.55, 4.40, 3.55, color=GREEN, note="独立验证 TPS 换算")
    add_callout(slide, "方向正确", "6/6 个实质压力点，TPS 升降方向全部预测正确。", 8.48, 5.80, 3.75, 0.76, GREEN, PALE_GREEN)
    footer(slide, 11, "strict online pre-intervention holdout")
    notes.append("结果一：用三组指标拆分验证 latency 模型、TPS 换算和端到端误差。")


def stage_actions(prs, notes):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(slide)
    add_title(slide, "五阶段推荐结果：模型推导出与验收一致的动作序列", "推荐阶段不读取实际混合 TPS；下表是推荐冻结后与真实阶段对齐展示。", "11 实验结果二")
    report = read_json(STAGE_REPORT)
    descriptions = ["低 TP，单 AP", "新增 AP", "TP 达饱和", "饱和且新 AP 到达", "TP demand 突增"]
    actions = ["保持富裕", "SB 8192→4096", "AP work_mem→256", "阻塞新 AP", "SB 4096→8192"]
    data = [["阶段", "负载变化", "推荐 SB", "逐 Query work_mem", "关键动作"]]
    for desc, action, row in zip(descriptions, actions, report["stages"]):
        wm = "；".join(f"{key.upper()}={value}" for key, value in row["recommended_work_mem"].items())
        data.append([row["stage"], desc, f"{row['recommended_sb_mb']}MB", wm, action])
    add_table(slide, data, 0.52, 1.56, 12.28, 3.88, widths=[0.68, 2.15, 1.35, 5.12, 2.98], font_size=8.8)
    add_callout(slide, "动作解释", "S2 从 AP-first 上界得到降低 SB；S3 在 TP 饱和后限制 AP grant；S4 无可行新增 AP 候选，执行准入控制；S5 TP demand 突增后提高 SB。", 0.82, 5.77, 11.68, 0.82, TEAL, PALE_TEAL)
    footer(slide, 12, "five-stage recommendation + post-decision validation")
    notes.append("结果二：按验收动作讲因果，不只报配置表。")


def stability_result(prs, notes):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(slide)
    add_title(slide, "端到端结果：S3-S5 保护 TP TPS 波动 2.48%", "所有推荐配置均实际应用；所有已运行 AP 正常自然结束。", "12 实验结果三")
    slide.shapes.add_picture(str(CHART_STAGE), Inches(0.56), Inches(1.48), width=Inches(8.30))
    checks = read_json(STAGE_REPORT)["checks"]
    add_kpi(slide, f"{checks['protected_tp_variation_s3_s5_percent']:.2f}%", "S3-S5 保护 TPS 波动", 9.15, 1.62, 3.08, color=GREEN, note="验收门槛 <5%")
    add_callout(slide, "动作一致性", "✓ S2 降 SB\n✓ S3 降 AP work_mem\n✓ S4 阻塞新 AP\n✓ S5 提高 SB", 9.02, 3.05, 3.35, 1.46, TEAL, PALE_TEAL)
    add_callout(slide, "自然结束", "S4/S5 只阻塞新 AP；已运行 SQL 不取消，无 AP failure。", 9.02, 4.79, 3.35, 0.86, BLUE, PALE_TEAL)
    add_callout(slide, "S5 含义", "保护 TP 4308→4416 TPS；新增 demand 另贡献约 298 TPS。", 9.02, 5.91, 3.35, 0.72, ORANGE, PALE_ORANGE)
    footer(slide, 13, "validation_report.json")
    notes.append("结果三：2.48% 是保护 TP 的阶段间波动；S5 总 TPS 还包含新增 TP demand。")


def engineering(prs, notes):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(slide)
    add_title(slide, "工程落地：新机器可一键完成标定、冻结、盲测和推荐", "不是只在当前实验目录中手工运行；流程已经收敛为配置驱动、可恢复的控制器。", "13 工程化")
    add_flow_box(slide, "machine.json", "设备与 CPU mask\nopenGauss 路径\nTP command / terminals\ncandidate replay command", 0.62, 1.80, 2.82, 2.05, BLUE, fill=PALE_TEAL)
    add_arrow(slide, 3.49, 2.82, 4.04, 2.82, TEAL, 2.0)
    add_flow_box(slide, "modelctl run-all", "doctor / prepare\nstorage calibration\npath anchors / freeze\nunseen holdout / predict", 4.12, 1.80, 2.82, 2.05, TEAL, fill=PALE_GREEN)
    add_arrow(slide, 6.99, 2.82, 7.54, 2.82, TEAL, 2.0)
    add_flow_box(slide, "模型产物", "machine inventory\nfrozen surface / model\nholdout report\nrecommendations.csv", 7.62, 1.80, 2.82, 2.05, ORANGE, fill=PALE_ORANGE)
    add_arrow(slide, 10.49, 2.82, 11.04, 2.82, TEAL, 2.0)
    add_flow_box(slide, "执行控制", "阶段重启应用 SB\nsession work_mem\nAP admission\n持续监控 TP", 11.12, 1.80, 1.62, 2.05, GREEN, fill=PALE_GREEN)
    add_callout(slide, "可靠性设计", "workspace 互斥锁、持久日志、state.json 断点续跑、产物大小/哈希检查、数据库停止后 finally 恢复、负载自然结束。", 0.88, 4.54, 11.58, 0.86, TEAL, PALE_TEAL)
    add_text(slide, "bin/run_portable_model.sh /absolute/path/to/machine.json", 1.12, 5.78, 11.1, 0.42, size=17, color=INK, bold=True, align=PP_ALIGN.CENTER, font="Consolas")
    add_callout(slide, "完整性门槛", "缺少候选 replay 或公式盲测未通过时，流程拒绝输出完整推荐。", 1.20, 6.08, 10.92, 0.64, RED, PALE_RED)
    footer(slide, 14, "huawei6_modelctl.py / PORTABLE_MODEL_BOOTSTRAP.md")
    notes.append("工程页：说明这不是一次性实验脚本，而是新机器可执行流程。")


def contributions(prs, notes):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(slide)
    add_title(slide, "项目贡献与创新点", "贡献不在单一预测精度，而在将 plan、内存、I/O 和 TPS 串成可验证、可执行的闭环。", "14 总结")
    items = [
        ("1", "Plan-aware Replay", "从一次负载与源码/算子行为推导不同 work_mem 下的 plan、动态峰值和 spill，支持逐 Query grant。", BLUE, PALE_TEAL),
        ("2", "因果 I/O→TPS 模型", "分离设备曲面、数据库路径和事务响应时间；新机器重标定，TPS 无拟合乘数。", ORANGE, PALE_ORANGE),
        ("3", "双向联合搜索", "TP-first 给 SB 性能下界，AP-first 给内存上界，在交集内联合选择 SB/work_mem/AP-cap。", TEAL, PALE_GREEN),
        ("4", "可执行稳定控制", "在原版 openGauss 上通过阶段重启、逐 Query work_mem 和 AP 准入实现五阶段动作与 <5% 稳定目标。", GREEN, PALE_GREEN),
    ]
    for index, (number, title, body, color, fill) in enumerate(items):
        x = 0.68 + (index % 2) * 6.24
        y = 1.72 + (index // 2) * 2.32
        add_rect(slide, x, y, 5.70, 1.80, fill, line=color, radius=True, width=1.4)
        add_rect(slide, x + 0.18, y + 0.20, 0.56, 0.56, color, radius=True)
        add_text(slide, number, x + 0.18, y + 0.32, 0.56, 0.22, size=15, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, title, x + 0.92, y + 0.22, 4.45, 0.32, size=16, color=color, bold=True)
        add_text(slide, body, x + 0.92, y + 0.67, 4.46, 0.86, size=10.5, color=INK)
    add_callout(slide, "量化证据", "未见 QD latency MAPE 1.53%，TPS MAPE 1.81%；五阶段动作全部匹配，保护 TPS 波动 2.48%。", 1.08, 6.04, 11.12, 0.68, TEAL, PALE_TEAL)
    footer(slide, 15)
    notes.append("贡献页：四点分别对应模型能力、可迁移性、搜索方法和工程控制。")


def limitations(prs, notes):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(slide)
    add_title(slide, "当前边界与下一步工作", "答辩中主动区分“已严格验证”和“仍需扩展”，避免把当前负载结论泛化为所有数据库场景。", "15 边界")
    add_flow_box(slide, "已验证", "8KiB buffered TP 读\n128KiB random-read AP 压力\n当前 terminals 与 QD 有效域\n未见 QD latency/TPS holdout", 0.70, 1.74, 3.50, 2.16, GREEN, fill=PALE_GREEN)
    add_flow_box(slide, "尚未完全解决", "新 plan 的 spill→qAP 仍依赖 replay/锚点\n写型、顺序型和不同块大小 I/O\nCPU/锁/OS cache 跨负载泛化\n更大规模未见 Query 与机器", 4.92, 1.74, 3.50, 2.16, RED, fill=PALE_RED)
    add_flow_box(slide, "下一步", "增加多 I/O class 设备曲面\n建立 plan/spill→qAP 未见 Query holdout\n联合 CPU/锁与 I/O 响应时间\n扩大跨机器、跨 scale factor 验证", 9.14, 1.74, 3.50, 2.16, BLUE, fill=PALE_TEAL)
    add_callout(slide, "当前可以主张", "在已记录的机器与 I/O domain 内，模型能够从 Trace/观测产生与五阶段要求一致的配置动作，并将保护 TP TPS 波动控制在 5% 内。", 0.96, 4.62, 11.44, 0.88, TEAL, PALE_TEAL)
    add_callout(slide, "当前不能主张", "尚不能声称任意新 SQL、任意 I/O 类型和任意机器都无需额外标定即可达到同样误差。超出 domain 的候选会被拒绝，而不是静默外推。", 0.96, 5.79, 11.44, 0.86, RED, PALE_RED)
    footer(slide, 16)
    notes.append("边界页：明确当前结论的有效域和下一阶段研究计划。")


def conclusion(prs, notes):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(slide, INK)
    add_text(slide, "结论", 0.92, 0.78, 2.0, 0.38, size=18, color=TEAL, bold=True)
    add_text(slide, "从“预测命中率”升级为\n“预测配置动作并保护 TP TPS”", 0.90, 1.38, 11.4, 1.18, size=34, color=WHITE, bold=True)
    items = [
        ("方法", "Plan-aware Replay + 机器因果模型 + TP/AP 双向约束搜索"),
        ("结果", "未见 QD：latency 1.53%，TPS 1.81%；五阶段保护 TPS 波动 2.48%"),
        ("落地", "原版 openGauss、阶段重启、逐 Query work_mem、AP 准入、一键新机器标定"),
    ]
    for index, (title, body) in enumerate(items):
        y = 3.08 + index * 0.92
        add_rect(slide, 0.96, y, 1.26, 0.56, TEAL if index == 0 else (GREEN if index == 1 else ORANGE), radius=True)
        add_text(slide, title, 0.96, y + 0.14, 1.26, 0.22, size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, body, 2.48, y + 0.12, 9.75, 0.32, size=15, color=WHITE)
    add_text(slide, "谢谢", 0.94, 6.38, 2.0, 0.48, size=26, color=TEAL, bold=True)
    add_text(slide, "Q & A", 10.12, 6.40, 2.0, 0.42, size=22, color=GRAY, bold=True, align=PP_ALIGN.RIGHT)
    notes.append("结尾：重申从 hit ratio 到 TPS 与控制动作的升级，给出三项量化证据。")


def build() -> None:
    make_charts()
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    notes: list[str] = []
    builders = (
        cover, background, objective, challenges, architecture, work_replay,
        work_causal_model, work_search, controller, protocol, formula_result,
        stage_actions, stability_result, engineering, contributions, limitations,
        conclusion,
    )
    for builder in builders:
        builder(prs, notes)
    prs.core_properties.title = "Huawei6 项目答辩：混合负载内存联合优化"
    prs.core_properties.subject = "Trace Replay、I/O latency→TPS 因果模型与双向联合搜索"
    prs.core_properties.author = "Huawei6 project"
    prs.save(OUT)
    shutil.copy2(OUT, ROOT_COPY)
    OUTLINE.write_text(
        "# Huawei6 项目答辩讲稿\n\n" + "\n".join(
            f"## 第 {index} 页\n{note}\n" for index, note in enumerate(notes, 1)
        ),
        encoding="utf-8",
    )
    print(json.dumps({"pptx": str(OUT), "root_copy": str(ROOT_COPY), "slides": len(prs.slides), "notes": str(OUTLINE)}, ensure_ascii=False, indent=2))


if __name__ == "__main__":
    build()
