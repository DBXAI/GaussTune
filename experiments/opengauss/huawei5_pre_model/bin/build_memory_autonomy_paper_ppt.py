#!/usr/bin/env python3
"""Build a paper-style presentation for the Huawei5 memory-autonomy project."""

from __future__ import annotations

import csv
import json
from pathlib import Path

import matplotlib.pyplot as plt
import numpy as np
from matplotlib import font_manager
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
LATEST = ROOT / "artifacts" / "00_latest"
LATEST.mkdir(parents=True, exist_ok=True)

OUT = LATEST / "Huawei5_memory_autonomy_paper_report_20260727.pptx"
OUTLINE = LATEST / "Huawei5_memory_autonomy_paper_report_20260727_outline.md"
CHART_MODEL = LATEST / "paper_model_validation_summary_20260727.png"
CHART_STAGE = LATEST / "paper_five_stage_tps_regret_20260727.png"
CHART_CROSS_STAGE = LATEST / "paper_cross_stage_tps_stability_20260727.png"
CHART_KERNEL = LATEST / "paper_kernel_resize_acceptance_20260727.png"

FONT = "Microsoft YaHei"
MONO = "Consolas"

INK = RGBColor(30, 43, 53)
TEAL = RGBColor(16, 117, 127)
GREEN = RGBColor(43, 139, 92)
ORANGE = RGBColor(222, 132, 40)
RED = RGBColor(194, 67, 62)
BLUE = RGBColor(52, 111, 166)
PURPLE = RGBColor(111, 82, 150)
LIGHT = RGBColor(244, 247, 248)
MID = RGBColor(213, 221, 225)
GRAY = RGBColor(96, 107, 115)
WHITE = RGBColor(255, 255, 255)
PALE_TEAL = RGBColor(226, 241, 242)
PALE_ORANGE = RGBColor(251, 239, 223)
PALE_GREEN = RGBColor(230, 243, 236)
PALE_RED = RGBColor(249, 232, 231)


def rgb_hex(color: RGBColor) -> str:
    return f"#{color[0]:02x}{color[1]:02x}{color[2]:02x}"


def configure_plot_font() -> None:
    candidates = [
        "Microsoft YaHei",
        "Noto Sans CJK SC",
        "Noto Sans CJK JP",
        "WenQuanYi Micro Hei",
        "SimHei",
        "DejaVu Sans",
    ]
    installed = {f.name for f in font_manager.fontManager.ttflist}
    selected = next((name for name in candidates if name in installed), "DejaVu Sans")
    plt.rcParams["font.sans-serif"] = [selected]
    plt.rcParams["axes.unicode_minus"] = False


def set_bg(slide, color=WHITE):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, x, y, w, h, fill, *, line=None, radius=False, width=1.0):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line or fill
    shape.line.width = Pt(width)
    if radius:
        shape.adjustments[0] = 0.08
    return shape


def add_text(
    slide,
    text,
    x,
    y,
    w,
    h,
    *,
    size=18,
    color=INK,
    bold=False,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    font=FONT,
    margin=0.03,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_bullets(slide, items, x, y, w, h, *, size=16, color=INK, spacing=8):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.03)
    tf.margin_top = Inches(0.02)
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {item}"
        p.font.name = FONT
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(spacing)
    return box


def add_title(slide, title, subtitle=None, section=None):
    add_rect(slide, 0, 0, 13.333, 0.11, TEAL)
    if section:
        add_text(slide, section, 0.62, 0.26, 2.8, 0.28, size=9, color=TEAL, bold=True)
    add_text(slide, title, 0.62, 0.56, 12.0, 0.52, size=27, bold=True)
    if subtitle:
        add_text(slide, subtitle, 0.64, 1.10, 11.9, 0.38, size=12, color=GRAY)


def add_footer(slide, page, source=None):
    add_rect(slide, 0.62, 7.17, 12.08, 0.012, MID)
    footer = "Huawei5 · Trace Replay 与 openGauss 在线内存自治 · 2026-07-27"
    add_text(slide, footer, 0.62, 7.22, 6.9, 0.18, size=7.5, color=GRAY)
    if source:
        add_text(slide, f"来源：{source}", 7.05, 7.22, 5.15, 0.18, size=7.0, color=GRAY, align=PP_ALIGN.RIGHT)
    add_text(slide, str(page), 12.28, 7.22, 0.42, 0.18, size=8, color=GRAY, align=PP_ALIGN.RIGHT)


def add_arrow(slide, x1, y1, x2, y2, color=GRAY, width=2.0):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    line.line.color.rgb = color
    line.line.width = Pt(width)
    line.line.end_arrowhead = True
    return line


def add_flow_box(slide, title, sub, x, y, w, h, color, *, fill=LIGHT):
    add_rect(slide, x, y, w, h, fill, line=color, radius=True, width=1.5)
    add_text(slide, title, x + 0.10, y + 0.16, w - 0.20, 0.33, size=15, color=color, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, sub, x + 0.12, y + 0.57, w - 0.24, h - 0.66, size=10, color=GRAY, align=PP_ALIGN.CENTER)


def add_kpi(slide, value, label, x, y, w, *, color=TEAL, note=None):
    add_text(slide, value, x, y, w, 0.52, size=29, color=color, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, label, x, y + 0.56, w, 0.32, size=12, bold=True, align=PP_ALIGN.CENTER)
    if note:
        add_text(slide, note, x, y + 0.91, w, 0.38, size=9, color=GRAY, align=PP_ALIGN.CENTER)


def add_table(slide, data, x, y, w, h, *, widths=None, font_size=10.5):
    rows, cols = len(data), len(data[0])
    table = slide.shapes.add_table(rows, cols, Inches(x), Inches(y), Inches(w), Inches(h)).table
    if widths:
        for i, width in enumerate(widths):
            table.columns[i].width = Inches(width)
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = str(data[r][c])
            cell.margin_left = Inches(0.04)
            cell.margin_right = Inches(0.04)
            cell.margin_top = Inches(0.03)
            cell.margin_bottom = Inches(0.03)
            cell.fill.solid()
            cell.fill.fore_color.rgb = INK if r == 0 else (LIGHT if r % 2 else WHITE)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER
                for run in p.runs:
                    run.font.name = FONT
                    run.font.size = Pt(font_size if r else font_size - 0.5)
                    run.font.bold = r == 0
                    run.font.color.rgb = WHITE if r == 0 else INK
    return table


def add_callout(slide, title, body, x, y, w, h, color, fill):
    add_rect(slide, x, y, w, h, fill, line=color, radius=True, width=1.2)
    add_text(slide, title, x + 0.16, y + 0.10, w - 0.32, 0.28, size=13.5, color=color, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, body, x + 0.16, y + 0.39, w - 0.32, max(0.22, h - 0.46), size=10.5, color=INK, align=PP_ALIGN.CENTER)


def make_model_chart() -> None:
    configure_plot_font()
    fig, axes = plt.subplots(1, 3, figsize=(13.2, 3.7))
    fig.patch.set_facecolor("white")

    labels = ["SB", "OS", "Combined"]
    values = [0.61, 1.82, 0.50]
    bars = axes[0].bar(labels, values, color=[rgb_hex(BLUE), rgb_hex(ORANGE), rgb_hex(TEAL)], width=0.62)
    axes[0].set_title("缓存命中率 held-out MAE", fontsize=13, fontweight="bold")
    axes[0].set_ylabel("百分点（越低越好）")
    axes[0].set_ylim(0, 2.15)
    for bar, value in zip(bars, values):
        axes[0].text(bar.get_x() + bar.get_width()/2, value + 0.06, f"{value:.2f} pp", ha="center", fontsize=11)

    labels = ["Plan family", "Spill 分类"]
    values = [100.0, 85.7]
    bars = axes[1].bar(labels, values, color=[rgb_hex(GREEN), rgb_hex(TEAL)], width=0.58)
    axes[1].set_title("单次负载 + 源码 replay 严格留出", fontsize=13, fontweight="bold")
    axes[1].set_ylabel("准确率（%）")
    axes[1].set_ylim(0, 112)
    for bar, value, frac in zip(bars, values, ["7/7", "6/7"]):
        axes[1].text(bar.get_x() + bar.get_width()/2, value + 2.5, f"{frac}\n{value:.1f}%", ha="center", fontsize=10)

    labels = ["首轮未见 SQL", "扩大未见 SQL", "一次锚点 v2"]
    values = [100.0, 71.4, 100.0]
    bars = axes[2].bar(labels, values, color=[rgb_hex(GREEN), rgb_hex(ORANGE), rgb_hex(BLUE)], width=0.62)
    axes[2].set_title("Spill 泛化：扩样本暴露边界", fontsize=13, fontweight="bold")
    axes[2].set_ylabel("分类准确率（%）")
    axes[2].set_ylim(0, 112)
    axes[2].tick_params(axis="x", labelrotation=10)
    for bar, value, frac in zip(bars, values, ["7/7", "15/21", "6/6"]):
        axes[2].text(bar.get_x() + bar.get_width()/2, value + 2.5, f"{frac}", ha="center", fontsize=10)

    for ax in axes:
        ax.spines[["top", "right"]].set_visible(False)
        ax.grid(axis="y", alpha=0.18)
    fig.tight_layout(w_pad=2.2)
    fig.savefig(CHART_MODEL, dpi=190, bbox_inches="tight")
    plt.close(fig)


def load_stage_rows():
    with (LATEST / "five_stage_saturated_tps_validation_20260726.csv").open(newline="", encoding="utf-8") as f:
        return list(csv.DictReader(f))


def make_stage_chart() -> None:
    configure_plot_font()
    rows = load_stage_rows()
    stages = [f"S{i+1}" for i in range(len(rows))]
    rec = [float(r["recommended_actual_tps"]) for r in rows]
    best = [float(r["best_challenger_tps"]) for r in rows]
    regret = [float(r["tps_regret_pct"]) for r in rows]
    x = np.arange(len(stages))
    width = 0.34
    fig, ax = plt.subplots(figsize=(11.8, 4.2))
    bars1 = ax.bar(x - width/2, rec, width, color=rgb_hex(TEAL), label="模型推荐配置的实测 TPS")
    bars2 = ax.bar(x + width/2, best, width, color="#a7b1b9", label="已测挑战点最高 TPS")
    ax.axhline(0, color="#333333", linewidth=0.8)
    ax.set_xticks(x, stages)
    ax.set_ylabel("饱和 TP TPS")
    ax.set_title("五阶段推荐配置：实测 TPS 距已测最高点均小于 5%", fontsize=15, fontweight="bold")
    ax.legend(frameon=False, ncol=2, loc="upper left")
    ax.grid(axis="y", alpha=0.18)
    ax.spines[["top", "right"]].set_visible(False)
    ax.set_ylim(0, max(best) * 1.20)
    for i, (b1, b2, reg) in enumerate(zip(bars1, bars2, regret)):
        top = max(b1.get_height(), b2.get_height())
        ax.text(i, top + 30, f"regret {reg:.2f}%", ha="center", color=rgb_hex(GREEN), fontweight="bold", fontsize=11)
    fig.tight_layout()
    fig.savefig(CHART_STAGE, dpi=190, bbox_inches="tight")
    plt.close(fig)


def make_cross_stage_chart() -> None:
    configure_plot_font()
    rows = load_stage_rows()
    stages = [f"S{i+1}" for i in range(len(rows))]
    tps = np.array([float(r["recommended_actual_tps"]) for r in rows])
    x = np.arange(len(stages))
    reference = tps[0]
    lower, upper = reference * 0.95, reference * 1.05
    fig, ax = plt.subplots(figsize=(10.6, 4.0))
    ax.axhspan(lower, upper, color=rgb_hex(PALE_GREEN), alpha=0.95, label="以 S1 为参照的 ±5% 带")
    ax.plot(x, tps, color=rgb_hex(TEAL), linewidth=2.8, marker="o", markersize=8, label="各阶段推荐配置实测 TPS")
    for i, value in enumerate(tps):
        ax.text(i, value + 14, f"{value:.1f}", ha="center", fontsize=10, fontweight="bold")
    for i in range(1, len(tps)):
        delta = (tps[i] / tps[i-1] - 1.0) * 100
        ax.text(i - 0.5, (tps[i] + tps[i-1]) / 2 - 22, f"{delta:+.2f}%", ha="center", color=rgb_hex(ORANGE), fontsize=9)
    ax.set_xticks(x, stages)
    ax.set_ylabel("饱和 TP TPS")
    ax.set_title("严格比较五阶段绝对 TPS：当前最大/最小差 19.71%，不满足 5%", fontsize=14.5, fontweight="bold")
    ax.set_ylim(min(lower, tps.min()) - 60, tps.max() + 75)
    ax.grid(axis="y", alpha=0.18)
    ax.spines[["top", "right"]].set_visible(False)
    ax.legend(frameon=False, loc="upper left")
    fig.tight_layout()
    fig.savefig(CHART_CROSS_STAGE, dpi=190, bbox_inches="tight")
    plt.close(fig)


def make_kernel_chart() -> None:
    configure_plot_font()
    base = ROOT / "results" / "kernel_online_resize_tps_20260726"
    paths = [
        base / "gated_8mb_read_only_aggregate" / "aggregate_summary.json",
        base / "gated_8mb_1s_read_write_aggregate" / "aggregate_summary.json",
    ]
    summaries = [json.loads(p.read_text(encoding="utf-8")) for p in paths]
    labels = ["只读 TP\n8MB / 0.5s", "读写 TP\n8MB / 1s"]
    one_s = [s["worst_1s_drop_pct"] for s in summaries]
    mean = [s["worst_migration_mean_drop_pct"] for s in summaries]
    post = [abs(s["worst_post_delta_pct"]) for s in summaries]
    x = np.arange(2)
    width = 0.22
    fig, ax = plt.subplots(figsize=(9.8, 4.2))
    groups = [
        ax.bar(x - width, one_s, width, color=rgb_hex(TEAL), label="最差 1 秒下降"),
        ax.bar(x, mean, width, color=rgb_hex(BLUE), label="迁移期均值下降"),
        ax.bar(x + width, post, width, color=rgb_hex(ORANGE), label="迁移后变化绝对值"),
    ]
    ax.axhline(3.0, color=rgb_hex(RED), linestyle="--", linewidth=1.8, label="3% 验收线")
    ax.set_title("128MB→64MB 在线缩容：三轮重复均通过严格单秒红线", fontsize=15, fontweight="bold", pad=42)
    ax.set_ylabel("TPS 下降（%）")
    ax.set_xticks(x, labels)
    ax.set_ylim(0, 3.5)
    ax.grid(axis="y", alpha=0.18)
    ax.spines[["top", "right"]].set_visible(False)
    ax.legend(frameon=False, ncol=4, loc="lower center", bbox_to_anchor=(0.5, 1.01), fontsize=9)
    for bars in groups:
        for bar in bars:
            ax.text(bar.get_x() + bar.get_width()/2, bar.get_height() + 0.07, f"{bar.get_height():.2f}%", ha="center", fontsize=10)
    fig.tight_layout()
    fig.savefig(CHART_KERNEL, dpi=190, bbox_inches="tight")
    plt.close(fig)


def add_notes(slide, text):
    try:
        slide.notes_slide.notes_text_frame.text = text
    except Exception:
        pass


def build() -> None:
    make_model_chart()
    make_stage_chart()
    make_cross_stage_chart()
    make_kernel_chart()

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    outline = []

    # 1. Cover
    slide = prs.slides.add_slide(blank)
    set_bg(slide, INK)
    add_rect(slide, 0, 0, 0.20, 7.5, TEAL)
    add_text(slide, "基于 Trace Replay 与内核在线扩缩的", 0.92, 1.13, 11.6, 0.58, size=24, color=RGBColor(190, 224, 226), bold=True)
    add_text(slide, "openGauss 混合负载内存自治", 0.90, 1.76, 11.7, 0.90, size=38, color=WHITE, bold=True)
    add_text(slide, "问题、方法、实现与实验验证", 0.93, 2.83, 9.0, 0.48, size=21, color=RGBColor(202, 210, 215))
    add_rect(slide, 0.93, 3.68, 5.45, 0.54, TEAL, radius=True)
    add_text(slide, "Huawei5 五阶段 TP/AP 混合负载", 1.08, 3.81, 5.15, 0.27, size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "核心目标：联合推荐 shared_buffers 与 work_mem，并在不中断数据库的条件下执行配置切换", 0.93, 5.34, 10.9, 0.74, size=16, color=RGBColor(211, 218, 222))
    add_text(slide, "研究进展汇报 · 2026-07-27", 0.93, 6.45, 4.3, 0.30, size=11, color=RGBColor(171, 183, 190))
    add_notes(slide, "开场先给出研究主线：预测层回答配多少，内核层回答如何在线切换。")
    outline.append("1. 封面：Trace Replay 与内核在线扩缩的 openGauss 混合负载内存自治。")

    # 2. Abstract
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, "摘要：从离线参数调优走向运行时内存自治", "一条主线解决“推荐什么”，另一条主线解决“如何无重启执行”。", "ABSTRACT")
    add_callout(slide, "问题", "固定 shared_buffers 无法同时适应 TP 热集、AP spill 与阶段性并发；work_mem 也不是独立旋钮。", 0.70, 1.72, 3.75, 2.00, RED, PALE_RED)
    add_callout(slide, "方法", "以一次负载 trace、EXPLAIN Plan 和 openGauss 算子源码为输入，联合 replay 缓存路径、算子内存与动态内存约束。", 4.78, 1.72, 3.78, 2.00, TEAL, PALE_TEAL)
    add_callout(slide, "实现", "控制器按阶段选择 SB/work_mem；内核以粒度化退休协议在线改变 active buffer 边界并释放物理页。", 8.89, 1.72, 3.76, 2.00, BLUE, LIGHT)
    add_rect(slide, 0.70, 4.15, 11.95, 1.72, INK, radius=True)
    add_text(slide, "主要结果", 0.96, 4.42, 1.45, 0.33, size=16, color=WHITE, bold=True)
    add_text(slide, "五阶段推荐在已测挑战网格上最大 TPS regret 4.55%；128→64MB 在线缩容三轮重复最差单秒 TPS 下降 2.76%，0 错误、0 重启。", 2.33, 4.36, 9.91, 0.72, size=16, color=WHITE, bold=True)
    add_text(slide, "结论边界：当前证明了 replay 推荐有效性和小规模内核路径可行性；尚未证明跨阶段归一化 TPS 波动 ≤5%，也未完成完整五阶段在线验收。", 0.97, 5.24, 11.13, 0.40, size=11, color=RGBColor(199, 210, 215))
    add_footer(slide, 2, "SATURATED_FIVE_STAGE_TPS_VALIDATION.md；kernel_online_resize_tps_20260726/README.md")
    add_notes(slide, "这一页相当于论文摘要，先讲结果，再强调结果边界。")
    outline.append("2. 摘要：问题、方法、实现、主要结果与结论边界。")

    # 3. Background
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, "问题背景：TP 与 AP 争用同一份物理内存", "SB 保护 TP 热页，work_mem 减少 AP spill，但二者同时挤压 OS page cache 与动态内存池。", "1 · BACKGROUND")
    add_flow_box(slide, "shared_buffers", "TP 热页驻留\n命中不足 → 数据页 I/O", 0.82, 1.77, 2.45, 1.42, TEAL, fill=PALE_TEAL)
    add_flow_box(slide, "OS page cache", "补偿小 SB\n受 refault/active-inactive 影响", 5.44, 1.77, 2.45, 1.42, BLUE, fill=LIGHT)
    add_flow_box(slide, "work_mem / 动态池", "Join / Sort / Agg\n不足 → spill 临时 I/O", 10.05, 1.77, 2.45, 1.42, ORANGE, fill=PALE_ORANGE)
    add_arrow(slide, 3.32, 2.48, 5.34, 2.48, GRAY, 2.0)
    add_arrow(slide, 9.98, 2.48, 7.97, 2.48, GRAY, 2.0)
    add_rect(slide, 3.68, 3.72, 5.98, 1.06, INK, radius=True)
    add_text(slide, "memory_target_max：必须共同满足的物理内存上限", 3.91, 4.02, 5.53, 0.38, size=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_bullets(slide, [
        "SB 增大不一定提升 TPS：Linux page cache 可能已经吸收部分缺页；过大反而压缩 AP 与 OS 空间。",
        "work_mem 是每个并发内存算子的额度，不是每个阶段的一次性总预算。",
        "最优配置随 AP Query、并发数、TP 热集与执行 Plan 改变，静态经验规则无法稳定覆盖。",
    ], 0.96, 5.18, 11.43, 1.37, size=14.5, spacing=6)
    add_footer(slide, 3, "RUNTIME_MEMORY_CONTROLLER_REPLAY.md；ONE_SHOT_SOURCE_PLAN_REPLAY.md")
    add_notes(slide, "这里解释为什么 SB 和 work_mem 必须联合预测，而不是分别找两个最大值。")
    outline.append("3. 问题背景：SB、OS cache 与动态算子内存竞争同一物理内存。")

    # 4. Research questions
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, "研究问题与验收约束", "评价重点不是拟合历史 TPS，而是能否用可解释的执行机制外推未测配置。", "2 · QUESTIONS")
    questions = [
        ("RQ1", "缓存路径", "只采一次访问轨迹，能否预测不同 SB 下 TP 的 SB/OS/combined 命中？", TEAL),
        ("RQ2", "算子内存", "只执行一次 SQL，能否推测 Plan 切换与每个 Plan 的 spill 边界？", ORANGE),
        ("RQ3", "联合推荐", "能否在动态池和物理内存约束下选择 SB × work_mem，而非两个独立拐点？", BLUE),
        ("RQ4", "在线执行", "能否不重启数据库完成 SB 切换，并把最差单秒 TPS 下降控制在 3% 内？", GREEN),
    ]
    y = 1.72
    for code, title, body, color in questions:
        add_rect(slide, 0.78, y, 1.05, 0.72, color, radius=True)
        add_text(slide, code, 0.82, y + 0.20, 0.97, 0.30, size=15, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, title, 2.10, y + 0.02, 1.72, 0.34, size=17, color=color, bold=True)
        add_text(slide, body, 3.75, y, 8.70, 0.62, size=14)
        y += 1.06
    add_rect(slide, 0.78, 6.11, 11.74, 0.55, PALE_RED, line=RED, radius=True)
    add_text(slide, "数据隔离原则：预测文件先冻结并记录哈希，验证 TPS、spill 标签不得回写模型；失败点必须保留。", 1.02, 6.25, 11.25, 0.27, size=13, color=RED, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, 4, "GENERALIZATION_VALIDATION.md；SATURATED_FIVE_STAGE_TPS_VALIDATION.md")
    add_notes(slide, "强调预注册和冻结预测，回应‘看答案后校准’的质疑。")
    outline.append("4. 研究问题：缓存、Plan/Spill、联合推荐和无重启执行；验证标签不回写。")

    # 5. Architecture
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, "总体框架：预测层、决策层、执行层闭环", "输入是 trace、Plan 与源码语义；TPS 只用于独立验收，不作为推荐器特征。", "3 · SYSTEM")
    add_flow_box(slide, "一次负载采集", "TP 页访问 trace\nAP 算子 trace\n并发与生命周期", 0.67, 1.80, 2.20, 1.46, BLUE)
    add_flow_box(slide, "候选 Plan 扫描", "不同 work_mem 下 EXPLAIN\nPlan family SHA\nrows / width", 3.18, 1.80, 2.20, 1.46, PURPLE)
    add_flow_box(slide, "双向联合 Replay", "缓存路径 + Plan/Spill\nOS cache + 动态峰值\n物理 I/O", 5.69, 1.80, 2.24, 1.46, TEAL, fill=PALE_TEAL)
    add_flow_box(slide, "约束搜索与控制器", "TPS 平台优先\n内存安全边界\n准入/排队/粒度化迁移", 8.24, 1.80, 2.23, 1.46, ORANGE, fill=PALE_ORANGE)
    add_flow_box(slide, "openGauss 内核", "active/desired/allocation\n页退休与物理释放\n在线生效", 10.78, 1.80, 1.90, 1.46, GREEN, fill=PALE_GREEN)
    for x in [2.89, 5.40, 7.95, 10.49]:
        add_arrow(slide, x, 2.53, x + 0.25, 2.53, GRAY, 1.8)
    add_rect(slide, 0.78, 3.78, 11.74, 1.12, INK, radius=True)
    add_text(slide, "推荐目标", 1.03, 4.00, 1.28, 0.31, size=16, color=WHITE, bold=True)
    add_text(slide, "先进入 TP-SB 命中平台，再最小化 TP 数据 I/O + AP spill I/O；同等性能下选择更小的动态峰值与总内存。", 2.25, 3.96, 9.95, 0.46, size=15, color=WHITE, bold=True)
    add_text(slide, "安全门：MemAvailable 保留量、max_dynamic_memory、并发算子生命周期、SB active 已提交状态。", 2.25, 4.47, 9.65, 0.29, size=10.5, color=RGBColor(198, 211, 216))
    add_callout(slide, "离线输出", "每阶段 SB、全局或 per-query work_mem、预测命中率、spill I/O、动态峰值和置信度。", 0.78, 5.34, 5.62, 1.02, TEAL, PALE_TEAL)
    add_callout(slide, "运行时输出", "目标配置、迁移速率、准入数/排队数；仅在 active 提交后转移动态内存额度。", 6.69, 5.34, 5.83, 1.02, GREEN, PALE_GREEN)
    add_footer(slide, 5, "JOINT_BIDIRECTIONAL_REPLAY.md；OPENGAUSS_RUNTIME_SB_KERNEL_DESIGN.md")
    add_notes(slide, "系统图是整场汇报的导航，后面四页依次拆解每个模块。")
    outline.append("5. 总体框架：一次采集、Plan 扫描、联合 replay、约束控制与内核执行。")

    # 6. SB replay
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, "方法一：执行路径感知的缓存 Trace Replay", "模型模拟页在 SB 与 Linux page cache 中的状态迁移，不直接回归 TPS。", "4 · CACHE REPLAY")
    add_flow_box(slide, "输入", "TP-only 页访问序列\n读写类型、时间与阶段\n多锚点用于校验路径", 0.78, 1.70, 2.42, 1.40, BLUE)
    add_arrow(slide, 3.25, 2.40, 3.72, 2.40, GRAY, 2.0)
    add_flow_box(slide, "SB 模拟", "固定容量 buffer pool\n冷热、淘汰、dirty/pin 约束\n得到 TP-SB hit", 3.78, 1.70, 2.42, 1.40, TEAL, fill=PALE_TEAL)
    add_arrow(slide, 6.25, 2.40, 6.72, 2.40, GRAY, 2.0)
    add_flow_box(slide, "Linux 模拟", "active / inactive\nrefault 距离与频繁页保护\n得到 OS conditional hit", 6.78, 1.70, 2.42, 1.40, ORANGE, fill=PALE_ORANGE)
    add_arrow(slide, 9.25, 2.40, 9.72, 2.40, GRAY, 2.0)
    add_flow_box(slide, "联合输出", "Combined = SB 命中\n+ SB miss 后 OS 命中\n+ TP 磁盘 miss", 9.78, 1.70, 2.72, 1.40, GREEN, fill=PALE_GREEN)
    add_rect(slide, 0.78, 3.53, 7.48, 2.38, LIGHT, line=MID, radius=True)
    add_text(slide, "为什么能外推不同 SB", 1.02, 3.78, 3.23, 0.35, size=18, color=TEAL, bold=True)
    add_bullets(slide, [
        "轨迹保留同一执行路径上的页重用距离，候选 SB 只改变容量与淘汰时机。",
        "TP 与 AP 均参与 replay，但评价缓存收益时单独统计 TP 页，避免 AP 扫描命中掩盖 TPS 拐点。",
        "Linux 模型显式保护频繁 refault 的 TP 页，修复小 SB 下 page cache 补偿被严重低估的问题。",
    ], 1.02, 4.20, 6.91, 1.46, size=13.2, spacing=5)
    add_rect(slide, 8.57, 3.53, 3.93, 2.38, INK, radius=True)
    add_text(slide, "held-out 命中率误差", 8.82, 3.83, 3.42, 0.34, size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_kpi(slide, "0.61 pp", "SB hit MAE", 8.84, 4.30, 1.55, color=BLUE)
    add_kpi(slide, "0.50 pp", "Combined MAE", 10.61, 4.30, 1.55, color=GREEN)
    add_text(slide, "OS conditional MAE：1.82 pp", 8.93, 5.47, 3.25, 0.26, size=10.5, color=RGBColor(202, 213, 218), align=PP_ALIGN.CENTER)
    add_footer(slide, 6, "artifacts/README.md；MODEL_NOTES.md")
    add_notes(slide, "先解释模型预测的是命中和 I/O，而不是凭 TPS 标签拟合曲线。")
    outline.append("6. 缓存 replay：TP-only 统计、SB 淘汰、Linux active/inactive/refault 与 combined。")

    # 7. Work_mem replay
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, "方法二：Plan 感知的算子内存 Replay", "每条 AP SQL 只执行一次；未执行 Plan 由 EXPLAIN 结构与 openGauss 源码规则合成。", "5 · OPERATOR REPLAY")
    add_flow_box(slide, "256MB 一次锚点", "实际 rows / width\n算子生命周期\n已发生的 spill", 0.70, 1.70, 2.20, 1.36, BLUE)
    add_flow_box(slide, "work_mem 扫描", "只执行 EXPLAIN\n识别 Plan family 区间\n提取 operatorMemKB", 3.16, 1.70, 2.20, 1.36, PURPLE)
    add_flow_box(slide, "源码语义重放", "Hash buckets / batches\nSortTuple / merge\nHashAgg group/context", 5.62, 1.70, 2.20, 1.36, TEAL, fill=PALE_TEAL)
    add_flow_box(slide, "生命周期聚合", "单 Query 峰值\n并发 Query 峰值求和\n动态池 grant 限制", 8.08, 1.70, 2.20, 1.36, ORANGE, fill=PALE_ORANGE)
    add_flow_box(slide, "输出", "Plan 切换区间\nspill/no-spill 边界\ntemp I/O 与置信度", 10.54, 1.70, 2.12, 1.36, GREEN, fill=PALE_GREEN)
    for x in [2.92, 5.38, 7.84, 10.30]:
        add_arrow(slide, x, 2.38, x + 0.20, 2.38, GRAY, 1.7)
    add_rect(slide, 0.70, 3.54, 7.38, 2.23, LIGHT, line=MID, radius=True)
    add_text(slide, "work_mem 的分配语义", 0.96, 3.80, 2.88, 0.34, size=18, color=ORANGE, bold=True)
    add_bullets(slide, [
        "DOP=1 时，一个 Query 中每个同时存活的 Sort/Hash/HashAggregate 都可能获得一个 work_mem。",
        "阶段动态峰值 = 每条 Query 的算子生命周期峰值之和；不能简单用 work_mem × Query 数。",
        "若新 Plan 没有同 Plan trace，则使用自身 Plan 节点和源码公式，trace 只校准行数误差、行宽和 allocator 开销。",
    ], 0.96, 4.20, 6.81, 1.37, size=12.8, spacing=5)
    add_rect(slide, 8.39, 3.54, 4.27, 2.23, INK, radius=True)
    add_text(slide, "严格留出：一次负载模式", 8.68, 3.84, 3.70, 0.34, size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_kpi(slide, "7 / 7", "Plan family", 8.73, 4.35, 1.52, color=GREEN)
    add_kpi(slide, "6 / 7", "Spill 分类", 10.75, 4.35, 1.52, color=ORANGE)
    add_text(slide, "唯一失败：运行时动态内存保护提前拒绝 Sort 扩容", 8.73, 5.43, 3.58, 0.29, size=9.7, color=RGBColor(203, 214, 219), align=PP_ALIGN.CENTER)
    add_footer(slide, 7, "ONE_SHOT_SOURCE_PLAN_REPLAY.md")
    add_notes(slide, "说明 Plan 不是 replay 猜出来的：候选配置只跑 EXPLAIN；执行行为由源码语义重放。")
    outline.append("7. 算子 replay：单次锚点、Plan 扫描、源码规则、生命周期与动态池。")

    # 8. Joint recommendation
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, "方法三：SB × work_mem 双向联合搜索", "联合不是两个独立拐点取交集，而是每个二维候选都重新计算缓存、spill、动态峰值与内存安全。", "6 · JOINT SEARCH")
    add_rect(slide, 0.76, 1.66, 4.12, 4.62, LIGHT, line=MID, radius=True)
    add_text(slide, "对每个候选 (SB, work_mem)", 1.02, 1.93, 3.62, 0.36, size=18, color=TEAL, bold=True, align=PP_ALIGN.CENTER)
    steps = [
        ("1", "Replay TP 页路径", "得到 SB/OS 命中与 TP 磁盘 miss", BLUE),
        ("2", "Replay AP 算子", "Plan、spill I/O 与动态峰值", ORANGE),
        ("3", "反向修正 OS cache", "SB 与动态内存共同挤压 page cache", PURPLE),
        ("4", "检查安全边界", "MemAvailable、动态池、并发准入", RED),
    ]
    y = 2.52
    for n, title, body, color in steps:
        add_rect(slide, 1.03, y, 0.45, 0.45, color, radius=True)
        add_text(slide, n, 1.08, y + 0.09, 0.35, 0.24, size=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, title, 1.65, y - 0.02, 2.77, 0.31, size=14, color=color, bold=True)
        add_text(slide, body, 1.65, y + 0.32, 2.82, 0.35, size=10.3, color=GRAY)
        y += 0.83
    add_rect(slide, 5.22, 1.66, 7.38, 2.16, INK, radius=True)
    add_text(slide, "选择规则", 5.55, 1.95, 1.46, 0.33, size=17, color=WHITE, bold=True)
    add_text(slide, "① 进入 TP-SB 饱和平台（差 ≤ 0.001 absolute）", 5.59, 2.43, 6.42, 0.34, size=15, color=WHITE)
    add_text(slide, "② 最小化 TP 数据 I/O + AP spill I/O", 5.59, 2.81, 6.42, 0.34, size=15, color=WHITE)
    add_text(slide, "③ 并列时选择动态峰值与总内存更小的配置", 5.59, 3.19, 6.42, 0.34, size=15, color=WHITE)
    add_rect(slide, 5.22, 4.15, 7.38, 2.13, PALE_TEAL, line=TEAL, radius=True)
    add_text(slide, "S5 数值例：为什么不是简单取最大值", 5.53, 4.43, 6.77, 0.35, size=17, color=TEAL, bold=True)
    add_text(slide, "候选必须同时回答：TP 是否已进入命中平台？AP 是否仍 spill？动态峰值是否挤压 OS cache？总内存是否可部署？", 5.53, 4.91, 6.61, 0.65, size=13.5)
    add_text(slide, "因此 SB 与 work_mem 会相互影响：增大 work_mem 可能降低 spill，却减少 OS cache；缩小 SB 可换动态额度，却可能增加 TP 数据 I/O。", 5.53, 5.59, 6.62, 0.45, size=11, color=GRAY)
    add_footer(slide, 8, "JOINT_BIDIRECTIONAL_REPLAY.md；SATURATED_FIVE_STAGE_TPS_VALIDATION.md")
    add_notes(slide, "用这一页回答 SB 和 work_mem 到底有什么关系：资源约束和 OS cache 让二者双向耦合。")
    outline.append("8. 双向联合搜索：每个二维候选重算缓存、spill、OS cache 与安全边界。")

    # 9. Kernel
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, "执行层：openGauss shared_buffers 在线扩缩原型", "启动时保留最大地址空间，运行时改变 active 边界；缩容必须安全处理 clean、dirty 与 pinned 页。", "7 · KERNEL")
    add_flow_box(slide, "max", "启动期固定\n数组与地址空间上限", 0.82, 1.73, 2.18, 1.22, BLUE)
    add_flow_box(slide, "desired", "控制器请求的目标\nSIGHUP 可更新", 3.25, 1.73, 2.18, 1.22, PURPLE)
    add_flow_box(slide, "allocation", "正在退休/激活的边界\n按粒度推进", 5.68, 1.73, 2.18, 1.22, ORANGE, fill=PALE_ORANGE)
    add_flow_box(slide, "active", "前台可分配的边界\n提交后才生效", 8.11, 1.73, 2.18, 1.22, TEAL, fill=PALE_TEAL)
    add_flow_box(slide, "physical", "MADV_REMOVE\n释放物理页", 10.54, 1.73, 1.94, 1.22, GREEN, fill=PALE_GREEN)
    for x in [3.02, 5.45, 7.88, 10.31]:
        add_arrow(slide, x, 2.34, x + 0.19, 2.34, GRAY, 1.7)
    protocols = [
        ("Clean", "映射锁下重检 → 从 BufTable 删除 → 标记退休", GREEN),
        ("Dirty", "先写回并确认状态，再执行同样退休协议", ORANGE),
        ("Pinned / I/O", "跳过并重试；不可提前释放额度或物理页", RED),
    ]
    y = 3.48
    for title, body, color in protocols:
        add_rect(slide, 0.85, y, 1.37, 0.53, color, radius=True)
        add_text(slide, title, 0.89, y + 0.13, 1.29, 0.25, size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, body, 2.51, y + 0.04, 6.10, 0.39, size=13.5)
        y += 0.77
    add_rect(slide, 8.91, 3.48, 3.57, 2.08, INK, radius=True)
    add_text(slide, "原型接口", 9.19, 3.77, 3.01, 0.31, size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "shared_buffers_target\nshared_buffers_resize_granule\nshared_buffers_resize_interval", 9.22, 4.25, 2.95, 0.86, size=14, color=RGBColor(207, 223, 226), font=MONO, align=PP_ALIGN.CENTER)
    add_rect(slide, 8.91, 5.76, 3.57, 0.53, PALE_RED, line=RED, radius=True)
    add_text(slide, "尚缺：WLM quota / graceful debt", 9.04, 5.91, 3.31, 0.25, size=11, color=RED, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, 9, "OPENGAUSS_RUNTIME_SB_KERNEL_DESIGN.md")
    add_notes(slide, "解释为什么在线缩容不能只改一个 GUC：必须有退休协议与提交边界。")
    outline.append("9. 内核实现：max/desired/allocation/active 边界与安全页退休。")

    # 10. Experiment design
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, "实验设计：分层验证，避免把同一批标签用于训练与验收", "模型准确性、配置 regret 和在线切换稳定性分别使用不同指标。", "8 · EXPERIMENT")
    data = [
        ["层次", "输入 / 隔离", "评价指标", "当前规模"],
        ["缓存 replay", "锚点 trace 与测试 SB 分离", "SB/OS/combined hit MAE", "held-out 配置"],
        ["Plan / Spill", "冻结 Plan/预测 CSV 后执行", "Plan SHA、spill 分类、temp I/O", "7 点严格留出；21 点外部"],
        ["五阶段推荐", "32 terminal 饱和 TP；AP=1/1/2/4/4", "配置 regret；跨阶段归一化波动", "前者完成；后者待连续在线实跑"],
        ["在线缩容", "独立 openGauss；每轮重建表", "最差单秒 TPS 下降、错误、重启", "12 terminal；2 负载×3 轮"],
    ]
    add_table(slide, data, 0.72, 1.72, 11.90, 3.58, widths=[1.70, 4.15, 3.35, 2.70], font_size=10.5)
    add_callout(slide, "公平性控制", "每个 (SB, profile) 独立重启、清 Linux page cache、预热 90 秒；S2/S3 追加 120 秒长测。", 0.72, 5.62, 5.70, 0.76, TEAL, PALE_TEAL)
    add_callout(slide, "严格红线", "在线缩容在第 20 秒触发；验收关注每一秒样本，最差下降必须 <3%，且错误数和重启数均为 0。", 6.72, 5.62, 5.90, 0.76, RED, PALE_RED)
    add_footer(slide, 10, "GENERALIZATION_VALIDATION.md；SATURATED_FIVE_STAGE_TPS_VALIDATION.md；kernel README")
    add_notes(slide, "这一页对应论文的实验设置，明确每个数字的分母和协议。")
    outline.append("10. 实验设计：分层指标、冻结预测、独立重启清 cache 与三轮在线缩容。")

    # 11. Model effects
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, "效果一：模型在已知路径上准确，未见 SQL 泛化仍是主要风险", "扩大测试集后准确率下降，促使模型从跨 Query 比例迁移升级为结构约束的一次锚点 replay。", "9 · RESULTS")
    slide.shapes.add_picture(str(CHART_MODEL), Inches(0.64), Inches(1.63), width=Inches(12.10), height=Inches(3.52))
    add_callout(slide, "可以确认", "缓存命中子模型误差低；一次负载 + 源码 replay 可识别留出 Plan family，并在 7 个点中判断对 6 个 spill 状态。", 0.76, 5.42, 5.72, 0.87, GREEN, PALE_GREEN)
    add_callout(slide, "不能过度宣称", "21 个完全未见 Query 的 spill 分类仅 15/21；Plan family 由执行前 EXPLAIN 提供，不是 replay 自己预测。", 6.78, 5.42, 5.78, 0.87, RED, PALE_RED)
    add_footer(slide, 11, "GENERALIZATION_VALIDATION.md；ONE_SHOT_REPLAY_V2.md")
    add_notes(slide, "主动展示失败结果。v2 的 6/6 是新构造 Query 的一次锚点验证，不覆盖任意 SQL。")
    outline.append("11. 模型效果：命中率 MAE、Plan/Spill 留出结果与未见 SQL 泛化边界。")

    # 12. Stage results
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, "效果二：五阶段联合推荐在已测挑战网格上全部低于 5% regret", "regret = (同阶段已测最高 TPS − 推荐配置实测 TPS) / 已测最高 TPS。", "10 · RESULTS")
    slide.shapes.add_picture(str(CHART_STAGE), Inches(0.68), Inches(1.56), width=Inches(8.24), height=Inches(3.88))
    rows = load_stage_rows()
    table = [["阶段", "推荐 SB", "推荐 work_mem", "TPS regret"]]
    for i, row in enumerate(rows, 1):
        table.append([f"S{i}", f"{int(row['recommended_sb_mb'])}MB", f"{int(row['recommended_work_mem_mb'])}MB", f"{float(row['tps_regret_pct']):.3f}%"])
    add_table(slide, table, 9.18, 1.71, 3.42, 3.48, widths=[0.55, 1.02, 1.18, 0.67], font_size=8.8)
    add_rect(slide, 0.80, 5.66, 11.74, 0.67, INK, radius=True)
    add_text(slide, "结论：配置 regret 5/5 通过，最大 4.547%；这不等于五阶段之间 TPS 波动 ≤5%。S4 仍有 spill。", 1.03, 5.84, 11.27, 0.31, size=13.5, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, 12, "five_stage_saturated_tps_validation_20260726.csv")
    add_notes(slide, "强调这是已测挑战网格，不是连续二维空间的全局最优数学证明。")
    outline.append("12. 五阶段效果：推荐点 TPS 与已测最高点对比，5/5 regret <5%。")

    # 13. Cross-stage stability gap
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, "新增验收项：五阶段之间 TPS 相对波动 ≤5%", "该指标不同于配置 regret；当前绝对 TPS 结果不通过，归一化口径尚未完成连续在线验证。", "11 · ACCEPTANCE GAP")
    slide.shapes.add_picture(str(CHART_CROSS_STAGE), Inches(0.70), Inches(1.55), width=Inches(7.52), height=Inches(3.72))
    add_rect(slide, 8.52, 1.70, 4.02, 1.62, PALE_RED, line=RED, radius=True)
    add_text(slide, "严格绝对 TPS 口径", 8.80, 1.96, 3.46, 0.32, size=16, color=RED, bold=True, align=PP_ALIGN.CENTER)
    add_kpi(slide, "19.71%", "最大 / 最小相对差", 8.84, 2.38, 3.38, color=RED, note="当前不满足 5%")
    add_rect(slide, 8.52, 3.60, 4.02, 1.67, PALE_TEAL, line=TEAL, radius=True)
    add_text(slide, "建议的归一化口径", 8.80, 3.85, 3.46, 0.32, size=16, color=TEAL, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "normalized TPS = 当前 TP TPS /\n同一 TP 输入强度下的无 AP 基线 TPS", 8.80, 4.27, 3.46, 0.58, size=12, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "验收：每阶段稳态降幅 ≤5%；切换窗口仍执行单秒抖动 ≤3%。", 8.80, 4.89, 3.46, 0.28, size=9.6, color=GRAY, align=PP_ALIGN.CENTER)
    add_callout(slide, "为什么不能直接比较绝对 TPS", "S1-S4 是低 TP 输入，S5 明确包含 TP 突增；若强制绝对 TPS 同带，只能人为限速高负载阶段，失去保护 TP 吞吐的意义。", 0.82, 5.60, 5.72, 0.83, ORANGE, PALE_ORANGE)
    add_callout(slide, "要补的实验", "在一个不中断的五阶段运行中记录逐秒 TP TPS；分别采集低 TP 与高 TP 的无 AP 基线，并冻结控制策略后验证归一化波动。", 6.82, 5.60, 5.72, 0.83, BLUE, LIGHT)
    add_footer(slide, 13, "five_stage_saturated_tps_validation_20260726.csv；新增验收分析")
    add_notes(slide, "明确纠正之前口径：regret 小于 5% 不能替代跨阶段稳定性。")
    outline.append("13. 新增验收缺口：绝对 TPS 最大/最小差 19.71%；应按同输入强度无 AP 基线归一化。")

    # 14. Kernel effects
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, "效果三：在线缩容原型满足 3% 单秒 TPS 红线", "隔离 openGauss 5.1.0：128MB→64MB，无 postmaster 重启；读写负载采用较慢的 8MB/s。", "12 · RESULTS")
    slide.shapes.add_picture(str(CHART_KERNEL), Inches(0.72), Inches(1.52), width=Inches(8.02), height=Inches(3.93))
    add_rect(slide, 9.00, 1.70, 3.57, 3.39, INK, radius=True)
    add_text(slide, "三轮重复验收", 9.26, 1.99, 3.06, 0.34, size=17, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_kpi(slide, "2.76%", "最差单秒下降", 9.30, 2.55, 2.98, color=GREEN, note="只读 Run 2；仍低于 3%")
    add_kpi(slide, "0 / 0", "错误 / 重启", 9.30, 3.89, 2.98, color=TEAL, note="只读与读写共 6 轮")
    add_callout(slide, "节流是必要的", "读写负载使用 8MB/0.5s 时最差单秒下降达到 4.59%；改为 8MB/1s 后降至 2.65%。", 0.84, 5.57, 5.72, 0.84, ORANGE, PALE_ORANGE)
    add_callout(slide, "适用边界", "当前只证明 128MB 隔离原型；尚未验证 256MB 生产粒度、checkpoint 重叠、DMS/huge page 与完整五阶段。", 6.82, 5.57, 5.72, 0.84, RED, PALE_RED)
    add_footer(slide, 14, "kernel_online_resize_tps_20260726/README.md")
    add_notes(slide, "这页证明内核路径可行，但不要说完整生产验收已完成。")
    outline.append("14. 内核效果：128→64MB 三轮在线缩容，最差单秒下降 2.76%，0 错误/重启。")

    # 15. Conclusion
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, "结论与下一步", "当前已经形成可解释的预测-决策-执行链路，但生产闭环仍需完成三项工程验证。", "13 · CONCLUSION")
    add_callout(slide, "贡献 1 · 可解释预测", "用页访问、算子生命周期和源码容量规则替代纯 TPS 回归；可以定位误差来自 cache、Plan、grant 还是 spill。", 0.74, 1.65, 3.76, 1.62, TEAL, PALE_TEAL)
    add_callout(slide, "贡献 2 · 联合配置", "在每个二维候选上共同计算 TP 数据 I/O、AP spill、OS cache、动态峰值与安全边界。", 4.78, 1.65, 3.76, 1.62, BLUE, LIGHT)
    add_callout(slide, "贡献 3 · 在线执行", "实现 active/desired/allocation 边界与 clean/dirty/pin 退休协议，小规模实测满足 3% 红线。", 8.82, 1.65, 3.76, 1.62, GREEN, PALE_GREEN)
    add_text(slide, "下一步优先级", 0.82, 3.82, 2.25, 0.38, size=20, color=INK, bold=True)
    next_steps = [
        ("01", "打通 WLM 动态额度交换", "接入 maxChunksPerProcess runtime quota、admission queue 和 running AP graceful debt。", RED),
        ("02", "补齐内核可观测性与故障验证", "状态视图、resize generation、checkpoint/dirty/pin 超时、crash recovery。", ORANGE),
        ("03", "完成生产规模五阶段闭环", "256MB 粒度、大 SB、连续 TP+AP 五阶段；归一化稳态 TPS 波动 ≤5%，切换单秒抖动 ≤3%。", TEAL),
    ]
    y = 4.34
    for code, title, body, color in next_steps:
        add_rect(slide, 0.84, y, 0.70, 0.56, color, radius=True)
        add_text(slide, code, 0.88, y + 0.15, 0.62, 0.25, size=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, title, 1.80, y - 0.02, 3.33, 0.31, size=15, color=color, bold=True)
        add_text(slide, body, 5.05, y - 0.01, 7.29, 0.44, size=12)
        y += 0.72
    add_footer(slide, 15, "本项目最新结果汇总，2026-07-27")
    add_notes(slide, "收束时分清已完成和待完成：预测推荐已可演示，生产级运行时闭环仍在工程化。")
    outline.append("15. 结论与下一步：WLM 额度交换、故障验证、生产规模五阶段闭环。")

    prs.save(OUT)
    OUTLINE.write_text(
        "# Huawei5 内存自治论文式汇报提纲\n\n"
        + "\n".join(f"- {line}" for line in outline)
        + "\n\n## 口径说明\n\n"
        + "- 五阶段结论限定为已测挑战网格，不能表述为连续参数空间的全局最优证明。\n"
        + "- 在线缩容结论限定为 128MB 隔离原型，不能表述为完整五阶段生产内核验收。\n"
        + "- 完全未见 Query 的 21 点 spill 分类为 15/21；v2 的 6/6 是一次新 Query 锚点协议。\n",
        encoding="utf-8",
    )
    print(OUT)
    print(OUTLINE)


if __name__ == "__main__":
    build()
