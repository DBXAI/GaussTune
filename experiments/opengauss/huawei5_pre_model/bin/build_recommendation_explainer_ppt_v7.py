#!/usr/bin/env python3
"""Build the text-first Huawei5 recommendation-mechanics presentation."""

from __future__ import annotations

import csv
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches

from build_model_progress_ppt import (
    BLUE,
    GRAY,
    GREEN,
    INK,
    LIGHT,
    MID,
    ORANGE,
    RED,
    TEAL,
    WHITE,
    add_bullets,
    add_kpi,
    add_rect,
    add_table,
    add_text,
    add_title,
    set_bg,
)


ROOT = Path(__file__).resolve().parents[1]
RESULTS = ROOT / "results"
ARTIFACTS = ROOT / "artifacts"
LATEST = ARTIFACTS / "00_latest"
FIGURES = ARTIFACTS / "01_current_joint_model/figures/five_stage_plan_aware"
OUT = LATEST / "Huawei5_plan_aware_recommendation_v10_20260724.pptx"
OUTLINE = LATEST / "Huawei5_plan_aware_recommendation_v10_20260724_outline.md"
QUERY_ACCURACY_FIGURE = ARTIFACTS / "01_current_joint_model/figures/all_query_workmem_prediction_vs_actual.png"
PLAN_VALIDATION_FIGURE = ARTIFACTS / "01_current_joint_model/figures/plan_aware_heldout_spill_validation.png"


def read_csv(path: Path) -> list[dict[str, str]]:
    with path.open(newline="", encoding="utf-8") as fh:
        return list(csv.DictReader(fh))


def footer(slide, page: int) -> None:
    add_rect(slide, 0.62, 7.17, 12.08, 0.012, MID)
    add_text(slide, "Huawei5 | Plan-aware SB/work_mem replay | 2026-07-24", 0.62, 7.22, 7.3, 0.18, size=8, color=GRAY)
    add_text(slide, str(page), 12.2, 7.22, 0.5, 0.18, size=8, color=GRAY, align=PP_ALIGN.RIGHT)


def conclusion(slide, text: str, color=TEAL, y: float = 1.58) -> None:
    add_rect(slide, 0.68, y, 11.96, 0.64, color, radius=True)
    add_text(slide, text, 0.94, y + 0.16, 11.44, 0.29, size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)


def info_box(slide, title: str, body: str, x: float, y: float, w: float, h: float, color) -> None:
    add_rect(slide, x, y, w, h, LIGHT, line=color, radius=True)
    add_text(slide, title, x + 0.18, y + 0.18, w - 0.36, 0.36, size=18, color=color, bold=True)
    add_text(slide, body, x + 0.18, y + 0.68, w - 0.36, h - 0.83, size=13, color=INK)


def build() -> None:
    LATEST.mkdir(parents=True, exist_ok=True)
    recs = read_csv(RESULTS / "plan_aware_replay_20260724/replay_expanded/stage_joint_recommendations.csv")
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    outline: list[str] = []

    # 1 Cover
    slide = prs.slides.add_slide(blank)
    set_bg(slide, INK)
    add_rect(slide, 0, 0, 0.18, 7.5, TEAL)
    add_text(slide, "Huawei5 配置推荐模型", 0.94, 1.34, 10.8, 0.72, size=39, color=WHITE, bold=True)
    add_text(slide, "SB 和 work_mem 的推荐值是怎么得到的？", 0.96, 2.34, 11.2, 0.66, size=29, color=RGBColor(181, 219, 222), bold=True)
    add_text(slide, "模型工作流程 · 二维联合选择 · 推荐值含义 · 五阶段结果", 0.97, 3.42, 10.2, 0.38, size=18, color=WHITE)
    add_rect(slide, 0.97, 4.20, 4.35, 0.52, TEAL, radius=True)
    add_text(slide, "核心：推荐最小有效配置，而不是最大内存", 1.10, 4.33, 4.08, 0.25, size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "2026-07-24", 0.97, 6.23, 2.0, 0.25, size=12, color=RGBColor(194, 201, 206))
    outline.append("1. 封面：SB/work_mem 推荐值如何得到。")

    # 2 Meaning
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, "先理解推荐值的含义", "推荐值不是“越大越好”，而是达到目标后不再浪费内存的最小配置。", "01 / Meaning")
    conclusion(slide, "模型寻找的是性能平台起点和可接受 I/O 下的最小内存点。", GREEN)
    info_box(slide, "推荐 SB", "TP 热页在数据库缓存中的命中率已经进入平台；继续增大 SB，TPS 或 TP-SB hit 不再明显提高。", 0.76, 2.62, 3.72, 2.18, TEAL)
    info_box(slide, "推荐 work_mem", "关键 Join/Agg/Sort 的 spill 已消除，或在无法全 no-spill 时达到可部署的受控 spill 点。", 4.81, 2.62, 3.72, 2.18, ORANGE)
    info_box(slide, "推荐组合", "同时满足 TP 性能、spill I/O、系统内存安全；同等性能下选择 SB + 动态内存占用更小的组合。", 8.86, 2.62, 3.72, 2.18, GREEN)
    add_bullets(slide, ["它表示模型认为最合算的配置，不等于把所有内存参数设到最大", "推荐点之外仍可能有相同性能的平台点，但会消耗更多内存", "只有完成真实二维 TPS 复跑后，才能称为已验证的全局最优组合"], 1.00, 5.24, 11.20, 1.24, size=15, spacing=9)
    footer(slide, 2)
    outline.append("2. 推荐值含义：平台起点、最小 no-spill、Pareto 最小内存。")

    # 3 Inputs and complete flow
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, "模型依靠什么数据工作？", "它是 trace replay，不是读取实际最优配置训练 TPS 回归器。", "02 / Inputs")
    data = [
        ["输入", "记录内容", "在模型中的作用"],
        ["页面访问 trace", "TP/AP 访问页、时间、BulkRead 策略", "重放不同 SB 下的数据库缓存命中"],
        ["算子内存 trace", "Join/Agg/Sort 分配、batch、spill、生命周期", "预测不同 work_mem 的峰值和 spill"],
        ["EXPLAIN 计划族", "不同 work_mem 下的执行路径", "防止跨计划错误外推"],
        ["系统内存实验", "SB、动态内存与 MemAvailable/RSS", "计算 OS cache 余量和内存安全"],
    ]
    add_table(slide, 5, 3, data, 0.70, 1.72, 11.92, 3.35, widths=[2.20, 4.55, 5.17], font_size=12)
    add_rect(slide, 0.90, 5.46, 11.54, 0.94, LIGHT, line=INK, radius=True)
    add_text(slide, "每个候选 (SB, work_mem) 都重新经过：计划匹配 → 算子 replay → 动态内存/spill → SB/OS replay → TP miss/I/O → 推荐规则", 1.20, 5.72, 10.94, 0.42, size=16, color=INK, bold=True, align=PP_ALIGN.CENTER)
    footer(slide, 3)
    outline.append("3. 模型输入：页面 trace、算子 trace、计划族和系统内存实验。")

    # 4 SB recommendation
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, "SB 推荐值是怎么得到的？", "mixed trace 参与状态变化，但只对 TP 页计分。", "03 / SB Recommendation")
    conclusion(slide, "对每个 SB 重放同一段访问轨迹，选择达到最大 TP-SB hit 99% 的最小 SB。", TEAL)
    steps = [
        ("1", "候选 SB", "128 / 256 / 512 / 1024 / 1504 / 2048 / 4096 / 8192MB"),
        ("2", "重放 mixed trace", "TP 与 AP 都进入 SB replacement 状态；BulkRead 使用 ring 行为。"),
        ("3", "只统计 TP", "输出 TP-SB hit、进入 OS 的 TP miss，以及 Linux refault 结果。"),
        ("4", "寻找平台起点", "满足 TP-SB hit ≥ 99% × 本阶段最大值的最小 SB。"),
    ]
    for index, (number, title, body) in enumerate(steps):
        y = 2.58 + index * 0.83
        add_rect(slide, 0.94, y, 0.52, 0.52, TEAL, radius=True)
        add_text(slide, number, 0.95, y + 0.12, 0.50, 0.24, size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, title, 1.72, y - 0.01, 2.35, 0.34, size=17, color=TEAL, bold=True)
        add_text(slide, body, 4.15, y - 0.01, 7.85, 0.52, size=14, color=INK)
    add_rect(slide, 1.08, 6.08, 11.10, 0.52, INK, radius=True)
    add_text(slide, "例如 S5：512MB TP-SB hit=97.06%，1024MB=99.39%；1024MB 首次进入平台，因此推荐 1024MB。", 1.34, 6.22, 10.58, 0.25, size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    footer(slide, 4)
    outline.append("4. SB 推荐：TP-only trace replay 和 99% 平台起点。")

    # 5 work_mem recommendation
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, "work_mem 推荐值是怎么得到的？", "work_mem 是算子级上限，不等于整条查询只使用这一份内存。", "04 / work_mem Recommendation")
    conclusion(slide, "重放每个算子的内存增长和生命周期，寻找关键算子不再 spill 的最小候选值。", ORANGE)
    operator_table = [
        ["算子", "trace/replay 内容", "work_mem 不足时"],
        ["Hash Join", "tuple bytes、bucket、batch", "batch 增加并写读临时分区"],
        ["HashAggregate", "group 数、每组分配", "未容纳 group spill"],
        ["Sort", "tuple chunk、memtuple、merge pass", "external merge 多轮 I/O"],
    ]
    add_table(slide, 4, 3, operator_table, 0.72, 2.55, 7.42, 2.62, widths=[1.72, 2.95, 2.75], font_size=12)
    add_rect(slide, 8.48, 2.55, 4.05, 2.62, LIGHT, line=ORANGE, radius=True)
    add_text(slide, "推荐过程", 8.78, 2.83, 3.45, 0.34, size=20, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
    add_bullets(slide, ["检查该 work_mem 的计划族", "按 operator start/end 计算并发峰值", "计算 spill 临时数据与读写 I/O", "选择首个 no-spill 或受控 spill 点"], 8.82, 3.33, 3.32, 1.50, size=13, spacing=7)
    add_rect(slide, 1.08, 5.73, 11.10, 0.72, INK, radius=True)
    add_text(slide, "例如 S5：1137MB 仍预测 1.67GiB spill；1150MB 变为 0；1208MB 同样为 0 但多占内存，因此推荐 1150MB。", 1.34, 5.94, 10.58, 0.31, size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    footer(slide, 5)
    outline.append("5. work_mem 推荐：算子级 replay、生命周期峰值、首个 no-spill 点。")

    # 6 coupling
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, "SB 和 work_mem 为什么必须联合预测？", "二者没有直接等式，但共同争用 RAM，并通过 OS cache 和 spill 相互影响。", "05 / Coupling")
    conclusion(slide, "模型不是分别算完再取交集，而是对每一个二维组合重新计算系统状态。", GREEN)
    info_box(slide, "SB 增大", "数据库缓存命中提高；同时固定内存占用增加，留给 OS page cache 和动态内存的空间减少。", 0.72, 2.52, 3.72, 1.68, TEAL)
    info_box(slide, "work_mem 增大", "Join/Agg/Sort spill 减少；同时并发动态内存峰值增加，也会压缩 OS page cache。", 4.80, 2.52, 3.72, 1.68, ORANGE)
    info_box(slide, "spill 的反馈", "临时文件增加物理 I/O，并以 streaming 页进入 Linux cache，可能淘汰普通页面。", 8.88, 2.52, 3.72, 1.68, RED)
    add_rect(slide, 0.90, 4.62, 11.54, 1.12, LIGHT, line=INK, radius=True)
    add_text(slide, "对候选 (SB, W) 计算", 1.18, 4.89, 2.50, 0.34, size=17, color=INK, bold=True)
    add_text(slide, "DynamicPeak(W)", 3.70, 4.89, 1.75, 0.34, size=15, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "+", 5.48, 4.89, 0.35, 0.34, size=18, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "OSCapacity(SB,W)", 5.85, 4.89, 2.10, 0.34, size=15, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "→", 8.03, 4.89, 0.42, 0.34, size=18, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "TP disk I/O + spill I/O + memory safety", 8.47, 4.89, 3.58, 0.34, size=14, color=GREEN, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "MemAvailable = 23546.38 - 0.29220 × SB - 0.41804 × dynamic_peak  (MB)", 1.12, 6.07, 11.10, 0.34, size=15, color=GRAY, bold=True, align=PP_ALIGN.CENTER)
    footer(slide, 6)
    outline.append("6. 联合关系：RAM 竞争、OS cache 和 spill 反馈。")

    # 7 selection
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, "二维推荐算法怎样从 192 个候选中选出一个？", "每一步都减少候选集合，最后选择最小有效组合。", "06 / Selection")
    selection = [
        ["步骤", "规则", "目的"],
        ["1", "排除没有同计划族 trace 的点", "不跨执行路径外推"],
        ["2", "排除 MemAvailable 低于 3.2GB 的点", "防止回收、OOM 和主机抖动"],
        ["3", "保留 TP-SB hit 达到阶段最大值 99% 的点", "进入 TP 性能平台"],
        ["4", "在物理 I/O 与内存占用上求 Pareto 前沿", "平衡 TP miss、spill 和 RAM"],
        ["5", "近最小 I/O 候选中选择内存占用最小者", "得到推荐配置"],
    ]
    add_table(slide, 6, 3, selection, 0.80, 1.78, 11.70, 4.50, widths=[1.15, 6.20, 4.35], font_size=13)
    add_text(slide, "结果：模型推荐的是“当前证据范围内的 Pareto 最优”，不是保证所有未知计划路径上的数学全局最优。", 1.00, 6.48, 11.35, 0.38, size=14, color=RED, bold=True, align=PP_ALIGN.CENTER)
    footer(slide, 7)
    outline.append("7. 二维选择规则：计划、内存安全、TP 平台、Pareto、最小内存。")

    # 8 concrete example
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, "S5 算例：为什么最终是 1024MB + 1150MB？", "不是机械取交集，而是比较同一张二维候选表。", "07 / Example")
    example = [
        ["SB + work_mem", "TP-SB hit", "spill I/O", "预测物理 I/O", "总内存占用", "判定"],
        ["512 + 1137", "97.06%", "1.67GiB", "4.50GiB", "4.80GiB", "TP 未到平台"],
        ["1024 + 1137", "99.39%", "1.67GiB", "4.50GiB", "5.30GiB", "仍有 spill"],
        ["1024 + 1150", "99.39%", "0", "2.83GiB", "5.31GiB", "推荐"],
        ["2048 + 1150", "99.80%", "0", "2.83GiB", "6.31GiB", "I/O 不降，多占 1GB"],
    ]
    add_table(slide, 5, 6, example, 0.62, 1.78, 12.08, 3.50, widths=[1.78, 1.52, 1.48, 1.85, 1.83, 3.62], font_size=11)
    add_rect(slide, 0.92, 5.70, 11.54, 0.75, GREEN, radius=True)
    add_text(slide, "1024 + 1150：已经进入 TP 平台，spill 从 1.67GiB 降为 0；继续加到 2048MB SB 不降 I/O，只多占 1GB。", 1.18, 5.92, 11.02, 0.33, size=15, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    footer(slide, 8)
    outline.append("8. S5 候选表逐项推出 1024MB + 1150MB。")

    # 9 stage recommendations and meaning
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, "五阶段最终推荐值及其含义", "每个数字都对应一个具体的 trace replay 结论。", "08 / Output")
    names = {"stage1_memory_rich": "S1", "stage2_reach_limit": "S2", "stage3_protect_tp": "S3", "stage4_backpressure": "S4", "stage5_tp_surge": "S5"}
    meaning = {
        "S1": "256MB 进入 TP 平台；1MB 已无 spill",
        "S2": "256MB 进入 TP 平台；Q3 在 1150MB 首次 no-spill",
        "S3": "512MB 进入 TP 平台；1083MB 覆盖关键算子",
        "S4": "256MB 达 TP 目标；6500MB 是 Q9 换 Plan 前的最低 I/O 点",
        "S5": "1024MB 命中最高 TPS 平台；1150MB 首次消除关键 spill",
    }
    table = [["阶段", "推荐 SB", "推荐 work_mem", "推荐值的实际含义"]]
    for row in recs:
        short = names[row["stage"]]
        table.append([short, f"{row['recommended_sb_mb']}MB", f"{row['recommended_work_mem_mb']}MB", meaning[short]])
    add_table(slide, 6, 4, table, 0.70, 1.75, 11.92, 4.45, widths=[1.05, 1.70, 2.20, 6.97], font_size=13)
    add_text(slide, "S4 已扩展到实测可部署上界 7140MB；6500MB 后 Q9 换 Plan，spill 反而上升。", 1.00, 6.42, 11.35, 0.38, size=14, color=GREEN, bold=True, align=PP_ALIGN.CENTER)
    footer(slide, 9)
    outline.append("9. 五阶段推荐及每个值的含义。")

    # 10 evidence
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, "这些推荐目前有多少验证证据？", "把命中率、内存边界和 TPS 验证分开看。", "09 / Evidence")
    add_kpi(slide, "0.50pp", "combined hit held-out MAE", 0.78, 1.75, 2.55, color=TEAL, note="SB 0.61pp / OS 1.82pp")
    add_kpi(slide, "7 / 7", "Plan-aware held-out 命中", 3.67, 1.75, 2.55, color=ORANGE, note="Plan 和 spill 分类均正确")
    add_kpi(slide, "1024MB", "S5 预测与实测 TPS 平台", 6.56, 1.75, 2.55, color=GREEN, note="TPS regret 0%")
    add_kpi(slide, "6500MB", "S4 完整域推荐", 9.45, 1.75, 2.55, color=GREEN, note="7141MB 新 Plan 实测引擎失败")
    add_bullets(slide, ["Q5/Q7/Q9/Q21：7 个 held-out 点的 Plan 与 spill 分类全部命中", "Q7/Q9/Q21：有 spill 点的 temp I/O 预测与实测接近重合", "S4：6500MB 后 Q9 切换 Plan，stage spill 从约 15.0GiB 回升到 23.6GiB", "Q21：7140MB 实测成功，7141MB 新 Plan 实测触发 2GiB 非法分配"], 1.00, 3.52, 11.15, 2.15, size=16, spacing=13)
    add_rect(slide, 1.02, 6.08, 11.24, 0.52, INK, radius=True)
    add_text(slide, "后续最关键实验：对推荐组合及相邻组合做同并发、同阶段的真实二维 TPS 验证。", 1.28, 6.22, 10.72, 0.25, size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    footer(slide, 10)
    outline.append("10. 验证证据和尚缺的真实二维 TPS 复跑。")

    # 11 Plan-aware held-out validation
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, "Plan 改变后还能预测 spill 吗？", "锚点和验证点严格分开；蓝柱为预测，绿柱为实际 temp read + write。", "10 / Plan-aware Validation")
    slide.shapes.add_picture(str(PLAN_VALIDATION_FIGURE), Inches(0.67), Inches(1.48), width=Inches(12.00), height=Inches(5.05))
    add_text(slide, "Q5/Q7/Q9/Q21：Plan 7/7、spill 分类 7/7；有 spill 点的 I/O 柱基本重合。", 1.00, 6.55, 11.35, 0.34, size=13, color=GREEN, bold=True, align=PP_ALIGN.CENTER)
    footer(slide, 11)
    outline.append("11. 新 Plan 族 held-out 验证：Plan、spill 分类和 I/O 数量。")

    # 12 All-query work_mem accuracy
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, "所有 AP Query：预测 work_mem 与实测值差多少？", "蓝点为预测，绿方块为实测；连线长度表示误差。", "10 / Query Accuracy")
    slide.shapes.add_picture(str(QUERY_ACCURACY_FIGURE), Inches(0.67), Inches(1.55), width=Inches(12.00), height=Inches(4.98))
    add_text(slide, "Q18/Q21 没有可实现的实际 no-spill 值，因此标为 host/engine infeasible，不计入 5/6 准确率分母。", 1.00, 6.57, 11.35, 0.34, size=12, color=RED, bold=True, align=PP_ALIGN.CENTER)
    footer(slide, 12)
    outline.append("12. 所有 AP Query 的预测/实测最小 no-spill work_mem 对比。")

    # 13-17 figures
    figure_slides = [
        ("S1", "推荐 256MB SB + 1MB work_mem；TPS 被限速，SB 最优无法识别；1MB 已足够无 spill。"),
        ("S2", "推荐 256MB SB + 1150MB work_mem；TPS 被限速；Q3 的 1150MB no-spill 边界逐 MB 验证。"),
        ("S3", "推荐 512MB SB + 1083MB work_mem；TPS 被限速但最低 P95 在 512MB；work_mem 边界验证通过。"),
        ("S4", "推荐 256MB SB + 6500MB work_mem；6500MB 后 Q9 换 Plan，spill 回升；7141MB 新 Plan 实测不可部署。"),
        ("S5", "推荐 1024MB SB + 1150MB work_mem；1024MB 命中实测最高 TPS；1150MB 为首个零 spill 支持点。"),
    ]
    for page, (short, explanation) in enumerate(figure_slides, start=13):
        slide = prs.slides.add_slide(blank)
        set_bg(slide)
        add_title(slide, f"{short} 推荐结果", "左侧：真实 TPS 与推荐 SB；右侧：预测 spill 与推荐 work_mem。", "Result")
        image = FIGURES / f"{short.lower()}_simple_tps_workmem_recommendation.png"
        slide.shapes.add_picture(str(image), Inches(0.62), Inches(1.52), width=Inches(12.10), height=Inches(4.28))
        add_rect(slide, 0.90, 6.08, 11.54, 0.58, INK, radius=True)
        add_text(slide, explanation, 1.14, 6.23, 11.06, 0.27, size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        footer(slide, page)
        outline.append(f"{page}. {short} TPS/spill 推荐图。")

    prs.save(OUT)
    OUTLINE.write_text("# 推荐模型说明版 PPT 大纲\n\n" + "\n".join(f"- {line}" for line in outline) + "\n", encoding="utf-8")
    print(OUT)
    print(OUTLINE)


if __name__ == "__main__":
    build()
