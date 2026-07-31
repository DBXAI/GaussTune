#!/usr/bin/env python3
"""Build a concise presentation focused on model mechanics and current errors."""

from __future__ import annotations

import csv
import json
from collections import defaultdict
from pathlib import Path

import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches

from build_model_progress_ppt import (
    ARTIFACTS,
    BLUE,
    FONT,
    GRAY,
    GREEN,
    INK,
    LIGHT,
    MID,
    ORANGE,
    RED,
    TEAL,
    WHITE,
    add_arrow,
    add_bullets,
    add_flow_box,
    add_footer,
    add_kpi,
    add_rect,
    add_table,
    add_text,
    add_title,
    rgb_hex,
    set_bg,
)


ROOT = Path(__file__).resolve().parents[1]
VALIDATION = ROOT / "results" / "sb_recommendation_validation_20260711_234537"
OUT = ARTIFACTS / "Huawei5_original_load_performance_validation_v6_tp_only_20260716.pptx"
OUTLINE = ARTIFACTS / "Huawei5_original_load_performance_validation_v6_tp_only_20260716_outline.md"
S5_RESULT_DIR = ROOT / "results" / "s5_tp_protected_full_binary_20260716"
S5_UPGRADE = S5_RESULT_DIR / "s5_tp_protected_predictions.csv"
S5_METRICS = S5_RESULT_DIR / "s5_tp_protected_metrics.json"
ORIGINAL_LOAD_PERF = VALIDATION / "original_load_stage_tp_latency.csv"
TP_ONLY_ALIGNMENT = ROOT / "results" / "tp_only_performance_alignment_20260716"
TP_ONLY_SUMMARY = TP_ONLY_ALIGNMENT / "tp_only_performance_summary.csv"
TP_ONLY_CHART = ARTIFACTS / "tp_only_replay_performance_alignment_20260716.png"


def read_csv(path: Path):
    with path.open(newline="", encoding="utf-8") as fh:
        return list(csv.DictReader(fh))


def make_error_chart(rows, s5_metrics, s5_rows, path: Path):
    by_stage = defaultdict(list)
    for row in rows:
        by_stage[row["stage"]].append(row)
    order = [
        "stage1_memory_rich",
        "stage2_reach_limit",
        "stage3_protect_tp",
        "stage4_backpressure",
        "stage5_tp_surge",
    ]
    sb = [sum(abs(float(r["sb_err_pp"])) for r in by_stage[s]) / len(by_stage[s]) for s in order]
    os_hit = [sum(abs(float(r["os_err_pp"])) for r in by_stage[s]) / len(by_stage[s]) for s in order]
    sb[-1] = sum(
        abs(float(row["protected_sb"]) - float(row["actual_sb"])) * 100
        for row in s5_rows
    ) / len(s5_rows)
    old_s5_os = os_hit[-1]
    os_hit[-1] = float(s5_metrics["protected_os_mae_pp"])
    xs = list(range(5))
    fig, ax = plt.subplots(figsize=(9.6, 4.0))
    ax.bar([x - 0.18 for x in xs], sb, width=0.36, color=rgb_hex(TEAL), label="SB hit MAE")
    ax.bar([x + 0.18 for x in xs], os_hit, width=0.36, color=rgb_hex(ORANGE), label="OS hit MAE (S5 upgraded)")
    ax.scatter([xs[-1] + 0.18], [old_s5_os], marker="x", s=90, linewidths=2.2, color=rgb_hex(RED), label="S5 OS MAE before upgrade", zorder=4)
    ax.set_xticks(xs, ["S1", "S2", "S3", "S4", "S5"])
    ax.set_ylabel("mean absolute error (percentage points)")
    ax.set_ylim(0, 27)
    ax.grid(axis="y", alpha=0.2)
    ax.legend(frameon=False, ncol=2, loc="upper left")
    for x, value in zip(xs, sb):
        ax.text(x - 0.18, value + 0.45, f"{value:.2f}", ha="center", fontsize=9)
    for x, value in zip(xs, os_hit):
        ax.text(x + 0.18, value + 0.45, f"{value:.2f}", ha="center", fontsize=9)
    ax.annotate(
        f"old {old_s5_os:.2f}",
        xy=(xs[-1] + 0.18, old_s5_os),
        xytext=(xs[-1] - 0.62, old_s5_os - 2.2),
        arrowprops={"arrowstyle": "->", "color": rgb_hex(RED)},
        color=rgb_hex(RED),
        fontsize=9,
    )
    fig.tight_layout()
    fig.savefig(path, dpi=180)
    plt.close(fig)
    return sum(sb) / len(sb), sum(os_hit) / len(os_hit)


def make_stage_accuracy_chart(rows, path: Path):
    rows = sorted(rows, key=lambda row: int(row["sb_mb"]))
    sbs = [int(row["sb_mb"]) for row in rows]
    fig, axes = plt.subplots(1, 3, figsize=(13.0, 4.15))
    series = [
        ("SB hit rate", "actual_sb", "pred_sb", TEAL),
        ("OS conditional hit rate", "actual_os", "pred_os", ORANGE),
        ("Combined hit rate", "actual_combined", "pred_combined", GREEN),
    ]
    metrics = {}
    for ax, (title, actual_key, pred_key, color) in zip(axes, series):
        actual = [float(row[actual_key]) * 100 for row in rows]
        predicted = [float(row[pred_key]) * 100 for row in rows]
        mae = sum(abs(a - p) for a, p in zip(actual, predicted)) / len(rows)
        metrics[title] = mae
        ax.plot(sbs, actual, color=rgb_hex(INK), marker="o", linewidth=2.2, label="Actual")
        ax.plot(sbs, predicted, color=rgb_hex(color), marker="s", linewidth=2.0, linestyle="--", label="Predicted")
        ax.set_xscale("log", base=2)
        ax.set_xticks(sbs, [str(value) for value in sbs], rotation=35)
        low = min(actual + predicted)
        high = max(actual + predicted)
        padding = max(1.5, (high - low) * 0.14)
        ax.set_ylim(max(0, low - padding), min(100, high + padding))
        ax.set_title(f"{title}\nMAE {mae:.2f}pp", fontsize=11, fontweight="bold")
        ax.set_xlabel("shared_buffers (MB)")
        ax.set_ylabel("hit rate (%)")
        ax.grid(alpha=0.2)
        ax.legend(frameon=False, fontsize=8, loc="best")

    actual_best = max(rows, key=lambda row: float(row["actual_combined"]))
    pred_best = max(rows, key=lambda row: float(row["pred_combined"]))
    axes[2].axvline(int(pred_best["sb_mb"]), color=rgb_hex(TEAL), linewidth=1.3, linestyle="--", alpha=0.7)
    axes[2].axvline(int(actual_best["sb_mb"]), color=rgb_hex(GREEN), linewidth=1.3, linestyle=":", alpha=0.9)
    axes[2].text(
        0.03,
        0.04,
        f"pred best: {pred_best['sb_mb']}MB\nactual best: {actual_best['sb_mb']}MB",
        transform=axes[2].transAxes,
        fontsize=8,
        color=rgb_hex(INK),
        bbox={"boxstyle": "round,pad=0.3", "facecolor": "white", "edgecolor": rgb_hex(MID), "alpha": 0.9},
    )
    fig.tight_layout()
    fig.savefig(path, dpi=190)
    plt.close(fig)
    metrics["pred_best_sb_mb"] = int(pred_best["sb_mb"])
    metrics["actual_best_sb_mb"] = int(actual_best["sb_mb"])
    metrics["regret_pp"] = (
        float(actual_best["actual_combined"]) - float(pred_best["actual_combined"])
    ) * 100
    return metrics


def make_regret_chart(summary, path: Path):
    stages = ["S1", "S2", "S3", "S4", "S5"]
    regret = [
        (float(row["actual_best_combined"]) - float(row["actual_combined_at_pred_best"])) * 100
        for row in summary
    ]
    fig, ax = plt.subplots(figsize=(8.4, 2.35))
    bars = ax.barh(stages[::-1], regret[::-1], color=[rgb_hex(RED), rgb_hex(TEAL), rgb_hex(TEAL), rgb_hex(RED), rgb_hex(ORANGE)])
    ax.set_xlabel("actual combined-hit loss at recommended SB (pp)")
    ax.set_xlim(0, 3.8)
    ax.grid(axis="x", alpha=0.2)
    for bar, value in zip(bars, regret[::-1]):
        ax.text(value + 0.05, bar.get_y() + bar.get_height() / 2, f"{value:.2f}", va="center", fontsize=9)
    fig.tight_layout()
    fig.savefig(path, dpi=180)
    plt.close(fig)


def make_tps_chart(rows, hit_best, path: Path):
    by_stage = defaultdict(list)
    for row in rows:
        by_stage[row["stage"]].append(row)
    stages = [
        ("stage1_memory_rich", "S1", TEAL),
        ("stage2_reach_limit", "S2", BLUE),
        ("stage3_protect_tp", "S3", GREEN),
        ("stage4_backpressure", "S4", ORANGE),
        ("stage5_tp_surge", "S5", RED),
    ]
    fig, axes = plt.subplots(2, 3, figsize=(12.2, 6.0))
    for ax, (stage, short, color) in zip(axes.flat, stages):
        stage_rows = sorted(by_stage[stage], key=lambda row: int(row["sb_mb"]))
        sbs = [int(row["sb_mb"]) for row in stage_rows]
        tps = [float(row["tps"]) for row in stage_rows]
        best_idx = max(range(len(stage_rows)), key=lambda idx: tps[idx])
        best_sb = sbs[best_idx]
        ax.plot(sbs, tps, color=rgb_hex(color), marker="o", linewidth=2.2)
        ax.axvline(hit_best[stage], color=rgb_hex(TEAL), linestyle="--", linewidth=1.4, label="combined recommendation")
        ax.axvline(best_sb, color=rgb_hex(INK), linestyle=":", linewidth=1.5, label="TPS best")
        ax.set_xscale("log", base=2)
        ax.set_xticks(sbs, [str(value) for value in sbs], rotation=35)
        ax.set_title(f"{short}: TPS best {best_sb}MB", fontsize=11, fontweight="bold")
        ax.set_xlabel("shared_buffers (MB)")
        ax.set_ylabel("TPS")
        ax.grid(alpha=0.2)
        ax.legend(frameon=False, fontsize=7, loc="best")
    axes.flat[-1].axis("off")
    fig.tight_layout()
    fig.savefig(path, dpi=190)
    plt.close(fig)


def make_original_latency_chart(rows, hit_best, path: Path):
    by_stage = defaultdict(list)
    for row in rows:
        by_stage[row["stage"]].append(row)
    stages = [
        ("stage1_memory_rich", "S1", TEAL),
        ("stage2_reach_limit", "S2", BLUE),
        ("stage3_protect_tp", "S3", GREEN),
        ("stage4_backpressure", "S4", ORANGE),
        ("stage5_tp_surge", "S5 low stream", RED),
    ]
    fig, axes = plt.subplots(2, 3, figsize=(12.2, 6.0))
    for ax, (stage, short, color) in zip(axes.flat, stages):
        stage_rows = sorted(by_stage[stage], key=lambda row: int(row["sb_mb"]))
        sbs = [int(row["sb_mb"]) for row in stage_rows]
        p95 = [float(row["latency_p95_ms"]) for row in stage_rows]
        best_idx = min(range(len(stage_rows)), key=lambda idx: p95[idx])
        best_sb = sbs[best_idx]
        ax.plot(sbs, p95, color=rgb_hex(color), marker="o", linewidth=2.2)
        ax.axvline(hit_best[stage], color=rgb_hex(TEAL), linestyle="--", linewidth=1.4, label="combined recommendation")
        ax.axvline(best_sb, color=rgb_hex(INK), linestyle=":", linewidth=1.5, label="lowest p95")
        ax.set_xscale("log", base=2)
        ax.set_xticks(sbs, [str(value) for value in sbs], rotation=35)
        ax.set_title(f"{short}: lowest p95 at {best_sb}MB", fontsize=11, fontweight="bold")
        ax.set_xlabel("shared_buffers (MB)")
        ax.set_ylabel("TPCC p95 latency (ms)")
        ax.grid(alpha=0.2)
        ax.legend(frameon=False, fontsize=7, loc="best")
    axes.flat[-1].axis("off")
    fig.tight_layout()
    fig.savefig(path, dpi=190)
    plt.close(fig)


def build():
    rows = read_csv(VALIDATION / "recommendation_validation_by_stage.csv")
    summary = read_csv(VALIDATION / "recommendation_validation_summary.csv")
    s5_upgrade = read_csv(S5_UPGRADE)
    s5_metrics = json.loads(S5_METRICS.read_text(encoding="utf-8"))
    upgraded_s5 = max(s5_upgrade, key=lambda row: float(row["protected_combined"]))
    actual_s5 = max(s5_upgrade, key=lambda row: float(row["actual_combined"]))
    for row in summary:
        if row["stage"] != "stage5_tp_surge":
            continue
        row["pred_best_sb_mb"] = upgraded_s5["sb_mb"]
        row["pred_best_combined"] = upgraded_s5["protected_combined"]
        row["actual_best_sb_mb"] = actual_s5["sb_mb"]
        row["actual_best_combined"] = actual_s5["actual_combined"]
        row["actual_combined_at_pred_best"] = upgraded_s5["actual_combined"]
        row["matched"] = "yes" if upgraded_s5["sb_mb"] == actual_s5["sb_mb"] else "no"
    original_perf_rows = read_csv(ORIGINAL_LOAD_PERF)
    tp_only_summary = read_csv(TP_ONLY_SUMMARY)
    hit_best = {row["stage"]: int(row["pred_best_sb_mb"]) for row in summary}
    stage_defs = [
        ("stage1_memory_rich", "S1", "内存充足阶段 · Q1 · AP×1"),
        ("stage2_reach_limit", "S2", "接近缓存容量边界 · Q3 · AP×1"),
        ("stage3_protect_tp", "S3", "TP 热页保护阶段 · Q5/Q7 · AP×2"),
        ("stage4_backpressure", "S4", "高并发与 backpressure · Q9/Q13/Q18/Q21 · AP×4"),
        ("stage5_tp_surge", "S5", "TP surge · Q1/Q3/Q5/Q7 · AP×4"),
    ]
    current_stage_rows = {
        stage: [dict(row) for row in rows if row["stage"] == stage]
        for stage, _short, _description in stage_defs[:-1]
    }
    current_stage_rows["stage5_tp_surge"] = [
        {
            "stage": "stage5_tp_surge",
            "sb_mb": row["sb_mb"],
            "actual_sb": row["actual_sb"],
            "pred_sb": row["protected_sb"],
            "actual_os": row["actual_os"],
            "pred_os": row["protected_os"],
            "actual_combined": row["actual_combined"],
            "pred_combined": row["protected_combined"],
        }
        for row in s5_upgrade
    ]

    error_chart = ARTIFACTS / "ppt_v3_hit_rate_errors_with_s5_upgrade.png"
    regret_chart = ARTIFACTS / "ppt_v3_recommendation_regret.png"
    latency_chart = ARTIFACTS / "ppt_v5_original_load_p95_latency_by_stage.png"
    sb_overall_mae, upgraded_os_overall_mae = make_error_chart(rows, s5_metrics, s5_upgrade, error_chart)
    stage_charts = {}
    stage_metrics = {}
    for stage, short, _description in stage_defs:
        chart = ARTIFACTS / f"ppt_v3_{short.lower()}_prediction_vs_actual.png"
        stage_charts[stage] = chart
        stage_metrics[stage] = make_stage_accuracy_chart(current_stage_rows[stage], chart)
    make_regret_chart(summary, regret_chart)
    make_original_latency_chart(original_perf_rows, hit_best, latency_chart)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # 1. Cover
    slide = prs.slides.add_slide(blank)
    set_bg(slide, INK)
    add_rect(slide, 0, 0, 0.18, 7.5, TEAL)
    add_text(slide, "Huawei5 现有 SB 预测模型", 0.92, 1.45, 10.9, 0.75, size=36, color=WHITE, bold=True)
    add_text(slide, "工作原理、两类命中率误差与配置推荐效果", 0.94, 2.38, 10.8, 0.55, size=23, color=TEAL)
    add_text(slide, "阶段性结果 · 基于 8 个有效 shared_buffers 实测点", 0.94, 3.4, 7.2, 0.4, size=14, color=WHITE)
    add_text(slide, "2026-07-16", 0.94, 6.25, 2.5, 0.3, size=12, color=GRAY)

    # 2. Inputs and goal
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, "模型做什么：从一次真实访问轨迹预测多个 SB 配置", "现有模型不是监督学习模型，而是缓存机制 replay。", "01 / Model")
    add_rect(slide, 0.7, 1.72, 4.05, 4.78, LIGHT, line=GREEN)
    add_text(slide, "输入", 0.98, 2.0, 1.2, 0.4, size=22, color=GREEN, bold=True)
    add_bullets(slide, [
        "1504MB 下运行完整五阶段负载",
        "页面标识：relation + block number",
        "访问时间、backend、BufferAccessStrategy",
        "各阶段起止边界与可用 OS 内存",
        "候选 SB：128MB–24576MB",
    ], 0.98, 2.62, 3.35, 2.7, size=15)
    add_arrow(slide, 4.85, 4.05, 5.75, 4.05, TEAL)
    add_rect(slide, 5.75, 1.72, 2.95, 4.78, LIGHT, line=TEAL)
    add_text(slide, "Replay", 6.0, 2.0, 2.45, 0.4, size=22, color=TEAL, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "对每个候选 SB\n重新模拟同一页面访问序列", 6.02, 2.82, 2.4, 0.95, size=18, color=INK, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "DB Buffer\n+\nLinux Page Cache", 6.05, 4.25, 2.35, 1.15, size=18, color=TEAL, bold=True, align=PP_ALIGN.CENTER)
    add_arrow(slide, 8.78, 4.05, 9.62, 4.05, ORANGE)
    add_rect(slide, 9.62, 1.72, 3.0, 4.78, INK)
    add_text(slide, "输出", 9.88, 2.0, 2.48, 0.4, size=22, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_kpi(slide, "SB hit", "数据库缓存命中率", 9.92, 2.7, 2.35, color=GREEN)
    add_kpi(slide, "OS hit", "SB miss 后 OS 命中率", 9.92, 4.05, 2.35, color=ORANGE)
    add_text(slide, "最终选择 combined hit 最大的 SB", 9.92, 5.6, 2.35, 0.42, size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, 2)

    # 3. Stage and replay process
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, "模型如何工作：连续重放五个阶段", "阶段边界只重置统计计数，不重置模拟缓存状态。", "01 / Model")
    stages = [
        ("S1", "Q1\nAP×1", GREEN),
        ("S2", "Q3\nAP×1", BLUE),
        ("S3", "Q5/Q7\nAP×2", TEAL),
        ("S4", "Q9/Q13/Q18/Q21\nAP×4", ORANGE),
        ("S5", "Q1/Q3/Q5/Q7\nAP×4 + TP surge", RED),
    ]
    x = 0.72
    for idx, (name, detail, color) in enumerate(stages):
        add_flow_box(slide, name, x, 1.82, 2.15, 1.32, color, sub=detail)
        if idx < 4:
            add_arrow(slide, x + 2.17, 2.48, x + 2.42, 2.48, MID, 1.5)
        x += 2.48
    add_rect(slide, 0.72, 3.7, 12.0, 2.35, LIGHT, line=MID)
    process = [
        ("1", "读取访问", "按时间读取 page access"),
        ("2", "模拟 SB", "命中或淘汰页面"),
        ("3", "模拟 OS", "SB miss 进入 page cache"),
        ("4", "阶段统计", "分别累计两类 hit rate"),
    ]
    x = 0.98
    for idx, (num, title, detail) in enumerate(process):
        add_rect(slide, x, 4.08, 2.55, 1.35, WHITE, line=TEAL, radius=True)
        add_text(slide, num, x + 0.12, 4.22, 0.42, 0.4, size=20, color=TEAL, bold=True)
        add_text(slide, title, x + 0.55, 4.22, 1.7, 0.32, size=16, bold=True)
        add_text(slide, detail, x + 0.18, 4.75, 2.18, 0.34, size=11, color=GRAY, align=PP_ALIGN.CENTER)
        if idx < 3:
            add_arrow(slide, x + 2.55, 4.75, x + 2.82, 4.75, TEAL, 1.5)
        x += 2.9
    add_footer(slide, 3)

    # 4. Formula and recommendation
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, "模型如何形成最终推荐", "两个缓存层分别预测，再计算 combined hit。", "01 / Model")
    add_rect(slide, 0.72, 1.75, 5.82, 4.72, LIGHT, line=TEAL)
    add_text(slide, "Shared Buffer", 1.0, 2.02, 2.8, 0.4, size=21, color=TEAL, bold=True)
    add_bullets(slide, [
        "容量 = candidate SB / 8KB",
        "普通访问近似 clock sweep",
        "BulkRead 使用 ring 控制淘汰",
        "输出 pred_SB_hit",
    ], 1.0, 2.62, 4.85, 2.0, size=15)
    add_text(slide, "Linux Page Cache（S5 升级）", 1.0, 4.78, 4.2, 0.4, size=19, color=ORANGE, bold=True)
    add_text(slide, "首次访问只进入 inactive；二次访问或短距离 refault 才晋升 active。已知 BulkRead 路径保持在 streaming inactive，回收时优先淘汰。", 1.02, 5.22, 4.95, 0.9, size=12, color=INK, bold=True)
    add_rect(slide, 6.82, 1.75, 5.82, 4.72, INK)
    add_text(slide, "combined = 1 − (1−SB) × (1−OS)", 7.08, 2.35, 5.3, 0.5, size=20, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "对每个阶段、每个候选 SB 计算 combined hit", 7.25, 3.25, 4.95, 0.4, size=15, color=WHITE, align=PP_ALIGN.CENTER)
    add_arrow(slide, 9.72, 3.9, 9.72, 4.48, ORANGE)
    add_rect(slide, 7.62, 4.48, 4.2, 0.88, TEAL, radius=True)
    add_text(slide, "推荐 combined hit 最大的可行 SB", 7.78, 4.72, 3.88, 0.32, size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "注意：目前推荐目标是命中率，不是直接预测 TPS。", 7.35, 5.78, 4.75, 0.4, size=12, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, 4)

    # 5. Two hit-rate errors
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, "两种命中率的预测误差（已纳入 S5 升级）", "每阶段均使用 8 个 SB 实测点；红色叉号标出升级前的 S5 OS 误差。", "02 / Error")
    add_kpi(slide, f"{sb_overall_mae:.2f} pp", "SB hit 总体 MAE", 0.72, 1.58, 2.45, color=TEAL)
    add_kpi(slide, f"{upgraded_os_overall_mae:.2f} pp", "升级后 OS hit 总体 MAE", 3.22, 1.58, 2.45, color=ORANGE)
    add_kpi(slide, f"{s5_metrics['protected_os_mae_pp']:.2f} pp", "升级后 Stage 5 OS MAE", 5.72, 1.58, 2.55, color=RED)
    slide.shapes.add_picture(str(error_chart), Inches(0.7), Inches(2.82), width=Inches(8.25), height=Inches(3.45))
    add_rect(slide, 9.28, 1.7, 3.35, 4.58, LIGHT, line=MID)
    add_text(slide, "主要观察", 9.55, 2.0, 2.8, 0.38, size=20, color=INK, bold=True, align=PP_ALIGN.CENTER)
    add_bullets(slide, [
        "S1、S5 的 SB hit 本身较准",
        "S2、S4 的 SB hit 被明显高估",
        "S5 OS 误差由 24.68pp 降至 12.67pp",
        "升级后总体 OS MAE 由 7.27pp 降至 4.87pp",
        "SB 与 OS 误差可能在 combined 中抵消",
    ], 9.55, 2.56, 2.72, 2.72, size=12, spacing=7)
    add_text(slide, "S5 仍是 OS 层误差最大的阶段，但方向和推荐排序已经明显改善。", 9.58, 5.48, 2.66, 0.56, size=12, color=RED, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, 5)

    # 6-10. Per-stage prediction curves
    for idx, (stage, short, description) in enumerate(stage_defs):
        slide = prs.slides.add_slide(blank)
        set_bg(slide)
        model_note = "当前 replay 模型"
        if stage == "stage5_tp_surge":
            model_note = "升级后的 active/inactive/refault replay"
        add_title(
            slide,
            f"{short}：不同 SB 下预测值与真实值",
            f"{description}；黑色为真实值，彩色虚线为{model_note}。",
            "03 / Stage Curves",
        )
        slide.shapes.add_picture(
            str(stage_charts[stage]),
            Inches(0.62),
            Inches(1.58),
            width=Inches(12.1),
            height=Inches(3.86),
        )
        metrics = stage_metrics[stage]
        add_rect(slide, 0.72, 5.72, 12.0, 0.88, LIGHT, line=TEAL)
        add_text(
            slide,
            (
                f"SB MAE {metrics['SB hit rate']:.2f}pp    |    "
                f"OS MAE {metrics['OS conditional hit rate']:.2f}pp    |    "
                f"Combined MAE {metrics['Combined hit rate']:.2f}pp    |    "
                f"推荐 {metrics['pred_best_sb_mb']}MB / 实际最优 {metrics['actual_best_sb_mb']}MB    |    "
                f"推荐点损失 {metrics['regret_pp']:.2f}pp"
            ),
            0.92,
            5.98,
            11.6,
            0.34,
            size=13,
            color=INK,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        add_footer(slide, 6 + idx)

    # 11. Original-load latency curves
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(
        slide,
        "原负载验证：不同 SB 下的 TP 响应时间",
        "保持原始 2/12 terminals 与 40/180 TPS 限速；固定吞吐下比较 p95 延迟。",
        "04 / Performance Validation",
    )
    slide.shapes.add_picture(
        str(latency_chart),
        Inches(0.62),
        Inches(1.55),
        width=Inches(12.1),
        height=Inches(5.35),
    )
    add_footer(slide, 11)

    # 12. Hit recommendation versus original-load latency
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(
        slide,
        "Combined 推荐能否命中原负载性能最优点？",
        "S1-S4 的 TP 被固定在 40 TPS，因此使用 p95 延迟判断性能。",
        "04 / Performance Validation",
    )
    perf_by_stage = defaultdict(list)
    for row in original_perf_rows:
        perf_by_stage[row["stage"]].append(row)
    perf_table = [["阶段", "Combined 推荐", "最低 p95", "最优 p95", "推荐点 p95", "延迟差"]]
    for stage, short, _description in stage_defs:
        stage_rows = perf_by_stage[stage]
        best = min(stage_rows, key=lambda row: float(row["latency_p95_ms"]))
        recommended = next(row for row in stage_rows if int(row["sb_mb"]) == hit_best[stage])
        best_p95 = float(best["latency_p95_ms"])
        recommended_p95 = float(recommended["latency_p95_ms"])
        perf_table.append([
            short,
            f"{hit_best[stage]} MB",
            f"{best['sb_mb']} MB",
            f"{best_p95:.1f} ms",
            f"{recommended_p95:.1f} ms",
            f"{recommended_p95 - best_p95:+.1f} ms",
        ])
    table = add_table(
        slide,
        6,
        6,
        perf_table,
        0.72,
        1.68,
        8.2,
        4.3,
        widths=[0.8, 1.55, 1.35, 1.35, 1.55, 1.35],
        font_size=11,
    )
    for row_idx in range(1, 6):
        for col_idx in (2, 5):
            for paragraph in table.cell(row_idx, col_idx).text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = RED
                    run.font.bold = True
    add_rect(slide, 9.22, 1.68, 3.4, 4.3, LIGHT, line=ORANGE)
    add_text(slide, "当前结论", 9.52, 1.98, 2.8, 0.4, size=20, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
    add_bullets(slide, [
        "S1-S4 实际 TPS 基本固定为 40",
        "S1-S4 的最低 p95 都出现在 512MB",
        "S4 的 combined 推荐 512MB 与延迟最优一致",
        "S1-S3 推荐点延迟只高约 2–4ms",
        "S5 图仅覆盖原低压 TP 流，不能代表 12-terminal 高压流的完整延迟",
    ], 9.48, 2.65, 2.9, 2.55, size=12, spacing=7)
    add_text(slide, "原配置下不能用“TPS 最高点”判断 S1-S4，因为负载发生器已经把 TPS 固定。", 9.48, 5.28, 2.9, 0.52, size=11, color=RED, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, 12)

    # 13. Recommendation gap and reasons
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, "命中率推荐配置：原验证结论", "S1-S4 沿用原完整验证；S5 已替换为 active/inactive/refault 完整历史 replay。", "05 / Summary")
    table_data = [["阶段", "模型推荐", "实际最优", "配置倍数", "实际损失"]]
    for row in summary:
        pred = int(row["pred_best_sb_mb"])
        actual = int(row["actual_best_sb_mb"])
        ratio = pred / actual
        regret = (float(row["actual_best_combined"]) - float(row["actual_combined_at_pred_best"])) * 100
        table_data.append([
            row["stage"].split("_")[0].replace("stage", "S"),
            f"{pred} MB",
            f"{actual} MB",
            f"{ratio:.1f}×",
            f"{regret:.2f} pp",
        ])
    table = add_table(slide, 6, 5, table_data, 0.72, 1.68, 7.25, 3.45, widths=[0.85, 1.55, 1.55, 1.35, 1.95], font_size=12)
    for r in (1, 4, 5):
        for c in (3, 4):
            for p in table.cell(r, c).text_frame.paragraphs:
                for run in p.runs:
                    run.font.color.rgb = RED if r in (4, 5) else ORANGE
                    run.font.bold = True
    slide.shapes.add_picture(str(regret_chart), Inches(0.72), Inches(5.28), width=Inches(7.25), height=Inches(1.32))
    add_rect(slide, 8.28, 1.68, 4.35, 4.92, LIGHT, line=MID)
    add_text(slide, "结论", 8.58, 1.98, 3.75, 0.4, size=21, color=INK, bold=True, align=PP_ALIGN.CENTER)
    add_bullets(slide, [
        "S2、S3 推荐准确",
        "S1 配置相差 4×，但命中率只损失 0.83pp：曲线较平，实际影响不大",
        "S4 相差 4×，损失 3.42pp：并发与 backpressure 路径未建模",
        "S5 推荐由 1504MB 移到 128MB；实际最优 256MB，配置相差 2×但只损失 0.32pp",
    ], 8.58, 2.65, 3.55, 2.75, size=13, spacing=8)
    exact_count = sum(row["matched"] == "yes" for row in summary)
    average_regret = sum(
        float(row["actual_best_combined"]) - float(row["actual_combined_at_pred_best"])
        for row in summary
    ) / len(summary) * 100
    add_kpi(slide, f"{exact_count} / 5", "推荐完全命中", 8.82, 5.55, 1.55, color=RED)
    add_kpi(slide, f"{average_regret:.2f} pp", "平均实际损失", 10.45, 5.55, 1.7, color=ORANGE)
    add_text(slide, "S5 配置没有精确命中，但已经落在实际最优附近的平台区；当前主要剩余问题转为 S4。", 8.58, 5.12, 3.55, 0.5, size=10, color=TEAL, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, 13)

    # 14. TP-only performance signal
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(
        slide,
        "新增：混合回放，但只统计 TP 命中",
        "AP 访问仍参与缓存竞争和淘汰；只有 TP 页面访问进入性能评分。",
        "06 / TP-only Validation",
    )
    add_rect(slide, 0.72, 1.72, 3.45, 4.78, LIGHT, line=BLUE)
    add_text(slide, "完整输入轨迹", 0.98, 2.02, 2.9, 0.38, size=21, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
    add_bullets(slide, [
        "TPCC 与 TPCH 按原时间顺序共同 replay",
        "保留 AP 对 SB/OS 缓存状态的污染",
        "用 relation 标识区分 TP 与 AP 页面",
        "系统目录访问仍参与状态，但不计入 TP 分数",
    ], 1.02, 2.68, 2.82, 2.48, size=13, spacing=8)
    add_arrow(slide, 4.28, 4.05, 5.18, 4.05, TEAL)
    add_rect(slide, 5.18, 1.72, 3.0, 4.78, INK)
    add_text(slide, "每个候选 SB", 5.45, 2.02, 2.48, 0.38, size=21, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_kpi(slide, "Replay all", "TP + AP 共同改变缓存", 5.52, 2.72, 2.32, color=BLUE)
    add_kpi(slide, "Score TP", "只累计 TP hit/miss", 5.52, 4.12, 2.32, color=TEAL)
    add_arrow(slide, 8.3, 4.05, 9.18, 4.05, ORANGE)
    add_rect(slide, 9.18, 1.72, 3.45, 4.78, LIGHT, line=ORANGE)
    add_text(slide, "性能推荐信号", 9.45, 2.02, 2.9, 0.38, size=21, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
    add_bullets(slide, [
        "不再使用 TP combined 最大值",
        "使用 TP SB hit 的 99% 饱和拐点",
        "原因：OS 会补偿 SB miss，使 combined 近似常数",
        "该信号描述 TP 工作集进入数据库缓存的程度",
    ], 9.48, 2.68, 2.82, 2.48, size=13, spacing=8)
    add_text(slide, "关键约束：这是性能相关排序信号，不是由实测 TPS 回归得到的校准曲线。", 9.5, 5.58, 2.78, 0.55, size=11, color=RED, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, 14)

    # 15. TP-only curves versus performance
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(
        slide,
        "TP-only replay 与真实 TP 性能逐点对齐",
        "S5 使用 12-terminal 高压流 + 2-terminal 低压流的总 TPS；S1-S4 因限速使用 p95。",
        "06 / TP-only Validation",
    )
    slide.shapes.add_picture(
        str(TP_ONLY_CHART),
        Inches(0.55),
        Inches(1.43),
        width=Inches(12.25),
        height=Inches(5.78),
    )
    add_footer(slide, 15)

    # 16. TP-only result summary
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(
        slide,
        "TP-only 统计能否预测 TPS 最优点？",
        "答案：S5 吞吐平台精确命中；其余阶段当前实验不能形成可识别的 TPS 最优点。",
        "06 / TP-only Validation",
    )
    table_data = [["阶段", "TP-SB 拐点", "实测目标点", "验证口径", "预测点代价"]]
    for row in tp_only_summary:
        if row["stage"] == "stage5_tp_surge":
            basis = "总 TP TPS 99% 平台"
            cost = f"{float(row['tps_at_predicted_knee']):.1f} TPS"
        else:
            basis = "最低 TP p95"
            cost = f"p95 +{float(row['p95_regret_pct']):.1f}%"
        table_data.append([
            row["stage_short"],
            f"{row['predicted_tp_sb_knee_mb']} MB",
            f"{row['actual_target_sb_mb']} MB",
            basis,
            cost,
        ])
    table = add_table(
        slide,
        6,
        5,
        table_data,
        0.72,
        1.68,
        7.55,
        3.72,
        widths=[0.85, 1.55, 1.55, 2.05, 1.55],
        font_size=11,
    )
    for row_idx in (3, 5):
        for col_idx in (1, 2, 4):
            for paragraph in table.cell(row_idx, col_idx).text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = GREEN
                    run.font.bold = True
    add_rect(slide, 8.58, 1.68, 4.05, 4.92, LIGHT, line=TEAL)
    add_text(slide, "结论", 8.88, 1.98, 3.45, 0.38, size=21, color=TEAL, bold=True, align=PP_ALIGN.CENTER)
    add_kpi(slide, "1024 MB", "S5 预测 / 实测 TPS 平台", 9.05, 2.62, 3.1, color=GREEN)
    add_kpi(slide, "0.990", "S5 TP-SB 与总 TPS Pearson", 9.05, 3.92, 3.1, color=TEAL)
    add_bullets(slide, [
        "S5 总 TPS：131 → 170 → 216 → 220",
        "1024MB 后继续增大 SB，吞吐几乎不再增加",
        "五阶段 TP combined 全部为平线，不能用于 TPS 排序",
        "S1-S4 需解除 40 TPS 限速后，才能严格验证 TPS 最优点",
    ], 8.96, 5.02, 3.2, 1.25, size=11, spacing=5)
    add_footer(slide, 16)

    prs.save(OUT)
    OUTLINE.write_text(
        "\n".join(
            [
                "# Huawei5 current model - concise outline",
                "",
                "1. Model input/output: one 1504MB trace replays many candidate SBs.",
                "2. Five stages are replayed continuously; cache state is retained.",
                "3. SB and OS hit are simulated separately; recommendation maximizes combined hit.",
                "4. After the S5 upgrade, overall OS MAE is 4.87pp and Stage5 OS MAE is 12.67pp.",
                "5-9. One slide per stage compares predicted and actual SB hit, OS hit, and combined hit across 8 SB sizes.",
                "10. Original 2/12-terminal rate-limited runs are used for performance validation; concurrency is not changed.",
                "11. S1-S4 stay near 40 TPS, so p95 latency rather than TPS identifies the better configuration.",
                "12. S1-S4 lowest p95 occurs at 512MB; S4 recommendation matches, while S1-S3 add only about 2-4ms.",
                "13. TP-only replay keeps TP and AP in the same cache state but scores only TP relation accesses.",
                "14. TP combined is flat; TP SB hit is used as the performance-oriented signal.",
                "15. Stage5 TP-SB knee and actual total-TP TPS plateau both occur at 1024MB; Pearson correlation is 0.990.",
            ]
        )
        + "\n",
        encoding="utf-8",
    )
    print(OUT)
    print(OUTLINE)


if __name__ == "__main__":
    build()
