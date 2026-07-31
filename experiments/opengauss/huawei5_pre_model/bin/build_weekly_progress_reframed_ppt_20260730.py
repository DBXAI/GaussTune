#!/usr/bin/env python3
"""Build the problem-first Huawei5 weekly progress presentation."""

from __future__ import annotations

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches

from build_memory_autonomy_paper_ppt import (
    BLUE,
    GRAY,
    GREEN,
    INK,
    LIGHT,
    MID,
    ORANGE,
    PALE_GREEN,
    PALE_ORANGE,
    PALE_RED,
    PALE_TEAL,
    PURPLE,
    RED,
    TEAL,
    WHITE,
    add_arrow,
    add_bullets,
    add_callout,
    add_flow_box,
    add_kpi,
    add_notes,
    add_rect,
    add_table,
    add_text,
    add_title,
    set_bg,
)


ROOT = Path(__file__).resolve().parents[1]
LATEST = ROOT / "artifacts" / "00_latest"
OUT = LATEST / "Huawei5_weekly_progress_reframed_20260730.pptx"
OUTLINE = LATEST / "Huawei5_weekly_progress_reframed_20260730_outline.md"

MODEL_SUMMARY = LATEST / "paper_model_validation_summary_20260727.png"
STAGE_VALIDATION = LATEST / "five_stage_saturated_tps_validation_20260726.png"
KERNEL_ACCEPTANCE = LATEST / "kernel_online_resize_read_write_acceptance_20260727.png"
WEEKLY_VALIDATION = LATEST / "weekly_tp_validation_20260729.png"
V8_TIMELINE = LATEST / "weekly_v8_controller_timeline_20260729.png"
WORKMEM_VALIDATION = ROOT / "artifacts/01_current_joint_model/figures/all_query_workmem_prediction_vs_actual.png"
S5_JOINT = ROOT / "artifacts/01_current_joint_model/figures/joint_bidirectional_effect_s5_20260722.png"
S5_SB_TPS = ROOT / "artifacts/02_validation_figures/s5_20260716_17/s5_tp_sb_hit_vs_total_tps_20260716.png"


def footer(slide, page: int, source: str = "") -> None:
    add_rect(slide, 0.62, 7.17, 12.08, 0.012, MID)
    add_text(
        slide,
        "Huawei5 · TP 优先的 Trace Replay 与在线资源闭环 · 2026-07-30",
        0.62,
        7.22,
        7.2,
        0.18,
        size=7.5,
        color=GRAY,
    )
    if source:
        add_text(slide, f"来源：{source}", 7.18, 7.22, 4.95, 0.18, size=7, color=GRAY, align=PP_ALIGN.RIGHT)
    add_text(slide, str(page), 12.28, 7.22, 0.42, 0.18, size=8, color=GRAY, align=PP_ALIGN.RIGHT)


def add_image_contain(slide, path: Path, x: float, y: float, w: float, h: float):
    with Image.open(path) as image:
        iw, ih = image.size
    scale = min(w / iw, h / ih)
    draw_w, draw_h = iw * scale, ih * scale
    draw_x = x + (w - draw_w) / 2
    draw_y = y + (h - draw_h) / 2
    return slide.shapes.add_picture(str(path), Inches(draw_x), Inches(draw_y), width=Inches(draw_w), height=Inches(draw_h))


def add_section_label(slide, text: str, x: float, y: float, w: float, color=TEAL) -> None:
    add_rect(slide, x, y, 0.08, 0.38, color)
    add_text(slide, text, x + 0.17, y - 0.01, w - 0.17, 0.40, size=15, color=color, bold=True)


def build() -> None:
    for path in [MODEL_SUMMARY, STAGE_VALIDATION, KERNEL_ACCEPTANCE, WEEKLY_VALIDATION, V8_TIMELINE, WORKMEM_VALIDATION, S5_JOINT, S5_SB_TPS]:
        if not path.exists():
            raise FileNotFoundError(path)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    outline: list[str] = []

    # 1. Cover
    slide = prs.slides.add_slide(blank)
    set_bg(slide, INK)
    add_rect(slide, 0, 0, 0.22, 7.5, TEAL)
    add_text(slide, "问题、模型、实验效果与不足", 0.94, 1.04, 10.8, 0.48, size=23, color=PALE_TEAL, bold=True)
    add_text(slide, "让 AP/TP 混合负载中的 TP TPS 保持稳定", 0.92, 1.70, 11.65, 1.12, size=36, color=WHITE, bold=True)
    add_text(slide, "基于 Trace Replay、源码级 Plan/Spill 推演与在线资源闭环", 0.94, 2.98, 11.3, 0.48, size=18, color=MID)
    add_rect(slide, 0.94, 3.78, 5.92, 0.56, TEAL, radius=True)
    add_text(slide, "Huawei5 五阶段 TP/AP 混合负载", 1.08, 3.94, 5.64, 0.25, size=12.5, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "目标：TP 相对基线保持率 ≥95%，AP 不取消并最终自然完成", 0.94, 5.30, 11.1, 0.42, size=16, color=MID)
    add_text(slide, "一周研发进展重构版 · 2026-07-30", 0.94, 6.46, 5.2, 0.24, size=10, color=MID)
    add_notes(slide, "开场先讲项目要解决的运行时问题，而不是先报开发事项。")
    outline.append("1. 封面：让混合负载中的 TP TPS 保持稳定。")

    # 2. Problem
    slide = prs.slides.add_slide(blank); set_bg(slide)
    add_title(slide, "当前问题：AP 执行会通过四条路径压低 TP TPS", "问题不是找一个静态最大值，而是在 AP 不断变化时持续保护 TP。", "01 · PROBLEM")
    add_flow_box(slide, "TP 事务负载", "32 terminals\n固定 offered 800 TPS\n对延迟和吞吐敏感", 0.62, 2.06, 2.38, 1.62, TEAL, fill=PALE_TEAL)
    add_flow_box(slide, "AP 复杂查询", "Join / Sort / Aggregate\n执行数分钟到数小时\n必须自然完成", 10.32, 2.06, 2.38, 1.62, ORANGE, fill=PALE_ORANGE)
    center = [
        ("CPU", "算子计算与并发 worker", BLUE),
        ("I/O", "数据扫描、临时 spill", RED),
        ("Cache", "扫描污染 page cache", PURPLE),
        ("Memory", "work_mem 与 SB 争用", GREEN),
    ]
    for i, (title, body, color) in enumerate(center):
        x = 3.46 + (i % 2) * 3.34
        y = 1.70 + (i // 2) * 1.66
        add_callout(slide, title, body, x, y, 2.86, 1.22, color, LIGHT)
    add_arrow(slide, 3.13, 2.87, 3.42, 2.87, color=GRAY)
    add_arrow(slide, 10.20, 2.87, 9.91, 2.87, color=GRAY)
    add_rect(slide, 0.82, 4.98, 11.70, 1.04, INK, radius=True)
    add_text(slide, "直接后果", 1.06, 5.25, 1.32, 0.30, size=16, color=WHITE, bold=True)
    add_text(slide, "AP 的 Plan、spill 和执行阶段变化，使同一 SB/work_mem 或固定 CPU/I/O 配额在前期安全、后期却可能让 TP 跌破 95%。", 2.35, 5.17, 9.82, 0.52, size=15, color=WHITE, bold=True)
    add_text(slide, "因此，项目必须同时回答：该怎么配置、配置何时切换、切换后 TP 是否真的恢复。", 0.92, 6.28, 11.50, 0.40, size=17, color=RED, bold=True, align=PP_ALIGN.CENTER)
    footer(slide, 2, "项目验收目标；V8_EXPLORATION.md")
    add_notes(slide, "把问题分成 CPU、I/O、cache 和动态内存四条竞争路径。")
    outline.append("2. 当前问题：AP 通过 CPU、I/O、Cache 和 Memory 四条路径压低 TP。")

    # 3. Goal and acceptance
    slide = prs.slides.add_slide(blank); set_bg(slide)
    add_title(slide, "项目目标：TP 稳定优先，但不能靠杀掉或永久饿死 AP", "静态推荐是起点；最终验收是完整自然生命周期中的闭环稳定性。", "01 · TARGET")
    add_kpi(slide, "≥95%", "TP 相对基线保持率", 0.74, 1.72, 2.62, color=TEAL, note="固定 800 TPS 参考；15 秒窗口")
    add_kpi(slide, "5/5", "五阶段都要满足", 3.36, 1.72, 2.62, color=BLUE, note="同一并发与同一口径")
    add_kpi(slide, "0", "AP Query 被取消", 5.98, 1.72, 2.62, color=GREEN, note="保留 session、Plan 与状态")
    add_kpi(slide, "自然完成", "AP 最终进展约束", 8.60, 1.72, 3.34, color=ORANGE, note="阶段计时结束不等于 cancel")
    data = [
        ["验收层次", "回答的问题", "当前结论"],
        ["预测准确性", "命中率、Plan、spill 边界能否外推？", "多数留出点有效，扩大未见 Query 仍有边界"],
        ["配置推荐", "推荐 SB × work_mem 是否接近已测最高 TPS？", "已测挑战网格 5/5 regret <5%"],
        ["在线执行", "SB 能否无重启切换且不伤 TP？", "128→64MB 三轮最差单秒下降 2.65%"],
        ["闭环稳定", "AP 自然结束全过程 TP 是否始终 ≥95%？", "S2 v8 未通过；10/222 个窗口越界"],
    ]
    add_table(slide, data, 0.72, 3.35, 11.92, 2.68, widths=[2.05, 5.02, 4.85], font_size=10.8)
    add_text(slide, "本汇报严格区分“推荐 regret <5%”与“运行期 TPS 波动 <5%”，两者不是同一个验收指标。", 0.82, 6.36, 11.72, 0.38, size=14.5, color=RED, bold=True, align=PP_ALIGN.CENTER)
    footer(slide, 3, "SATURATED_FIVE_STAGE_TPS_VALIDATION.md；V8_RESULT.md")
    outline.append("3. 目标与验收：TP ≥95%、五阶段、0 cancel、AP 自然完成。")

    # 4. Model I/O
    slide = prs.slides.add_slide(blank); set_bg(slide)
    add_title(slide, "模型总览：输入执行证据，输出候选配置与安全策略", "它不是用历史最优 TPS 标签训练的黑盒回归，而是执行行为 Replay + 在线反馈。", "02 · MODEL I/O")
    add_section_label(slide, "输入", 0.66, 1.62, 2.45, BLUE)
    add_flow_box(slide, "一次锚点负载", "TP 页访问序列\nAP 算子 trace\n实际 rows / width / lifecycle", 0.66, 2.04, 2.48, 1.52, BLUE, fill=LIGHT)
    add_flow_box(slide, "Plan 与源码", "不同 work_mem 的 EXPLAIN\nPlan family / operatorMemKB\nHash/Sort/Agg 容量规则", 0.66, 3.78, 2.48, 1.52, PURPLE, fill=LIGHT)
    add_flow_box(slide, "系统约束与反馈", "内存上限 / 并发 / DOP\n真实 TP TPS / AP 进展\nCPU / I/O / freeze", 0.66, 5.52, 2.48, 1.14, GREEN, fill=LIGHT)
    add_section_label(slide, "内部建模", 3.55, 1.62, 5.84, TEAL)
    add_flow_box(slide, "缓存 Replay", "SB 淘汰 + Linux active/inactive/refault\n单独统计 TP-SB、TP-OS 与物理 miss", 3.55, 2.04, 2.60, 1.52, TEAL, fill=PALE_TEAL)
    add_flow_box(slide, "Plan/Spill Replay", "work_mem 改变 Plan 与算子容量\n按生命周期计算 spill 和动态峰值", 6.38, 2.04, 2.60, 1.52, ORANGE, fill=PALE_ORANGE)
    add_flow_box(slide, "SB × work_mem 联合搜索", "spill I/O 反向污染 OS cache\nSB 又压缩动态池；每个组合重新计算", 3.55, 4.02, 5.43, 1.42, PURPLE, fill=LIGHT)
    add_flow_box(slide, "在线 TP SLO 闭环", "预测候选 → 试用 → 观测 TP/AP\n升档、降档、冻结归因与安全回退", 3.55, 5.70, 5.43, 0.96, GREEN, fill=PALE_GREEN)
    add_section_label(slide, "输出", 9.40, 1.62, 3.24, ORANGE)
    add_flow_box(slide, "阶段配置", "shared_buffers\nglobal / per-query work_mem\n预测命中、spill 与峰值", 9.40, 2.04, 3.24, 1.52, ORANGE, fill=PALE_ORANGE)
    add_flow_box(slide, "运行策略", "AP CPU / I/O 安全上界\nQuery 边界 grant\nfreeze / rollback 条件", 9.40, 3.82, 3.24, 1.52, GREEN, fill=PALE_GREEN)
    add_flow_box(slide, "可信边界", "Plan/spill 置信度\n不可部署点与未知区域", 9.40, 5.60, 3.24, 1.06, RED, fill=PALE_RED)
    footer(slide, 4, "ONE_SHOT_SOURCE_PLAN_REPLAY.md；JOINT_BIDIRECTIONAL_REPLAY.md")
    add_notes(slide, "这一页回答模型输入、内部状态和输出，后面三页分别拆开。")
    outline.append("4. 模型总览：输入、内部建模、输出和在线反馈。")

    # 5. Data protocol
    slide = prs.slides.add_slide(blank); set_bg(slide)
    add_title(slide, "数据如何进入模型：一次真实执行 + 未执行配置推演", "目标是减少反复跑完整负载，同时保证验证集不参与候选生成。", "02 · DATA")
    steps = [
        ("1", "真实执行一次", "保留 TP 页访问、AP Plan、算子 rows/width、生命周期和 spill。", BLUE),
        ("2", "只扫 EXPLAIN", "对候选 work_mem 生成 Plan family 区间，不执行这些候选 SQL。", PURPLE),
        ("3", "源码级合成", "未执行 Plan 使用自身节点结构与 Hash/Sort/Agg 源码容量规则。", TEAL),
        ("4", "联合回放", "在所有 SB × work_mem 点重算缓存、spill、动态峰值和物理 I/O。", ORANGE),
        ("5", "冻结预测", "保存预测文件与哈希，再独立重启、清 cache、跑 TPS/spill 验证。", GREEN),
    ]
    y = 1.65
    for code, title, body, color in steps:
        add_rect(slide, 0.76, y, 0.64, 0.55, color, radius=True)
        add_text(slide, code, 0.80, y + 0.15, 0.56, 0.24, size=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, title, 1.66, y + 0.02, 2.08, 0.34, size=16, color=color, bold=True)
        add_text(slide, body, 3.70, y, 8.72, 0.58, size=13.5)
        y += 0.96
    add_callout(slide, "模型会用的数据", "trace、Plan、源码规则、容量/并发约束，以及锚点中的行数误差和 allocator 开销。", 0.78, 6.10, 5.70, 0.78, TEAL, PALE_TEAL)
    add_callout(slide, "模型不会用的数据", "目标配置的实测最优 TPS、验证集 spill 标签，以及验证后人为校准的答案。", 6.76, 6.10, 5.76, 0.78, RED, PALE_RED)
    footer(slide, 5, "严格留出与预注册实验协议")
    outline.append("5. 数据协议：一次执行、EXPLAIN 扫描、源码合成、冻结后独立验证。")

    # 6. Cache replay
    slide = prs.slides.add_slide(blank); set_bg(slide)
    add_title(slide, "子模型一：用页访问轨迹推演不同 SB 的 TP 缓存路径", "AP 与 TP 都参与回放，但评价时单独统计 TP 页，避免 AP 扫描命中掩盖 TPS 拐点。", "03 · CACHE REPLAY")
    add_flow_box(slide, "页访问序列", "page id / read-write\n时间与阶段\nTP/AP 标签", 0.72, 1.78, 2.20, 1.46, BLUE, fill=LIGHT)
    add_arrow(slide, 2.98, 2.51, 3.47, 2.51, color=GRAY)
    add_flow_box(slide, "SB 状态机", "容量、冷热、淘汰\ndirty / pin 约束\n得到 TP-SB hit", 3.54, 1.78, 2.38, 1.46, TEAL, fill=PALE_TEAL)
    add_arrow(slide, 5.98, 2.51, 6.47, 2.51, color=GRAY)
    add_flow_box(slide, "Linux page cache", "active / inactive\nrefault 距离\n保护频繁访问 TP 页", 6.54, 1.78, 2.38, 1.46, ORANGE, fill=PALE_ORANGE)
    add_arrow(slide, 8.98, 2.51, 9.47, 2.51, color=GRAY)
    add_flow_box(slide, "输出", "TP-SB hit\nTP-OS conditional hit\nTP 物理磁盘 miss", 9.54, 1.78, 2.98, 1.46, GREEN, fill=PALE_GREEN)
    add_callout(slide, "修复过的关键问题", "原模型严重低估小 SB 下 Linux page cache 的补偿。升级 active/inactive 与 refault 后，频繁访问 TP 页能够被保护。", 0.82, 3.84, 5.68, 1.40, ORANGE, PALE_ORANGE)
    add_callout(slide, "为什么和 TPS 有关系", "对饱和 TP，TP 物理 miss 降低通常对应 TPS 上升并形成平台；模型用该机制筛选 SB 候选，而不是直接拟合 TPS 数值。", 6.78, 3.84, 5.68, 1.40, TEAL, PALE_TEAL)
    add_kpi(slide, "0.61 pp", "SB hit held-out MAE", 1.20, 5.55, 3.10, color=BLUE)
    add_kpi(slide, "1.82 pp", "OS hit held-out MAE", 5.12, 5.55, 3.10, color=ORANGE)
    add_kpi(slide, "0.50 pp", "Combined held-out MAE", 9.02, 5.55, 3.10, color=GREEN)
    footer(slide, 6, "缓存 replay held-out 验证")
    outline.append("6. 缓存 Replay：SB 淘汰、Linux refault 保护和 TP 独立命中统计。")

    # 7. Operator replay
    slide = prs.slides.add_slide(blank); set_bg(slide)
    add_title(slide, "子模型二：不同 work_mem 下同时推演 Plan 与 spill", "work_mem 不是阶段总预算；每个同时存活的内存算子都可能获得自己的额度。", "03 · PLAN / SPILL")
    add_flow_box(slide, "锚点 Plan trace", "实际 rows / width\n算子生命周期\n一次真实 spill", 0.64, 1.75, 2.10, 1.42, BLUE, fill=LIGHT)
    add_flow_box(slide, "候选 Plan family", "EXPLAIN 扫描\n去成本 SHA\noperatorMemKB", 3.00, 1.75, 2.10, 1.42, PURPLE, fill=LIGHT)
    add_flow_box(slide, "源码容量公式", "Hash bucket/batch\nSortTuple/merge\nHashAgg groups", 5.36, 1.75, 2.10, 1.42, TEAL, fill=PALE_TEAL)
    add_flow_box(slide, "Query 生命周期", "同一 Query 并发算子\n峰值与 spill 时刻\nPlan 切换区间", 7.72, 1.75, 2.10, 1.42, ORANGE, fill=PALE_ORANGE)
    add_flow_box(slide, "阶段聚合", "并发 Query 峰值求和\nmax_dynamic_memory\nglobal/per-query grant", 10.08, 1.75, 2.58, 1.42, GREEN, fill=PALE_GREEN)
    for x in [2.77, 5.13, 7.49, 9.85]:
        add_arrow(slide, x, 2.46, x + 0.19, 2.46, color=GRAY, width=1.6)
    add_callout(slide, "可预测的内容", "每个 Plan family 的 work_mem 区间、spill/no-spill 边界、临时 I/O、单 Query 与阶段动态峰值。", 0.78, 3.70, 5.74, 1.32, TEAL, PALE_TEAL)
    add_callout(slide, "尚不精确的内容", "新 Plan 的 spill I/O 数量受全局运行时内存状态影响；Materialize、WindowAgg、SetOp 尚未完整源码化。", 6.80, 3.70, 5.74, 1.32, RED, PALE_RED)
    add_kpi(slide, "7/7", "严格留出 Plan family", 1.00, 5.44, 2.82, color=GREEN)
    add_kpi(slide, "6/7", "严格留出 spill 分类", 4.05, 5.44, 2.82, color=TEAL)
    add_kpi(slide, "15/21", "扩大未见 Query spill", 7.10, 5.44, 2.82, color=ORANGE)
    add_kpi(slide, "48.8%", "未见 Query spill I/O MAPE", 10.15, 5.44, 2.82, color=RED)
    footer(slide, 7, "ONE_SHOT_SOURCE_PLAN_REPLAY.md；STRICT_UNSEEN_QUERY_VALIDATION.md")
    outline.append("7. Plan/Spill Replay：Plan family、源码容量、生命周期和多 Query 动态池。")

    # 8. Joint interaction
    slide = prs.slides.add_slide(blank); set_bg(slide)
    add_title(slide, "联合模型：SB 与 work_mem 必须在每个候选点双向重算", "两者共同消耗内存，并通过 spill 与 page cache 相互影响。", "03 · JOINT MODEL")
    add_flow_box(slide, "增大 SB", "TP-SB hit ↑\n但 OS cache 和动态池空间 ↓", 0.82, 1.82, 2.58, 1.42, TEAL, fill=PALE_TEAL)
    add_arrow(slide, 3.48, 2.53, 5.12, 2.53, color=GRAY)
    add_flow_box(slide, "动态内存 grant", "可部署 work_mem 上界改变\n并发 Query 峰值受限", 5.22, 1.82, 2.78, 1.42, ORANGE, fill=PALE_ORANGE)
    add_arrow(slide, 8.10, 2.53, 9.72, 2.53, color=GRAY)
    add_flow_box(slide, "spill 与物理 I/O", "AP spill 可能 ↓ 或 ↑\nOS cache 污染也随之改变", 9.82, 1.82, 2.70, 1.42, PURPLE, fill=LIGHT)
    add_arrow(slide, 10.20, 3.52, 3.00, 3.52, color=RED, width=2.0)
    add_text(slide, "反向作用：spill I/O 会挤占 TP I/O 并污染 OS cache", 3.72, 3.57, 5.84, 0.30, size=13, color=RED, bold=True, align=PP_ALIGN.CENTER)
    data = [
        ["候选评价顺序", "意义"],
        ["1. 进入 TP-SB 命中平台", "优先保证 TP 热集不再因 SB 不足持续下降"],
        ["2. 满足总内存与动态池边界", "排除 work_mem 峰值不可部署、SB 过大的组合"],
        ["3. 最小化 TP 数据 I/O + AP spill I/O", "在平台内选真正降低共享 I/O 压力的组合"],
        ["4. 同等性能选更小内存", "保留 OS cache 和运行时调整空间"],
    ]
    add_table(slide, data, 0.82, 4.12, 7.26, 2.26, widths=[3.12, 4.14], font_size=10.5)
    add_callout(slide, "模型输出的含义", "推荐值是满足约束后的性能平台起点，不是“把 SB 和 work_mem 都设到最大”，也不是两个独立拐点的简单交集。", 8.42, 4.12, 4.10, 2.26, GREEN, PALE_GREEN)
    footer(slide, 8, "JOINT_BIDIRECTIONAL_REPLAY.md")
    outline.append("8. 联合模型：SB、动态池、spill 和 OS cache 双向反馈。")

    # 9. Online loop
    slide = prs.slides.add_slide(blank); set_bg(slide)
    add_title(slide, "从离线推荐到在线闭环：预测缩小空间，反馈保证 TP", "控制器不凭一次预测永久固定资源；它观测真实 TP/AP 后执行升档、降档和回退。", "04 · ONLINE CONTROL")
    add_flow_box(slide, "Replay 初始候选", "阶段 SB / work_mem\nAP CPU / I/O 安全起点", 0.62, 1.82, 2.18, 1.46, PURPLE, fill=LIGHT)
    add_flow_box(slide, "稳定门", "TP 连续窗口 ≥98%\n才允许 AP 提交或升档", 3.00, 1.82, 2.18, 1.46, TEAL, fill=PALE_TEAL)
    add_flow_box(slide, "90 秒试用", "AP 进展必须改善\n期间 TP 不低于 95%", 5.38, 1.82, 2.18, 1.46, BLUE, fill=LIGHT)
    add_flow_box(slide, "因果确认", "串行改变 SB 与 I/O\nfreezer 暂停同一 SQL", 7.76, 1.82, 2.18, 1.46, ORANGE, fill=PALE_ORANGE)
    add_flow_box(slide, "安全回退", "80→40→20→10→5\n保留已学习路径上界", 10.14, 1.82, 2.52, 1.46, GREEN, fill=PALE_GREEN)
    for x in [2.82, 5.20, 7.58, 9.96]:
        add_arrow(slide, x, 2.55, x + 0.14, 2.55, color=GRAY, width=1.5)
    add_callout(slide, "可以做", "在 AP 严重影响 TP 时降低 AP CPU/I/O、缩小后续 Query 的 work_mem、必要时有界 freeze；动态调整 shared_buffers target。", 0.80, 4.06, 5.78, 1.50, TEAL, PALE_TEAL)
    add_callout(slide, "不能做", "180 秒结束直接 cancel SQL；假装运行中的 work_mem 已缩小；为得到好看曲线而降低 TPS 基线或忽略自然收尾期。", 6.78, 4.06, 5.78, 1.50, RED, PALE_RED)
    add_text(slide, "最终控制目标：TP 保持率达标，同时 AP 在资源受控条件下持续取得进展并自然完成。", 0.90, 6.03, 11.56, 0.48, size=18, color=INK, bold=True, align=PP_ALIGN.CENTER)
    footer(slide, 9, "TP SLO controller；动态 SB 内核；AP cgroup/freezer")
    outline.append("9. 在线闭环：稳定门、试用、因果确认、回退和自然完成。")

    # 10. Evidence map
    slide = prs.slides.add_slide(blank); set_bg(slide)
    add_title(slide, "效果如何证明：五类实验分别验证模型链路", "每张图只回答一个问题，避免用单一准确率概括整个系统。", "05 · EXPERIMENT MAP")
    data = [
        ["验证对象", "实验协议", "关键结果", "能证明什么"],
        ["缓存 Replay", "训练锚点与 held-out SB 点隔离", "SB 0.61pp；OS 1.82pp；Combined 0.50pp", "缓存路径可外推"],
        ["Plan / Spill", "Plan/Query 严格留出；预测冻结后执行", "Plan 7/7；spill 6/7；扩大样本 15/21", "边界较准，I/O 量仍有限"],
        ["联合推荐", "32 terminals；独立重启、清 cache、预热", "五阶段已测网格最大 regret 4.55%", "推荐点接近已测 TPS 最高点"],
        ["在线 SB 内核", "128→64MB；三轮读写 TP；单秒检查", "最差单秒下降 2.65%；0 错误", "配置可无重启执行"],
        ["在线 TP SLO", "固定 800 TPS；AP 自然结束全生命周期", "TP-only 0 越界；v8 为 10/222 越界", "闭环可诊断，但尚未验收"],
    ]
    add_table(slide, data, 0.56, 1.58, 12.22, 4.72, widths=[2.05, 3.28, 3.40, 3.49], font_size=9.5)
    add_text(slide, "下面 6 页依次展示这些实验的原始效果图与结论边界。", 0.84, 6.51, 11.64, 0.35, size=16, color=TEAL, bold=True, align=PP_ALIGN.CENTER)
    footer(slide, 10, "本周及历史已冻结实验结果")
    outline.append("10. 实验证据地图：缓存、Plan/Spill、联合推荐、内核执行、在线闭环。")

    # 11. Model accuracy evidence
    slide = prs.slides.add_slide(blank); set_bg(slide)
    add_title(slide, "效果一：缓存与 Plan/Spill 的留出验证", "缓存误差较小；Plan family 稳定，但扩大未见 Query 后 spill 泛化仍暴露边界。", "05 · EFFECT 1")
    add_image_contain(slide, MODEL_SUMMARY, 0.62, 1.56, 12.10, 4.36)
    add_callout(slide, "结论", "缓存 replay 可以作为 SB 候选信号；Plan 切换判断较稳；spill/no-spill 比 spill I/O 数量可靠。", 0.80, 6.02, 5.78, 0.78, GREEN, PALE_GREEN)
    add_callout(slide, "边界", "扩大未见 Query 为 15/21，说明当前算子覆盖和运行时全局内存状态仍不足，不能写成全泛化。", 6.78, 6.02, 5.78, 0.78, ORANGE, PALE_ORANGE)
    footer(slide, 11, "paper_model_validation_summary_20260727.png")
    outline.append("11. 效果一：缓存与 Plan/Spill 留出准确性。")

    # 12. Work_mem evidence
    slide = prs.slides.add_slide(blank); set_bg(slide)
    add_title(slide, "效果二：逐 Query 最小 no-spill work_mem 预测", "模型能够识别多数可观测边界，也能标记在当前主机或引擎约束下不可达到的点。", "05 · EFFECT 2")
    add_image_contain(slide, WORKMEM_VALIDATION, 0.62, 1.52, 12.10, 4.80)
    add_callout(slide, "命中", "Q1/Q3/Q7/Q9/Q13 一致；可观测边界 5/6 正确。", 0.80, 6.08, 5.78, 0.78, GREEN, PALE_GREEN)
    add_callout(slide, "错误与不可达", "Q5 在 305MB 切换 Plan；Q18/Q21 无可部署 no-spill 值。", 6.78, 6.08, 5.78, 0.78, ORANGE, PALE_ORANGE)
    footer(slide, 12, "all_query_workmem_prediction_vs_actual.png")
    outline.append("12. 效果二：逐 Query work_mem 边界与不可部署点。")

    # 13. Joint S5 evidence
    slide = prs.slides.add_slide(blank); set_bg(slide)
    add_title(slide, "效果三：S5 中 SB 信号与 TPS 平台一致", "联合模型再用 work_mem 的 spill I/O 和总内存边界筛选同一平台内的配置。", "05 · EFFECT 3")
    add_image_contain(slide, S5_SB_TPS, 0.58, 1.58, 6.20, 4.68)
    add_image_contain(slide, S5_JOINT, 6.88, 1.58, 5.88, 4.68)
    add_text(slide, "SB：预测和实测 TPS 首个平台均从 1024MB 开始", 0.72, 6.30, 5.92, 0.38, size=14, color=TEAL, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "work_mem：1150MB 为首个 zero-spill 边界；联合点必须再检查总内存", 6.94, 6.30, 5.70, 0.38, size=14, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
    footer(slide, 13, "S5 TP hit/TPS 与 joint bidirectional replay")
    outline.append("13. 效果三：S5 的 SB TPS 平台与 work_mem spill 联合约束。")

    # 14. Five-stage recommendation
    slide = prs.slides.add_slide(blank); set_bg(slide)
    add_title(slide, "效果四：五阶段推荐点在已测挑战网格上均接近最高 TPS", "预测在验证 TPS 产生前冻结；每个配置独立重启、清 Linux cache，并使用 32-terminal 饱和 TP。", "05 · EFFECT 4")
    add_image_contain(slide, STAGE_VALIDATION, 0.60, 1.52, 12.12, 4.78)
    add_kpi(slide, "5/5", "阶段 regret <5%", 0.88, 6.16, 2.72, color=GREEN)
    add_kpi(slide, "4.55%", "最大 regret：S3", 3.62, 6.16, 2.72, color=ORANGE)
    add_text(slide, "边界：这是“推荐配置距已测最高点”的 regret，不代表五个阶段之间的绝对 TPS 波动已经小于 5%。", 6.62, 6.25, 5.72, 0.45, size=12.8, color=RED, bold=True, align=PP_ALIGN.CENTER)
    footer(slide, 14, "five_stage_saturated_tps_validation_20260726.png")
    outline.append("14. 效果四：五阶段已测挑战网格推荐 regret。")

    # 15. Runtime feasibility
    slide = prs.slides.add_slide(blank); set_bg(slide)
    add_title(slide, "效果五：配置可以在线执行，800 TPS 基线也可长期维持", "先分别证明执行器不破坏 TP、以及没有 AP 时固定基线本身可持续。", "05 · EFFECT 5")
    add_image_contain(slide, KERNEL_ACCEPTANCE, 0.62, 1.48, 12.10, 4.70)
    add_text(slide, "在线执行：128→64MB 三轮最差单秒下降 2.65%，低于 3% 红线", 0.68, 6.12, 5.84, 0.50, size=13.5, color=GREEN, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "TP-only 1 小时基线：平均 100.34%，最低 97.17%，0 个窗口低于 95%", 6.72, 6.12, 5.94, 0.50, size=13.5, color=TEAL, bold=True, align=PP_ALIGN.CENTER)
    footer(slide, 15, "kernel resize acceptance；TP-only summary.json")
    outline.append("15. 效果五：动态 SB 内核和 TP-only 长稳态基线。")

    # 16. Closed loop evidence
    slide = prs.slides.add_slide(blank); set_bg(slide)
    add_title(slide, "效果六：闭环能发现路径变化并回退，但全生命周期仍未达标", "Q3 不取消并自然完成；控制器识别 80MiB/s 后期不安全并回退到 40/20。", "05 · EFFECT 6")
    add_image_contain(slide, V8_TIMELINE, 0.58, 1.50, 12.18, 4.88)
    add_kpi(slide, "5535.6s", "Q3 自然完成", 0.82, 6.13, 2.52, color=BLUE)
    add_kpi(slide, "0", "前 180 秒越界", 3.34, 6.13, 2.52, color=GREEN)
    add_kpi(slide, "84.50%", "全生命周期最低保持率", 5.86, 6.13, 3.10, color=RED)
    add_kpi(slide, "10/222", "全生命周期越界窗口", 9.04, 6.13, 3.10, color=RED)
    footer(slide, 16, "V8_RESULT.md；weekly_v8_controller_timeline_20260729.png")
    outline.append("16. 效果六：在线回退有效，但 S2 全生命周期验收失败。")

    # 17. Limitations and next steps
    slide = prs.slides.add_slide(blank); set_bg(slide)
    add_title(slide, "当前不足与下一步", "模型链路已经形成，但最终项目目标仍取决于预防式策略和完整五阶段验收。", "06 · LIMITATIONS")
    gaps = [
        ["不足", "为什么会影响最终目标", "下一步处理"],
        ["在线探索会先污染 page cache", "发现 80MiB/s 不安全时，TP 热页损伤已发生", "持久化阶段/Query/路径安全上界，生产从低档启动"],
        ["运行中 work_mem 不能真正热缩", "已启动算子的内存债务和 spill 路径仍存在", "只在 Query 边界发放 grant，运行中记录 graceful debt"],
        ["扩大未见 Query spill 仅 15/21", "未覆盖算子和全局内存保护导致边界错误", "补充 Materialize/WindowAgg/SetOp 源码与新 Query 留出"],
        ["v8 缺少 10/5MiB/s 后备档", "20MiB/s 后期仍可能压低 TP", "完整回退集合 20→10→5，freeze 仅作短期因果确认"],
        ["五阶段全生命周期尚未验收", "当前只有 S2 探索失败证据", "先完成 S2 exploit-only，再按同口径扩展 S1-S5"],
    ]
    add_table(slide, gaps, 0.58, 1.56, 12.18, 4.62, widths=[3.02, 4.30, 4.86], font_size=9.6)
    add_rect(slide, 0.76, 6.34, 11.82, 0.66, PALE_RED, line=RED, radius=True)
    add_text(slide, "最终验收线：所有 AP Query 自然完成、任何阶段不截断前 180 秒、五阶段 TP 归一化保持率均在 ±5% 内。", 0.94, 6.52, 11.46, 0.28, size=13.5, color=RED, bold=True, align=PP_ALIGN.CENTER)
    footer(slide, 17, "当前代码、严格留出结果与下一轮实验计划")
    add_notes(slide, "结尾明确哪些已经完成、哪些还是验收缺口，不把探索运行包装成最终推荐。")
    outline.append("17. 当前不足：探索污染、work_mem 生效边界、泛化、低档回退和五阶段验收。")

    prs.save(OUT)
    OUTLINE.write_text(
        "# Huawei5 一周进展重构版提纲（2026-07-30）\n\n"
        + "\n".join(f"- {item}" for item in outline)
        + "\n\n## 汇报主线\n\n"
        + "1. 先说明 AP/TP 混合负载为何导致 TP TPS 波动，以及最终验收目标。\n"
        + "2. 再说明模型的输入、内部 Replay、输出和在线闭环。\n"
        + "3. 用六组独立实验图验证缓存、Plan/Spill、联合推荐、内核与闭环效果。\n"
        + "4. 最后明确：推荐 regret 已通过已测网格，但五阶段全生命周期 TPS ±5% 尚未验收。\n",
        encoding="utf-8",
    )
    print(OUT)
    print(OUTLINE)


if __name__ == "__main__":
    build()
