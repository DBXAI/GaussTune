#!/usr/bin/env python3
"""
generate_ppt.py — 生成 PostgreSQL 内存性能实验汇报 PPT
背景使用 数据库监控.pptx 第一页模板（去除文字）
输出: 实验汇报.pptx
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

FIGURES  = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'figures')
OUTPUT   = os.path.join(os.path.dirname(os.path.abspath(__file__)), '实验汇报.pptx')

# ── 颜色主题 ─────────────────────────────────────────────────────────────
BLUE_DARK  = RGBColor(0x1B, 0x3A, 0x5C)
BLUE_MID   = RGBColor(0x2C, 0x5F, 0x8A)
BLUE_LIGHT = RGBColor(0x3A, 0x7C, 0xBD)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_DARK  = RGBColor(0x33, 0x33, 0x33)
GRAY_MID   = RGBColor(0x66, 0x66, 0x66)
ACCENT_RED = RGBColor(0xC0, 0x39, 0x2B)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SLD_W, SLD_H = prs.slide_width, prs.slide_height

# ══════════════════════════════════════════════════════════════════════════
#  工具函数
# ══════════════════════════════════════════════════════════════════════════

def inject_bg(slide):
    pass  # no-op, using solid color backgrounds

def add_rect(slide, left, top, width, height, color, alpha_pct=100):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def add_textbox(slide, left, top, width, height, text, size=18,
                color=GRAY_DARK, bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text; p.font.size = Pt(size); p.font.color.rgb = color
    p.font.bold = bold; p.font.name = 'Microsoft YaHei'; p.alignment = align
    return tb

def add_bullets(slide, left, top, width, height, items, size=16,
                color=GRAY_DARK, spacing=1.35):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item; p.font.size = Pt(size); p.font.color.rgb = color
        p.font.name = 'Microsoft YaHei'
        p.space_after = Pt(size * (spacing - 1))
    return tb

def add_img(slide, name, left, top, width=None, height=None):
    path = os.path.join(FIGURES, name)
    if not os.path.exists(path):
        print(f"  [warn] missing: {name}"); return None
    kw = {'left': left, 'top': top}
    if width:  kw['width']  = width
    if height: kw['height'] = height
    return slide.shapes.add_picture(path, **kw)


# ── 页面模板 ─────────────────────────────────────────────────────────────

def title_slide(title, subtitle=""):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(s, 0, 0, SLD_W, SLD_H, BLUE_DARK)
    add_rect(s, 0, Inches(3.0), SLD_W, Inches(0.08), BLUE_LIGHT)
    add_textbox(s, Inches(1), Inches(1.5), Inches(11), Inches(1.5),
                title, size=40, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    if subtitle:
        add_textbox(s, Inches(1), Inches(3.5), Inches(11), Inches(1.5),
                    subtitle, size=22, color=RGBColor(0xAA,0xCC,0xEE),
                    align=PP_ALIGN.CENTER)
    return s

def section_slide(title):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(s, 0, 0, SLD_W, SLD_H, BLUE_MID)
    add_textbox(s, Inches(1), Inches(2.8), Inches(11), Inches(2),
                title, size=36, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    return s

def content_slide(title):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.background; bg.fill.solid(); bg.fill.fore_color.rgb = WHITE
    add_rect(s, 0, 0, SLD_W, Inches(0.85), BLUE_DARK)
    add_textbox(s, Inches(0.5), Inches(0.1), Inches(12), Inches(0.65),
                title, size=24, color=WHITE, bold=True)
    return s


# ══════════════════════════════════════════════════════════════════════════
#  幻灯片内容
# ══════════════════════════════════════════════════════════════════════════

# ── 1. 封面 ──────────────────────────────────────────────────────────────
title_slide("PostgreSQL 内存分配与缓存性能实验",
            "内存 → 缓存命中率 → 执行时间 的量化分析\n2026 年 5 月")

# ── 2. 实验概述 ──────────────────────────────────────────────────────────
s = content_slide("实验概述")
add_bullets(s, Inches(0.5), Inches(1.2), Inches(6), Inches(5.5), [
    "目标：量化 PostgreSQL 内存分配对查询性能的影响",
    "建立完整因果链：内存 → 缓存命中率 → SQL 执行时间",
    "",
    "环境：",
    "  ● PostgreSQL 12, Linux 5.4, NVMe SSD",
    "  ● 通过 cgroup v1 精确控制进程内存上限",
    "",
    "数据集：",
    "  ● TPC-H SF10 (~14GB) — OLAP 分析负载",
    "  ● TPC-C 100 warehouse (~10GB) — OLTP 事务负载",
    "  ● sysbench 10×2GB 表 (~20GB) — 通用负载",
], size=18)
add_bullets(s, Inches(6.8), Inches(1.4), Inches(6), Inches(5), [
    "四个实验层层递进：",
    "",
    "Exp1: 测量单次磁盘读延迟",
    "  → 每次 cache miss 的代价是多少？",
    "",
    "Exp2: 测量当前 cache miss 率",
    "  → 现在有多少 miss？TP vs AP？",
    "",
    "Exp3: 预测不同内存下的 miss 率",
    "  → 增大内存能减少多少 miss？",
    "",
    "Exp4: 直接测量执行时间变化",
    "  → 实际能快多少？验证预测",
], size=17, color=GRAY_MID)

# ── 3. 实验关联 ──────────────────────────────────────────────────────────
s = content_slide("实验逻辑关系")
boxes = [
    ("Exp1\n磁盘读延迟",      Inches(1.5), Inches(2.0), BLUE_LIGHT),
    ("Exp2\nCache Miss 基线", Inches(1.5), Inches(4.5), RGBColor(0x27,0xAE,0x60)),
    ("Exp3\nSBPX 预测",       Inches(5.5), Inches(3.0), RGBColor(0xF3,0x9C,0x12)),
    ("Exp4\n实测执行时间",    Inches(9.5), Inches(3.0), ACCENT_RED),
]
for label, x, y, clr in boxes:
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y,
                            Inches(2.5), Inches(1.5))
    sh.fill.solid(); sh.fill.fore_color.rgb = clr; sh.line.fill.background()
    tf = sh.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = label; p.font.size = Pt(16)
    p.font.color.rgb = WHITE; p.font.bold = True
    p.font.name = 'Microsoft YaHei'; p.alignment = PP_ALIGN.CENTER
add_textbox(s, Inches(4.0), Inches(2.2), Inches(1.5), Inches(0.5),
            "→ 延迟输入", size=13, color=GRAY_MID)
add_textbox(s, Inches(4.0), Inches(4.5), Inches(1.5), Inches(0.5),
            "→ miss 基线", size=13, color=GRAY_MID)
add_textbox(s, Inches(8.0), Inches(2.8), Inches(1.5), Inches(0.5),
            "→ 对比验证", size=13, color=GRAY_MID)
add_bullets(s, Inches(0.5), Inches(6.2), Inches(12), Inches(1), [
    "每个实验都在 TPC-H/C 和 sysbench 两套数据集上运行，确保结论的普适性"
], size=15, color=GRAY_MID)

# ══════════════════════════════════════════════════════════════════════════
#  Exp1
# ══════════════════════════════════════════════════════════════════════════
section_slide("实验一：Buffer Victimization + 磁盘读延迟")

s = content_slide("Exp1: 磁盘读延迟 — 测什么？为什么测？")
add_bullets(s, Inches(0.5), Inches(1.2), Inches(5.8), Inches(5.5), [
    "问题：每次 cache miss 到底要付出多大代价？",
    "",
    "PostgreSQL 读数据的过程：",
    "  1. 先查 shared_buffers（内存缓存）",
    "  2. 如果命中 → 直接返回，耗时 ~0",
    "  3. 如果未命中（cache miss）→ 调用 pread64()",
    "     从磁盘读取 8KB 的数据页",
    "",
    "Exp1 就是测量这个 pread64() 每次花多久",
    "→ 这就是 \"每次 cache miss 的代价\"",
    "",
    "为什么重要：",
    "  后续实验要回答 \"减少 N 次 miss 能省多少时间\"",
    "  答案 = N × 本实验测到的单次延迟",
], size=17)
add_bullets(s, Inches(6.8), Inches(1.2), Inches(6), Inches(5.5), [
    "怎么测：",
    "  ● strace -e pread64 挂载到 PG 进程",
    "  ● 记录每次 pread64() 调用的耗时（微秒）",
    "  ● 清除 OS page cache 后执行全表扫描",
    "    → 强制产生大量磁盘读",
    "",
    "分两组跑：",
    "  ● TPC-C: stock 表扫描（~3.6GB）",
    "    → 数据分布随机，产生随机 I/O",
    "  ● TPC-H: lineitem 表扫描（~10GB）",
    "    → 数据顺序存储，产生顺序 I/O",
    "",
    "  对比：随机 vs 顺序 I/O 的延迟差异",
], size=16)

s = content_slide("Exp1: TPC-C — 随机 I/O 延迟分布")
add_bullets(s, Inches(0.3), Inches(1.0), Inches(12), Inches(2), [
    "stock 全表扫描 (3.6GB, 随机布局), 54,636 次磁盘读  |  左: 延迟直方图 (大部分 <10μs 命中 OS cache)  右: 累积分布",
    "解读: 50% 的读在 6μs 内完成 (命中 OS cache), 但最慢的 1% 需要 6ms+ (真正的磁盘寻道)",
], size=14)
add_img(s, 'exp1_tpcc_pread64.png', Inches(0.2), Inches(2.8), width=Inches(12.8))

s = content_slide("Exp1: TPC-H — 顺序 I/O 延迟分布")
add_bullets(s, Inches(0.3), Inches(1.0), Inches(12), Inches(2), [
    "lineitem 全表扫描 (10GB, 顺序布局), 147,478 次磁盘读  |  顺序读让 OS 提前预读 (readahead)",
    "解读: p99=5411μs, 比 TPC-C 的 6277μs 低 14% — 顺序扫描的尾部延迟更低",
], size=14)
add_img(s, 'exp1_tpch_pread64.png', Inches(0.2), Inches(2.8), width=Inches(12.8))

s = content_slide("Exp1: 随机 vs 顺序 I/O 延迟对比")
add_bullets(s, Inches(0.3), Inches(1.0), Inches(12), Inches(2), [
    "左: CDF 叠加 — 两者 p50 相同 (6μs), 但 TPC-C 的曲线在尾部更靠右 (慢读更多)",
    "右: 百分位对比 — p99 差异最明显: TPC-C 随机 I/O 比 TPC-H 顺序 I/O 慢 ~16%",
    "结论: cache miss 的代价取决于 I/O 模式, OLAP 顺序扫描的单次 miss 代价反而更低",
], size=14)
add_img(s, 'exp1_tpcc_vs_tpch.png', Inches(0.2), Inches(3.2), width=Inches(12.8))

s = content_slide("Exp1: 磁盘读延迟 — sysbench 结果")
add_bullets(s, Inches(0.3), Inches(1.0), Inches(5), Inches(2), [
    "151,808 次 pread64 (sbtest1 全表扫描, cold start)",
    "p50=9μs, p75=12μs, p90=16μs, p99=2100μs",
    "NVMe SSD + OS page cache 回填极快",
    "真实磁盘延迟集中在 p99 (~2ms)",
], size=14)
add_img(s, 'sysbench_exp1_pread64.png', Inches(0.2), Inches(3.0), width=Inches(12.8))

# ══════════════════════════════════════════════════════════════════════════
#  Exp2
# ══════════════════════════════════════════════════════════════════════════
section_slide("实验二：Cache Miss 率测量")

s = content_slide("Exp2: Cache Miss 率 — 目标与方法")
add_bullets(s, Inches(0.5), Inches(1.2), Inches(6), Inches(5.5), [
    "目标：量化当前 128MB 配置下的 cache miss 严重程度",
    "  ● 分别测量 TP (OLTP) 和 AP (OLAP) 负载的 miss 率",
    "  ● 评估混合负载下 TP/AP 的相互干扰",
    "",
    "Cache Miss 定义：",
    "  miss_rate = blks_read / (blks_hit + blks_read)",
    "  blks_hit = shared_buffers 命中",
    "  blks_read = 需从 OS/磁盘读入",
], size=18)
add_bullets(s, Inches(6.8), Inches(1.2), Inches(6), Inches(5.5), [
    "TPC 版本：三阶段测量",
    "  Phase 1: TPC-C only (60s)",
    "  Phase 2: TPC-H only (60s)",
    "  Phase 3: 混合 (120s)",
    "",
    "sysbench 版本：五阶段测量",
    "  Phase 1: oltp_read_only (30s)",
    "  Phase 2: oltp_write_only (30s)",
    "  Phase 3: 混合读写 (30s)",
    "  Phase 4: 全表扫描 AP 风格 (cold)",
    "  Phase 5: 扫描后 OLTP 恢复",
], size=16)

s = content_slide("Exp2: Cache Miss — TPC-H/C 结果")
add_bullets(s, Inches(0.3), Inches(1.0), Inches(5), Inches(1.8), [
    "TPC-C (TP): miss 率仅 4.87% — 热点数据驻留",
    "TPC-H (AP): miss 率 100% — 工作集远超 buffer pool",
    "混合负载: TP miss 仅上升 0.10pp — usagecount 保护",
], size=14)
add_img(s, 'fig3_cachemiss_by_phase.png', Inches(0.5), Inches(2.8), width=Inches(12))

s = content_slide("Exp2: Cache Miss — sysbench 结果")
add_bullets(s, Inches(0.3), Inches(1.0), Inches(6), Inches(1.8), [
    "OLTP 读: 8.04% | OLTP 写: 8.56% | 混合: 6.29%",
    "全表扫描 (cold): 95.09% — 与 TPC-H 一致",
    "扫描后 OLTP 恢复: 5.48% — buffer pool 快速回热",
    "sysbench miss 率 (8%) > TPC-C (5%): 10 张表访问更分散",
], size=14)
add_img(s, 'sysbench_exp2_cachemiss.png', Inches(0.5), Inches(2.8), width=Inches(12))

# ══════════════════════════════════════════════════════════════════════════
#  Exp3
# ══════════════════════════════════════════════════════════════════════════
section_slide("实验三：SBPX 预测")

s = content_slide("Exp3: SBPX — 目标与方法")
add_bullets(s, Inches(0.5), Inches(1.2), Inches(6), Inches(5.5), [
    "目标：不实际增大内存，预测增大后的 miss 率变化",
    "",
    "核心理论：Stack Distance（重用距离）",
    "  ● 每次 buffer 请求，计算自上次访问同一 page",
    "    以来有多少不同 page 被访问过",
    "  ● Miss Ratio Curve: miss_rate(B) = P(距离 > B)",
    "  ● 只需一次 trace → 预测任意 buffer size 的 miss 率",
    "",
    "SHARDS 近似算法（1% 采样，误差 < 1%）：",
    "  ● FNV-1a hash 采样，O(N log N) 计算",
    "  ● 内存开销降低 100x，适合生产环境",
], size=17)
add_bullets(s, Inches(6.8), Inches(1.2), Inches(6), Inches(5.5), [
    "预测公式：",
    "  saved_time = Δmiss_count × avg_disk_latency",
    "  （磁盘延迟来自 Exp1 测量值）",
    "",
    "采集方式：",
    "  ● pg_buffercache 轮询（500ms 间隔）",
    "  ● 比较相邻快照推断 hit/miss",
    "",
    "验证方式：",
    "  ● 与 Exp4 实测数据对比",
], size=16)

s = content_slide("Exp3: SBPX — TPC-H/C 结果")
add_bullets(s, Inches(0.3), Inches(1.0), Inches(5.5), Inches(1.8), [
    "TPC-C: miss 从 41%(64MB) → 36%(128MB) → 34%(256MB+) 趋平",
    "TPC-H: miss 始终 ~99% — 增大 buffer 对全表扫描无效",
    "TPC-C 工作集 ~1130MB | 4GB 可节省约 4.2s",
], size=14)
add_img(s, 'fig5_sbpx_mrc.png', Inches(0.3), Inches(2.8), width=Inches(12.5))

s = content_slide("Exp3: SBPX — sysbench 结果")
add_bullets(s, Inches(0.3), Inches(1.0), Inches(6), Inches(1.8), [
    "4.3M page access 采集 (60s oltp_read_only)",
    "128MB → 32% miss, 256MB → 26% miss, 512MB+ → 25.7% (趋平)",
    "工作集 ~8.6GB, 当前 128MB 仅覆盖 1.5%",
    "256MB 以上收益递减 — OLTP 热点已驻留",
], size=14)
add_img(s, 'sysbench_exp3_sbpx.png', Inches(0.3), Inches(2.8), width=Inches(12.5))

# ══════════════════════════════════════════════════════════════════════════
#  Exp4
# ══════════════════════════════════════════════════════════════════════════
section_slide("实验四：SQL 执行时间 vs 内存分配")

s = content_slide("Exp4: 执行时间 vs 内存 — 目标与方法")
add_bullets(s, Inches(0.5), Inches(1.2), Inches(6), Inches(5), [
    "目标：直接回答——给 PG 多分配 X MB 内存，",
    "      查询实际能快多少？",
    "",
    "方法：cgroup 精确控制总内存",
    "  1. 修改 shared_buffers (= total_mem / 4)",
    "  2. 重启 PG, 加入 cgroup",
    "  3. 清除 OS page cache (cold start)",
    "  4. 运行查询, \\timing 计时",
    "  5. 采集 blks_hit / blks_read 计算 miss 率",
], size=17)
add_bullets(s, Inches(6.8), Inches(1.2), Inches(6), Inches(5.5), [
    "TPC 实验规模：",
    "  ● 20 个内存档位 (256MB → 24GB)",
    "  ● 12 条查询 (TPC-H Q1/Q3/Q5/Q6/Q9/Q12/Q13/Q18",
    "    + TPC-C STOCK/WH/OL/CUST)",
    "  ● 每配置 3 次重复",
    "",
    "sysbench 实验规模：",
    "  ● 28 个内存档位 (128MB → 24GB, 64MB 步进)",
    "  ● 16 条查询 (全表聚合/排序/JOIN/窗口函数/",
    "    CTE/子查询/DISTINCT 等)",
    "  ● 每配置 5 次重复 (cold/warm/hot)",
    "",
    "拟合：T(M) = a × M^(-b) + c（幂律衰减）",
], size=15)

s = content_slide("Exp4: TPC-H 曲线拟合 — T(M) = a·M⁻ᵇ + c")
add_bullets(s, Inches(0.3), Inches(1.0), Inches(12), Inches(1.2), [
    "每条查询的执行时间 vs 内存分配散点图 + 幂律拟合曲线，阴影区间 = min-max 范围",
], size=14)
add_img(s, 'curve_fit_tpch.png', Inches(0.3), Inches(2.2), width=Inches(12.5))

s = content_slide("Exp4: TPC-C 曲线拟合 — T(M) = a·M⁻ᵇ + c")
add_bullets(s, Inches(0.3), Inches(1.0), Inches(12), Inches(1.2), [
    "TPC-C 查询以全表扫描/排序为主，执行时间随内存增大显著下降，R² > 0.95",
], size=14)
add_img(s, 'curve_fit_tpcc.png', Inches(0.3), Inches(2.2), width=Inches(12.5))

s = content_slide("Exp4: 加速比总览 — TPC-H/C")
add_bullets(s, Inches(0.3), Inches(1.0), Inches(12), Inches(1.2), [
    "以最小内存为基准，各查询的加速比随内存增长曲线 — I/O 密集查询可达 2-5x 加速",
], size=14)
add_img(s, 'speedup_all.png', Inches(1.0), Inches(2.2), width=Inches(11))

s = content_slide("Exp4: sysbench 曲线拟合 (28 档, 16 查询)")
add_bullets(s, Inches(0.3), Inches(1.0), Inches(12), Inches(1.2), [
    "sysbench 数据集验证：均匀分布数据下执行时间同样符合幂律衰减，b 值反映内存敏感度",
], size=14)
add_img(s, 'sysbench_curve_fit.png', Inches(0.1), Inches(2.2), width=Inches(13))

s = content_slide("Exp4: sysbench 加速比")
add_bullets(s, Inches(0.3), Inches(1.0), Inches(12), Inches(1.2), [
    "全表扫描/多表 JOIN 加速显著，点查/索引查询不受内存影响 — 结论与 TPC 一致",
], size=14)
add_img(s, 'sysbench_speedup.png', Inches(0.5), Inches(2.2), width=Inches(12))

s = content_slide("Exp4: Cache Miss 率 vs 内存分配")
add_img(s, 'miss_rate_vs_memory.png',  Inches(0.2), Inches(1.2), width=Inches(6.3))
add_img(s, 'sysbench_miss_rate.png',   Inches(6.7), Inches(1.2), width=Inches(6.3))
add_bullets(s, Inches(0.3), Inches(5.8), Inches(12.5), Inches(1.5), [
    "左: TPC-H/C — TPC-H miss 从 99% 降至 20% (16GB), TPC-C 从 99% 快速下降",
    "右: sysbench — miss 率从 99% 逐步下降, 存在明显\"膝盖点\" (~30-50% 工作集覆盖时)",
], size=14)

# ══════════════════════════════════════════════════════════════════════════
#  核心发现
# ══════════════════════════════════════════════════════════════════════════
s = content_slide("核心发现总结")
findings = [
    ("1. 幂律衰减",
     "I/O 密集型查询的执行时间随内存符合 T(M) = a·M⁻ᵇ+c，R² > 0.95"),
    ("2. 收益递减",
     "覆盖 ~50% 工作集后，继续增加内存的边际收益快速下降"),
    ("3. 查询类型决定收益",
     "全表扫描/大表 JOIN: 2-5x 加速 | 索引查询: < 1.1x，几乎无影响"),
    ("4. SBPX 预测有效",
     "预测 miss 率与实测吻合，可用于生产环境的内存规划"),
    ("5. TP/AP 缓存竞争有限",
     "TPC-C miss 率在混合负载下仅上升 0.1pp，clock-sweep 保护有效"),
    ("6. 结论跨数据集一致",
     "TPC-H/C 和 sysbench 两套数据集得出相同结论"),
]
y = Inches(1.3)
for title, desc in findings:
    add_rect(s, Inches(0.4), y, Inches(12.5), Inches(0.85),
             RGBColor(0xF5,0xF7,0xFA))
    add_textbox(s, Inches(0.6), y + Inches(0.05), Inches(4), Inches(0.4),
                title, size=16, color=BLUE_DARK, bold=True)
    add_textbox(s, Inches(0.6), y + Inches(0.4), Inches(12), Inches(0.4),
                desc, size=14, color=GRAY_MID)
    y += Inches(0.95)

# ── 下一步 ───────────────────────────────────────────────────────────────
s = content_slide("下一步计划")
add_bullets(s, Inches(0.5), Inches(1.3), Inches(12), Inches(5.5), [
    "1. exp4_mem_sweep：使用 EXPLAIN ANALYZE 获取精确 buffer 计数",
    "   + benchbase 测量 TPC-C TPS + 混合负载下的吞吐量衰减",
    "",
    "2. 结合 SBPX 预测与实测数据，构建内存分配推荐模型",
    "   → 输入工作集大小与 SLA 要求，输出推荐 shared_buffers",
    "",
    "3. 验证 adaptive buffer pool sizing 策略",
    "   → 根据实时 miss rate 动态调整 shared_buffers",
    "",
    "4. 多租户场景下的 buffer pool 隔离实验",
    "   → TP/AP 隔离 vs 共享，量化干扰与隔离的 trade-off",
], size=18)

# ── 结尾 ─────────────────────────────────────────────────────────────────
title_slide("谢谢", "Q & A")

# ── 保存 ─────────────────────────────────────────────────────────────────
prs.save(OUTPUT)
n = len(prs.slides)
sz = os.path.getsize(OUTPUT) / 1024 / 1024
print(f"PPT saved: {OUTPUT}")
print(f"Slides: {n}  |  Size: {sz:.1f} MB")
