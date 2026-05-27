#!/usr/bin/env python3
"""
generate_validation_ppt.py — 预测验证分析 PPT
验证 Exp1+2+3 预测是否与 Exp4 实测吻合
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

FIGURES = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'figures')
OUTPUT  = os.path.join(os.path.dirname(os.path.abspath(__file__)), '预测验证分析.pptx')

BLUE_DARK  = RGBColor(0x1B, 0x3A, 0x5C)
BLUE_LIGHT = RGBColor(0x3A, 0x7C, 0xBD)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_DARK  = RGBColor(0x33, 0x33, 0x33)
GRAY_MID   = RGBColor(0x66, 0x66, 0x66)
ACCENT_RED = RGBColor(0xC0, 0x39, 0x2B)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SLD_W, SLD_H = prs.slide_width, prs.slide_height

def add_rect(s, l, t, w, h, c):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = c
    sh.line.fill.background(); sh.shadow.inherit = False; return sh

def tb(s, l, t, w, h, txt, sz=18, c=GRAY_DARK, b=False, a=PP_ALIGN.LEFT):
    box = s.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame; tf.word_wrap = True; p = tf.paragraphs[0]
    p.text = txt; p.font.size = Pt(sz); p.font.color.rgb = c
    p.font.bold = b; p.font.name = 'Microsoft YaHei'; p.alignment = a
    return box

def bl(s, l, t, w, h, items, sz=16, c=GRAY_DARK):
    box = s.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item; p.font.size = Pt(sz); p.font.color.rgb = c
        p.font.name = 'Microsoft YaHei'; p.space_after = Pt(sz * 0.35)
    return box

def img(s, name, l, t, w=None):
    path = os.path.join(FIGURES, name)
    if not os.path.exists(path): return None
    if w:
        return s.shapes.add_picture(path, left=l, top=t, width=w)
    return s.shapes.add_picture(path, left=l, top=t)

def title_slide(title, subtitle=""):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(s, 0, 0, SLD_W, SLD_H, BLUE_DARK)
    add_rect(s, 0, Inches(3.0), SLD_W, Inches(0.08), BLUE_LIGHT)
    tb(s, Inches(1), Inches(1.5), Inches(11), Inches(1.5),
       title, sz=40, c=WHITE, b=True, a=PP_ALIGN.CENTER)
    if subtitle:
        tb(s, Inches(1), Inches(3.5), Inches(11), Inches(1.5),
           subtitle, sz=22, c=RGBColor(0xAA,0xCC,0xEE), a=PP_ALIGN.CENTER)
    return s

def content_slide(title):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.background; bg.fill.solid(); bg.fill.fore_color.rgb = WHITE
    add_rect(s, 0, 0, SLD_W, Inches(0.85), BLUE_DARK)
    tb(s, Inches(0.5), Inches(0.1), Inches(12), Inches(0.65),
       title, sz=24, c=WHITE, b=True)
    return s


# ═══════════════════════════════════════════════════════════════════════
# Slide 1: 封面
# ═══════════════════════════════════════════════════════════════════════
title_slide("预测验证分析",
            "Exp1+2+3 的预测结果 vs Exp4 的实测结果\n模型准不准？哪里不准？为什么？")

# ═══════════════════════════════════════════════════════════════════════
# Slide 2: 预测公式
# ═══════════════════════════════════════════════════════════════════════
s = content_slide("预测公式：三个实验各贡献什么")
bl(s, Inches(0.5), Inches(1.2), Inches(12), Inches(5.5), [
    "核心公式：",
    "",
    "    预测执行时间(M) = 基准时间 - 节省时间",
    "",
    "    节省时间 = (基准miss率 - 预测miss率) × 总访问次数 × 每次miss代价",
    "",
    "三个实验各贡献一个参数：",
    "",
    "    Exp1 → 每次 miss 代价（μs/miss）",
    "        TPC-C: 375.8 μs    TPC-H: 285.9 μs    sysbench: 88.6 μs",
    "",
    "    Exp2 → 基准 miss 率（当前 128MB 下）",
    "        TPC-C: 99.9%    TPC-H: 98.8%    sysbench: 8.0%",
    "",
    "    Exp3 (SBPX) → 预测 miss 率（增大内存后会变成多少）",
    "",
    "    Exp4 → 实际测量执行时间 → 和预测比，看准不准",
], sz=17)

# ═══════════════════════════════════════════════════════════════════════
# Slide 3: TPC-H 预测 vs 实测
# ═══════════════════════════════════════════════════════════════════════
s = content_slide("TPC-H：预测 vs 实测（偏离严重）")
bl(s, Inches(0.3), Inches(1.0), Inches(12), Inches(1.2), [
    "左: 蓝=实测, 红=预测, 红色阴影=预测误差 | 右: 每个内存档位的误差百分比 (绿<10%, 黄<30%, 红>30%)",
], sz=14)
img(s, 'validation_tpch.png', Inches(0.2), Inches(2.0), Inches(12.8))

# ═══════════════════════════════════════════════════════════════════════
# Slide 4: TPC-C 预测 vs 实测
# ═══════════════════════════════════════════════════════════════════════
s = content_slide("TPC-C：预测 vs 实测（几乎完美）")
bl(s, Inches(0.3), Inches(1.0), Inches(12), Inches(1.2), [
    "左: 两条线完全重叠 (Y 轴放大到 ±5% 才能看出差异) | 右: 误差全在 ±1% 以内",
], sz=14)
img(s, 'validation_tpcc.png', Inches(0.2), Inches(2.0), Inches(12.8))

# ═══════════════════════════════════════════════════════════════════════
# Slide 5: 误差根因
# ═══════════════════════════════════════════════════════════════════════
s = content_slide("根因：每次 miss 的代价不是常数")
bl(s, Inches(0.3), Inches(1.0), Inches(12), Inches(1.2), [
    "蓝虚线 = Exp1 测的常数代价 (286μs) | 红线 = 从 Exp4 反推出的真实代价 → 内存越大，代价越高",
], sz=14)
img(s, 'validation_cost_not_constant.png', Inches(1.5), Inches(2.0), Inches(10))

# ═══════════════════════════════════════════════════════════════════════
# Slide 6: 结果解读
# ═══════════════════════════════════════════════════════════════════════
s = content_slide("结果解读")
bl(s, Inches(0.5), Inches(1.2), Inches(12), Inches(5.5), [
    "TPC-C：预测误差 < 1%，几乎完美",
    "  ● 原因：miss 率始终 99.9%，内存变化对它没影响",
    "  ● 公式里的 Δmiss 约等于 0，所以预测时间 ≈ 基准时间 → 自然准确",
    "  ● 但这也说明：TPC-C 的查询对 shared_buffers 不敏感",
    "",
    "TPC-H：小内存时准确（<5% 误差），大内存时严重高估收益",
    "  ● 8GB 时预测误差 -75%：预测说应该快很多，但实际只快了一点",
    "  ● 12GB 时预测时间甚至变成负数 — 模型已经崩溃了",
    "",
    "问题出在哪里？",
], sz=17)

# ═══════════════════════════════════════════════════════════════════════
# Slide 7: 为什么预测不准
# ═══════════════════════════════════════════════════════════════════════
s = content_slide("为什么 TPC-H 大内存时预测不准？")
bl(s, Inches(0.5), Inches(1.2), Inches(5.8), Inches(5.5), [
    "根本原因：每次 miss 的代价不是常数",
    "",
    "Exp1 测的 cost_per_miss = 285.9 μs",
    "这是小内存(128MB)下的平均值，包含了：",
    "  ● 64% 命中 OS cache → ~10 μs",
    "  ● 36% 真正读磁盘   → ~920 μs",
    "",
    "但当 shared_buffers 从 128MB 增大到 8GB 后：",
    "  ● 容易命中的\"热页\"已经被 buffer pool 接住了",
    "  ● 剩下还 miss 的都是\"冷页\"",
    "  ● 冷页在 OS cache 里也没有",
    "  ● → 100% 打磁盘，每次 ~920 μs",
], sz=17)
bl(s, Inches(6.8), Inches(1.2), Inches(6), Inches(5.5), [
    "也就是说：",
    "",
    "小内存时：miss 的页有很多是\"温\"的，",
    "  OS cache 能兜底，代价低（~286 μs）",
    "",
    "大内存时：miss 的页都是\"冰冷\"的，",
    "  OS cache 也没有，代价高（~920 μs）",
    "",
    "模型的假设 \"cost_per_miss 是常数\" 不成立",
    "→ 用小内存下测的低代价去乘大内存下",
    "   减少的 miss 数 → 高估了节省的时间",
    "",
    "修正方向：cost_per_miss 应该是 miss 率的函数",
    "  miss 率越低 → 剩余 miss 越\"冷\" → 代价越高",
], sz=16, c=GRAY_MID)

# ═══════════════════════════════════════════════════════════════════════
# Slide 8: 总结
# ═══════════════════════════════════════════════════════════════════════
s = content_slide("预测验证总结")
y = Inches(1.3)
findings = [
    ("公式", "预测时间 = 基准时间 - Δmiss × 总请求 × 代价/miss"),
    ("TPC-C", "预测误差 < 1% — miss 率不变，预测自然准确（但无意义）"),
    ("TPC-H 小内存", "预测误差 < 5% — OS cache 行为稳定，代价近似常数"),
    ("TPC-H 大内存", "预测严重偏离 — cost_per_miss 不是常数，随 miss 率下降而上升"),
    ("模型局限", "\"每次 miss 代价相同\" 的假设在大内存时不成立"),
    ("修正方向", "需要建立 cost_per_miss = f(miss_rate) 的动态模型"),
]
for title_text, desc in findings:
    add_rect(s, Inches(0.4), y, Inches(12.5), Inches(0.85),
             RGBColor(0xF5,0xF7,0xFA))
    tb(s, Inches(0.6), y + Inches(0.05), Inches(3.5), Inches(0.4),
       title_text, sz=16, c=BLUE_DARK, b=True)
    tb(s, Inches(0.6), y + Inches(0.4), Inches(12), Inches(0.4),
       desc, sz=14, c=GRAY_MID)
    y += Inches(0.95)


# ── 保存 ─────────────────────────────────────────────────────────────
prs.save(OUTPUT)
n = len(prs.slides)
sz = os.path.getsize(OUTPUT) / 1024 / 1024
print(f"PPT saved: {OUTPUT}")
print(f"Slides: {n}  |  Size: {sz:.1f} MB")
